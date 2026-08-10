import os
import difflib
import shutil
import subprocess
import re
import csv
from pathlib import Path
from datetime import datetime

from llm_local import ask_llm
from tools import find_project


# ==========================
# CONFIG
# ==========================
TEXT_EXTENSIONS = {
    ".py", ".js", ".jsx", ".ts", ".tsx",
    ".html", ".css", ".scss", ".sass",
    ".json", ".md", ".txt", ".yml", ".yaml",
    ".env", ".ini", ".cfg", ".toml",
    ".xml", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".php", ".rb", ".go", ".rs",
    ".sql", ".bat", ".ps1", ".sh",
    ".dockerfile", ".gitignore", ".log"
}

SKIP_DIRS = {
    "node_modules",
    "venv",
    ".venv",
    "jarvis-env",
    "__pycache__",
    ".git",
    ".idea",
    ".vscode",
    "dist",
    "build",
    ".next",
    ".cache",
    "site-packages"
}

MAX_READ_CHARS = 12000
MAX_LLM_CHARS = 10000
MAX_SEARCH_RESULTS = 30
BACKUP_ROOT = "file_backups"

HUD_PROJECT_FILE = "hud_project.txt"
HUD_CURRENT_FILE = "hud_current_file.txt"
HUD_ACTION_FILE = "hud_action.txt"
HUD_AI_STATUS_FILE = "hud_ai_status.txt"
HUD_RESULT_FILE = "hud_result.txt"


def write_hud_file(path, value):
    try:
        with open(path, "w", encoding="utf-8", errors="ignore") as f:
            f.write(str(value))
    except Exception:
        pass


def update_hud_context(
    project=None,
    file_path=None,
    action=None,
    ai_status=None,
    result=None
):
    if project is not None:
        write_hud_file(HUD_PROJECT_FILE, project)

    if file_path is not None:
        write_hud_file(HUD_CURRENT_FILE, file_path)

    if action is not None:
        write_hud_file(HUD_ACTION_FILE, action)

    if ai_status is not None:
        write_hud_file(HUD_AI_STATUS_FILE, ai_status)

    if result is not None:
        text = str(result).replace("\n", " ").strip()
        write_hud_file(HUD_RESULT_FILE, text[:220])




# ==========================
# HELPERS
# ==========================
def normalize_name(name):
    return "".join(
        ch for ch in str(name).lower()
        if ch.isalnum()
    )


def get_project(project_name):
    project = find_project(project_name)

    if not project:
        return None, f"Project not found: {project_name}"

    path = project.get("path")

    if not path or not os.path.exists(path):
        return None, f"Project path not found: {project_name}"

    return project, None


def is_skipped_path(path):
    parts = Path(path).parts

    for part in parts:
        if part.lower() in SKIP_DIRS:
            return True

    return False


def is_text_file(path):
    name = os.path.basename(path).lower()
    ext = os.path.splitext(path)[1].lower()

    if name in {
        "dockerfile",
        ".gitignore",
        ".env",
        "makefile",
        "readme"
    }:
        return True

    return ext in TEXT_EXTENSIONS


def collect_project_files(project_path):
    results = []

    for root, dirs, files in os.walk(project_path):
        dirs[:] = [
            d for d in dirs
            if d.lower() not in SKIP_DIRS
        ]

        for file in files:
            full_path = os.path.join(root, file)

            if is_skipped_path(full_path):
                continue

            rel_path = os.path.relpath(
                full_path,
                project_path
            )

            results.append({
                "name": file,
                "relative_path": rel_path,
                "full_path": full_path
            })

    return results


def find_project_file(project_name, file_query):
    project, error = get_project(project_name)

    if error:
        return None, error

    project_path = project["path"]
    files = collect_project_files(project_path)

    if not files:
        return None, "No files found in this project."

    query_norm = normalize_name(file_query)

    normalized_query_path = str(file_query).replace("/", os.sep).replace("\\", os.sep).lower()

    # 1. Exact relative path match
    for item in files:
        normalized_rel = item["relative_path"].replace("/", os.sep).replace("\\", os.sep).lower()

        if normalized_query_path == normalized_rel:
            return item, None

    # 2. Exact filename match
    for item in files:
        if file_query.lower() == item["name"].lower():
            return item, None

    # 3. Partial path / filename match
    partial_matches = []

    for item in files:
        rel_norm = normalize_name(item["relative_path"])
        name_norm = normalize_name(item["name"])

        if query_norm in rel_norm or query_norm in name_norm:
            partial_matches.append(item)

    if len(partial_matches) == 1:
        return partial_matches[0], None

    if len(partial_matches) > 1:
        preview = "\n".join(
            item["relative_path"]
            for item in partial_matches[:15]
        )

        return None, (
            "Multiple files matched. Be more specific:\n"
            f"{preview}"
        )

    # 4. Fuzzy match
    best_item = None
    best_score = 0

    for item in files:
        candidates = [
            normalize_name(item["name"]),
            normalize_name(item["relative_path"])
        ]

        for candidate in candidates:
            score = difflib.SequenceMatcher(
                None,
                query_norm,
                candidate
            ).ratio()

            if score > best_score:
                best_score = score
                best_item = item

    if best_item and best_score >= 0.60:
        return best_item, None

    return None, f"File not found: {file_query}"


def read_file_content(path, max_chars=MAX_READ_CHARS):
    if not os.path.exists(path):
        return None, "File not found."

    if not is_text_file(path):
        return None, (
            "This file does not look like a readable text/code file."
        )

    try:
        with open(
            path,
            "r",
            encoding="utf-8",
            errors="ignore"
        ) as f:
            return f.read(max_chars), None

    except Exception as e:
        return None, f"Could not read file: {e}"


def ensure_backup_dir():
    os.makedirs(BACKUP_ROOT, exist_ok=True)
    return BACKUP_ROOT


def create_project_file_backup(project_name, file_query):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    path = item["full_path"]

    if not os.path.exists(path):
        return f"Cannot create backup. File not found: {path}"

    ensure_backup_dir()

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_project = normalize_name(project_name) or "project"
    safe_file = normalize_name(os.path.basename(path)) or "file"
    backup_name = f"{safe_project}_{safe_file}_{timestamp}.bak"
    backup_path = os.path.join(BACKUP_ROOT, backup_name)

    try:
        shutil.copy2(path, backup_path)

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action="Backup created",
            ai_status="READY",
            result=backup_path
        )

        return (
            "Backup created successfully.\n"
            f"Original: {path}\n"
            f"Backup: {os.path.abspath(backup_path)}"
        )

    except Exception as e:
        return f"Could not create backup: {e}"


def safe_write_project_file(project_name, file_query, new_content):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    path = item["full_path"]

    if not str(new_content).strip():
        return "Refused: new content is empty."

    backup_result = create_project_file_backup(
        project_name,
        file_query
    )

    if not backup_result.startswith("Backup created successfully."):
        return backup_result

    try:
        with open(
            path,
            "w",
            encoding="utf-8",
            errors="ignore"
        ) as f:
            f.write(str(new_content))

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action="Safe write completed",
            ai_status="READY",
            result=path
        )

        return (
            "Safe write completed.\n"
            f"{backup_result}\n"
            f"Updated file: {path}"
        )

    except Exception as e:
        return f"Backup was created, but file write failed: {e}"


def restore_latest_project_file_backup(project_name, file_query):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    path = item["full_path"]

    if not os.path.isdir(BACKUP_ROOT):
        return "No backups folder found."

    safe_project = normalize_name(project_name) or "project"
    safe_file = normalize_name(os.path.basename(path)) or "file"
    prefix = f"{safe_project}_{safe_file}_"

    candidates = [
        os.path.join(BACKUP_ROOT, name)
        for name in os.listdir(BACKUP_ROOT)
        if name.startswith(prefix) and name.endswith(".bak")
    ]

    if not candidates:
        return f"No backup found for: {path}"

    candidates.sort(reverse=True)
    latest_backup = candidates[0]

    try:
        shutil.copy2(latest_backup, path)

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action="Backup restored",
            ai_status="READY",
            result=latest_backup
        )

        return (
            "Latest backup restored.\n"
            f"Backup: {os.path.abspath(latest_backup)}\n"
            f"Restored file: {path}"
        )

    except Exception as e:
        return f"Could not restore backup: {e}"


def list_project_file_backups():
    if not os.path.isdir(BACKUP_ROOT):
        return "No backups folder found."

    backups = [
        name for name in os.listdir(BACKUP_ROOT)
        if name.endswith(".bak")
    ]

    if not backups:
        return "No backups found."

    backups.sort(reverse=True)

    output = [f"Backups found: {len(backups)}"]

    for name in backups[:50]:
        output.append(f" - {name}")

    if len(backups) > 50:
        output.append(f"... and {len(backups) - 50} more")

    return "\n".join(output)



# ==========================
# RESOLVE PROJECT FILE WITH CONTEXT
# ==========================
def resolve_project_file_with_content(project_name, file_query, max_chars=MAX_LLM_CHARS):
    item, error = find_project_file(project_name, file_query)

    if error:
        update_hud_context(
            project=project_name,
            file_path=file_query,
            action="File resolve failed",
            ai_status="ERROR",
            result=error
        )
        return None, None, error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=max_chars
    )

    if read_error:
        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action="File read failed",
            ai_status="ERROR",
            result=read_error
        )
        return item, None, read_error

    update_hud_context(
        project=project_name,
        file_path=item["relative_path"],
        action="Project file resolved",
        ai_status="READY",
        result=item["full_path"]
    )

    return item, content, None


# ==========================
# OPEN PROJECT FILE
# ==========================
def open_project_file(project_name, file_query):
    item, error = find_project_file(
        project_name,
        file_query
    )

    if error:
        return error

    try:
        os.startfile(item["full_path"])

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action="Opening project file",
            ai_status="READY",
            result=item["full_path"]
        )

        return (
            f"Opening file:\n"
            f"{item['relative_path']}"
        )

    except Exception as e:
        return f"Could not open file: {e}"


# ==========================
# READ PROJECT FILE
# ==========================
def read_project_file(project_name, file_query):
    item, error = find_project_file(
        project_name,
        file_query
    )

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"]
    )

    if read_error:
        return read_error

    return (
        f"FILE: {item['relative_path']}\n"
        f"PATH: {item['full_path']}\n\n"
        f"{content}"
    )


# ==========================
# PREVIEW PROJECT FILE
# ==========================
def preview_project_file(project_name, file_query, max_chars=3000):
    item, error = find_project_file(
        project_name,
        file_query
    )

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=max_chars
    )

    if read_error:
        return read_error

    return (
        "Safe preview only. No changes made.\n"
        f"FILE: {item['relative_path']}\n"
        f"PATH: {item['full_path']}\n\n"
        f"{content}"
    )


# ==========================
# SEARCH PROJECT FILES
# ==========================
def search_project_files(project_name, keyword):
    project, error = get_project(project_name)

    if error:
        return error

    project_path = project["path"]
    keyword_lower = keyword.lower()

    results = []

    files = collect_project_files(project_path)

    for item in files:
        rel = item["relative_path"]

        if keyword_lower in rel.lower():
            results.append(
                f"{rel}  [path match]"
            )
            continue

        if not is_text_file(item["full_path"]):
            continue

        content, err = read_file_content(
            item["full_path"],
            max_chars=8000
        )

        if err or not content:
            continue

        if keyword_lower in content.lower():
            results.append(
                f"{rel}  [content match]"
            )

        if len(results) >= MAX_SEARCH_RESULTS:
            break

    if not results:
        return "No matching files found."

    return "\n".join(results)




# ==========================
# DEVELOPER LINE OPERATIONS
# ==========================
def read_project_file_lines(project_name, file_query, start_line=None, end_line=None, with_numbers=True):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return read_error

    lines = content.splitlines()

    if not lines:
        return f"File is empty: {item['relative_path']}"

    if start_line is None:
        start_line = 1

    if end_line is None:
        end_line = len(lines)

    try:
        start_line = max(1, int(start_line))
        end_line = min(len(lines), int(end_line))
    except Exception:
        return "Invalid line range."

    if start_line > end_line:
        return "Invalid line range."

    selected = lines[start_line - 1:end_line]

    if with_numbers:
        width = len(str(end_line))
        selected_text = "\n".join(
            f"{str(index).rjust(width)} | {line}"
            for index, line in enumerate(selected, start=start_line)
        )
    else:
        selected_text = "\n".join(selected)

    update_hud_context(
        project=project_name,
        file_path=item["relative_path"],
        action=f"Read lines {start_line}-{end_line}",
        ai_status="READY",
        result=item["relative_path"]
    )

    return (
        f"FILE: {item['relative_path']}\n"
        f"PATH: {item['full_path']}\n"
        f"LINES: {start_line}-{end_line}\n\n"
        f"{selected_text}"
    )


def extract_project_file_lines(project_name, file_query, start_line, end_line):
    item, error = find_project_file(project_name, file_query)

    if error:
        return None, error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return None, read_error

    lines = content.splitlines()

    try:
        start_line = max(1, int(start_line))
        end_line = min(len(lines), int(end_line))
    except Exception:
        return None, "Invalid line range."

    if start_line > end_line:
        return None, "Invalid line range."

    return "\n".join(lines[start_line - 1:end_line]), None


def replace_project_file_lines(project_name, file_query, start_line, end_line, new_code):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return read_error

    lines = content.splitlines()

    try:
        start_line = max(1, int(start_line))
        end_line = min(len(lines), int(end_line))
    except Exception:
        return "Invalid line range."

    if start_line > end_line:
        return "Invalid line range."

    backup_result = create_project_file_backup(project_name, file_query)

    if not backup_result.startswith("Backup created successfully."):
        return backup_result

    replacement_lines = str(new_code).splitlines()

    updated_lines = (
        lines[:start_line - 1]
        + replacement_lines
        + lines[end_line:]
    )

    try:
        with open(item["full_path"], "w", encoding="utf-8", errors="ignore") as f:
            f.write("\n".join(updated_lines) + "\n")

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action=f"Replaced lines {start_line}-{end_line}",
            ai_status="READY",
            result=item["relative_path"]
        )

        return (
            "Lines replaced safely.\n"
            f"File: {item['relative_path']}\n"
            f"Lines: {start_line}-{end_line}\n"
            f"{backup_result}"
        )

    except Exception as e:
        return f"Backup created, but line replacement failed: {e}"


def insert_project_file_lines(project_name, file_query, insert_at_line, new_code):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return read_error

    lines = content.splitlines()

    try:
        insert_at_line = max(1, int(insert_at_line))
    except Exception:
        return "Invalid insert line."

    insert_at_line = min(insert_at_line, len(lines) + 1)

    backup_result = create_project_file_backup(project_name, file_query)

    if not backup_result.startswith("Backup created successfully."):
        return backup_result

    new_lines = str(new_code).splitlines()

    updated_lines = (
        lines[:insert_at_line - 1]
        + new_lines
        + lines[insert_at_line - 1:]
    )

    try:
        with open(item["full_path"], "w", encoding="utf-8", errors="ignore") as f:
            f.write("\n".join(updated_lines) + "\n")

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action=f"Inserted code at line {insert_at_line}",
            ai_status="READY",
            result=item["relative_path"]
        )

        return (
            "Lines inserted safely.\n"
            f"File: {item['relative_path']}\n"
            f"Insert line: {insert_at_line}\n"
            f"{backup_result}"
        )

    except Exception as e:
        return f"Backup created, but line insert failed: {e}"


def copy_project_file_lines_between_projects(
    source_project,
    source_file,
    source_start,
    source_end,
    target_project,
    target_file,
    target_start,
    target_end=None,
    mode="replace"
):
    extracted, error = extract_project_file_lines(
        source_project,
        source_file,
        source_start,
        source_end
    )

    if error:
        return error

    if mode == "insert" or target_end is None:
        return insert_project_file_lines(
            target_project,
            target_file,
            target_start,
            extracted
        )

    return replace_project_file_lines(
        target_project,
        target_file,
        target_start,
        target_end,
        extracted
    )


def open_project_file_in_app(project_name, file_query, app_name):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    try:
        from tools import open_file_in_app
        return open_file_in_app(item["full_path"], app_name)
    except Exception:
        try:
            subprocess.Popen([app_name, item["full_path"]], shell=True)
            return f"Opening file in {app_name}: {item['relative_path']}"
        except Exception as e:
            return f"Could not open file in {app_name}: {e}"


def review_project_file_lines(project_name, file_query, start_line, end_line):
    code, error = extract_project_file_lines(
        project_name,
        file_query,
        start_line,
        end_line
    )

    if error:
        return error

    prompt = f"""
You are JARVIS, a senior software engineer and cybersecurity reviewer.

Review ONLY this code fragment.
Do not invent context outside the fragment.
If something depends on missing context, say so.

Project:
{project_name}

File:
{file_query}

Lines:
{start_line}-{end_line}

Code:
{code}

Return:
1. What this fragment does
2. Bugs or risky logic
3. Security risks
4. Maintainability issues
5. Concrete improvements
6. A better version of the fragment if useful
"""

    return run_project_file_llm(
        project_name,
        file_query,
        prompt,
        f"Review lines {start_line}-{end_line}"
    )


def improve_project_file_lines(project_name, file_query, start_line, end_line):
    code, error = extract_project_file_lines(
        project_name,
        file_query,
        start_line,
        end_line
    )

    if error:
        return error

    prompt = f"""
You are JARVIS, a senior developer.

Improve ONLY this code fragment.
Preserve behavior unless there is an obvious bug.
Return:
1. Main problems
2. Improved code fragment in a fenced code block
3. Short explanation

Project:
{project_name}

File:
{file_query}

Lines:
{start_line}-{end_line}

Code:
{code}
"""

    return run_project_file_llm(
        project_name,
        file_query,
        prompt,
        f"Improve lines {start_line}-{end_line}"
    )


def explain_project_file_lines(project_name, file_query, start_line, end_line):
    code, error = extract_project_file_lines(
        project_name,
        file_query,
        start_line,
        end_line
    )

    if error:
        return error

    prompt = f"""
You are JARVIS, a clear programming tutor.

Explain ONLY this code fragment in a practical way.

Project:
{project_name}
File:
{file_query}
Lines:
{start_line}-{end_line}

Code:
{code}
"""

    return run_project_file_llm(
        project_name,
        file_query,
        prompt,
        f"Explain lines {start_line}-{end_line}"
    )



# ==========================================================
# LOCAL REPORT GENERATOR
# Replaces old dependency on developer_assistant.py.
# Supports: md, txt, html, docx/doc, pdf, pptx/ppt, xlsx/xls/csv.
# ==========================================================
REPORT_DIR = "developer_reports"


def ensure_report_dir():
    os.makedirs(REPORT_DIR, exist_ok=True)
    return REPORT_DIR


def report_timestamp():
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def safe_report_name(value):
    value = str(value).strip()
    value = re.sub(r"[^A-Za-z0-9_.-]+", "_", value)
    return value.strip("_") or "report"


def open_report_path(path):
    try:
        os.startfile(os.path.abspath(path))
        return True
    except Exception:
        return False


def write_text_report(path, content):
    with open(path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(str(content))


def create_local_markdown_report(title, content, open_after=True):
    ensure_report_dir()

    path = os.path.join(
        REPORT_DIR,
        f"{safe_report_name(title)}_{report_timestamp()}.md"
    )

    body = (
        f"# {title}\n\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        f"{content}\n"
    )

    write_text_report(path, body)

    if open_after:
        open_report_path(path)

    return path


def create_local_html_report(title, content, open_after=True):
    import html

    ensure_report_dir()

    path = os.path.join(
        REPORT_DIR,
        f"{safe_report_name(title)}_{report_timestamp()}.html"
    )

    body = f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>{html.escape(str(title))}</title>
<style>
body {{
    font-family: Arial, sans-serif;
    background: #0b1020;
    color: #e7f6ff;
    margin: 40px;
}}
h1 {{ color: #58d8ff; }}
pre {{
    white-space: pre-wrap;
    background: #111a33;
    border: 1px solid #263a70;
    border-radius: 12px;
    padding: 20px;
}}
</style>
</head>
<body>
<h1>{html.escape(str(title))}</h1>
<p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
<pre>{html.escape(str(content))}</pre>
</body>
</html>"""

    write_text_report(path, body)

    if open_after:
        open_report_path(path)

    return path


def create_local_docx_report(title, content, open_after=True):
    ensure_report_dir()

    try:
        from docx import Document
    except Exception:
        return create_local_markdown_report(
            title + "_docx_fallback",
            "python-docx is not installed. Markdown fallback created.\n\n" + str(content),
            open_after=open_after
        )

    path = os.path.join(
        REPORT_DIR,
        f"{safe_report_name(title)}_{report_timestamp()}.docx"
    )

    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    for block in str(content).split("\n\n"):
        doc.add_paragraph(block)

    doc.save(path)

    if open_after:
        open_report_path(path)

    return path


def create_local_pdf_report(title, content, open_after=True):
    ensure_report_dir()

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception:
        return create_local_markdown_report(
            title + "_pdf_fallback",
            "reportlab is not installed. Markdown fallback created.\n\n" + str(content),
            open_after=open_after
        )

    path = os.path.join(
        REPORT_DIR,
        f"{safe_report_name(title)}_{report_timestamp()}.pdf"
    )

    c = canvas.Canvas(path, pagesize=A4)
    width, height = A4
    x = 42
    y = height - 50

    c.setFont("Helvetica-Bold", 15)
    c.drawString(x, y, str(title)[:90])
    y -= 24

    c.setFont("Helvetica", 8)

    for raw_line in str(content).splitlines():
        line = str(raw_line)

        chunks = [line[i:i + 115] for i in range(0, len(line), 115)] or [""]

        for chunk in chunks:
            if y < 45:
                c.showPage()
                c.setFont("Helvetica", 8)
                y = height - 45

            c.drawString(x, y, chunk)
            y -= 10

    c.save()

    if open_after:
        open_report_path(path)

    return path


def create_local_pptx_report(title, content, open_after=True):
    ensure_report_dir()

    try:
        from pptx import Presentation
    except Exception:
        return create_local_markdown_report(
            title + "_pptx_fallback",
            "python-pptx is not installed. Markdown fallback created.\n\n" + str(content),
            open_after=open_after
        )

    path = os.path.join(
        REPORT_DIR,
        f"{safe_report_name(title)}_{report_timestamp()}.pptx"
    )

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = str(title)[:80]
    slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

    sections = str(content).split("\n\n")

    for section in sections[:15]:
        section = section.strip()

        if not section:
            continue

        lines = section.splitlines()
        slide_title = lines[0][:70] if lines else "Report Section"
        slide_body = "\n".join(lines[1:]) if len(lines) > 1 else section

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        slide.placeholders[1].text = slide_body[:900]

    prs.save(path)

    if open_after:
        open_report_path(path)

    return path


def create_local_excel_report(title, content, project_name="", file_query="", open_after=True):
    ensure_report_dir()

    path = os.path.join(
        REPORT_DIR,
        f"{safe_report_name(title)}_{report_timestamp()}.xlsx"
    )

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill

        wb = Workbook()

        ws = wb.active
        ws.title = "Summary"
        ws["A1"] = title
        ws["A1"].font = Font(bold=True, size=16)
        ws["A3"] = "Project"
        ws["B3"] = project_name
        ws["A4"] = "File"
        ws["B4"] = file_query
        ws["A5"] = "Generated"
        ws["B5"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        ws["A7"] = "Report preview"
        ws["A7"].font = Font(bold=True)

        lines = str(content).splitlines()

        for row_index, line in enumerate(lines[:500], start=8):
            ws.cell(row=row_index, column=1).value = line

        ws.column_dimensions["A"].width = 120
        ws.column_dimensions["B"].width = 40

        ws2 = wb.create_sheet("Report Lines")
        ws2.append(["Line", "Text"])

        for cell in ws2[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="D9EAF7")

        for index, line in enumerate(lines, start=1):
            ws2.append([index, line])

        ws2.column_dimensions["A"].width = 12
        ws2.column_dimensions["B"].width = 140

        ws3 = wb.create_sheet("Findings")
        ws3.append(["Type", "Line"])

        for cell in ws3[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="FCE4D6")

        keywords = ["critical", "high", "medium", "low", "risk", "bug", "issue", "security", "todo", "fixme"]

        for line in lines:
            lower = line.lower()

            if any(keyword in lower for keyword in keywords):
                if "critical" in lower:
                    kind = "CRITICAL"
                elif "high" in lower:
                    kind = "HIGH"
                elif "medium" in lower:
                    kind = "MEDIUM"
                elif "low" in lower:
                    kind = "LOW"
                elif "security" in lower:
                    kind = "SECURITY"
                elif "bug" in lower:
                    kind = "BUG"
                else:
                    kind = "INFO"

                ws3.append([kind, line])

        ws3.column_dimensions["A"].width = 18
        ws3.column_dimensions["B"].width = 140

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    cell.alignment = Alignment(vertical="top", wrap_text=True)

        wb.save(path)

    except Exception:
        path = path.replace(".xlsx", ".csv")
        with open(path, "w", encoding="utf-8", errors="ignore", newline="") as f:
            writer = csv.writer(f)
            writer.writerow(["project", "file", "line_number", "text"])

            for index, line in enumerate(str(content).splitlines(), start=1):
                writer.writerow([project_name, file_query, index, line])

    if open_after:
        open_report_path(path)

    return path


def create_local_report(title, content, format_type="md", open_after=True, project_name="", file_query=""):
    fmt = str(format_type).lower().strip()

    if fmt in {"word", "doc", "docx"}:
        path = create_local_docx_report(title, content, open_after=open_after)
    elif fmt == "pdf":
        path = create_local_pdf_report(title, content, open_after=open_after)
    elif fmt in {"ppt", "pptx", "powerpoint", "presentation"}:
        path = create_local_pptx_report(title, content, open_after=open_after)
    elif fmt in {"excel", "xls", "xlsx", "spreadsheet"}:
        path = create_local_excel_report(
            title,
            content,
            project_name=project_name,
            file_query=file_query,
            open_after=open_after
        )
    elif fmt in {"html", "web"}:
        path = create_local_html_report(title, content, open_after=open_after)
    elif fmt in {"txt", "text"}:
        ensure_report_dir()
        path = os.path.join(REPORT_DIR, f"{safe_report_name(title)}_{report_timestamp()}.txt")
        write_text_report(path, content)

        if open_after:
            open_report_path(path)
    else:
        path = create_local_markdown_report(title, content, open_after=open_after)

    return f"Report created:\n{path}"

def generate_file_review_report(project_name, file_query, format_type="md"):
    content = analyze_project_file(project_name, file_query)

    return create_local_report(
        title=f"JARVIS File Review - {file_query}",
        content=content,
        format_type=format_type,
        open_after=True,
        project_name=project_name,
        file_query=file_query
    )


def generate_file_lines_review_report(project_name, file_query, start_line, end_line, format_type="md"):
    content = review_project_file_lines(project_name, file_query, start_line, end_line)

    return create_local_report(
        title=f"JARVIS Code Fragment Review - {file_query} lines {start_line}-{end_line}",
        content=content,
        format_type=format_type,
        open_after=True,
        project_name=project_name,
        file_query=file_query
    )


def project_file_command_help():
    return """
Project File Assistant commands supported:
- read_project_file_lines(project, file, start, end)
- review_project_file_lines(project, file, start, end)
- improve_project_file_lines(project, file, start, end)
- copy_project_file_lines_between_projects(source_project, source_file, source_start, source_end, target_project, target_file, target_start, target_end)
- generate_file_review_report(project, file, format_type)
- generate_file_lines_review_report(project, file, start, end, format_type)
"""


# ==========================
# LLM PROMPT HELPER
# ==========================
def build_project_file_prompt(project_name, file_query, role, task):
    item, content, error = resolve_project_file_with_content(
        project_name,
        file_query,
        max_chars=MAX_LLM_CHARS
    )

    if error:
        return None, error

    prompt = f"""
You are JARVIS, {role}.

Use ONLY the code below.
Do not invent files, frameworks, functions, APIs, routes, databases, or features.
If something is not visible in this file, say: "Not visible in this file."
Do not claim you changed code unless a separate safe apply command was used.

Project:
{project_name}

File:
{item['relative_path']}

Path:
{item['full_path']}

Code:
{content}

Task:
{task}
"""

    update_hud_context(
        project=project_name,
        file_path=item["relative_path"],
        action="AI prompt prepared",
        ai_status="THINKING",
        result="Sending code to local AI."
    )

    return prompt, None


# ==========================
# SAFE LLM CALL HELPER
# ==========================
def run_project_file_llm(project_name, file_query, prompt, action):
    update_hud_context(
        project=project_name,
        file_path=file_query,
        action=action,
        ai_status="THINKING",
        result="Local AI is processing..."
    )

    result = ask_llm(prompt)

    update_hud_context(
        project=project_name,
        file_path=file_query,
        action=f"{action} completed",
        ai_status="READY",
        result=result
    )

    return result


# ==========================
# ANALYZE PROJECT FILE
# ==========================
def analyze_project_file(project_name, file_query):
    prompt, error = build_project_file_prompt(
        project_name,
        file_query,
        "a senior software engineer and cybersecurity reviewer",
        """
Analyze this project file.

Return:
1. File purpose
2. Code quality score from 1 to 10
3. Security issues
4. Performance issues
5. Bugs or risky logic
6. Maintainability issues
7. Concrete improvements
8. Final recommendation

Be practical, specific, and concise.
"""
    )

    if error:
        return error

    return run_project_file_llm(project_name, file_query, prompt, "Analyze project file")


# ==========================
# IMPROVE PROJECT FILE
# ==========================
def improve_project_file(project_name, file_query):
    prompt, error = build_project_file_prompt(
        project_name,
        file_query,
        "a senior developer",
        """
Improve this code without changing its main behavior.

Return:
1. Main weaknesses
2. Better architecture
3. Safer/security improvements
4. Performance improvements
5. Cleaner code suggestions
6. A corrected/improved code version if possible
7. Short explanation of what changed

Do not invent external files that are not necessary.
Keep the answer useful and practical.
"""
    )

    if error:
        return error

    return run_project_file_llm(project_name, file_query, prompt, "Improve project file")


# ==========================
# SAFE PATCH SUGGESTION
# ==========================
def suggest_safe_project_file_patch(project_name, file_query):
    prompt, error = build_project_file_prompt(
        project_name,
        file_query,
        "a strict senior developer and code safety reviewer",
        """
Analyze this file and return ONLY a safe edit plan.

Return:
1. What should be changed
2. Why it should be changed
3. Exact functions/sections affected
4. Risks
5. Backup recommendation
6. Patch-style proposal

Important:
- Do not say the patch was applied.
- Do not modify files.
- Do not invent unrelated code.
"""
    )

    if error:
        return error

    return run_project_file_llm(project_name, file_query, prompt, "Suggest safe patch")


# ==========================
# OPTIMIZE PROJECT FILE
# ==========================
def optimize_project_file(project_name, file_query):
    prompt, error = build_project_file_prompt(
        project_name,
        file_query,
        "a performance-focused software engineer",
        """
Optimize this file.

Return:
1. Current performance problems
2. Memory usage concerns
3. Unnecessary complexity
4. Optimization recommendations
5. Improved code snippets
6. Risks of each optimization

Be concrete and avoid vague advice.
"""
    )

    if error:
        return error

    return run_project_file_llm(project_name, file_query, prompt, "Optimize project file")



# ==========================
# SECURITY PROJECT FILE
# ==========================
def security_project_file(project_name, file_query):
    prompt, error = build_project_file_prompt(
        project_name,
        file_query,
        "a cybersecurity-focused code reviewer",
        """
Perform a strict security review.

Return:
1. Security risk level: Low / Medium / High
2. Vulnerabilities found
3. Authentication/authorization issues
4. Input validation issues
5. Secret/key handling issues
6. Logging/privacy issues
7. Exact safe fixes
8. Final recommendation

Be concrete and do not invent code that is not visible.
"""
    )

    if error:
        return error

    return run_project_file_llm(
        project_name,
        file_query,
        prompt,
        "Security review project file"
    )


def explain_project_file(project_name, file_query):
    prompt, error = build_project_file_prompt(
        project_name,
        file_query,
        "a clear software engineering tutor",
        """
Explain this file clearly.

Return:
1. What this file does
2. Important functions/classes
3. How data flows through it
4. What can break
5. What to learn from it

Keep it understandable and practical.
"""
    )

    if error:
        return error

    return run_project_file_llm(
        project_name,
        file_query,
        prompt,
        "Explain project file"
    )


# ==========================
# SAFE FULL REPLACEMENT
# ==========================
def apply_safe_project_file_replacement(project_name, file_query, new_content):
    """
    Safely replaces an entire project file.
    A backup is always created first.
    Call this only after explicit user confirmation.
    """

    return safe_write_project_file(
        project_name,
        file_query,
        new_content
    )


# ==========================
# AUTO SAFE APPLY PATCH
# ==========================
def extract_code_block(text):
    """
    Extracts the first fenced code block from an LLM response.
    If no fenced block exists, returns the full text only if it looks like code.
    """

    text = str(text).strip()

    if not text:
        return None

    if "```" in text:
        parts = text.split("```")

        # Usually: text, language, code, text...
        if len(parts) >= 3:
            code_part = parts[1]

            # Remove optional language line
            lines = code_part.splitlines()

            if lines and len(lines[0].strip()) <= 20 and not lines[0].strip().startswith(("#", "import", "from", "def", "class", "{", "<")):
                code_part = "\n".join(lines[1:])

            return code_part.strip()

    # Fallback: allow direct code only if it looks like code.
    code_indicators = [
        "import ",
        "from ",
        "def ",
        "class ",
        "const ",
        "let ",
        "function ",
        "export ",
        "#include",
        "public class",
        "<?php",
        "<html",
    ]

    lowered = text.lower()

    if any(indicator in lowered for indicator in code_indicators):
        return text

    return None


def validate_generated_replacement(original_content, new_content):
    """
    Basic safety checks before replacing a file.
    """

    original_content = str(original_content or "")
    new_content = str(new_content or "")

    if not new_content.strip():
        return "Generated content is empty."

    if len(new_content.strip()) < 20:
        return "Generated content is too short."

    if original_content and len(new_content) < max(30, int(len(original_content) * 0.25)):
        return (
            "Generated content is much shorter than the original file. "
            "Refusing automatic replacement."
        )

    dangerous_markers = [
        "rm -rf",
        "del /f /s /q",
        "format c:",
        "shutdown /s",
        "os.remove(",
        "shutil.rmtree(",
    ]

    lowered = new_content.lower()

    for marker in dangerous_markers:
        if marker in lowered:
            return f"Generated content contains risky marker: {marker}"

    return None


def generate_safe_full_replacement(project_name, file_query, mode="improve"):
    """
    Asks the LLM for a complete improved file.
    Does NOT write anything by itself.
    """

    item, error = find_project_file(project_name, file_query)

    if error:
        return None, error

    original_content, read_error = read_file_content(
        item["full_path"],
        max_chars=MAX_LLM_CHARS
    )

    if read_error:
        return None, read_error

    prompt = f"""
You are JARVIS, a strict senior developer.

Task:
Improve the file safely and return the COMPLETE updated file content.

Patch mode:
{mode}

Mode instructions:
- improve: improve readability, safety and maintainability.
- fix: focus on bugs and risky logic.
- secure: focus on security hardening without breaking behavior.
- optimize: focus on performance and clarity without changing behavior.

Rules:
- Return ONLY one fenced code block.
- Do not add explanations outside the code block.
- Do not invent missing files.
- Keep the same language/framework.
- Preserve the main behavior.
- Improve readability, safety and maintainability.
- Do not delete important logic.
- Do not include destructive commands.

Project:
{project_name}

File:
{item['relative_path']}

Original code:
{original_content}
"""

    update_hud_context(
        project=project_name,
        file_path=item["relative_path"],
        action=f"Generating AI patch: {mode}",
        ai_status="THINKING",
        result="Local AI is generating full replacement."
    )

    response = ask_llm(prompt)

    update_hud_context(
        project=project_name,
        file_path=item["relative_path"],
        action="AI patch generated",
        ai_status="READY",
        result="Validating generated code."
    )

    generated_code = extract_code_block(response)

    if not generated_code:
        return None, (
            "The model did not return a valid code block. "
            "No changes were made."
        )

    validation_error = validate_generated_replacement(
        original_content,
        generated_code
    )

    if validation_error:
        return None, validation_error

    return generated_code, None


def apply_ai_safe_project_file_patch(project_name, file_query, mode="improve"):
    """
    Generates a full improved replacement with LLM and applies it safely.
    A backup is always created first.
    """

    update_hud_context(
        project=project_name,
        file_path=file_query,
        action=f"Starting safe AI patch: {mode}",
        ai_status="THINKING",
        result="Resolving file and preparing backup."
    )

    generated_code, error = generate_safe_full_replacement(
        project_name,
        file_query,
        mode=mode
    )

    if error:
        update_hud_context(
            project=project_name,
            file_path=file_query,
            action="AI patch failed",
            ai_status="ERROR",
            result=error
        )
        return error

    result = apply_safe_project_file_replacement(
        project_name,
        file_query,
        generated_code
    )

    update_hud_context(
        project=project_name,
        file_path=file_query,
        action=f"Safe AI patch completed: {mode}",
        ai_status="READY",
        result=result
    )

    return (
        f"AI safe patch attempted. Mode: {mode}\n"
        f"{result}"
    )


# ==========================
# AI PATCH CONVENIENCE MODES
# ==========================
def fix_project_file(project_name, file_query):
    return apply_ai_safe_project_file_patch(
        project_name,
        file_query,
        mode="fix"
    )


def secure_project_file(project_name, file_query):
    return apply_ai_safe_project_file_patch(
        project_name,
        file_query,
        mode="secure"
    )


def auto_optimize_project_file(project_name, file_query):
    return apply_ai_safe_project_file_patch(
        project_name,
        file_query,
        mode="optimize"
    )

# ==========================
# COMPATIBILITY ALIASES
# ==========================
# These names make the module easier to call from jarvis_agent.py and developer_assistant.py.
read_file_lines = read_project_file_lines
extract_file_lines = extract_project_file_lines
replace_file_lines = replace_project_file_lines
insert_file_lines = insert_project_file_lines
copy_lines_between_project_files = copy_project_file_lines_between_projects
review_file_lines = review_project_file_lines
improve_file_lines = improve_project_file_lines
explain_file_lines = explain_project_file_lines
open_file_in_app = open_project_file_in_app
generate_review_report = generate_file_review_report
generate_lines_review_report = generate_file_lines_review_report

# ==========================================================
# ADVANCED DEVELOPER FILE TOOLS
# Added for stronger J.A.R.V.I.S developer commands.
# Safe by default: read/search/report operations do not modify files.
# Write operations always create backups first.
# ==========================================================
def _safe_symbol_pattern(symbol_name, symbol_type="symbol"):
    symbol = re.escape(str(symbol_name).strip())

    if symbol_type == "function":
        return [
            rf"^\s*def\s+{symbol}\b",
            rf"^\s*async\s+def\s+{symbol}\b",
            rf"^\s*function\s+{symbol}\b",
            rf"^\s*const\s+{symbol}\s*=",
            rf"^\s*let\s+{symbol}\s*=",
            rf"^\s*var\s+{symbol}\s*=",
            rf"^\s*export\s+function\s+{symbol}\b",
            rf"^\s*public\s+.*\s+{symbol}\s*\(",
            rf"^\s*private\s+.*\s+{symbol}\s*\(",
            rf"^\s*protected\s+.*\s+{symbol}\s*\(",
        ]

    if symbol_type == "class":
        return [
            rf"^\s*class\s+{symbol}\b",
            rf"^\s*export\s+class\s+{symbol}\b",
            rf"^\s*public\s+class\s+{symbol}\b",
        ]

    if symbol_type == "import":
        return [
            rf"^\s*import\s+.*{symbol}",
            rf"^\s*from\s+.*{symbol}.*\s+import",
            rf"^\s*using\s+.*{symbol}",
            rf"^\s*#include\s+.*{symbol}",
            rf"^\s*require\(.+{symbol}.+\)",
        ]

    return [rf"\b{symbol}\b"]


def find_project_file_symbol(project_name, file_query, symbol_name, symbol_type="symbol", context_lines=3):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return read_error

    lines = content.splitlines()
    patterns = _safe_symbol_pattern(symbol_name, symbol_type)

    matches = []

    for index, line in enumerate(lines, start=1):
        if any(re.search(pattern, line, flags=re.IGNORECASE) for pattern in patterns):
            start = max(1, index - context_lines)
            end = min(len(lines), index + context_lines)

            snippet = "\n".join(
                f"{str(i).rjust(len(str(end)))} | {lines[i - 1]}"
                for i in range(start, end + 1)
            )

            matches.append(
                f"Match at line {index}:\n{snippet}"
            )

    update_hud_context(
        project=project_name,
        file_path=item["relative_path"],
        action=f"Find {symbol_type}: {symbol_name}",
        ai_status="READY",
        result=f"{len(matches)} matches"
    )

    if not matches:
        return f"{symbol_type.title()} not found: {symbol_name}"

    return (
        f"Found {len(matches)} match(es) for {symbol_type} '{symbol_name}'\n"
        f"FILE: {item['relative_path']}\n"
        f"PATH: {item['full_path']}\n\n"
        + "\n\n".join(matches[:30])
    )


def find_project_function(project_name, file_query, function_name):
    return find_project_file_symbol(
        project_name,
        file_query,
        function_name,
        symbol_type="function"
    )


def find_project_class(project_name, file_query, class_name):
    return find_project_file_symbol(
        project_name,
        file_query,
        class_name,
        symbol_type="class"
    )


def find_project_import(project_name, file_query, import_name):
    return find_project_file_symbol(
        project_name,
        file_query,
        import_name,
        symbol_type="import"
    )


def find_project_symbol(project_name, file_query, symbol_name):
    return find_project_file_symbol(
        project_name,
        file_query,
        symbol_name,
        symbol_type="symbol"
    )


def find_project_todos(project_name, file_query=None, limit=80):
    project, error = get_project(project_name)

    if error:
        return error

    files = []

    if file_query:
        item, file_error = find_project_file(project_name, file_query)

        if file_error:
            return file_error

        files = [item]
    else:
        files = collect_project_files(project["path"])

    markers = ["TODO", "FIXME", "HACK", "BUG", "XXX"]
    results = []

    for item in files:
        if not is_text_file(item["full_path"]):
            continue

        content, read_error = read_file_content(
            item["full_path"],
            max_chars=10_000_000
        )

        if read_error:
            continue

        for index, line in enumerate(content.splitlines(), start=1):
            upper = line.upper()

            if any(marker in upper for marker in markers):
                results.append(
                    f"{item['relative_path']}:{index} -> {line.strip()}"
                )

                if len(results) >= limit:
                    break

        if len(results) >= limit:
            break

    if not results:
        return "No TODO/FIXME/HACK/BUG markers found."

    return (
        f"TODO/FIXME markers found: {len(results)}\n\n"
        + "\n".join(results)
    )


def find_project_text(project_name, keyword, file_query=None, limit=80):
    project, error = get_project(project_name)

    if error:
        return error

    keyword_lower = str(keyword).lower()
    results = []

    if file_query:
        item, file_error = find_project_file(project_name, file_query)

        if file_error:
            return file_error

        files = [item]
    else:
        files = collect_project_files(project["path"])

    for item in files:
        if not is_text_file(item["full_path"]):
            continue

        content, read_error = read_file_content(
            item["full_path"],
            max_chars=10_000_000
        )

        if read_error:
            continue

        for index, line in enumerate(content.splitlines(), start=1):
            if keyword_lower in line.lower():
                results.append(
                    f"{item['relative_path']}:{index} -> {line.strip()}"
                )

                if len(results) >= limit:
                    break

        if len(results) >= limit:
            break

    if not results:
        return f"No matches found for: {keyword}"

    return (
        f"Matches for '{keyword}': {len(results)}\n\n"
        + "\n".join(results)
    )


def replace_project_file_text(project_name, file_query, old_text, new_text, max_replacements=1):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return read_error

    if old_text not in content:
        return "Old text was not found. No changes made."

    backup_result = create_project_file_backup(project_name, file_query)

    if not backup_result.startswith("Backup created successfully."):
        return backup_result

    updated = content.replace(
        old_text,
        new_text,
        int(max_replacements)
    )

    try:
        with open(item["full_path"], "w", encoding="utf-8", errors="ignore") as f:
            f.write(updated)

        update_hud_context(
            project=project_name,
            file_path=item["relative_path"],
            action="Text replacement completed",
            ai_status="READY",
            result=item["relative_path"]
        )

        return (
            "Text replaced safely.\n"
            f"File: {item['relative_path']}\n"
            f"Replacements: {max_replacements}\n"
            f"{backup_result}"
        )

    except Exception as e:
        return f"Backup created, but replacement failed: {e}"


def create_unified_diff_for_project_file(project_name, file_query, new_content, context=3):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    old_content, read_error = read_file_content(
        item["full_path"],
        max_chars=10_000_000
    )

    if read_error:
        return read_error

    old_lines = old_content.splitlines(keepends=True)
    new_lines = str(new_content).splitlines(keepends=True)

    diff = difflib.unified_diff(
        old_lines,
        new_lines,
        fromfile=f"original/{item['relative_path']}",
        tofile=f"new/{item['relative_path']}",
        n=context
    )

    diff_text = "".join(diff)

    if not diff_text.strip():
        return "No differences."

    return diff_text


def compare_project_files(project_a, file_a, project_b, file_b, context=3):
    item_a, error_a = find_project_file(project_a, file_a)

    if error_a:
        return error_a

    item_b, error_b = find_project_file(project_b, file_b)

    if error_b:
        return error_b

    content_a, read_error_a = read_file_content(
        item_a["full_path"],
        max_chars=10_000_000
    )

    if read_error_a:
        return read_error_a

    content_b, read_error_b = read_file_content(
        item_b["full_path"],
        max_chars=10_000_000
    )

    if read_error_b:
        return read_error_b

    diff = difflib.unified_diff(
        content_a.splitlines(keepends=True),
        content_b.splitlines(keepends=True),
        fromfile=f"{project_a}/{item_a['relative_path']}",
        tofile=f"{project_b}/{item_b['relative_path']}",
        n=context
    )

    diff_text = "".join(diff)

    if not diff_text.strip():
        return "Files are identical."

    return diff_text[:20000]


def open_project_file_in_vscode(project_name, file_query, line=None):
    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    commands = [
        "code",
        "code.cmd",
        os.path.expandvars(r"%LOCALAPPDATA%\Programs\Microsoft VS Code\Code.exe"),
        r"C:\Program Files\Microsoft VS Code\Code.exe",
    ]

    target = item["full_path"]

    if line:
        target = f"{target}:{int(line)}"

    for command in commands:
        try:
            if os.path.exists(command) or shutil.which(command):
                subprocess.Popen(
                    [command, "-g", target],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    shell=False
                )

                update_hud_context(
                    project=project_name,
                    file_path=item["relative_path"],
                    action="Open file in VS Code",
                    ai_status="READY",
                    result=item["relative_path"]
                )

                return f"Opening in VS Code: {item['relative_path']}"
        except Exception:
            continue

    return "VS Code command was not found. Add VS Code to PATH or install it."


def open_project_file_in_ide(project_name, file_query, ide_name="VS Code", line=None):
    ide_lower = str(ide_name).lower().strip()

    if ide_lower in {"vs code", "vscode", "visual studio code", "code"}:
        return open_project_file_in_vscode(project_name, file_query, line=line)

    item, error = find_project_file(project_name, file_query)

    if error:
        return error

    try:
        from tools import open_file_in_app
        return open_file_in_app(item["full_path"], ide_name)
    except Exception:
        pass

    try:
        subprocess.Popen(
            [ide_name, item["full_path"]],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            shell=True
        )

        return f"Opening file in {ide_name}: {item['relative_path']}"

    except Exception as e:
        return f"Could not open file in {ide_name}: {e}"


def _report_lines_to_csv(path, project_name, file_query, content):
    with open(path, "w", encoding="utf-8", errors="ignore", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["project", "file", "line_number", "text"])

        for index, line in enumerate(str(content).splitlines(), start=1):
            writer.writerow([project_name, file_query, index, line])


def generate_file_report_excel(project_name, file_query, content=None, open_after=True):
    if content is None:
        content = analyze_project_file(project_name, file_query)

    os.makedirs("developer_reports", exist_ok=True)

    safe_project = normalize_name(project_name) or "project"
    safe_file = normalize_name(file_query) or "file"
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = os.path.join(
        "developer_reports",
        f"{safe_project}_{safe_file}_file_report_{timestamp}.xlsx"
    )

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill

        wb = Workbook()

        ws = wb.active
        ws.title = "Summary"
        ws["A1"] = "J.A.R.V.I.S File Report"
        ws["A1"].font = Font(bold=True, size=16)
        ws["A3"] = "Project"
        ws["B3"] = project_name
        ws["A4"] = "File"
        ws["B4"] = file_query
        ws["A5"] = "Generated"
        ws["B5"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        ws["A7"] = "Report preview"
        ws["A7"].font = Font(bold=True)

        for row_index, line in enumerate(str(content).splitlines()[:400], start=8):
            ws.cell(row=row_index, column=1).value = line

        ws.column_dimensions["A"].width = 120
        ws.column_dimensions["B"].width = 35

        ws2 = wb.create_sheet("Report Lines")
        ws2.append(["Line", "Text"])

        for cell in ws2[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="D9EAF7")

        for index, line in enumerate(str(content).splitlines(), start=1):
            ws2.append([index, line])

        ws2.column_dimensions["A"].width = 12
        ws2.column_dimensions["B"].width = 140

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    cell.alignment = Alignment(vertical="top", wrap_text=True)

        wb.save(path)

    except Exception:
        path = path.replace(".xlsx", ".csv")
        _report_lines_to_csv(path, project_name, file_query, content)

    if open_after:
        try:
            os.startfile(os.path.abspath(path))
        except Exception:
            pass

    return f"Report created:\n{path}"


def generate_file_report(project_name, file_query, format_type="md"):
    fmt = str(format_type).lower().strip()
    content = analyze_project_file(project_name, file_query)

    if fmt in {"excel", "xls", "xlsx", "spreadsheet"}:
        return generate_file_report_excel(
            project_name,
            file_query,
            content=content,
            open_after=True
        )

    try:
        from project_review_assistant import export_report_content
        return export_report_content(
            project_name,
            content,
            format_type=fmt,
            report_kind=f"file_review_{normalize_name(file_query)}",
            open_after=True
        )
    except Exception:
        return generate_file_review_report(
            project_name,
            file_query,
            format_type=fmt
        )


def generate_file_lines_report(project_name, file_query, start_line, end_line, format_type="md"):
    fmt = str(format_type).lower().strip()
    content = review_project_file_lines(project_name, file_query, start_line, end_line)

    if fmt in {"excel", "xls", "xlsx", "spreadsheet"}:
        return generate_file_report_excel(
            project_name,
            f"{file_query}_lines_{start_line}_{end_line}",
            content=content,
            open_after=True
        )

    try:
        from project_review_assistant import export_report_content
        return export_report_content(
            project_name,
            content,
            format_type=fmt,
            report_kind=f"file_lines_{normalize_name(file_query)}_{start_line}_{end_line}",
            open_after=True
        )
    except Exception:
        return generate_file_lines_review_report(
            project_name,
            file_query,
            start_line,
            end_line,
            format_type=fmt
        )


def project_file_advanced_help():
    return """
Advanced Project File Assistant commands/functions:
- read_project_file_lines(project, file, start, end)
- find_project_function(project, file, function_name)
- find_project_class(project, file, class_name)
- find_project_import(project, file, import_name)
- find_project_symbol(project, file, symbol_name)
- find_project_todos(project, optional_file)
- find_project_text(project, keyword, optional_file)
- replace_project_file_lines(project, file, start, end, new_code)
- insert_project_file_lines(project, file, line, new_code)
- replace_project_file_text(project, file, old_text, new_text)
- compare_project_files(project_a, file_a, project_b, file_b)
- open_project_file_in_ide(project, file, ide_name, optional_line)
- generate_file_report(project, file, format_type)
- generate_file_lines_report(project, file, start, end, format_type)
"""


# Extra compatibility aliases
find_function_in_file = find_project_function
find_class_in_file = find_project_class
find_import_in_file = find_project_import
find_symbol_in_file = find_project_symbol
find_todos = find_project_todos
find_text = find_project_text
replace_text = replace_project_file_text
compare_files = compare_project_files
open_file_in_ide = open_project_file_in_ide
generate_file_report_any = generate_file_report
generate_file_lines_report_any = generate_file_lines_report
