import os
import re
import ast
import json
import subprocess
from collections import Counter, defaultdict
from datetime import datetime

from llm_local import ask_llm
from deep_project_memory import find_deep_project


MAX_FILES_FOR_CONTEXT = 70
MAX_CHARS_PER_FILE = 4500


# ==========================
# BASIC HELPERS
# ==========================
def get_project(project_name):
    project = find_deep_project(project_name)

    if not project:
        return None, f"Project not found in deep memory: {project_name}"

    return project, None


def get_code_files(project):
    files = []

    for item in project.get("files", []):
        path = item.get("relative_path", "")
        content = item.get("content", "")
        ext = item.get("extension", "")

        if not content.strip():
            continue

        files.append({
            "path": path,
            "extension": ext,
            "content": content
        })

    return files


# ==========================
# SMART CONTEXT SELECTION
# ==========================
IMPORTANT_FILE_NAMES = {
    "package.json",
    "requirements.txt",
    "pyproject.toml",
    "docker-compose.yml",
    "docker-compose.yaml",
    "dockerfile",
    "vite.config.ts",
    "vite.config.js",
    "tailwind.config.js",
    "postcss.config.js",
    "tsconfig.json",
    "angular.json",
    "next.config.js",
    "next.config.ts",
    "readme.md",
    "readme.en.md",
    ".env",
    ".env.example",
    "config.py",
    "settings.py",
    "app.py",
    "api.py",
    "main.py",
    "server.py",
    "__init__.py",
}

IMPORTANT_PATH_TOKENS = [
    "routes/",
    "api/",
    "schemas/",
    "models/",
    "utils/",
    "middleware",
    "auth",
    "token",
    "jwt",
    "security",
    "permission",
    "admin",
    "dashboard",
    "database",
    "db",
    "crypto",
    "encrypt",
    "backup",
    "scanner",
    "logger",
    "audit",
    "frontend",
    "components",
    "src/",
    "main.",
    "app.",
    "scanform",
    "filescanlog",
    "test",
    "tests",
    "docs/",
    ".github/",
]


def file_priority_score(item):
    path = normalize_path(item.get("path", ""))
    name = os.path.basename(path)
    ext = item.get("extension", "").lower()
    content = item.get("content", "").lower()[:2500]

    score = 0

    if name in IMPORTANT_FILE_NAMES:
        score += 80

    for token in IMPORTANT_PATH_TOKENS:
        if token in path:
            score += 18

    if ext in {".py", ".ts", ".tsx", ".js", ".jsx"}:
        score += 20

    if ext in {".json", ".yml", ".yaml", ".toml", ".md"}:
        score += 12

    security_keywords = [
        "fastapi",
        "flask",
        "apirouter",
        "blueprint",
        "jwt",
        "bcrypt",
        "passlib",
        "cryptcontext",
        "oauth2",
        "token",
        "secret_key",
        "auth",
        "permission",
        "admin",
        "sqlite",
        "sqlalchemy",
        "firebase",
        "axios",
        "upload",
        "multipart",
        "scan",
        "audit",
        "logger",
        "docker",
        "vite",
        "react",
        "tailwind",
    ]

    for keyword in security_keywords:
        if keyword in content or keyword in path:
            score += 4

    # Avoid making lock files dominate context.
    if "package-lock.json" in path or "yarn.lock" in path or "pnpm-lock" in path:
        score -= 60

    # Keep tests, but lower priority than source files.
    if is_test_path(path):
        score -= 12

    return score


def select_context_files(files, limit=MAX_FILES_FOR_CONTEXT):
    scored = []

    for item in files:
        scored.append(
            (
                file_priority_score(item),
                item.get("path", "").lower(),
                item
            )
        )

    scored.sort(
        key=lambda row: (
            -row[0],
            row[1]
        )
    )

    selected = [
        item for score, path, item in scored[:limit]
    ]

    return selected


def extract_json_dependencies(content):
    try:
        data = json.loads(content)
    except Exception:
        return {}

    deps = {}

    for key in [
        "dependencies",
        "devDependencies",
        "peerDependencies"
    ]:
        if isinstance(data.get(key), dict):
            deps[key] = sorted(data[key].keys())

    scripts = data.get("scripts", {})

    if isinstance(scripts, dict):
        deps["scripts"] = scripts

    return deps


def build_project_manifest(project, files=None):
    if files is None:
        files = get_code_files(project)

    selected = select_context_files(files)

    key_files = [
        item["path"]
        for item in selected[:40]
    ]

    dependencies = {}

    for item in files:
        path = normalize_path(item["path"])

        if path.endswith("package.json"):
            dependencies[item["path"]] = extract_json_dependencies(
                item["content"]
            )

        elif path.endswith("requirements.txt"):
            dependencies[item["path"]] = [
                line.strip()
                for line in item["content"].splitlines()
                if line.strip() and not line.strip().startswith("#")
            ][:80]

    route_files = [
        item["path"]
        for item in files
        if "routes/" in normalize_path(item["path"])
        or "dashboard" in normalize_path(item["path"])
        or "admin" in normalize_path(item["path"])
        or "auth" in normalize_path(item["path"])
    ][:60]

    frontend_files = [
        item["path"]
        for item in files
        if item["extension"] in {".tsx", ".jsx", ".ts", ".js", ".css"}
        and (
            "src/" in normalize_path(item["path"])
            or "components/" in normalize_path(item["path"])
            or "frontend" in normalize_path(item["path"])
            or "file-scan-log-app" in normalize_path(item["path"])
        )
    ][:60]

    return (
        "REAL PROJECT MANIFEST\n"
        f"Project name: {project.get('name')}\n"
        f"Project path: {project.get('path')}\n"
        f"Files indexed: {project.get('files_count')}\n"
        f"Detected tech stack: {', '.join(project.get('tech_stack', []))}\n\n"
        f"Key files selected for analysis:\n{key_files}\n\n"
        f"Route/API/Auth/Admin related files:\n{route_files}\n\n"
        f"Frontend/UI related files:\n{frontend_files}\n\n"
        f"Dependencies/scripts extracted from package/requirements files:\n{dependencies}\n"
    )


def build_project_context(project, files=None):
    if files is None:
        files = get_code_files(project)

    selected_files = select_context_files(files)

    context = build_project_manifest(project, files)

    context += "\n\nREAL PROJECT FILE CONTENTS USED AS EVIDENCE:\n"

    for item in selected_files:
        context += (
            "\n\n==============================\n"
            f"FILE: {item['path']}\n"
            f"PRIORITY_SCORE: {file_priority_score(item)}\n"
            "==============================\n"
        )

        context += item["content"][:MAX_CHARS_PER_FILE]

    return context


def project_overview(project):
    files = get_code_files(project)

    extensions = Counter(
        item["extension"]
        for item in files
    )

    folders = Counter()

    for item in files:
        folder = os.path.dirname(item["path"])

        if not folder:
            folder = "root"

        folders[folder] += 1

    selected = select_context_files(files, limit=25)

    key_files = [
        item["path"]
        for item in selected
    ]

    return (
        f"Project name: {project.get('name')}\n"
        f"Project path: {project.get('path')}\n"
        f"Files indexed: {project.get('files_count')}\n"
        f"Tech stack: {', '.join(project.get('tech_stack', []))}\n\n"
        f"Extensions:\n{dict(extensions)}\n\n"
        f"Main folders:\n{dict(folders.most_common(15))}\n\n"
        f"Highest-priority files for this review:\n{key_files}"
    )


def strict_prompt_prefix():
    return """
IMPORTANT GROUNDING RULES:

- Use ONLY the real indexed project files and manifests shown below.
- Do NOT invent files, services, classes, functions, databases, APIs, cloud providers, ML modules, teams, timelines, or features.
- If something is not visible in the indexed files, write exactly: "Not visible in indexed files."
- Mention exact file paths for every concrete claim.
- Prefer evidence from package.json, requirements.txt, routes, auth, config, frontend/src, Docker, tests, and README.
- Do NOT describe this as a generic cybersecurity platform. Describe only what is visible in the code.
- When giving risks or improvements, tie them to exact files and snippets when possible.
- Be practical, technical, and realistic.
"""


# ==========================
# PROJECT REVIEW
# ==========================
def review_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    context = build_project_context(project, files)

    prompt = f"""
You are JARVIS, a senior software architect.

{strict_prompt_prefix()}

Review this entire project based ONLY on the real file manifest and file contents below. Cite exact file paths in every section.

PROJECT OVERVIEW:
{project_overview(project)}

PROJECT FILES:
{context}

Return:

1. Project purpose
2. Architecture overview
3. Strong parts
4. Weak parts
5. Code quality score from 1 to 10
6. Maintainability score from 1 to 10
7. Security score from 1 to 10
8. Main risks
9. Priority improvements
10. Final recommendation
"""

    return ask_llm(prompt)


# ==========================
# SECURITY REVIEW
# ==========================
def find_security_issues(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    context = build_project_context(project, files)

    prompt = f"""
You are JARVIS, a cybersecurity reviewer.

{strict_prompt_prefix()}

Analyze this project for security issues based ONLY on the real file manifest and file contents below. Cite exact file paths for every issue.

PROJECT OVERVIEW:
{project_overview(project)}

PROJECT FILES:
{context}

Return:

1. Overall security risk: Low / Medium / High
2. Authentication issues
3. Authorization issues
4. JWT/token issues
5. Password and secret handling issues
6. Input validation issues
7. File upload risks
8. Logging/privacy risks
9. Dependency/configuration risks
10. Concrete fixes by priority
11. Files that should be fixed first

Mention exact file names when visible.
"""

    return ask_llm(prompt)


# ==========================
# DEAD CODE HELPERS
# ==========================
def normalize_path(path):
    return path.replace("\\", "/").lower()


def file_stem(path):
    return os.path.splitext(
        os.path.basename(path)
    )[0].lower()


def remove_comments_and_strings_python(code):
    try:
        tree = ast.parse(code)
        names = set()

        for node in ast.walk(tree):
            if isinstance(node, ast.Name):
                names.add(node.id.lower())
            elif isinstance(node, ast.Attribute):
                names.add(node.attr.lower())

        return names

    except Exception:
        # Fallback simple token extraction
        return set(
            re.findall(
                r"[A-Za-z_][A-Za-z0-9_]*",
                code.lower()
            )
        )


def extract_python_symbols(path, content):
    result = {
        "functions": [],
        "classes": [],
        "imports": []
    }

    try:
        tree = ast.parse(content)

        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef):
                result["functions"].append(node.name)
            elif isinstance(node, ast.AsyncFunctionDef):
                result["functions"].append(node.name)
            elif isinstance(node, ast.ClassDef):
                result["classes"].append(node.name)
            elif isinstance(node, ast.Import):
                for alias in node.names:
                    result["imports"].append(alias.name)
            elif isinstance(node, ast.ImportFrom):
                module = node.module or ""
                result["imports"].append(module)

    except Exception:
        pass

    return result


def extract_js_symbols(content):
    functions = []
    classes = []

    function_patterns = [
        r"function\s+([A-Za-z_][A-Za-z0-9_]*)\s*\(",
        r"const\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*\(",
        r"const\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*async\s*\(",
        r"export\s+default\s+function\s+([A-Za-z_][A-Za-z0-9_]*)",
        r"export\s+const\s+([A-Za-z_][A-Za-z0-9_]*)",
    ]

    for pattern in function_patterns:
        functions.extend(
            re.findall(pattern, content)
        )

    classes.extend(
        re.findall(
            r"class\s+([A-Za-z_][A-Za-z0-9_]*)",
            content
        )
    )

    return {
        "functions": sorted(set(functions)),
        "classes": sorted(set(classes))
    }


def build_reference_text(files):
    text_parts = []

    for item in files:
        text_parts.append(
            normalize_path(item["path"])
        )
        text_parts.append(
            item["content"].lower()
        )

    return "\n".join(text_parts)


# ==========================
# DEAD CODE CLASSIFICATION
# ==========================
def is_test_path(path):
    lower = normalize_path(path)

    return (
        lower.startswith("test/")
        or lower.startswith("tests/")
        or "/test/" in lower
        or "/tests/" in lower
        or lower.endswith("_test.py")
        or lower.startswith("__tests__/")
        or ".spec." in lower
        or ".test." in lower
    )


def is_documentation_or_config_path(path):
    lower = normalize_path(path)
    name = os.path.basename(lower)

    doc_config_tokens = [
        "readme",
        "license",
        "changelog",
        "deployment",
        "security_overview",
        "docs/",
        ".github/",
        "package-lock.json",
        "yarn.lock",
        "pnpm-lock.yaml",
        "requirements.txt",
        "dockerfile",
        "docker-compose",
        "cloudflare",
        "waf",
        "config",
        "settings",
        ".env",
        ".yml",
        ".yaml",
        ".md",
        ".json",
        ".toml",
        ".ini",
        ".cfg",
        ".sh",
        ".bat",
        ".ps1",
    ]

    if any(token in lower for token in doc_config_tokens):
        return True

    if name in {
        "package.json",
        "angular.json",
        "vite.config.ts",
        "vite.config.js",
        "tailwind.config.js",
        "postcss.config.js",
        "tsconfig.json",
        "pyproject.toml",
    }:
        return True

    return False


def is_entrypoint_or_framework_file(path):
    lower = normalize_path(path)
    name = os.path.basename(lower)

    entry_tokens = [
        "__init__.py",
        "main.py",
        "app.py",
        "api.py",
        "server.py",
        "index.js",
        "index.ts",
        "main.tsx",
        "main.ts",
        "app.tsx",
        "app.jsx",
        "app.component.ts",
        "manage.py",
        "routes/",
        "schemas/",
        "models/",
        "middleware",
    ]

    return any(token in lower or name == token for token in entry_tokens)


def classify_dead_code_path(path):
    if is_test_path(path):
        return "test"

    if is_documentation_or_config_path(path):
        return "docs_config"

    return "production"


def risk_label_for_file(path):
    if is_test_path(path):
        return "LOW RISK TO DELETE / archive only if tests are obsolete"

    if is_documentation_or_config_path(path):
        return "MEDIUM RISK - documentation/config/deployment may be used outside code"

    if is_entrypoint_or_framework_file(path):
        return "HIGH RISK - may be framework entrypoint, route, model, schema, or dynamic import"

    return "MEDIUM RISK - verify references and runtime usage first"


def risk_label_for_symbol(path):
    if is_test_path(path):
        return "LOW RISK - test function, used by test runner"

    if is_entrypoint_or_framework_file(path):
        return "HIGH RISK - may be called by framework/router/dynamic import"

    return "MEDIUM RISK - verify references before removing"


def append_limited_section(output, title, items, limit=40):
    output.append(f"\n{title}")

    if not items:
        output.append("None found.")
        return

    for item in items[:limit]:
        output.append(item)

    if len(items) > limit:
        output.append(f"... and {len(items) - limit} more")



# ==========================
# IMPROVED DEAD CODE
# ==========================
def find_dead_code(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    all_text = build_reference_text(files)

    unused_files = {
        "production": [],
        "test": [],
        "docs_config": []
    }

    unused_functions = {
        "production": [],
        "test": [],
        "docs_config": []
    }

    unused_classes = {
        "production": [],
        "test": [],
        "docs_config": []
    }

    unused_imports = {
        "production": [],
        "test": [],
        "docs_config": []
    }

    for item in files:
        path = item["path"]
        ext = item["extension"]
        content = item["content"]

        stem = file_stem(path)
        lower_path = normalize_path(path)

        category = classify_dead_code_path(path)

        # Count references to file/module name.
        reference_count = all_text.count(stem)

        # Keep important project/framework files out of unused-file list.
        should_skip_file_candidate = (
            is_entrypoint_or_framework_file(path)
            or stem in {
                "readme",
                "package",
                "requirements",
                "config",
                "settings",
                "main",
                "index",
                "app",
            }
        )

        if reference_count <= 1 and not should_skip_file_candidate:
            unused_files[category].append(
                f"{path} [{risk_label_for_file(path)}]"
            )

        if ext == ".py":
            symbols = extract_python_symbols(path, content)
            current_tokens = remove_comments_and_strings_python(content)

            for func in symbols["functions"]:
                if func.startswith("_"):
                    continue

                # Ignore pytest-style functions in tests from production warnings.
                count = all_text.count(func.lower())

                if count <= 1:
                    unused_functions[category].append(
                        f"{path} -> function {func} [{risk_label_for_symbol(path)}]"
                    )

            for cls in symbols["classes"]:
                count = all_text.count(cls.lower())

                if count <= 1:
                    unused_classes[category].append(
                        f"{path} -> class {cls} [{risk_label_for_symbol(path)}]"
                    )

            for imp in symbols["imports"]:
                if not imp:
                    continue

                short = imp.split(".")[-1].lower()

                # Avoid marking framework/test imports as definitely unused.
                if short and short not in current_tokens:
                    unused_imports[category].append(
                        f"{path} -> import {imp} [{risk_label_for_symbol(path)}]"
                    )

        elif ext in {".js", ".jsx", ".ts", ".tsx"}:
            symbols = extract_js_symbols(content)

            for func in symbols["functions"]:
                count = all_text.count(func.lower())

                if count <= 1:
                    unused_functions[category].append(
                        f"{path} -> function/component {func} [{risk_label_for_symbol(path)}]"
                    )

            for cls in symbols["classes"]:
                count = all_text.count(cls.lower())

                if count <= 1:
                    unused_classes[category].append(
                        f"{path} -> class {cls} [{risk_label_for_symbol(path)}]"
                    )

    output = [
        "Dead code / unused code scan",
        "Important: These are heuristic results, not delete commands.",
        "Production code is separated from tests and documentation/config files.",
    ]

    append_limited_section(
        output,
        "Possible unused production files:",
        unused_files["production"],
        40
    )

    append_limited_section(
        output,
        "Possible unused test files:",
        unused_files["test"],
        30
    )

    append_limited_section(
        output,
        "Possible unused documentation/config/deployment files:",
        unused_files["docs_config"],
        40
    )

    append_limited_section(
        output,
        "Possible unused production functions/components:",
        unused_functions["production"],
        60
    )

    append_limited_section(
        output,
        "Possible unused test functions:",
        unused_functions["test"],
        40
    )

    append_limited_section(
        output,
        "Possible unused production classes:",
        unused_classes["production"],
        40
    )

    append_limited_section(
        output,
        "Possible unused test classes:",
        unused_classes["test"],
        20
    )

    append_limited_section(
        output,
        "Possible unused production imports:",
        unused_imports["production"],
        60
    )

    append_limited_section(
        output,
        "Possible unused test imports:",
        unused_imports["test"],
        30
    )

    append_limited_section(
        output,
        "Possible unused docs/config imports:",
        unused_imports["docs_config"],
        30
    )

    output.append(
        "\nRecommendation:\n"
        "Do not delete automatically. First search the full project, run tests, "
        "check dynamic imports/routes/config usage, and verify deployment scripts manually."
    )

    return "\n".join(output)


# ==========================
# DUPLICATE LOGIC HELPERS
# ==========================
def is_import_or_boilerplate_line(line):
    stripped = line.strip()

    if not stripped:
        return True

    ignore_prefixes = (
        "import ",
        "from ",
        "export ",
        "using ",
        "#include",
        "package ",
        "namespace ",
        "//",
        "#",
        "/*",
        "*",
        "*/",
    )

    if stripped.startswith(ignore_prefixes):
        return True

    if stripped in {
        "{",
        "}",
        ");",
        ");",
        "};",
        ")",
        "(",
        "return;",
    }:
        return True

    return False


def normalize_code_line(line):
    line = line.strip()

    # Remove inline comments
    line = re.sub(r"#.*$", "", line)
    line = re.sub(r"//.*$", "", line)

    # Normalize string literals
    line = re.sub(r'"[^"]*"', '"STR"', line)
    line = re.sub(r"'[^']*'", "'STR'", line)

    # Normalize numbers
    line = re.sub(r"\b\d+\b", "NUM", line)

    # Normalize whitespace
    line = re.sub(r"\s+", " ", line)

    return line.strip()


def extract_logic_blocks(content, min_lines=8):
    raw_lines = content.splitlines()

    cleaned = []

    for line in raw_lines:
        normalized = normalize_code_line(line)

        if not normalized:
            continue

        if is_import_or_boilerplate_line(normalized):
            continue

        cleaned.append(normalized)

    blocks = []

    for i in range(0, len(cleaned) - min_lines + 1):
        block_lines = cleaned[i:i + min_lines]

        block = "\n".join(block_lines)

        # Ignore tiny or mostly punctuation blocks
        if len(block) < 180:
            continue

        # Require at least some logic keywords/operators
        logic_markers = [
            "if ",
            "for ",
            "while ",
            "try",
            "except",
            "catch",
            "return ",
            "await ",
            "=>",
            "=",
            "raise ",
            "throw ",
            ".append",
            ".push",
        ]

        if not any(marker in block for marker in logic_markers):
            continue

        blocks.append(block)

    return blocks


# ==========================
# IMPROVED DUPLICATES
# ==========================
def find_duplicate_code(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    block_map = defaultdict(set)

    for item in files:
        path = item["path"]
        content = item["content"]

        blocks = extract_logic_blocks(
            content,
            min_lines=8
        )

        for block in blocks:
            block_map[block].add(path)

    duplicates = []

    for block, paths in block_map.items():
        unique_paths = sorted(paths)

        if len(unique_paths) > 1:
            duplicates.append(
                (block, unique_paths)
            )

    if not duplicates:
        return (
            "No obvious duplicated business logic blocks found.\n"
            "Import-only duplicates and boilerplate were ignored."
        )

    duplicates.sort(
        key=lambda item: (
            len(item[1]),
            len(item[0])
        ),
        reverse=True
    )

    output = [
        "Possible duplicated business logic blocks found:",
        "Import-only duplicates and boilerplate were ignored.\n"
    ]

    for index, (block, paths) in enumerate(
        duplicates[:10],
        start=1
    ):
        output.append(f"\nDuplicate #{index}")
        output.append("Files:")
        output.extend(paths)
        output.append("Code sample:")
        output.append(block[:900])

    return "\n".join(output)


# ==========================
# ARCHITECTURE REPORT
# ==========================
def generate_architecture_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    context = build_project_context(project, files)

    prompt = f"""
You are JARVIS, a senior software architect.

{strict_prompt_prefix()}

Generate an architecture report for this project based ONLY on the real file manifest and file contents below. Cite exact file paths in every section.

PROJECT OVERVIEW:
{project_overview(project)}

PROJECT FILES:
{context}

Return:

1. High-level architecture
2. Backend structure
3. Frontend structure
4. Database layer
5. Security layer
6. API/routes structure
7. Data flow
8. Important dependencies
9. Architecture risks
10. Recommended architecture improvements
"""

    return ask_llm(prompt)


# ==========================
# IMPROVEMENT ROADMAP
# ==========================
def generate_improvement_roadmap(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    context = build_project_context(project, files)

    prompt = f"""
You are JARVIS, a senior engineering manager.

{strict_prompt_prefix()}

Create an improvement roadmap for this project based ONLY on the real file manifest and file contents below. Every task must mention affected files.

PROJECT OVERVIEW:
{project_overview(project)}

PROJECT FILES:
{context}

Return:

1. Immediate fixes
2. Security improvements
3. Performance improvements
4. Refactoring plan
5. Testing improvements
6. Documentation improvements
7. Deployment/DevOps improvements
8. 7-day action plan
9. 30-day action plan
10. Final priority list
"""

    return ask_llm(prompt)


# ==========================
# PROJECT OPTIMIZATION
# ==========================
def optimize_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    context = build_project_context(project, files)

    prompt = f"""
You are JARVIS, a performance and maintainability expert.

{strict_prompt_prefix()}

Analyze this project for optimization opportunities based ONLY on the real file manifest and file contents below. Every optimization must mention affected files.

PROJECT OVERVIEW:
{project_overview(project)}

PROJECT FILES:
{context}

Return:

1. Performance bottlenecks
2. Unnecessary complexity
3. Repeated logic
4. Heavy dependencies
5. Frontend optimization opportunities
6. Backend optimization opportunities
7. Database optimization opportunities
8. Recommended optimizations by priority
9. Risks of optimization
"""

    return ask_llm(prompt)


# ==========================
# STRUCTURE ANALYSIS
# ==========================
def analyze_project_structure(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    folders = defaultdict(list)

    for item in files:
        folder = os.path.dirname(item["path"])

        if not folder:
            folder = "root"

        folders[folder].append(item["path"])

    output = [
        f"Project: {project.get('name')}",
        f"Path: {project.get('path')}",
        f"Files indexed: {project.get('files_count')}",
        f"Tech stack: {', '.join(project.get('tech_stack', []))}",
        "\nFolders:"
    ]

    for folder, paths in sorted(
        folders.items(),
        key=lambda item: item[0]
    ):
        output.append(
            f"\n{folder} ({len(paths)} files)"
        )

        for path in paths[:20]:
            output.append(f" - {path}")

    return "\n".join(output)


# Friendly aliases
def find_unused_files(project_name):
    return find_dead_code(project_name)


def find_unused_functions(project_name):
    return find_dead_code(project_name)


def find_duplicate_logic(project_name):
    return find_duplicate_code(project_name)


# ==========================
# PROJECT DOCUMENTATION
# ==========================
def generate_project_documentation(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    context = build_project_context(project, files)

    prompt = f"""
You are JARVIS, a senior technical writer and software architect.

{strict_prompt_prefix()}

Generate professional project documentation based ONLY on the real file manifest and file contents below. Mention exact files for features, frontend, backend, API, security, and deployment.

PROJECT OVERVIEW:
{project_overview(project)}

PROJECT FILES:
{context}

Return:

1. Executive summary
2. Project purpose
3. Main features
4. Tech stack
5. Folder structure overview
6. Backend overview
7. Frontend overview
8. Security overview
9. API overview
10. Deployment overview
11. Known risks
12. Future improvements
"""

    return ask_llm(prompt)


# Friendly alias
def generate_documentation(project_name):
    return generate_project_documentation(project_name)


# ==========================
# GROUNDED PROJECT EVIDENCE
# ==========================
def project_evidence_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    selected = select_context_files(files, limit=50)

    output = [
        "Grounded project evidence report",
        f"Project: {project.get('name')}",
        f"Path: {project.get('path')}",
        f"Files indexed: {project.get('files_count')}",
        f"Tech stack: {', '.join(project.get('tech_stack', []))}",
        "",
        "Most important files used for project reviews:"
    ]

    for item in selected:
        output.append(
            f" - {item['path']} [priority {file_priority_score(item)}]"
        )

    output.append("")
    output.append("Manifest summary:")
    output.append(build_project_manifest(project, files))

    return "\n".join(output)


def evidence_project(project_name):
    return project_evidence_report(project_name)


# ==========================
# PROJECT EVIDENCE + GROUNDED REVIEW
# ==========================
def build_evidence_context(project, files=None, limit=35):
    if files is None:
        files = get_code_files(project)

    selected = select_context_files(files, limit=limit)

    context = build_project_manifest(project, files)

    context += "\n\nSTRICT EVIDENCE FILE CONTENTS:\n"

    for item in selected:
        content = item["content"][:MAX_CHARS_PER_FILE]

        context += (
            "\n\n------------------------------\n"
            f"EVIDENCE FILE: {item['path']}\n"
            f"PRIORITY_SCORE: {file_priority_score(item)}\n"
            "------------------------------\n"
            f"{content}\n"
        )

    return context


def grounded_prompt_rules():
    return """
STRICT GROUNDED MODE:

- You MUST use only the EVIDENCE FILES shown below.
- You MUST mention exact file paths in every major point.
- You MUST NOT invent products, teams, timelines, cloud services, ML systems, dashboards, APIs, databases, or features.
- If something is not present in the evidence, write exactly: "Not visible in indexed files."
- Do not use generic cybersecurity marketing language.
- Prefer short, technical, verifiable statements.
- Every recommendation must name the file(s) it affects.
- If the evidence is insufficient, say what exact files are missing.
"""


def grounded_review_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    evidence = build_evidence_context(project, files)

    prompt = f"""
You are JARVIS, a strict senior codebase reviewer.

{grounded_prompt_rules()}

Perform a grounded codebase review.

PROJECT OVERVIEW:
{project_overview(project)}

EVIDENCE:
{evidence}

Return:

1. What this project actually contains
2. Real architecture based on files
3. Backend evidence with exact files
4. Frontend evidence with exact files
5. Security evidence with exact files
6. Database/config evidence with exact files
7. Strong parts with exact files
8. Weak parts with exact files
9. Concrete risks with exact files
10. Priority fixes with exact files
11. Final grounded recommendation

Do not include claims that are not directly supported by file evidence.
"""

    return ask_llm(prompt)


def grounded_security_review_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    evidence = build_evidence_context(project, files)

    prompt = f"""
You are JARVIS, a strict cybersecurity auditor.

{grounded_prompt_rules()}

Perform a grounded security audit.

PROJECT OVERVIEW:
{project_overview(project)}

EVIDENCE:
{evidence}

Return:

1. Overall security risk
2. Authentication findings with exact files
3. Authorization findings with exact files
4. Token/JWT findings with exact files
5. Password hashing / secret handling findings with exact files
6. File upload findings with exact files
7. Logging/privacy findings with exact files
8. Dependency/configuration findings with exact files
9. False assumptions to avoid
10. Priority security fixes with exact files

Do not invent vulnerabilities that are not visible in the evidence.
"""

    return ask_llm(prompt)


def grounded_architecture_review_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    evidence = build_evidence_context(project, files)

    prompt = f"""
You are JARVIS, a strict software architect.

{grounded_prompt_rules()}

Create a grounded architecture review.

PROJECT OVERVIEW:
{project_overview(project)}

EVIDENCE:
{evidence}

Return:

1. Actual architecture visible in indexed files
2. Backend modules and exact files
3. Frontend modules and exact files
4. API/routes and exact files
5. Data/config/storage layer and exact files
6. Security layer and exact files
7. Build/deployment evidence and exact files
8. Architecture gaps
9. Recommended architecture improvements with exact files

Do not describe anything not supported by the evidence.
"""

    return ask_llm(prompt)


def grounded_documentation_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    evidence = build_evidence_context(project, files)

    prompt = f"""
You are JARVIS, a strict technical documentation writer.

{grounded_prompt_rules()}

Generate documentation based only on evidence.

PROJECT OVERVIEW:
{project_overview(project)}

EVIDENCE:
{evidence}

Return:

1. Project summary based on actual files
2. Tech stack based on actual files
3. Folder structure
4. Backend components with exact files
5. Frontend components with exact files
6. API/routes with exact files
7. Security-related components with exact files
8. Configuration/deployment files
9. Tests present
10. Known limitations
11. Next improvements

Do not add fictional modules or features.
"""

    return ask_llm(prompt)


def grounded_review(project_name):
    return grounded_review_project(project_name)


def grounded_security_review(project_name):
    return grounded_security_review_project(project_name)


def grounded_architecture_review(project_name):
    return grounded_architecture_review_project(project_name)


def grounded_docs(project_name):
    return grounded_documentation_project(project_name)


# ==========================
# STRICT GROUNDED ANALYZER
# Rule-based analyzer. No LLM. No speculation.
# ==========================
def _contains_any(text, keywords):
    lower = text.lower()

    return any(
        keyword.lower() in lower
        for keyword in keywords
    )


def _files_matching(files, path_tokens=None, content_tokens=None, extensions=None):
    path_tokens = path_tokens or []
    content_tokens = content_tokens or []
    extensions = extensions or []

    matches = []

    for item in files:
        path = normalize_path(item.get("path", ""))
        content = item.get("content", "").lower()
        ext = item.get("extension", "").lower()

        if extensions and ext not in extensions:
            continue

        path_hit = (
            not path_tokens
            or any(token.lower() in path for token in path_tokens)
        )

        content_hit = (
            not content_tokens
            or any(token.lower() in content for token in content_tokens)
        )

        if path_hit and content_hit:
            matches.append(item)

    return matches


def _format_file_list(title, items, limit=30):
    output = [f"\n{title}"]

    if not items:
        output.append("None detected.")
        return output

    for item in items[:limit]:
        output.append(f" - {item['path']}")

    if len(items) > limit:
        output.append(f"... and {len(items) - limit} more")

    return output


def _extract_python_routes(item):
    content = item.get("content", "")
    path = item.get("path", "")

    routes = []

    patterns = [
        r"@router\.(get|post|put|delete|patch)\(\s*[\"']([^\"']+)[\"']",
        r"@app\.(get|post|put|delete|patch|route)\(\s*[\"']([^\"']+)[\"']",
        r"@[\w_]+\.route\(\s*[\"']([^\"']+)[\"']",
    ]

    for pattern in patterns:
        for match in re.findall(pattern, content):
            if isinstance(match, tuple):
                if len(match) == 2:
                    method, route = match
                    routes.append(f"{path} -> {method.upper()} {route}")
                elif len(match) == 1:
                    routes.append(f"{path} -> {match[0]}")
            else:
                routes.append(f"{path} -> {match}")

    return routes


def _extract_js_api_calls(item):
    content = item.get("content", "")
    path = item.get("path", "")

    calls = []

    patterns = [
        r"axios\.(get|post|put|delete|patch)\(\s*[\"']([^\"']+)[\"']",
        r"fetch\(\s*[\"']([^\"']+)[\"']",
    ]

    for pattern in patterns:
        for match in re.findall(pattern, content):
            if isinstance(match, tuple):
                if len(match) == 2:
                    method, endpoint = match
                    calls.append(f"{path} -> {method.upper()} {endpoint}")
            else:
                calls.append(f"{path} -> FETCH {match}")

    return calls


def _extract_package_data(files):
    packages = {}

    for item in files:
        path = normalize_path(item["path"])

        if path.endswith("package.json"):
            packages[item["path"]] = extract_json_dependencies(
                item["content"]
            )

    return packages


def _extract_requirements(files):
    requirements = {}

    for item in files:
        path = normalize_path(item["path"])

        if path.endswith("requirements.txt"):
            requirements[item["path"]] = [
                line.strip()
                for line in item["content"].splitlines()
                if line.strip() and not line.strip().startswith("#")
            ]

    return requirements


def strict_project_facts(project_name):
    project, error = get_project(project_name)

    if error:
        return None, None, error

    files = get_code_files(project)

    facts = {
        "project_name": project.get("name"),
        "project_path": project.get("path"),
        "files_count": project.get("files_count"),
        "tech_stack": project.get("tech_stack", []),
        "extensions": Counter(item["extension"] for item in files),
        "packages": _extract_package_data(files),
        "requirements": _extract_requirements(files),
        "python_files": _files_matching(files, extensions=[".py"]),
        "frontend_files": _files_matching(
            files,
            path_tokens=[
                "src/",
                "components/",
                "frontend",
                "file-scan-log-app"
            ],
            extensions=[".tsx", ".jsx", ".ts", ".js", ".css"]
        ),
        "fastapi_files": _files_matching(
            files,
            content_tokens=[
                "fastapi",
                "apirouter",
                "oauth2passwordrequestform"
            ]
        ),
        "flask_files": _files_matching(
            files,
            content_tokens=[
                "from flask",
                "flask(",
                ".route("
            ]
        ),
        "auth_files": _files_matching(
            files,
            path_tokens=[
                "auth",
                "token",
                "permission",
                "admin"
            ]
        ),
        "jwt_files": _files_matching(
            files,
            content_tokens=[
                "jwt",
                "create_access_token",
                "decode_jwt",
                "access_token"
            ]
        ),
        "password_hashing_files": _files_matching(
            files,
            content_tokens=[
                "bcrypt",
                "argon2",
                "passlib",
                "cryptcontext",
                "pwd_context"
            ]
        ),
        "secret_files": _files_matching(
            files,
            content_tokens=[
                "secret_key",
                "access_token_expire",
                "os.getenv",
                "load_dotenv",
                ".env"
            ]
        ),
        "database_files": _files_matching(
            files,
            content_tokens=[
                "sqlite",
                "sqlalchemy",
                "create_engine",
                "sessionmaker",
                "database",
                "db"
            ]
        ),
        "upload_files": _files_matching(
            files,
            content_tokens=[
                "multipart/form-data",
                "type=\"file\"",
                "upload",
                "formdata",
                "file"
            ]
        ),
        "logging_files": _files_matching(
            files,
            content_tokens=[
                "logging",
                "audit_log",
                "logger",
                "print("
            ]
        ),
        "docker_files": _files_matching(
            files,
            path_tokens=[
                "docker",
                "docker-compose"
            ]
        ),
        "ci_files": _files_matching(
            files,
            path_tokens=[
                ".github",
                "workflows"
            ]
        ),
        "test_files": [
            item for item in files
            if is_test_path(item["path"])
        ],
        "config_files": _files_matching(
            files,
            path_tokens=[
                "config",
                ".env",
                "settings",
                "requirements",
                "package.json",
                "vite.config",
                "tailwind.config",
                "postcss.config",
                "docker-compose"
            ]
        ),
    }

    routes = []
    api_calls = []

    for item in files:
        if item["extension"] == ".py":
            routes.extend(_extract_python_routes(item))

        if item["extension"] in {".js", ".jsx", ".ts", ".tsx"}:
            api_calls.extend(_extract_js_api_calls(item))

    facts["routes"] = routes
    facts["frontend_api_calls"] = api_calls

    return project, facts, None


def strict_grounded_analyzer_project(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    output = [
        "STRICT GROUNDED ANALYZER REPORT",
        "Mode: rule-based / no LLM / no speculation",
        "",
        f"Project: {facts['project_name']}",
        f"Path: {facts['project_path']}",
        f"Files indexed: {facts['files_count']}",
        f"Detected tech stack from deep memory: {', '.join(facts['tech_stack'])}",
        f"Extensions: {dict(facts['extensions'])}",
        "",
        "Dependency evidence:"
    ]

    if facts["packages"]:
        for path, data in facts["packages"].items():
            output.append(f" - {path}: {data}")

    if facts["requirements"]:
        for path, reqs in facts["requirements"].items():
            output.append(f" - {path}: {reqs}")

    if not facts["packages"] and not facts["requirements"]:
        output.append(" - No package.json or requirements.txt detected.")

    output.extend(
        _format_file_list(
            "FastAPI evidence files:",
            facts["fastapi_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Flask evidence files:",
            facts["flask_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Authentication / authorization related files:",
            facts["auth_files"]
        )
    )

    output.extend(
        _format_file_list(
            "JWT / token related files:",
            facts["jwt_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Password hashing evidence files:",
            facts["password_hashing_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Database / storage evidence files:",
            facts["database_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Frontend / UI evidence files:",
            facts["frontend_files"]
        )
    )

    output.extend(
        _format_file_list(
            "File upload / file scan evidence files:",
            facts["upload_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Logging / audit evidence files:",
            facts["logging_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Docker / deployment evidence files:",
            facts["docker_files"]
        )
    )

    output.extend(
        _format_file_list(
            "CI / workflow evidence files:",
            facts["ci_files"]
        )
    )

    output.extend(
        _format_file_list(
            "Test evidence files:",
            facts["test_files"]
        )
    )

    output.append("\nDetected backend routes:")
    if facts["routes"]:
        output.extend(f" - {route}" for route in facts["routes"][:80])
    else:
        output.append("None detected.")

    output.append("\nDetected frontend API calls:")
    if facts["frontend_api_calls"]:
        output.extend(f" - {call}" for call in facts["frontend_api_calls"][:80])
    else:
        output.append("None detected.")

    output.append(
        "\nConclusion:\n"
        "This report only lists evidence found directly in indexed files. "
        "It does not infer hidden services, cloud infrastructure, ML models, or databases unless present in files."
    )

    return "\n".join(output)


def strict_security_analyzer_project(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    findings = []

    for item in facts["auth_files"]:
        content = item["content"].lower()

        if "fake_users_db" in content:
            findings.append(
                f"MEDIUM/HIGH: {item['path']} contains fake_users_db / hardcoded demo users. Replace with real database-backed users before production."
            )

        if "cryptcontext" in content or "bcrypt" in content or "passlib" in content:
            findings.append(
                f"POSITIVE: {item['path']} contains password hashing evidence."
            )

        if "oauth2passwordrequestform" in content:
            findings.append(
                f"POSITIVE: {item['path']} uses OAuth2PasswordRequestForm-style login handling."
            )

    for item in facts["secret_files"]:
        content = item["content"]

        suspicious_defaults = [
            "your-secret-key",
            "changeme",
            "changeme123",
            "secret",
            "password123",
            "secure123"
        ]

        for secret in suspicious_defaults:
            if secret.lower() in content.lower():
                findings.append(
                    f"HIGH: {item['path']} contains weak/default secret or demo credential string: {secret}"
                )

        if "os.getenv" in content:
            findings.append(
                f"POSITIVE: {item['path']} reads configuration/secrets from environment variables."
            )

    for item in facts["jwt_files"]:
        content = item["content"].lower()

        if "jwt" in content:
            findings.append(
                f"INFO: {item['path']} contains JWT/token handling evidence."
            )

        if "expires_delta" in content or "exp" in content:
            findings.append(
                f"POSITIVE: {item['path']} appears to include token expiration handling."
            )

    for item in facts["upload_files"]:
        content = item["content"].lower()

        if "size >" in content or "5 * 1024 * 1024" in content:
            findings.append(
                f"POSITIVE: {item['path']} contains file size validation."
            )

        if "application/pdf" in content or "image/png" in content or "image/jpeg" in content:
            findings.append(
                f"POSITIVE: {item['path']} contains file type validation."
            )

        if "multipart/form-data" in content or "formdata" in content:
            findings.append(
                f"INFO: {item['path']} contains frontend upload/form-data logic."
            )

    for item in facts["logging_files"]:
        content = item["content"].lower()

        if "print(" in content and "audit" in content:
            findings.append(
                f"MEDIUM: {item['path']} uses print-based audit logging. Consider structured file/DB logging with retention."
            )

        if "ip" in content and "username" in content:
            findings.append(
                f"INFO: {item['path']} logs username/IP style data. Review privacy/retention requirements."
            )

    for item in facts["database_files"]:
        content = item["content"].lower()

        if "sqlite" in content:
            findings.append(
                f"INFO: {item['path']} contains SQLite evidence."
            )

        if "sqlalchemy" in content:
            findings.append(
                f"POSITIVE: {item['path']} contains SQLAlchemy ORM evidence."
            )

    output = [
        "STRICT SECURITY ANALYZER REPORT",
        "Mode: rule-based / no LLM / no speculation",
        "",
        f"Project: {facts['project_name']}",
        f"Files indexed: {facts['files_count']}",
        "",
        "Security findings:"
    ]

    if findings:
        output.extend(f" - {finding}" for finding in sorted(set(findings)))
    else:
        output.append(" - No explicit security findings detected by rules.")

    output.extend(
        _format_file_list(
            "Files to inspect first:",
            (
                facts["auth_files"]
                + facts["jwt_files"]
                + facts["secret_files"]
                + facts["upload_files"]
                + facts["logging_files"]
            ),
            limit=40
        )
    )

    output.append(
        "\nRecommended manual checks:\n"
        " - Confirm production secrets are not default/demo values.\n"
        " - Confirm authentication uses persistent users, not fake/demo dictionaries.\n"
        " - Confirm JWT signing key is strong and stored outside source code.\n"
        " - Confirm upload backend validates file type, size, and content, not only frontend.\n"
        " - Confirm logs do not expose sensitive personal data beyond what is necessary."
    )

    return "\n".join(output)


def strict_architecture_analyzer_project(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    output = [
        "STRICT ARCHITECTURE ANALYZER REPORT",
        "Mode: rule-based / no LLM / no speculation",
        "",
        f"Project: {facts['project_name']}",
        f"Path: {facts['project_path']}",
        "",
        "Architecture evidence summary:"
    ]

    if facts["fastapi_files"]:
        output.append(" - FastAPI evidence detected.")
    else:
        output.append(" - FastAPI evidence not detected by rules.")

    if facts["flask_files"]:
        output.append(" - Flask evidence detected.")
    else:
        output.append(" - Flask evidence not detected by rules.")

    if facts["frontend_files"]:
        output.append(" - Frontend/UI evidence detected.")
    else:
        output.append(" - Frontend/UI evidence not detected by rules.")

    if facts["database_files"]:
        output.append(" - Database/storage evidence detected.")
    else:
        output.append(" - Database/storage evidence not detected by rules.")

    if facts["docker_files"]:
        output.append(" - Docker/deployment evidence detected.")
    else:
        output.append(" - Docker/deployment evidence not detected by rules.")

    output.extend(
        _format_file_list(
            "Backend/API files:",
            facts["fastapi_files"] + facts["flask_files"] + facts["auth_files"],
            limit=50
        )
    )

    output.extend(
        _format_file_list(
            "Frontend files:",
            facts["frontend_files"],
            limit=50
        )
    )

    output.extend(
        _format_file_list(
            "Database/config files:",
            facts["database_files"] + facts["config_files"],
            limit=50
        )
    )

    output.extend(
        _format_file_list(
            "Deployment/CI files:",
            facts["docker_files"] + facts["ci_files"],
            limit=30
        )
    )

    output.append("\nRoutes detected:")
    if facts["routes"]:
        output.extend(f" - {route}" for route in facts["routes"][:80])
    else:
        output.append(" - None detected.")

    output.append("\nFrontend API calls detected:")
    if facts["frontend_api_calls"]:
        output.extend(f" - {call}" for call in facts["frontend_api_calls"][:80])
    else:
        output.append(" - None detected.")

    output.append(
        "\nArchitecture gaps detected by rules:\n"
        " - Check whether frontend API base URLs are centralized.\n"
        " - Check whether FastAPI/Flask structure is intentionally mixed or legacy.\n"
        " - Check whether config/secrets are consistently environment-based.\n"
        " - Check whether tests cover routes, permissions, upload, token generation, and dashboard access."
    )

    return "\n".join(output)


# Friendly aliases
def strict_grounded_analyzer(project_name):
    return strict_grounded_analyzer_project(project_name)


def strict_security_analyzer(project_name):
    return strict_security_analyzer_project(project_name)


def strict_architecture_analyzer(project_name):
    return strict_architecture_analyzer_project(project_name)


# ==========================
# CROSS PROJECT INTELLIGENCE
# Rule-based portfolio/project comparison. No LLM.
# ==========================
def _safe_count(items):
    return len(items) if items else 0


def _unique_paths(items):
    seen = set()
    result = []

    for item in items:
        path = item.get("path", "")

        if path and path not in seen:
            seen.add(path)
            result.append(path)

    return result


def _score_project_facts(facts):
    scores = {
        "security": 0,
        "architecture": 0,
        "maintainability": 0
    }

    reasons = {
        "security": [],
        "architecture": [],
        "maintainability": []
    }

    # Security score
    if facts["auth_files"]:
        scores["security"] += 15
        reasons["security"].append("Authentication/authorization files detected.")

    if facts["jwt_files"]:
        scores["security"] += 15
        reasons["security"].append("JWT/token handling evidence detected.")

    if facts["password_hashing_files"]:
        scores["security"] += 15
        reasons["security"].append("Password hashing evidence detected.")

    if facts["secret_files"]:
        scores["security"] += 8
        reasons["security"].append("Environment/config/secret handling evidence detected.")

    if facts["test_files"]:
        scores["security"] += 8
        reasons["security"].append("Tests detected.")

    if facts["ci_files"]:
        scores["security"] += 8
        reasons["security"].append("CI/workflow files detected.")

    if facts["upload_files"]:
        scores["security"] += 6
        reasons["security"].append("Upload/file handling evidence detected.")

    # Penalize weak/demo secrets
    weak_terms = [
        "your-secret-key",
        "changeme",
        "changeme123",
        "password123",
        "secure123"
    ]

    weak_hits = []

    for item in facts["secret_files"] + facts["auth_files"]:
        content = item.get("content", "").lower()

        for term in weak_terms:
            if term in content:
                weak_hits.append(f"{item['path']} -> {term}")

    if weak_hits:
        scores["security"] -= min(25, len(set(weak_hits)) * 5)
        reasons["security"].append(
            "Weak/demo secret or credential strings detected: "
            + ", ".join(sorted(set(weak_hits))[:8])
        )

    # Architecture score
    if facts["fastapi_files"] or facts["flask_files"]:
        scores["architecture"] += 15
        reasons["architecture"].append("Backend framework evidence detected.")

    if facts["frontend_files"]:
        scores["architecture"] += 15
        reasons["architecture"].append("Frontend/UI layer detected.")

    if facts["database_files"]:
        scores["architecture"] += 12
        reasons["architecture"].append("Database/storage layer detected.")

    if facts["routes"]:
        scores["architecture"] += 12
        reasons["architecture"].append("Backend routes detected.")

    if facts["frontend_api_calls"]:
        scores["architecture"] += 8
        reasons["architecture"].append("Frontend API calls detected.")

    if facts["docker_files"]:
        scores["architecture"] += 8
        reasons["architecture"].append("Docker/deployment evidence detected.")

    if facts["packages"] or facts["requirements"]:
        scores["architecture"] += 8
        reasons["architecture"].append("Dependency manifests detected.")

    # Penalize mixed backend frameworks when both are present.
    if facts["fastapi_files"] and facts["flask_files"]:
        scores["architecture"] -= 6
        reasons["architecture"].append(
            "Both FastAPI and Flask evidence detected; verify if this is intentional."
        )

    # Maintainability score
    if facts["test_files"]:
        scores["maintainability"] += 18
        reasons["maintainability"].append("Tests detected.")

    if facts["packages"] or facts["requirements"]:
        scores["maintainability"] += 12
        reasons["maintainability"].append("Dependency manifests detected.")

    if facts["ci_files"]:
        scores["maintainability"] += 10
        reasons["maintainability"].append("CI/workflow files detected.")

    if facts["docker_files"]:
        scores["maintainability"] += 8
        reasons["maintainability"].append("Docker/deployment files detected.")

    if facts["frontend_files"]:
        scores["maintainability"] += 8
        reasons["maintainability"].append("Frontend source organization detected.")

    if facts["python_files"]:
        scores["maintainability"] += 8
        reasons["maintainability"].append("Python source files detected.")

    if facts["files_count"]:
        if facts["files_count"] >= 20:
            scores["maintainability"] += 6
            reasons["maintainability"].append("Non-trivial project size detected.")

    # Clamp and convert to /10
    final_scores = {}

    for key, value in scores.items():
        value = max(0, min(100, value))
        final_scores[key] = round(value / 10, 1)

    final_scores["overall"] = round(
        (
            final_scores["security"]
            + final_scores["architecture"]
            + final_scores["maintainability"]
        ) / 3,
        1
    )

    return final_scores, reasons


def project_scorecard(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return None, error

    scores, reasons = _score_project_facts(facts)

    return {
        "name": facts["project_name"],
        "path": facts["project_path"],
        "files_count": facts["files_count"],
        "tech_stack": facts["tech_stack"],
        "scores": scores,
        "reasons": reasons,
        "evidence": {
            "fastapi_files": _unique_paths(facts["fastapi_files"]),
            "flask_files": _unique_paths(facts["flask_files"]),
            "frontend_files": _unique_paths(facts["frontend_files"]),
            "auth_files": _unique_paths(facts["auth_files"]),
            "jwt_files": _unique_paths(facts["jwt_files"]),
            "password_hashing_files": _unique_paths(facts["password_hashing_files"]),
            "database_files": _unique_paths(facts["database_files"]),
            "docker_files": _unique_paths(facts["docker_files"]),
            "ci_files": _unique_paths(facts["ci_files"]),
            "test_files": _unique_paths(facts["test_files"]),
            "routes": facts["routes"],
            "frontend_api_calls": facts["frontend_api_calls"],
        }
    }, None


def format_project_scorecard(card):
    scores = card["scores"]

    output = [
        f"Project: {card['name']}",
        f"Path: {card['path']}",
        f"Files indexed: {card['files_count']}",
        f"Tech stack: {', '.join(card['tech_stack'])}",
        "",
        "Scores:",
        f" - Security: {scores['security']}/10",
        f" - Architecture: {scores['architecture']}/10",
        f" - Maintainability: {scores['maintainability']}/10",
        f" - Overall: {scores['overall']}/10",
        "",
        "Reasons:"
    ]

    for category in ["security", "architecture", "maintainability"]:
        output.append(f"\n{category.capitalize()}:")

        for reason in card["reasons"][category]:
            output.append(f" - {reason}")

        if not card["reasons"][category]:
            output.append(" - No specific reasons detected.")

    evidence = card["evidence"]

    output.append("\nEvidence summary:")
    output.append(f" - Backend routes: {len(evidence['routes'])}")
    output.append(f" - Frontend API calls: {len(evidence['frontend_api_calls'])}")
    output.append(f" - Auth files: {len(evidence['auth_files'])}")
    output.append(f" - JWT files: {len(evidence['jwt_files'])}")
    output.append(f" - Frontend files: {len(evidence['frontend_files'])}")
    output.append(f" - Test files: {len(evidence['test_files'])}")
    output.append(f" - Docker files: {len(evidence['docker_files'])}")
    output.append(f" - CI files: {len(evidence['ci_files'])}")

    return "\n".join(output)


def score_project(project_name):
    card, error = project_scorecard(project_name)

    if error:
        return error

    return format_project_scorecard(card)


def _split_compare_command(text):
    # Supports: "A and B"
    parts = re.split(r"\s+and\s+", text, maxsplit=1, flags=re.IGNORECASE)

    if len(parts) != 2:
        return None, None, "Use format: compare projects <project A> and <project B>"

    return parts[0].strip(), parts[1].strip(), None


def compare_projects(projects_text):
    first, second, error = _split_compare_command(projects_text)

    if error:
        return error

    card_a, error_a = project_scorecard(first)

    if error_a:
        return error_a

    card_b, error_b = project_scorecard(second)

    if error_b:
        return error_b

    output = [
        "CROSS PROJECT COMPARISON",
        "Mode: rule-based / no LLM / no speculation",
        "",
        f"Project A: {card_a['name']}",
        f"Project B: {card_b['name']}",
        "",
        "Scores:",
        (
            f" - Security: {card_a['name']} {card_a['scores']['security']}/10 "
            f"vs {card_b['name']} {card_b['scores']['security']}/10"
        ),
        (
            f" - Architecture: {card_a['name']} {card_a['scores']['architecture']}/10 "
            f"vs {card_b['name']} {card_b['scores']['architecture']}/10"
        ),
        (
            f" - Maintainability: {card_a['name']} {card_a['scores']['maintainability']}/10 "
            f"vs {card_b['name']} {card_b['scores']['maintainability']}/10"
        ),
        (
            f" - Overall: {card_a['name']} {card_a['scores']['overall']}/10 "
            f"vs {card_b['name']} {card_b['scores']['overall']}/10"
        ),
        "",
        "Evidence counts:",
        (
            f" - Routes: {card_a['name']} {len(card_a['evidence']['routes'])} "
            f"vs {card_b['name']} {len(card_b['evidence']['routes'])}"
        ),
        (
            f" - Frontend API calls: {card_a['name']} {len(card_a['evidence']['frontend_api_calls'])} "
            f"vs {card_b['name']} {len(card_b['evidence']['frontend_api_calls'])}"
        ),
        (
            f" - Auth files: {card_a['name']} {len(card_a['evidence']['auth_files'])} "
            f"vs {card_b['name']} {len(card_b['evidence']['auth_files'])}"
        ),
        (
            f" - JWT files: {card_a['name']} {len(card_a['evidence']['jwt_files'])} "
            f"vs {card_b['name']} {len(card_b['evidence']['jwt_files'])}"
        ),
        (
            f" - Test files: {card_a['name']} {len(card_a['evidence']['test_files'])} "
            f"vs {card_b['name']} {len(card_b['evidence']['test_files'])}"
        ),
        (
            f" - Docker files: {card_a['name']} {len(card_a['evidence']['docker_files'])} "
            f"vs {card_b['name']} {len(card_b['evidence']['docker_files'])}"
        ),
        "",
        "Winner by category:"
    ]

    for category in ["security", "architecture", "maintainability", "overall"]:
        a_score = card_a["scores"][category]
        b_score = card_b["scores"][category]

        if a_score > b_score:
            winner = card_a["name"]
        elif b_score > a_score:
            winner = card_b["name"]
        else:
            winner = "Tie"

        output.append(f" - {category.capitalize()}: {winner}")

    output.append("\nProject A reasons:")
    output.append(format_project_scorecard(card_a))

    output.append("\nProject B reasons:")
    output.append(format_project_scorecard(card_b))

    return "\n".join(output)


def compare_security(projects_text):
    first, second, error = _split_compare_command(projects_text)

    if error:
        return error

    card_a, error_a = project_scorecard(first)

    if error_a:
        return error_a

    card_b, error_b = project_scorecard(second)

    if error_b:
        return error_b

    output = [
        "SECURITY COMPARISON",
        "Mode: rule-based / no LLM / no speculation",
        "",
        f"{card_a['name']}: {card_a['scores']['security']}/10",
        f"{card_b['name']}: {card_b['scores']['security']}/10",
        "",
        f"{card_a['name']} security reasons:"
    ]

    output.extend(f" - {reason}" for reason in card_a["reasons"]["security"])

    output.append(f"\n{card_b['name']} security reasons:")
    output.extend(f" - {reason}" for reason in card_b["reasons"]["security"])

    return "\n".join(output)


def compare_architecture(projects_text):
    first, second, error = _split_compare_command(projects_text)

    if error:
        return error

    card_a, error_a = project_scorecard(first)

    if error_a:
        return error_a

    card_b, error_b = project_scorecard(second)

    if error_b:
        return error_b

    output = [
        "ARCHITECTURE COMPARISON",
        "Mode: rule-based / no LLM / no speculation",
        "",
        f"{card_a['name']}: {card_a['scores']['architecture']}/10",
        f"{card_b['name']}: {card_b['scores']['architecture']}/10",
        "",
        f"{card_a['name']} architecture reasons:"
    ]

    output.extend(f" - {reason}" for reason in card_a["reasons"]["architecture"])

    output.append(f"\n{card_b['name']} architecture reasons:")
    output.extend(f" - {reason}" for reason in card_b["reasons"]["architecture"])

    return "\n".join(output)


def _load_all_deep_projects():
    # Import locally to avoid changing existing imports.
    from deep_project_memory import load_deep_projects

    data = load_deep_projects()

    latest = {}

    for item in data:
        name = item.get("name", "")

        if not name:
            continue

        latest[name.lower()] = item

    return list(latest.values())


def _rank_projects(metric):
    projects = _load_all_deep_projects()

    if not projects:
        return "No deep projects remembered."

    cards = []

    for project in projects:
        card, error = project_scorecard(project.get("name", ""))

        if not error:
            cards.append(card)

    if not cards:
        return "No projects could be scored."

    cards.sort(
        key=lambda card: card["scores"][metric],
        reverse=True
    )

    output = [
        f"PROJECT RANKING BY {metric.upper()}",
        "Mode: rule-based / no LLM / no speculation",
        ""
    ]

    for index, card in enumerate(cards, start=1):
        output.append(
            f"{index}. {card['name']} -> "
            f"{card['scores'][metric]}/10 "
            f"(overall {card['scores']['overall']}/10, files {card['files_count']})"
        )

    return "\n".join(output)


def rank_projects_by_security():
    return _rank_projects("security")


def rank_projects_by_architecture():
    return _rank_projects("architecture")


def rank_projects_by_maintainability():
    return _rank_projects("maintainability")


def best_project_in_memory():
    projects = _load_all_deep_projects()

    if not projects:
        return "No deep projects remembered."

    cards = []

    for project in projects:
        card, error = project_scorecard(project.get("name", ""))

        if not error:
            cards.append(card)

    if not cards:
        return "No projects could be scored."

    cards.sort(
        key=lambda card: card["scores"]["overall"],
        reverse=True
    )

    best = cards[0]

    return (
        "BEST PROJECT IN MEMORY\n"
        "Mode: rule-based / no LLM / no speculation\n\n"
        + format_project_scorecard(best)
    )


# Friendly aliases
def compare_project_pair(projects_text):
    return compare_projects(projects_text)


def compare_project_security(projects_text):
    return compare_security(projects_text)


def compare_project_architecture(projects_text):
    return compare_architecture(projects_text)


# ==========================
# PROJECT REPORT EXPORT
# Generates a Markdown report in reports/
# ==========================
def _safe_report_filename(name):
    cleaned = re.sub(r"[^A-Za-z0-9_.-]+", "_", name.strip())
    cleaned = cleaned.strip("_")

    if not cleaned:
        cleaned = "project"

    return cleaned


def _markdown_code_block(title, content):
    return (
        f"\n\n## {title}\n\n"
        "```text\n"
        f"{content}\n"
        "```\n"
    )


def export_project_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    os.makedirs("reports", exist_ok=True)

    safe_name = _safe_report_filename(project.get("name", project_name))
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    filename = f"{safe_name}_report_{timestamp}.md"
    path = os.path.join("reports", filename)

    sections = []

    sections.append(
        f"# Project Report: {project.get('name')}\n\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        f"Project path: `{project.get('path')}`\n\n"
        f"Files indexed: `{project.get('files_count')}`\n\n"
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`\n"
    )

    sections.append(
        _markdown_code_block(
            "Project Score",
            score_project(project_name)
        )
    )

    sections.append(
        _markdown_code_block(
            "Strict Grounded Analyzer",
            strict_grounded_analyzer_project(project_name)
        )
    )

    sections.append(
        _markdown_code_block(
            "Strict Security Analyzer",
            strict_security_analyzer_project(project_name)
        )
    )

    sections.append(
        _markdown_code_block(
            "Strict Architecture Analyzer",
            strict_architecture_analyzer_project(project_name)
        )
    )

    sections.append(
        _markdown_code_block(
            "Dead Code Scan",
            find_dead_code(project_name)
        )
    )

    sections.append(
        _markdown_code_block(
            "Duplicate Code Scan",
            find_duplicate_code(project_name)
        )
    )

    sections.append(
        "\n\n## Final Recommendation\n\n"
        "- Fix weak/demo secrets and credentials first.\n"
        "- Replace demo users with persistent database-backed authentication if needed.\n"
        "- Verify whether mixed FastAPI/Flask architecture is intentional.\n"
        "- Keep frontend validation, but confirm backend validation for uploads.\n"
        "- Run tests before deleting any dead-code candidates.\n"
    )

    content = "\n".join(sections)

    with open(path, "w", encoding="utf-8") as f:
        f.write(content)

    return (
        "Project report exported successfully.\n"
        f"Project: {project.get('name')}\n"
        f"File: {path}"
    )


def export_report(project_name):
    return export_project_report(project_name)


def export_project_markdown_report(project_name):
    return export_project_report(project_name)


# ==========================
# SMART DAILY PROJECT CHECK
# Daily overview across all deep-memory projects.
# Rule-based / no LLM / no speculation.
# ==========================
def _top_security_warnings_from_card(card, limit=5):
    warnings = []

    for reason in card.get("reasons", {}).get("security", []):
        lower = reason.lower()

        if (
            "weak/demo" in lower
            or "credential" in lower
            or "secret" in lower
            or "auth" in lower
            or "jwt" in lower
        ):
            warnings.append(reason)

    return warnings[:limit]


def _daily_priority_for_card(card):
    scores = card["scores"]
    evidence = card["evidence"]
    priorities = []

    if scores["security"] < 6:
        priorities.append(
            "Security score is below 6/10. Start with secrets, demo credentials, authentication, JWT, and upload validation."
        )

    if scores["architecture"] < 6:
        priorities.append(
            "Architecture score is below 6/10. Verify backend structure, routes, dependency manifests, and deployment files."
        )

    if scores["maintainability"] < 6:
        priorities.append(
            "Maintainability score is below 6/10. Add tests, CI, documentation, and reduce duplicated/dead code."
        )

    if not evidence.get("test_files"):
        priorities.append(
            "No test files detected. Add basic tests for authentication, API routes, and critical workflows."
        )

    if not evidence.get("docker_files"):
        priorities.append(
            "No Docker/deployment evidence detected. Add Dockerfile or docker-compose if deployment matters."
        )

    if not priorities:
        priorities.append(
            "Project looks stable by current rules. Focus on cleaning warnings and improving documentation."
        )

    return priorities


def daily_project_check():
    projects = _load_all_deep_projects()

    if not projects:
        return (
            "SMART DAILY PROJECT CHECK\n"
            "No deep projects remembered.\n\n"
            "Run: remember deep project <project name>"
        )

    cards = []

    for project in projects:
        card, error = project_scorecard(project.get("name", ""))

        if not error:
            cards.append(card)

    if not cards:
        return "No projects could be scored for daily check."

    cards.sort(
        key=lambda card: card["scores"]["overall"],
        reverse=True
    )

    best = cards[0]
    weakest_security = min(
        cards,
        key=lambda card: card["scores"]["security"]
    )
    weakest_maintainability = min(
        cards,
        key=lambda card: card["scores"]["maintainability"]
    )

    output = [
        "SMART DAILY PROJECT CHECK",
        "Mode: rule-based / no LLM / no speculation",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        "",
        f"Projects in deep memory: {len(cards)}",
        "",
        "Project ranking today:"
    ]

    for index, card in enumerate(cards, start=1):
        output.append(
            f"{index}. {card['name']} -> "
            f"overall {card['scores']['overall']}/10 "
            f"(security {card['scores']['security']}/10, "
            f"architecture {card['scores']['architecture']}/10, "
            f"maintainability {card['scores']['maintainability']}/10)"
        )

    output.extend([
        "",
        "Best project today:",
        (
            f" - {best['name']} with overall score "
            f"{best['scores']['overall']}/10"
        ),
        "",
        "Project needing most security attention:",
        (
            f" - {weakest_security['name']} with security score "
            f"{weakest_security['scores']['security']}/10"
        ),
        "",
        "Project needing most maintainability attention:",
        (
            f" - {weakest_maintainability['name']} with maintainability score "
            f"{weakest_maintainability['scores']['maintainability']}/10"
        ),
        "",
        "Top issues detected today:"
    ])

    issue_count = 0

    for card in cards:
        warnings = _top_security_warnings_from_card(card, limit=3)

        for warning in warnings:
            issue_count += 1
            output.append(f" - {card['name']}: {warning}")

    if issue_count == 0:
        output.append(" - No major security warnings detected by current rules.")

    output.extend([
        "",
        "Recommended action plan for today:"
    ])

    # Choose the most useful next task:
    focus_card = weakest_security

    if weakest_security["scores"]["security"] >= 6:
        focus_card = weakest_maintainability

    output.append(f"Focus project: {focus_card['name']}")

    for index, priority in enumerate(
        _daily_priority_for_card(focus_card)[:5],
        start=1
    ):
        output.append(f"{index}. {priority}")

    output.extend([
        "",
        "Suggested commands:",
        f" - strict security analyzer project {focus_card['name']}",
        f" - find dead code in project {focus_card['name']}",
        f" - find duplicates in project {focus_card['name']}",
        f" - export report {focus_card['name']}",
        "",
        "Daily recommendation:",
        "Fix the highest-risk security/configuration issues first, then improve tests and documentation."
    ])

    return "\n".join(output)


def smart_daily_check():
    return daily_project_check()


def daily_check():
    return daily_project_check()


# ==========================
# SUGGEST FIXES
# Rule-based fix suggestions. No automatic code changes.
# ==========================
def _add_fix(fixes, priority, file_path, issue, recommendation, why):
    fixes.append({
        "priority": priority,
        "file": file_path,
        "issue": issue,
        "recommendation": recommendation,
        "why": why
    })


def _priority_value(priority):
    order = {
        "CRITICAL": 0,
        "HIGH": 1,
        "MEDIUM": 2,
        "LOW": 3
    }

    return order.get(priority.upper(), 9)


def _format_fixes_report(project_name, fixes):
    if not fixes:
        return (
            "SUGGEST FIXES REPORT\n"
            "Mode: rule-based / no LLM / no automatic code changes\n\n"
            f"Project: {project_name}\n\n"
            "No concrete fixes were detected by the current rules.\n"
            "Run strict analyzers for deeper inspection."
        )

    fixes.sort(
        key=lambda item: (
            _priority_value(item["priority"]),
            item["file"].lower(),
            item["issue"].lower()
        )
    )

    output = [
        "SUGGEST FIXES REPORT",
        "Mode: rule-based / no LLM / no automatic code changes",
        "",
        f"Project: {project_name}",
        "",
        "Important:",
        "These are safe recommendations only. JARVIS does not modify files automatically.",
        "",
        "Priority fixes:"
    ]

    for index, fix in enumerate(fixes, start=1):
        output.extend([
            "",
            f"{index}. [{fix['priority']}] {fix['file']}",
            f"Issue: {fix['issue']}",
            f"Suggested fix: {fix['recommendation']}",
            f"Why: {fix['why']}"
        ])

    output.extend([
        "",
        "Recommended workflow:",
        "1. Fix HIGH priority items first.",
        "2. Run tests.",
        "3. Run strict security analyzer again.",
        "4. Export a new report.",
        "",
        f"Useful commands:",
        f" - strict security analyzer project {project_name}",
        f" - strict architecture analyzer project {project_name}",
        f" - find dead code in project {project_name}",
        f" - find duplicates in project {project_name}",
        f" - export report {project_name}"
    ])

    return "\n".join(output)


def suggest_fixes_for_project(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    fixes = []

    # Auth / demo users / weak credentials
    for item in facts["auth_files"]:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        if "fake_users_db" in lower:
            _add_fix(
                fixes,
                "HIGH",
                path,
                "Hardcoded demo user database detected: fake_users_db.",
                (
                    "Move users to a real database table/model, remove fake_users_db, "
                    "and load credentials only through registration/admin workflows."
                ),
                "Hardcoded users are unsafe and make authentication unrealistic for production."
            )

        if "password123" in lower or "secure123" in lower:
            _add_fix(
                fixes,
                "HIGH",
                path,
                "Weak/demo password string detected.",
                (
                    "Remove demo credentials from source code. Use environment variables "
                    "only for setup secrets and store real user passwords as hashes in the database."
                ),
                "Demo credentials can be reused accidentally and are easy to guess."
            )

        if "is_admin" in lower and "fake_users_db" in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Admin role is stored in demo dictionary.",
                (
                    "Move role/permission data to a database-backed user model and protect admin routes "
                    "with a reusable dependency such as require_admin_user."
                ),
                "Role checks should be centralized and persistent, not stored in local dictionaries."
            )

        if "httpexception(status_code=401" in lower and "www-authenticate" not in lower:
            _add_fix(
                fixes,
                "LOW",
                path,
                "401 response does not appear to include WWW-Authenticate header.",
                (
                    "Add headers={\"WWW-Authenticate\": \"Bearer\"} to authentication failures "
                    "if using OAuth2/Bearer token flows."
                ),
                "This improves standards compatibility for OAuth2 clients."
            )

    # Secrets / config
    for item in facts["secret_files"]:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        weak_terms = [
            "your-secret-key",
            "changeme",
            "changeme123",
            "secret",
            "password123",
            "secure123"
        ]

        for term in weak_terms:
            if term in lower:
                _add_fix(
                    fixes,
                    "HIGH",
                    path,
                    f"Weak/default secret or credential detected: {term}.",
                    (
                        "Replace default/demo values with required environment variables. "
                        "Fail fast at startup if the variable is missing or too weak."
                    ),
                    "Default secrets can allow token forgery, unauthorized access, or accidental exposure."
                )

        if "os.getenv" in content and "raise" not in lower and "runtimeerror" not in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Environment variables are used, but missing-value validation is not obvious.",
                (
                    "Add startup validation for required environment variables such as JWT secret, "
                    "database URL, email/API keys, and encryption keys."
                ),
                "A missing env var can silently fall back to insecure defaults or break runtime behavior."
            )

    # JWT / token handling
    for item in facts["jwt_files"]:
        path = item["path"]
        lower = item["content"].lower()

        if "jwt" in lower and ("secret_key" in lower or "algorithm" in lower):
            if "exp" not in lower and "expires_delta" not in lower:
                _add_fix(
                    fixes,
                    "HIGH",
                    path,
                    "JWT/token handling detected without obvious expiration handling.",
                    (
                        "Add exp claim to tokens and validate token expiration on every protected request."
                    ),
                    "Tokens without expiration remain valid too long if leaked."
                )

        if "decode" in lower and "except" not in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "JWT decode logic does not appear to have explicit exception handling.",
                (
                    "Wrap token decode in try/except and return clean 401 errors for expired, invalid, "
                    "or malformed tokens."
                ),
                "Bad tokens should not crash the application or leak internal errors."
            )

    # Upload / file scan
    for item in facts["upload_files"]:
        path = item["path"]
        lower = item["content"].lower()

        if "type=\"file\"" in lower or "formdata" in lower or "multipart/form-data" in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Frontend file upload logic detected.",
                (
                    "Keep frontend validation, but confirm the backend also validates file type, size, "
                    "extension, filename, and content before processing."
                ),
                "Frontend validation can be bypassed; backend validation is mandatory."
            )

        if "upload" in lower and "filename" in lower and "secure_filename" not in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Upload filename handling detected without obvious filename sanitization.",
                (
                    "Sanitize uploaded filenames and store files using generated safe names, "
                    "not user-controlled names."
                ),
                "Unsafe filenames may cause path traversal or overwrite issues."
            )

    # Logging / privacy
    for item in facts["logging_files"]:
        path = item["path"]
        lower = item["content"].lower()

        if "print(" in lower and ("audit" in lower or "logger" in lower):
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Print-based logging appears to be used for audit/security output.",
                (
                    "Replace print-based audit logs with structured logging, rotating files, or database audit events."
                ),
                "Security logs need timestamps, levels, retention, and consistent formatting."
            )

        if "username" in lower and "ip" in lower:
            _add_fix(
                fixes,
                "LOW",
                path,
                "Logs appear to include username/IP data.",
                (
                    "Review log retention and privacy rules. Avoid logging secrets, passwords, tokens, or full personal data."
                ),
                "Audit logs are useful, but they must not expose unnecessary personal or sensitive data."
            )

    # Database / architecture
    if facts["fastapi_files"] and facts["flask_files"]:
        _add_fix(
            fixes,
            "MEDIUM",
            "project architecture",
            "Both FastAPI and Flask evidence detected.",
            (
                "Decide whether both frameworks are intentionally used. If not, consolidate backend code "
                "around one framework to reduce complexity."
            ),
            "Mixed frameworks can make routing, auth, testing, and deployment harder to maintain."
        )

    if not facts["test_files"]:
        _add_fix(
            fixes,
            "MEDIUM",
            "tests/",
            "No test files detected.",
            (
                "Add tests for login/authentication, protected routes, token validation, upload validation, "
                "and key project workflows."
            ),
            "Tests improve maintainability and prevent regressions."
        )

    if not facts["ci_files"]:
        _add_fix(
            fixes,
            "LOW",
            ".github/workflows/",
            "No CI/workflow files detected.",
            (
                "Add a basic CI workflow that installs dependencies and runs tests/lint checks on push."
            ),
            "CI catches errors before code is merged or deployed."
        )

    if not facts["docker_files"]:
        _add_fix(
            fixes,
            "LOW",
            "Dockerfile / docker-compose.yml",
            "No Docker/deployment evidence detected.",
            (
                "Add Dockerfile or docker-compose.yml if you want reproducible local/prod deployment."
            ),
            "Reproducible environments make setup and deployment easier."
        )

    # Dead code and duplicate hints
    duplicate_result = find_duplicate_code(project_name)

    if "Duplicate #" in duplicate_result:
        _add_fix(
            fixes,
            "MEDIUM",
            "duplicated code candidates",
            "Duplicated business logic detected.",
            (
                "Run 'find duplicates in project <project>' and extract repeated logic into shared utilities/services."
            ),
            "Duplicated logic increases maintenance cost and bug risk."
        )

    dead_code_result = find_dead_code(project_name)

    if "Possible unused production" in dead_code_result:
        _add_fix(
            fixes,
            "LOW",
            "dead code candidates",
            "Possible unused code detected.",
            (
                "Run 'find dead code in project <project>', manually verify each candidate, then remove only after tests pass."
            ),
            "Removing verified dead code reduces complexity, but heuristic results must be checked manually."
        )

    return _format_fixes_report(
        project.get("name", project_name),
        fixes
    )


def suggest_project_fixes(project_name):
    return suggest_fixes_for_project(project_name)


def project_fixes(project_name):
    return suggest_fixes_for_project(project_name)


def suggest_fixes(project_name):
    return suggest_fixes_for_project(project_name)


# ==========================
# FALSE POSITIVE REDUCTION ENGINE
# Improves security findings and fix suggestions.
# Rule-based / no LLM / no automatic code changes.
# ==========================
IGNORED_SECURITY_PATH_PREFIXES = (
    "memory/",
    "reports/",
    ".git/",
    "node_modules/",
    "venv/",
    ".venv/",
    "jarvis-env/",
    "__pycache__/",
)

LOW_CONFIDENCE_SECRET_PATH_PREFIXES = (
    "docs/",
    "documentation/",
)

LOW_CONFIDENCE_SECRET_EXTENSIONS = {
    ".md",
    ".txt",
    ".json",
    ".lock",
}

REAL_SECRET_NAME_PATTERN = re.compile(
    r"\b("
    r"secret_key|jwt_secret|jwt_secret_key|access_token_secret|"
    r"api_key|private_key|encryption_key|hmac_secret|aes_secret|"
    r"password_pepper|password_salt|client_secret|token_secret|"
    r"openai_api_key|aws_secret_access_key"
    r")\b",
    re.IGNORECASE
)

ENV_READ_PATTERN = re.compile(
    r"os\.getenv\(\s*[\"']([A-Z0-9_]*(SECRET|KEY|TOKEN|PASSWORD|PEPPER|SALT)[A-Z0-9_]*)[\"']",
    re.IGNORECASE
)

ASSIGNMENT_SECRET_PATTERN = re.compile(
    r"(?i)\b([A-Z0-9_]*(SECRET|KEY|TOKEN|PASSWORD|PEPPER|SALT)[A-Z0-9_]*)\b\s*=\s*[\"']([^\"']+)[\"']"
)

WEAK_SECRET_VALUES = {
    "your-secret-key",
    "changeme",
    "changeme123",
    "password123",
    "secure123",
    "admin",
    "admin123",
    "test",
    "test123",
    "dev",
    "demo",
}


def _fp_path(path):
    return normalize_path(path or "")


def _is_ignored_security_file(path):
    lower = _fp_path(path)

    return any(
        lower.startswith(prefix)
        or f"/{prefix}" in lower
        for prefix in IGNORED_SECURITY_PATH_PREFIXES
    )


def _is_low_confidence_secret_file(path):
    lower = _fp_path(path)
    ext = os.path.splitext(lower)[1]

    if ext in LOW_CONFIDENCE_SECRET_EXTENSIONS:
        return True

    return any(
        lower.startswith(prefix)
        or f"/{prefix}" in lower
        for prefix in LOW_CONFIDENCE_SECRET_PATH_PREFIXES
    )


def _is_source_code_file(path):
    ext = os.path.splitext(_fp_path(path))[1]

    return ext in {
        ".py",
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".env",
        ".ini",
        ".cfg",
        ".yml",
        ".yaml",
        ".toml",
    }


def _security_relevant_items(items):
    filtered = []

    for item in items:
        path = item.get("path", "")

        if _is_ignored_security_file(path):
            continue

        filtered.append(item)

    return filtered


def _looks_like_real_secret_context(path, content):
    if _is_ignored_security_file(path):
        return False

    # For docs/json/md/txt, avoid simple word matches like "secret" in prose.
    if _is_low_confidence_secret_file(path):
        return bool(
            REAL_SECRET_NAME_PATTERN.search(content)
            or ENV_READ_PATTERN.search(content)
            or ASSIGNMENT_SECRET_PATTERN.search(content)
        )

    if not _is_source_code_file(path):
        return False

    return bool(
        REAL_SECRET_NAME_PATTERN.search(content)
        or ENV_READ_PATTERN.search(content)
        or ASSIGNMENT_SECRET_PATTERN.search(content)
    )


def _weak_secret_matches(path, content):
    if not _looks_like_real_secret_context(path, content):
        return []

    matches = []
    lower_content = content.lower()

    # Match obvious weak defaults, but avoid flagging the plain word "secret".
    for value in WEAK_SECRET_VALUES:
        if value in lower_content:
            matches.append(value)

    # Flag exact "secret" only when it appears as an assigned value for a secret-like variable.
    for match in ASSIGNMENT_SECRET_PATTERN.findall(content):
        variable_name = match[0]
        assigned_value = match[2].strip()

        if assigned_value.lower() == "secret":
            matches.append(f"{variable_name}=secret")

    return sorted(set(matches))


def _env_validation_missing(content):
    if "os.getenv" not in content:
        return False

    if "raise" in content.lower() or "runtimeerror" in content.lower() or "valueerror" in content.lower():
        return False

    return bool(ENV_READ_PATTERN.search(content))


def _security_finding_sort_key(text):
    order = {
        "CRITICAL": 0,
        "HIGH": 1,
        "MEDIUM/HIGH": 2,
        "MEDIUM": 3,
        "LOW": 4,
        "POSITIVE": 5,
        "INFO": 6,
    }

    label = text.split(":", 1)[0].strip().upper()
    return (order.get(label, 9), text.lower())


def _dedupe_keep_order(items):
    seen = set()
    result = []

    for item in items:
        if item in seen:
            continue

        seen.add(item)
        result.append(item)

    return result


# Override old strict security analyzer with false-positive reduction.
def strict_security_analyzer_project(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    findings = []

    auth_files = _security_relevant_items(facts["auth_files"])
    secret_files = _security_relevant_items(facts["secret_files"])
    jwt_files = _security_relevant_items(facts["jwt_files"])
    upload_files = _security_relevant_items(facts["upload_files"])
    logging_files = _security_relevant_items(facts["logging_files"])
    database_files = _security_relevant_items(facts["database_files"])

    for item in auth_files:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        if "fake_users_db" in lower:
            findings.append(
                f"HIGH: {path} contains fake_users_db / hardcoded demo users. Replace with real database-backed users before production."
            )

        if "password123" in lower or "secure123" in lower:
            findings.append(
                f"HIGH: {path} contains weak/demo password strings. Remove demo credentials from source code."
            )

        if "cryptcontext" in lower or "bcrypt" in lower or "passlib" in lower:
            findings.append(
                f"POSITIVE: {path} contains password hashing evidence."
            )

        if "oauth2passwordrequestform" in lower:
            findings.append(
                f"POSITIVE: {path} uses OAuth2PasswordRequestForm-style login handling."
            )

    for item in secret_files:
        path = item["path"]
        content = item["content"]

        weak_matches = _weak_secret_matches(path, content)

        for weak in weak_matches:
            findings.append(
                f"HIGH: {path} contains weak/default secret or demo credential string: {weak}"
            )

        if ENV_READ_PATTERN.search(content):
            findings.append(
                f"POSITIVE: {path} reads security-sensitive configuration from environment variables."
            )

        if _env_validation_missing(content):
            findings.append(
                f"MEDIUM: {path} reads security-sensitive environment variables, but missing-value validation is not obvious."
            )

    for item in jwt_files:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        if _is_ignored_security_file(path):
            continue

        if "jwt" in lower:
            findings.append(
                f"INFO: {path} contains JWT/token handling evidence."
            )

        if "expires_delta" in lower or "exp" in lower:
            findings.append(
                f"POSITIVE: {path} appears to include token expiration handling."
            )
        elif "jwt" in lower and ("encode" in lower or "decode" in lower):
            findings.append(
                f"HIGH: {path} contains JWT handling without obvious expiration evidence."
            )

        if "decode" in lower and "except" not in lower:
            findings.append(
                f"MEDIUM: {path} decodes tokens without obvious exception handling."
            )

    for item in upload_files:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        if _is_ignored_security_file(path):
            continue

        if "size >" in lower or "5 * 1024 * 1024" in lower:
            findings.append(
                f"POSITIVE: {path} contains file size validation."
            )

        if "application/pdf" in lower or "image/png" in lower or "image/jpeg" in lower:
            findings.append(
                f"POSITIVE: {path} contains file type validation."
            )

        if "multipart/form-data" in lower or "formdata" in lower or "type=\"file\"" in lower:
            findings.append(
                f"INFO: {path} contains upload/form-data logic. Confirm backend validation exists."
            )

        if "filename" in lower and "secure_filename" not in lower and _is_source_code_file(path):
            findings.append(
                f"MEDIUM: {path} references uploaded filenames without obvious filename sanitization."
            )

    for item in logging_files:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        if _is_ignored_security_file(path):
            continue

        if "print(" in lower and ("audit" in lower or "logger" in lower):
            findings.append(
                f"MEDIUM: {path} uses print-based audit/security logging. Consider structured logging with retention."
            )

        if "ip" in lower and "username" in lower and _is_source_code_file(path):
            findings.append(
                f"LOW: {path} logs username/IP style data. Review privacy and retention requirements."
            )

    for item in database_files:
        path = item["path"]
        content = item["content"].lower()

        if _is_ignored_security_file(path):
            continue

        if "sqlite" in content:
            findings.append(
                f"INFO: {path} contains SQLite evidence."
            )

        if "sqlalchemy" in content:
            findings.append(
                f"POSITIVE: {path} contains SQLAlchemy ORM evidence."
            )

    findings = _dedupe_keep_order(findings)
    findings.sort(key=_security_finding_sort_key)

    output = [
        "STRICT SECURITY ANALYZER REPORT",
        "Mode: rule-based / no LLM / no speculation",
        "False-positive reduction: enabled",
        "",
        f"Project: {facts['project_name']}",
        f"Files indexed: {facts['files_count']}",
        "",
        "Security findings:"
    ]

    if findings:
        output.extend(f" - {finding}" for finding in findings)
    else:
        output.append(" - No explicit security findings detected by rules.")

    first_files = (
        auth_files
        + jwt_files
        + secret_files
        + upload_files
        + logging_files
    )

    output.extend(
        _format_file_list(
            "Files to inspect first:",
            first_files,
            limit=40
        )
    )

    output.append(
        "\nRecommended manual checks:\n"
        " - Confirm production secrets are not default/demo values.\n"
        " - Ignore findings in docs/memory/reports unless they contain real active credentials.\n"
        " - Confirm authentication uses persistent users, not fake/demo dictionaries.\n"
        " - Confirm JWT signing key is strong and stored outside source code.\n"
        " - Confirm upload backend validates file type, size, filename, and content.\n"
        " - Confirm logs do not expose passwords, tokens, secrets, or unnecessary personal data."
    )

    return "\n".join(output)


# Override scoring to use reduced false positives.
def _score_project_facts(facts):
    scores = {
        "security": 0,
        "architecture": 0,
        "maintainability": 0
    }

    reasons = {
        "security": [],
        "architecture": [],
        "maintainability": []
    }

    auth_files = _security_relevant_items(facts["auth_files"])
    jwt_files = _security_relevant_items(facts["jwt_files"])
    password_hashing_files = _security_relevant_items(facts["password_hashing_files"])
    secret_files = _security_relevant_items(facts["secret_files"])
    upload_files = _security_relevant_items(facts["upload_files"])

    if auth_files:
        scores["security"] += 15
        reasons["security"].append("Authentication/authorization files detected.")

    if jwt_files:
        scores["security"] += 15
        reasons["security"].append("JWT/token handling evidence detected.")

    if password_hashing_files:
        scores["security"] += 15
        reasons["security"].append("Password hashing evidence detected.")

    if secret_files:
        scores["security"] += 8
        reasons["security"].append("Environment/config/secret handling evidence detected.")

    if facts["test_files"]:
        scores["security"] += 8
        reasons["security"].append("Tests detected.")

    if facts["ci_files"]:
        scores["security"] += 8
        reasons["security"].append("CI/workflow files detected.")

    if upload_files:
        scores["security"] += 6
        reasons["security"].append("Upload/file handling evidence detected.")

    weak_hits = []

    for item in secret_files + auth_files:
        for weak in _weak_secret_matches(item.get("path", ""), item.get("content", "")):
            weak_hits.append(f"{item['path']} -> {weak}")

    if weak_hits:
        scores["security"] -= min(25, len(set(weak_hits)) * 5)
        reasons["security"].append(
            "Weak/demo secret or credential strings detected: "
            + ", ".join(sorted(set(weak_hits))[:8])
        )

    if facts["fastapi_files"] or facts["flask_files"]:
        scores["architecture"] += 15
        reasons["architecture"].append("Backend framework evidence detected.")

    if facts["frontend_files"]:
        scores["architecture"] += 15
        reasons["architecture"].append("Frontend/UI layer detected.")

    if facts["database_files"]:
        scores["architecture"] += 12
        reasons["architecture"].append("Database/storage layer detected.")

    if facts["routes"]:
        scores["architecture"] += 12
        reasons["architecture"].append("Backend routes detected.")

    if facts["frontend_api_calls"]:
        scores["architecture"] += 8
        reasons["architecture"].append("Frontend API calls detected.")

    if facts["docker_files"]:
        scores["architecture"] += 8
        reasons["architecture"].append("Docker/deployment evidence detected.")

    if facts["packages"] or facts["requirements"]:
        scores["architecture"] += 8
        reasons["architecture"].append("Dependency manifests detected.")

    if facts["fastapi_files"] and facts["flask_files"]:
        scores["architecture"] -= 6
        reasons["architecture"].append(
            "Both FastAPI and Flask evidence detected; verify if this is intentional."
        )

    if facts["test_files"]:
        scores["maintainability"] += 18
        reasons["maintainability"].append("Tests detected.")

    if facts["packages"] or facts["requirements"]:
        scores["maintainability"] += 12
        reasons["maintainability"].append("Dependency manifests detected.")

    if facts["ci_files"]:
        scores["maintainability"] += 10
        reasons["maintainability"].append("CI/workflow files detected.")

    if facts["docker_files"]:
        scores["maintainability"] += 8
        reasons["maintainability"].append("Docker/deployment files detected.")

    if facts["frontend_files"]:
        scores["maintainability"] += 8
        reasons["maintainability"].append("Frontend source organization detected.")

    if facts["python_files"]:
        scores["maintainability"] += 8
        reasons["maintainability"].append("Python source files detected.")

    if facts["files_count"] and facts["files_count"] >= 20:
        scores["maintainability"] += 6
        reasons["maintainability"].append("Non-trivial project size detected.")

    final_scores = {}

    for key, value in scores.items():
        value = max(0, min(100, value))
        final_scores[key] = round(value / 10, 1)

    final_scores["overall"] = round(
        (
            final_scores["security"]
            + final_scores["architecture"]
            + final_scores["maintainability"]
        ) / 3,
        1
    )

    return final_scores, reasons


# Override suggest fixes with false-positive reduction and better prioritization.
def suggest_fixes_for_project(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    fixes = []

    auth_files = _security_relevant_items(facts["auth_files"])
    secret_files = _security_relevant_items(facts["secret_files"])
    jwt_files = _security_relevant_items(facts["jwt_files"])
    upload_files = _security_relevant_items(facts["upload_files"])
    logging_files = _security_relevant_items(facts["logging_files"])

    for item in auth_files:
        path = item["path"]
        content = item["content"]
        lower = content.lower()

        if "fake_users_db" in lower:
            _add_fix(
                fixes,
                "HIGH",
                path,
                "Hardcoded demo user database detected: fake_users_db.",
                (
                    "Move users to a real database table/model, remove fake_users_db, "
                    "and load credentials only through registration/admin workflows."
                ),
                "Hardcoded users are unsafe and make authentication unrealistic for production."
            )

        if "password123" in lower or "secure123" in lower:
            _add_fix(
                fixes,
                "HIGH",
                path,
                "Weak/demo password string detected.",
                (
                    "Remove demo credentials from source code. Store real user passwords only as hashes "
                    "in the database."
                ),
                "Demo credentials can be reused accidentally and are easy to guess."
            )

        if "is_admin" in lower and "fake_users_db" in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Admin role is stored in demo dictionary.",
                (
                    "Move role/permission data to a database-backed user model and protect admin routes "
                    "with a reusable dependency such as require_admin_user."
                ),
                "Role checks should be centralized and persistent, not stored in local dictionaries."
            )

        if "httpexception(status_code=401" in lower and "www-authenticate" not in lower:
            _add_fix(
                fixes,
                "LOW",
                path,
                "401 response does not appear to include WWW-Authenticate header.",
                (
                    "Add headers={\"WWW-Authenticate\": \"Bearer\"} to authentication failures "
                    "if using OAuth2/Bearer token flows."
                ),
                "This improves standards compatibility for OAuth2 clients."
            )

    for item in secret_files:
        path = item["path"]
        content = item["content"]

        for weak in _weak_secret_matches(path, content):
            _add_fix(
                fixes,
                "HIGH",
                path,
                f"Weak/default secret or credential detected: {weak}.",
                (
                    "Replace default/demo values with required environment variables. "
                    "Fail fast at startup if the variable is missing or too weak."
                ),
                "Default secrets can allow token forgery, unauthorized access, or accidental exposure."
            )

        if _env_validation_missing(content):
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Security-sensitive environment variables are used, but missing-value validation is not obvious.",
                (
                    "Add startup validation for required environment variables such as JWT secret, "
                    "database URL, email/API keys, and encryption keys."
                ),
                "A missing env var can silently fall back to insecure defaults or break runtime behavior."
            )

    for item in jwt_files:
        path = item["path"]
        lower = item["content"].lower()

        if _is_ignored_security_file(path):
            continue

        if "jwt" in lower and ("encode" in lower or "decode" in lower):
            if "exp" not in lower and "expires_delta" not in lower:
                _add_fix(
                    fixes,
                    "HIGH",
                    path,
                    "JWT/token handling detected without obvious expiration handling.",
                    "Add exp claim to tokens and validate token expiration on every protected request.",
                    "Tokens without expiration remain valid too long if leaked."
                )

        if "decode" in lower and "except" not in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "JWT decode logic does not appear to have explicit exception handling.",
                (
                    "Wrap token decode in try/except and return clean 401 errors for expired, invalid, "
                    "or malformed tokens."
                ),
                "Bad tokens should not crash the application or leak internal errors."
            )

    for item in upload_files:
        path = item["path"]
        lower = item["content"].lower()

        if _is_ignored_security_file(path):
            continue

        if "type=\"file\"" in lower or "formdata" in lower or "multipart/form-data" in lower:
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Upload/form-data logic detected.",
                (
                    "Keep frontend validation, but confirm the backend also validates file type, size, "
                    "extension, filename, and content before processing."
                ),
                "Frontend validation can be bypassed; backend validation is mandatory."
            )

        if "filename" in lower and "secure_filename" not in lower and _is_source_code_file(path):
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Upload filename handling detected without obvious filename sanitization.",
                (
                    "Sanitize uploaded filenames and store files using generated safe names, "
                    "not user-controlled names."
                ),
                "Unsafe filenames may cause path traversal or overwrite issues."
            )

    for item in logging_files:
        path = item["path"]
        lower = item["content"].lower()

        if _is_ignored_security_file(path):
            continue

        if "print(" in lower and ("audit" in lower or "logger" in lower):
            _add_fix(
                fixes,
                "MEDIUM",
                path,
                "Print-based logging appears to be used for audit/security output.",
                (
                    "Replace print-based audit logs with structured logging, rotating files, or database audit events."
                ),
                "Security logs need timestamps, levels, retention, and consistent formatting."
            )

        if "username" in lower and "ip" in lower and _is_source_code_file(path):
            _add_fix(
                fixes,
                "LOW",
                path,
                "Logs appear to include username/IP data.",
                (
                    "Review log retention and privacy rules. Avoid logging secrets, passwords, tokens, or full personal data."
                ),
                "Audit logs are useful, but they must not expose unnecessary personal or sensitive data."
            )

    if facts["fastapi_files"] and facts["flask_files"]:
        _add_fix(
            fixes,
            "MEDIUM",
            "project architecture",
            "Both FastAPI and Flask evidence detected.",
            (
                "Decide whether both frameworks are intentionally used. If not, consolidate backend code "
                "around one framework to reduce complexity."
            ),
            "Mixed frameworks can make routing, auth, testing, and deployment harder to maintain."
        )

    if not facts["test_files"]:
        _add_fix(
            fixes,
            "MEDIUM",
            "tests/",
            "No test files detected.",
            (
                "Add tests for login/authentication, protected routes, token validation, upload validation, "
                "and key project workflows."
            ),
            "Tests improve maintainability and prevent regressions."
        )

    if not facts["ci_files"]:
        _add_fix(
            fixes,
            "LOW",
            ".github/workflows/",
            "No CI/workflow files detected.",
            "Add a basic CI workflow that installs dependencies and runs tests/lint checks on push.",
            "CI catches errors before code is merged or deployed."
        )

    if not facts["docker_files"]:
        _add_fix(
            fixes,
            "LOW",
            "Dockerfile / docker-compose.yml",
            "No Docker/deployment evidence detected.",
            "Add Dockerfile or docker-compose.yml if you want reproducible local/prod deployment.",
            "Reproducible environments make setup and deployment easier."
        )

    duplicate_result = find_duplicate_code(project_name)

    if "Duplicate #" in duplicate_result:
        _add_fix(
            fixes,
            "MEDIUM",
            "duplicated code candidates",
            "Duplicated business logic detected.",
            "Run 'find duplicates in project <project>' and extract repeated logic into shared utilities/services.",
            "Duplicated logic increases maintenance cost and bug risk."
        )

    dead_code_result = find_dead_code(project_name)

    if "Possible unused production" in dead_code_result:
        _add_fix(
            fixes,
            "LOW",
            "dead code candidates",
            "Possible unused production code may exist.",
            "Run 'find dead code in project <project>', manually verify each candidate, then remove only after tests pass.",
            "Removing verified dead code reduces complexity, but heuristic results must be checked manually."
        )

    return _format_fixes_report(
        project.get("name", project_name),
        fixes
    )


def suggest_project_fixes(project_name):
    return suggest_fixes_for_project(project_name)


def project_fixes(project_name):
    return suggest_fixes_for_project(project_name)


def suggest_fixes(project_name):
    return suggest_fixes_for_project(project_name)

# ==========================
# STEP 6 - AI SECURITY AUDITOR
# Rule-based scans. No LLM. No speculation.
# ==========================
SECRET_VALUE_PATTERN = re.compile(
    r'''(?ix)
    (?P<key>
        api[_-]?key|
        secret[_-]?key|
        jwt[_-]?secret|
        token|
        access[_-]?token|
        refresh[_-]?token|
        private[_-]?key|
        client[_-]?secret|
        openai[_-]?api[_-]?key|
        github[_-]?token|
        firebase[_-]?key
    )
    \s*[:=]\s*
    (?P<quote>["\']?)
    (?P<value>[A-Za-z0-9_\-\.\/+=:]{8,})
    (?P=quote)
    '''
)

PASSWORD_VALUE_PATTERN = re.compile(
    r'''(?ix)
    (?P<key>
        password|
        passwd|
        pwd|
        admin_password|
        db_password|
        database_password
    )
    \s*[:=]\s*
    (?P<quote>["\']?)
    (?P<value>[^"'\n\r]{4,})
    (?P=quote)
    '''
)

API_KEY_SIGNATURES = [
    "sk-",
    "xoxb-",
    "ghp_",
    "github_pat_",
    "AIza",
    "AKIA",
    "BEGIN PRIVATE KEY",
    "Bearer ",
]

WEAK_SECRET_VALUES = {
    "secret",
    "changeme",
    "changeme123",
    "password",
    "password123",
    "admin",
    "admin123",
    "test",
    "dev",
    "your-secret-key",
    "secure123",
}

DANGEROUS_IMPORT_PATTERNS = [
    ("pickle", "Python pickle can deserialize unsafe data if input is untrusted."),
    ("subprocess", "subprocess can execute OS commands. Verify input sanitization."),
    ("os.system", "os.system executes shell commands and can cause command injection."),
    ("eval(", "eval executes dynamic code."),
    ("exec(", "exec executes dynamic code."),
    ("yaml.load(", "yaml.load may be unsafe without SafeLoader."),
    ("shell=True", "subprocess shell=True increases command injection risk."),
    ("child_process", "Node child_process can execute OS commands."),
    ("dangerouslySetInnerHTML", "React dangerouslySetInnerHTML can introduce XSS."),
]

SQLI_PATTERNS = [
    r'execute\s*\(\s*f["\']',
    r'execute\s*\(\s*["\'][^"\']*(select|insert|update|delete)[^"\']*["\']\s*\+',
    r'(select|insert|update|delete)[^"\']*\{[^}]+\}',
    r'cursor\.execute\s*\([^,]+%',
    r'\.raw\s*\(\s*f["\']',
]

XSS_PATTERNS = [
    "dangerouslySetInnerHTML",
    ".innerHTML",
    "document.write(",
    "insertAdjacentHTML",
    "v-html",
]


def _iter_security_scan_files(project_name):
    project, error = get_project(project_name)

    if error:
        return None, error

    return get_code_files(project), None


def _is_probably_placeholder(value):
    value_lower = str(value).strip().lower()

    if value_lower in WEAK_SECRET_VALUES:
        return True

    if "example" in value_lower or "placeholder" in value_lower:
        return True

    if value_lower.startswith("your_") or value_lower.startswith("your-"):
        return True

    return False


def _redact_secret(value):
    value = str(value)

    if len(value) <= 6:
        return "***"

    return value[:3] + "***" + value[-3:]


def _line_number(content, index):
    return content[:index].count("\n") + 1


def _format_security_findings(title, project_name, findings, recommendation):
    output = [
        title,
        "Mode: rule-based / no LLM / no automatic changes",
        f"Project: {project_name}",
        "",
        "Findings:"
    ]

    if findings:
        output.extend(f" - {item}" for item in findings)
    else:
        output.append(" - None detected by these rules.")

    output.append("")
    output.append("Recommendation:")
    output.append(recommendation)

    return "\n".join(output)


def find_api_keys(project_name):
    files, error = _iter_security_scan_files(project_name)

    if error:
        return error

    findings = []

    for item in files:
        path = item["path"]
        content = item["content"]

        for match in SECRET_VALUE_PATTERN.finditer(content):
            key = match.group("key")
            value = match.group("value")
            line = _line_number(content, match.start())

            severity = "MEDIUM" if _is_probably_placeholder(value) else "HIGH"

            findings.append(
                f"{severity}: {path}:{line} possible API/secret value {key}={_redact_secret(value)}"
            )

        for signature in API_KEY_SIGNATURES:
            index = content.find(signature)

            if index != -1:
                line = _line_number(content, index)
                findings.append(
                    f"HIGH: {path}:{line} contains API key/token signature '{signature}'"
                )

    return _format_security_findings(
        "API KEY / TOKEN SCAN",
        project_name,
        sorted(set(findings)),
        "Move real keys/tokens to environment variables, rotate exposed credentials, and keep only .env.example placeholders in source control."
    )


def find_passwords(project_name):
    files, error = _iter_security_scan_files(project_name)

    if error:
        return error

    findings = []

    for item in files:
        path = item["path"]
        content = item["content"]

        for match in PASSWORD_VALUE_PATTERN.finditer(content):
            key = match.group("key")
            value = match.group("value").strip()
            line = _line_number(content, match.start())

            if not value:
                continue

            severity = "MEDIUM" if _is_probably_placeholder(value) else "HIGH"

            findings.append(
                f"{severity}: {path}:{line} possible hardcoded password {key}={_redact_secret(value)}"
            )

        lower = content.lower()

        for weak in ["admin123", "password123", "changeme123", "secure123"]:
            index = lower.find(weak)

            if index != -1:
                line = _line_number(content, index)
                findings.append(
                    f"HIGH: {path}:{line} weak/demo password string detected: {weak}"
                )

    return _format_security_findings(
        "PASSWORD SCAN",
        project_name,
        sorted(set(findings)),
        "Remove hardcoded passwords, use environment variables/secrets manager, and ensure test credentials cannot be used in production."
    )


def find_hardcoded_secrets(project_name):
    api_report = find_api_keys(project_name)
    password_report = find_passwords(project_name)

    return (
        "HARDCODED SECRETS SCAN\n"
        "Mode: rule-based / no LLM / no automatic changes\n\n"
        f"{api_report}\n\n"
        f"{password_report}"
    )


def find_sql_injection(project_name):
    files, error = _iter_security_scan_files(project_name)

    if error:
        return error

    findings = []

    for item in files:
        path = item["path"]
        content = item["content"]

        for pattern in SQLI_PATTERNS:
            for match in re.finditer(pattern, content, flags=re.IGNORECASE):
                line = _line_number(content, match.start())
                sample = content[match.start():match.start() + 140].replace("\n", " ")
                findings.append(
                    f"HIGH: {path}:{line} possible dynamic SQL construction -> {sample}"
                )

    return _format_security_findings(
        "SQL INJECTION SCAN",
        project_name,
        sorted(set(findings)),
        "Use parameterized queries/ORM bindings only. Avoid string concatenation, f-strings, percent formatting, or raw SQL with user-controlled input."
    )


def find_xss_risks(project_name):
    files, error = _iter_security_scan_files(project_name)

    if error:
        return error

    findings = []

    for item in files:
        path = item["path"]
        content = item["content"]

        for marker in XSS_PATTERNS:
            index = content.find(marker)

            if index != -1:
                line = _line_number(content, index)
                findings.append(
                    f"HIGH: {path}:{line} possible XSS sink detected: {marker}"
                )

    return _format_security_findings(
        "XSS RISK SCAN",
        project_name,
        sorted(set(findings)),
        "Avoid HTML injection APIs. Sanitize user-controlled HTML with a trusted sanitizer and prefer safe text rendering."
    )


def find_dangerous_imports(project_name):
    files, error = _iter_security_scan_files(project_name)

    if error:
        return error

    findings = []

    for item in files:
        path = item["path"]
        content = item["content"]

        for marker, reason in DANGEROUS_IMPORT_PATTERNS:
            index = content.find(marker)

            if index != -1:
                line = _line_number(content, index)
                findings.append(
                    f"MEDIUM/HIGH: {path}:{line} uses '{marker}'. {reason}"
                )

    return _format_security_findings(
        "DANGEROUS IMPORT / API SCAN",
        project_name,
        sorted(set(findings)),
        "Review every dangerous API usage. Restrict inputs, avoid shell execution, and replace unsafe deserialization/dynamic execution where possible."
    )


def _count_high_medium_low(report):
    lower = report.lower()

    high = lower.count("high:")
    medium = lower.count("medium:")
    medium += lower.count("medium/high:")
    low = lower.count("low:")

    return high, medium, low


def full_security_audit(project_name):
    sections = [
        find_security_issues(project_name),
        strict_security_analyzer_project(project_name),
        find_api_keys(project_name),
        find_passwords(project_name),
        find_hardcoded_secrets(project_name),
        find_sql_injection(project_name),
        find_xss_risks(project_name),
        find_dangerous_imports(project_name),
    ]

    separator = "\n\n" + ("=" * 60) + "\n\n"
    report = separator.join(sections)

    high, medium, low = _count_high_medium_low(report)

    header = [
        "FULL SECURITY AUDIT",
        "Mode: mixed grounded LLM + rule-based scans / no automatic changes",
        f"Project: {project_name}",
        "",
        "Summary:",
        f" - High indicators: {high}",
        f" - Medium indicators: {medium}",
        f" - Low indicators: {low}",
        "",
        "Priority:",
        " - First: secrets, passwords, API keys, JWT signing keys.",
        " - Second: authentication/authorization and upload validation.",
        " - Third: SQL injection, XSS, dangerous imports/execution APIs.",
        " - Fourth: logging privacy and dependency/config hardening.",
        "",
        "=" * 60,
        ""
    ]

    return "\n".join(header) + report


def generate_security_roadmap(project_name):
    audit = full_security_audit(project_name)

    high, medium, low = _count_high_medium_low(audit)

    output = [
        "SECURITY ROADMAP",
        "Mode: generated from audit findings / no automatic changes",
        f"Project: {project_name}",
        "",
        "Risk summary:",
        f" - High indicators: {high}",
        f" - Medium indicators: {medium}",
        f" - Low indicators: {low}",
        "",
        "CRITICAL / DAY 0:",
        " - Rotate any real API keys, tokens, JWT secrets, passwords, or private keys found in source code.",
        " - Move secrets to environment variables or a secrets manager.",
        " - Confirm .env files are ignored and only .env.example placeholders are committed.",
        "",
        "HIGH / DAYS 1-3:",
        " - Fix authentication and authorization gaps first.",
        " - Replace demo users/default credentials with persistent database-backed users.",
        " - Harden JWT configuration: strong secret, expiration, issuer/audience checks where relevant.",
        " - Validate file uploads server-side: size, type, extension, and content scanning.",
        "",
        "MEDIUM / DAYS 4-7:",
        " - Replace unsafe SQL string construction with parameterized queries or ORM bindings.",
        " - Remove or guard dangerous imports and execution APIs such as eval, exec, os.system, subprocess shell=True.",
        " - Avoid unsafe HTML rendering; sanitize user-generated content before display.",
        "",
        "LOW / DAYS 7-14:",
        " - Improve structured audit logging and retention.",
        " - Review privacy impact of IP/username logging.",
        " - Add security tests for auth, upload, JWT, and sensitive routes.",
        "",
        "Evidence source:",
        "The detailed audit below is generated from indexed project files.",
        "",
        audit[:9000]
    ]

    return "\n".join(output)


# Friendly aliases
def scan_api_keys(project_name):
    return find_api_keys(project_name)


def scan_passwords(project_name):
    return find_passwords(project_name)


def scan_hardcoded_secrets(project_name):
    return find_hardcoded_secrets(project_name)


def scan_sql_injection(project_name):
    return find_sql_injection(project_name)


def scan_xss(project_name):
    return find_xss_risks(project_name)


def scan_dangerous_imports(project_name):
    return find_dangerous_imports(project_name)


def enterprise_audit(project_name):
    return full_security_audit(project_name)


def scan_entire_project(project_name):
    return full_security_audit(project_name)

# ==========================
# STEP 8 - AUTONOMOUS CODING ASSISTANT
# Safe autonomous reports. No automatic code changes.
# ==========================
def _autonomous_section(title, content, max_chars=9000):
    content = str(content).strip()

    if len(content) > max_chars:
        content = content[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 70
        + f"\n{title}\n"
        + "=" * 70
        + "\n"
        + content
    )


def _autonomous_project_header(project_name, mode):
    project, error = get_project(project_name)

    if error:
        return None, error

    header = [
        f"JARVIS AUTONOMOUS {mode}",
        "Mode: safe autonomous analysis / no automatic file changes",
        f"Project: {project.get('name')}",
        f"Path: {project.get('path')}",
        f"Files indexed: {project.get('files_count')}",
        f"Tech stack: {', '.join(project.get('tech_stack', []))}",
        "",
        "Safety:",
        " - This report does not modify code.",
        " - Use safe patch / apply patch commands only after reviewing backups and diffs.",
        " - Every recommendation should be verified with tests.",
    ]

    return "\n".join(header), None


def autonomous_review_project(project_name):
    header, error = _autonomous_project_header(
        project_name,
        "PROJECT REVIEW"
    )

    if error:
        return error

    sections = [
        header,
        _autonomous_section(
            "1. PROJECT SCORECARD",
            score_project(project_name)
        ),
        _autonomous_section(
            "2. GROUNDED REVIEW",
            strict_grounded_analyzer_project(project_name)
        ),
        _autonomous_section(
            "3. ARCHITECTURE ANALYSIS",
            strict_architecture_analyzer_project(project_name)
        ),
        _autonomous_section(
            "4. SECURITY ANALYSIS",
            strict_security_analyzer_project(project_name)
        ),
        _autonomous_section(
            "5. DEAD CODE SCAN",
            find_dead_code(project_name)
        ),
        _autonomous_section(
            "6. DUPLICATE CODE SCAN",
            find_duplicate_code(project_name)
        ),
        _autonomous_section(
            "7. IMPROVEMENT ROADMAP",
            generate_improvement_roadmap(project_name)
        ),
    ]

    return "\n".join(sections)


def autonomous_secure_project(project_name):
    header, error = _autonomous_project_header(
        project_name,
        "SECURITY HARDENING PLAN"
    )

    if error:
        return error

    sections = [
        header,
        _autonomous_section(
            "1. FULL SECURITY AUDIT",
            full_security_audit(project_name)
        ),
        _autonomous_section(
            "2. API KEY / TOKEN SCAN",
            find_api_keys(project_name)
        ),
        _autonomous_section(
            "3. PASSWORD SCAN",
            find_passwords(project_name)
        ),
        _autonomous_section(
            "4. HARDCODED SECRETS SCAN",
            find_hardcoded_secrets(project_name)
        ),
        _autonomous_section(
            "5. SQL INJECTION SCAN",
            find_sql_injection(project_name)
        ),
        _autonomous_section(
            "6. XSS SCAN",
            find_xss_risks(project_name)
        ),
        _autonomous_section(
            "7. DANGEROUS IMPORTS / EXECUTION SCAN",
            find_dangerous_imports(project_name)
        ),
        _autonomous_section(
            "8. SECURITY ROADMAP",
            generate_security_roadmap(project_name)
        ),
    ]

    return "\n".join(sections)


def autonomous_optimize_project(project_name):
    header, error = _autonomous_project_header(
        project_name,
        "OPTIMIZATION PLAN"
    )

    if error:
        return error

    sections = [
        header,
        _autonomous_section(
            "1. PROJECT SCORECARD",
            score_project(project_name)
        ),
        _autonomous_section(
            "2. OPTIMIZATION REPORT",
            optimize_project(project_name)
        ),
        _autonomous_section(
            "3. DUPLICATE LOGIC",
            find_duplicate_code(project_name)
        ),
        _autonomous_section(
            "4. DEAD CODE",
            find_dead_code(project_name)
        ),
        _autonomous_section(
            "5. ARCHITECTURE RISKS",
            strict_architecture_analyzer_project(project_name)
        ),
        _autonomous_section(
            "6. IMPROVEMENT ROADMAP",
            generate_improvement_roadmap(project_name)
        ),
    ]

    return "\n".join(sections)


def autonomous_fix_project(project_name):
    header, error = _autonomous_project_header(
        project_name,
        "FIX PLAN"
    )

    if error:
        return error

    sections = [
        header,
        _autonomous_section(
            "1. TOP PROJECT FIXES",
            suggest_fixes_for_project(project_name)
        ),
        _autonomous_section(
            "2. SECURITY FIXES",
            generate_security_roadmap(project_name)
        ),
        _autonomous_section(
            "3. DEAD CODE / CLEANUP FIXES",
            find_dead_code(project_name)
        ),
        _autonomous_section(
            "4. DUPLICATE CODE / REFACTORING FIXES",
            find_duplicate_code(project_name)
        ),
        _autonomous_section(
            "5. ARCHITECTURE FIXES",
            strict_architecture_analyzer_project(project_name)
        ),
        _autonomous_section(
            "6. SAFE EXECUTION PLAN",
            (
                "Recommended order:\n"
                "1. Refresh deep project memory.\n"
                "2. Backup the target files before any patch.\n"
                "3. Fix secrets/passwords/API keys first.\n"
                "4. Fix authentication/authorization and upload validation.\n"
                "5. Fix SQL injection/XSS/dangerous execution APIs.\n"
                "6. Remove duplicated/dead code only after tests pass.\n"
                "7. Run tests and re-run full security audit."
            )
        ),
    ]

    return "\n".join(sections)


def autonomous_improve_project(project_name):
    header, error = _autonomous_project_header(
        project_name,
        "IMPROVEMENT PLAN"
    )

    if error:
        return error

    sections = [
        header,
        _autonomous_section(
            "1. REVIEW EVERYTHING",
            autonomous_review_project(project_name),
            max_chars=12000
        ),
        _autonomous_section(
            "2. SECURITY HARDENING",
            autonomous_secure_project(project_name),
            max_chars=12000
        ),
        _autonomous_section(
            "3. OPTIMIZATION",
            autonomous_optimize_project(project_name),
            max_chars=12000
        ),
        _autonomous_section(
            "4. FIX PLAN",
            autonomous_fix_project(project_name),
            max_chars=12000
        ),
        _autonomous_section(
            "5. FINAL AUTONOMOUS RECOMMENDATION",
            (
                "Do not apply every change at once.\n"
                "Start with the highest-risk security findings, then architecture cleanup, "
                "then maintainability and optimization.\n"
                "After every patch: run tests, restart the app, and re-run JARVIS audit."
            )
        ),
    ]

    return "\n".join(sections)


def review_everything(project_name="CyberShield AI"):
    return autonomous_review_project(project_name)


def fix_project(project_name):
    return autonomous_fix_project(project_name)


def secure_project(project_name):
    return autonomous_secure_project(project_name)


def improve_project(project_name):
    return autonomous_improve_project(project_name)


def optimize_project_autonomous(project_name):
    return autonomous_optimize_project(project_name)



# ==========================
# STEP 10 - MEMORY AWARE REVIEWS
# ==========================
try:
    from deep_project_memory import (
        project_timeline,
        vulnerability_history,
        remembered_fixes,
        project_evolution,
        session_summary
    )
except Exception:
    pass


def memory_aware_review(project_name):
    review = review_project(project_name)

    try:
        timeline = project_timeline(project_name)
    except:
        timeline = "No timeline available."

    try:
        vulns = vulnerability_history(project_name)
    except:
        vulns = "No vulnerability history."

    try:
        fixes = remembered_fixes(project_name)
    except:
        fixes = "No remembered fixes."

    return (
        "MEMORY-AWARE PROJECT REVIEW\n\n"
        + review
        + "\n\n=== PROJECT TIMELINE ===\n"
        + timeline
        + "\n\n=== VULNERABILITY HISTORY ===\n"
        + vulns
        + "\n\n=== REMEMBERED FIXES ===\n"
        + fixes
    )


def memory_aware_security_review(project_name):
    report = find_security_issues(project_name)

    try:
        history = vulnerability_history(project_name)
    except:
        history = "No vulnerability history."

    return (
        "MEMORY-AWARE SECURITY REVIEW\n\n"
        + report
        + "\n\n=== HISTORICAL SECURITY FINDINGS ===\n"
        + history
    )


def project_evolution_report(project_name):
    try:
        return project_evolution(project_name)
    except Exception as e:
        return f"Evolution report error: {e}"


def engineering_session_summary():
    try:
        return session_summary()
    except Exception as e:
        return f"Session summary error: {e}"

# ==========================
# STEP 11 - AUTONOMOUS ENGINEERING PLANNER
# Senior Architect + Security Lead + Technical PM
# Safe planning only. No automatic code changes.
# ==========================
def _planner_header(project_name, mode):
    project, error = get_project(project_name)

    if error:
        return None, None, error

    card, card_error = project_scorecard(project_name)

    if card_error:
        card = None

    header = [
        f"JARVIS {mode}",
        "Mode: autonomous engineering planning / no automatic code changes",
        f"Project: {project.get('name')}",
        f"Path: {project.get('path')}",
        f"Files indexed: {project.get('files_count')}",
        f"Tech stack: {', '.join(project.get('tech_stack', []))}",
        "",
        "Safety rules:",
        " - This output is a plan, not an automatic patch.",
        " - Create backups before applying any fix.",
        " - Apply changes one file at a time.",
        " - Run tests after every change.",
    ]

    if card:
        header.extend([
            "",
            "Current scorecard:",
            f" - Security: {card['scores']['security']}/10",
            f" - Architecture: {card['scores']['architecture']}/10",
            f" - Maintainability: {card['scores']['maintainability']}/10",
            f" - Overall: {card['scores']['overall']}/10",
        ])

    return project, "\n".join(header), None


def _planner_section(title, content, max_chars=7000):
    content = str(content).strip()

    if len(content) > max_chars:
        content = content[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 70
        + f"\n{title}\n"
        + "=" * 70
        + "\n"
        + content
    )


def _extract_action_items_from_text(text, limit=20):
    lines = str(text).splitlines()
    actions = []

    keywords = [
        "fix",
        "replace",
        "move",
        "remove",
        "add",
        "validate",
        "sanitize",
        "harden",
        "refactor",
        "test",
        "backup",
        "rotate",
        "review",
        "centralize",
        "parameterized",
        "environment",
    ]

    for line in lines:
        clean = line.strip(" -\t")

        if not clean:
            continue

        lower = clean.lower()

        if any(keyword in lower for keyword in keywords):
            actions.append(clean)

        if len(actions) >= limit:
            break

    if not actions:
        return ["Run a full audit, inspect the highest-risk files, create backups, then apply fixes one by one."]

    return actions


def generate_project_roadmap(project_name):
    project, header, error = _planner_header(
        project_name,
        "PROJECT ROADMAP"
    )

    if error:
        return error

    sections = [
        header,
        _planner_section(
            "1. CURRENT PROJECT SCORE",
            score_project(project_name)
        ),
        _planner_section(
            "2. MEMORY-AWARE PROJECT CONTEXT",
            memory_aware_review(project_name)
            if "memory_aware_review" in globals()
            else autonomous_review_project(project_name),
            max_chars=10000
        ),
        _planner_section(
            "3. SECURITY ROADMAP",
            generate_security_roadmap(project_name)
        ),
        _planner_section(
            "4. IMPROVEMENT ROADMAP",
            generate_improvement_roadmap(project_name)
        ),
        _planner_section(
            "5. OPTIMIZATION PLAN",
            autonomous_optimize_project(project_name)
        ),
        _planner_section(
            "6. 30-DAY EXECUTION PLAN",
            (
                "Week 1:\n"
                " - Fix secrets/passwords/API keys and weak configuration.\n"
                " - Fix authentication, authorization, JWT, and upload validation issues.\n"
                " - Add basic tests for the highest-risk routes/files.\n\n"
                "Week 2:\n"
                " - Refactor duplicated logic and remove dead code only after tests pass.\n"
                " - Centralize configuration and API URLs.\n"
                " - Improve logging structure and privacy retention rules.\n\n"
                "Week 3:\n"
                " - Improve architecture boundaries and module separation.\n"
                " - Add documentation for setup, security, and deployment.\n"
                " - Add CI checks where missing.\n\n"
                "Week 4:\n"
                " - Run full security audit again.\n"
                " - Export project report.\n"
                " - Prepare release/deployment checklist."
            )
        )
    ]

    return "\n".join(sections)


def next_best_improvements(project_name):
    project, header, error = _planner_header(
        project_name,
        "NEXT BEST IMPROVEMENTS"
    )

    if error:
        return error

    security = full_security_audit(project_name)
    fixes = autonomous_fix_project(project_name)
    optimization = autonomous_optimize_project(project_name)

    actions = []
    actions.extend(_extract_action_items_from_text(security, limit=8))
    actions.extend(_extract_action_items_from_text(fixes, limit=8))
    actions.extend(_extract_action_items_from_text(optimization, limit=8))

    seen = set()
    unique = []

    for action in actions:
        key = action.lower()

        if key not in seen:
            seen.add(key)
            unique.append(action)

    output = [
        header,
        "",
        "Recommended next improvements, ordered safely:",
        ""
    ]

    priority_labels = [
        "Critical security",
        "High-risk auth/config",
        "Input validation",
        "Secrets and environment",
        "Tests",
        "Architecture cleanup",
        "Maintainability",
        "Performance",
        "Documentation",
        "Deployment readiness",
    ]

    for index, action in enumerate(unique[:10], start=1):
        label = priority_labels[index - 1] if index <= len(priority_labels) else "Improvement"

        output.append(
            f"{index}. [{label}] {action}"
        )

    output.append("")
    output.append(
        "Recommendation: apply only the first 1-3 items first, then run tests and re-audit."
    )

    return "\n".join(output)


def what_should_i_fix_next(project_name):
    return next_best_improvements(project_name)


def highest_risk_vulnerabilities(project_name):
    project, header, error = _planner_header(
        project_name,
        "HIGHEST RISK VULNERABILITIES"
    )

    if error:
        return error

    audit = full_security_audit(project_name)
    lines = audit.splitlines()

    high_lines = []

    for line in lines:
        lower = line.lower()

        if (
            "high:" in lower
            or "medium/high:" in lower
            or "secret" in lower
            or "password" in lower
            or "api key" in lower
            or "sql injection" in lower
            or "xss" in lower
            or "jwt" in lower
            or "auth" in lower
        ):
            high_lines.append(line)

    if not high_lines:
        high_lines = [
            "No high-risk vulnerabilities were detected by the current rule-based scan.",
            "Still run manual checks for secrets, auth, JWT, uploads, SQL injection, and XSS."
        ]

    output = [
        header,
        "",
        "Highest-risk findings:",
        ""
    ]

    output.extend(
        f"- {line.strip()}"
        for line in high_lines[:40]
        if line.strip()
    )

    output.append("")
    output.append("Fix order:")
    output.append("1. Secrets/API keys/passwords.")
    output.append("2. Authentication and authorization.")
    output.append("3. JWT/session handling.")
    output.append("4. Upload validation and input validation.")
    output.append("5. SQL injection/XSS/dangerous execution APIs.")

    return "\n".join(output)


def estimate_project_maturity(project_name):
    card, error = project_scorecard(project_name)

    if error:
        return error

    scores = card["scores"]
    overall = scores["overall"]

    if overall >= 8:
        level = "Advanced"
    elif overall >= 6:
        level = "Intermediate"
    elif overall >= 4:
        level = "Early / prototype"
    else:
        level = "Needs foundation work"

    output = [
        "PROJECT MATURITY ESTIMATE",
        "Mode: rule-based / no LLM / no speculation",
        "",
        format_project_scorecard(card),
        "",
        f"Maturity level: {level}",
        "",
        "Interpretation:",
        " - Security score shows how much security evidence exists in indexed files.",
        " - Architecture score shows visible structure, routes, UI/backend split, DB/config, deployment evidence.",
        " - Maintainability score shows tests, CI, documentation/config, project organization.",
        "",
        "Next step:",
        "Run: next best improvements <project>"
    ]

    return "\n".join(output)


def estimate_production_readiness(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    blockers = []
    warnings = []
    positives = []

    if not facts["test_files"]:
        blockers.append("No test files detected.")

    if not facts["ci_files"]:
        warnings.append("No CI/workflow files detected.")

    if not facts["docker_files"]:
        warnings.append("No Docker/deployment files detected.")

    if facts["auth_files"]:
        positives.append("Authentication/authorization files detected.")
    else:
        blockers.append("No clear authentication/authorization files detected.")

    if facts["jwt_files"]:
        positives.append("JWT/token handling evidence detected.")

    if facts["password_hashing_files"]:
        positives.append("Password hashing evidence detected.")
    else:
        warnings.append("No clear password hashing evidence detected.")

    security_report = full_security_audit(project_name).lower()

    if "high:" in security_report or "medium/high:" in security_report:
        blockers.append("High or medium/high security indicators detected by audit.")

    if not blockers and len(warnings) <= 1:
        readiness = "High"
    elif len(blockers) <= 1:
        readiness = "Medium"
    else:
        readiness = "Low"

    output = [
        "PRODUCTION READINESS ESTIMATE",
        "Mode: rule-based / no automatic changes",
        "",
        f"Project: {facts['project_name']}",
        f"Readiness: {readiness}",
        "",
        "Blockers:"
    ]

    output.extend(f" - {item}" for item in blockers) if blockers else output.append(" - None detected by rules.")

    output.append("")
    output.append("Warnings:")
    output.extend(f" - {item}" for item in warnings) if warnings else output.append(" - None detected by rules.")

    output.append("")
    output.append("Positive evidence:")
    output.extend(f" - {item}" for item in positives) if positives else output.append(" - None detected by rules.")

    output.append("")
    output.append("Before production:")
    output.append(" - Run tests.")
    output.append(" - Re-run full security audit.")
    output.append(" - Export project report.")
    output.append(" - Review secrets, logs, uploads, JWT, dependencies, and deployment config.")

    return "\n".join(output)


def generate_sprint_plan(project_name):
    project, header, error = _planner_header(
        project_name,
        "SPRINT PLAN"
    )

    if error:
        return error

    improvements = next_best_improvements(project_name)
    risks = highest_risk_vulnerabilities(project_name)

    return (
        header
        + _planner_section(
            "SPRINT GOAL",
            "Stabilize security, improve architecture, add test coverage, and prepare for release readiness."
        )
        + _planner_section(
            "BACKLOG INPUT",
            improvements,
            max_chars=8000
        )
        + _planner_section(
            "RISK INPUT",
            risks,
            max_chars=8000
        )
        + _planner_section(
            "7-DAY SPRINT PLAN",
            (
                "Day 1: Triage all high-risk security findings and create backups.\n"
                "Day 2: Fix secrets/config/auth/JWT risks.\n"
                "Day 3: Fix input validation, upload validation, SQL injection, XSS, dangerous APIs.\n"
                "Day 4: Add or improve tests for changed files.\n"
                "Day 5: Refactor duplicated/dead code after tests pass.\n"
                "Day 6: Improve documentation and deployment checklist.\n"
                "Day 7: Run full audit, export report, and decide release readiness."
            )
        )
    )


def generate_release_checklist(project_name):
    project, header, error = _planner_header(
        project_name,
        "RELEASE CHECKLIST"
    )

    if error:
        return error

    return (
        header
        + _planner_section(
            "RELEASE CHECKLIST",
            (
                "[ ] Project opens and runs locally.\n"
                "[ ] Dependencies install cleanly.\n"
                "[ ] Tests pass.\n"
                "[ ] Secrets are not committed.\n"
                "[ ] .env.example exists and .env is ignored.\n"
                "[ ] Authentication and authorization are tested.\n"
                "[ ] JWT/session expiration is configured.\n"
                "[ ] Upload validation is server-side, not only frontend.\n"
                "[ ] SQL injection and XSS risky patterns are checked.\n"
                "[ ] Dangerous imports/execution APIs are reviewed.\n"
                "[ ] Logs do not expose unnecessary personal/sensitive data.\n"
                "[ ] Error handling is user-safe and developer-useful.\n"
                "[ ] Build command works.\n"
                "[ ] Deployment configuration is documented.\n"
                "[ ] Full security audit re-run after final changes.\n"
                "[ ] Project report exported."
            )
        )
    )


def generate_deployment_checklist(project_name):
    project, header, error = _planner_header(
        project_name,
        "DEPLOYMENT CHECKLIST"
    )

    if error:
        return error

    return (
        header
        + _planner_section(
            "DEPLOYMENT CHECKLIST",
            (
                "[ ] Confirm production environment variables.\n"
                "[ ] Confirm database path/URL is production-safe.\n"
                "[ ] Confirm debug mode is disabled.\n"
                "[ ] Confirm CORS policy is restricted.\n"
                "[ ] Confirm HTTPS/reverse proxy settings.\n"
                "[ ] Confirm logs path and retention.\n"
                "[ ] Confirm backup/restore process.\n"
                "[ ] Confirm file upload size/type/content validation.\n"
                "[ ] Confirm dependency versions and vulnerability scan.\n"
                "[ ] Confirm frontend API base URL configuration.\n"
                "[ ] Confirm build artifacts are generated correctly.\n"
                "[ ] Confirm restart/recovery process.\n"
                "[ ] Confirm monitoring/health checks.\n"
                "[ ] Run smoke test after deployment."
            )
        )
    )


def become_project_architect(project_name):
    project, header, error = _planner_header(
        project_name,
        "PROJECT ARCHITECT MODE"
    )

    if error:
        return error

    return (
        header
        + _planner_section(
            "ARCHITECTURE VIEW",
            strict_architecture_analyzer_project(project_name)
        )
        + _planner_section(
            "SECURITY VIEW",
            highest_risk_vulnerabilities(project_name)
        )
        + _planner_section(
            "MATURITY VIEW",
            estimate_project_maturity(project_name)
        )
        + _planner_section(
            "PRODUCTION READINESS",
            estimate_production_readiness(project_name)
        )
        + _planner_section(
            "ARCHITECT DECISION",
            (
                "Recommended operating mode:\n"
                "1. Treat security as the first milestone.\n"
                "2. Stabilize architecture boundaries second.\n"
                "3. Add tests before large refactors.\n"
                "4. Release only after full audit and deployment checklist pass."
            )
        )
    )


# Friendly aliases for command routing
def generate_roadmap(project_name):
    return generate_project_roadmap(project_name)


def project_roadmap(project_name):
    return generate_project_roadmap(project_name)


def next_improvements(project_name):
    return next_best_improvements(project_name)


def fix_next(project_name):
    return what_should_i_fix_next(project_name)


def high_risk_vulnerabilities(project_name):
    return highest_risk_vulnerabilities(project_name)


def production_readiness(project_name):
    return estimate_production_readiness(project_name)


def project_maturity(project_name):
    return estimate_project_maturity(project_name)


def sprint_plan(project_name):
    return generate_sprint_plan(project_name)


def release_checklist(project_name):
    return generate_release_checklist(project_name)


def deployment_checklist(project_name):
    return generate_deployment_checklist(project_name)


def project_architect(project_name):
    return become_project_architect(project_name)

# ==========================
# STEP 16 - MULTI PROJECT INTELLIGENCE
# Portfolio intelligence / cross-project learning / roadmap
# Rule-based where possible. No automatic code changes.
# ==========================
def _portfolio_cards():
    projects = _load_all_deep_projects()

    cards = []

    for project in projects:
        name = project.get("name", "")

        if not name:
            continue

        card, error = project_scorecard(name)

        if not error and card:
            cards.append(card)

    cards.sort(
        key=lambda card: card["scores"]["overall"],
        reverse=True
    )

    return cards


def _portfolio_line(card):
    scores = card["scores"]

    return (
        f"{card['name']} | "
        f"overall {scores['overall']}/10 | "
        f"security {scores['security']}/10 | "
        f"architecture {scores['architecture']}/10 | "
        f"maintainability {scores['maintainability']}/10 | "
        f"files {card['files_count']} | "
        f"stack: {', '.join(card['tech_stack'])}"
    )


def compare_all_projects():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    output = [
        "MULTI PROJECT COMPARISON",
        "Mode: rule-based portfolio intelligence / no speculation",
        "",
        "Projects ranked by overall score:",
        ""
    ]

    for index, card in enumerate(cards, start=1):
        output.append(f"{index}. {_portfolio_line(card)}")

    output.append("")
    output.append("Best by category:")

    categories = [
        "overall",
        "security",
        "architecture",
        "maintainability"
    ]

    for category in categories:
        best = max(
            cards,
            key=lambda card: card["scores"][category]
        )

        output.append(
            f" - {category.capitalize()}: "
            f"{best['name']} ({best['scores'][category]}/10)"
        )

    output.append("")
    output.append("Weakest by category:")

    for category in categories:
        weakest = min(
            cards,
            key=lambda card: card["scores"][category]
        )

        output.append(
            f" - {category.capitalize()}: "
            f"{weakest['name']} ({weakest['scores'][category]}/10)"
        )

    return "\n".join(output)


def find_best_project():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    best = cards[0]

    return (
        "BEST PROJECT\n"
        "Mode: rule-based / no speculation\n\n"
        + format_project_scorecard(best)
    )


def find_weakest_project():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    weakest = min(
        cards,
        key=lambda card: card["scores"]["overall"]
    )

    return (
        "WEAKEST PROJECT\n"
        "Mode: rule-based / no speculation\n\n"
        + format_project_scorecard(weakest)
        + "\n\nRecommended next command:\n"
        + f"what should i fix next {weakest['name']}"
    )


def most_secure_project():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    best = max(
        cards,
        key=lambda card: card["scores"]["security"]
    )

    return (
        "MOST SECURE PROJECT\n"
        "Mode: rule-based / no speculation\n\n"
        + format_project_scorecard(best)
    )


def weakest_security_project():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    weakest = min(
        cards,
        key=lambda card: card["scores"]["security"]
    )

    return (
        "PROJECT WITH WEAKEST SECURITY\n"
        "Mode: rule-based / no speculation\n\n"
        + format_project_scorecard(weakest)
        + "\n\nRecommended next command:\n"
        + f"secure workflow {weakest['name']}"
    )


def most_production_ready_project():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    best = None
    best_score = -1
    best_text = ""

    for card in cards:
        text = estimate_production_readiness(card["name"])
        lower = text.lower()

        score = card["scores"]["overall"]

        if "readiness: high" in lower:
            score += 3
        elif "readiness: medium" in lower:
            score += 1
        elif "readiness: low" in lower:
            score -= 2

        if score > best_score:
            best_score = score
            best = card
            best_text = text

    if not best:
        return "No production readiness result could be generated."

    return (
        "MOST PRODUCTION-READY PROJECT\n"
        "Mode: rule-based estimate / verify manually before release\n\n"
        f"Selected project: {best['name']}\n\n"
        f"{best_text}"
    )


def project_needing_most_work():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    weakest = min(
        cards,
        key=lambda card: (
            card["scores"]["overall"],
            card["scores"]["security"],
            card["scores"]["maintainability"]
        )
    )

    return (
        "PROJECT NEEDING THE MOST WORK\n"
        "Mode: rule-based / no speculation\n\n"
        + format_project_scorecard(weakest)
        + "\n\nRecommended improvement path:\n"
        + "1. Run full security audit.\n"
        + "2. Fix highest-risk vulnerabilities.\n"
        + "3. Add or improve tests.\n"
        + "4. Generate sprint plan.\n"
        + "5. Re-score project.\n\n"
        + f"Suggested command: workflow project {weakest['name']}"
    )


def portfolio_summary():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    total = len(cards)

    avg_security = round(
        sum(card["scores"]["security"] for card in cards) / total,
        1
    )

    avg_architecture = round(
        sum(card["scores"]["architecture"] for card in cards) / total,
        1
    )

    avg_maintainability = round(
        sum(card["scores"]["maintainability"] for card in cards) / total,
        1
    )

    avg_overall = round(
        sum(card["scores"]["overall"] for card in cards) / total,
        1
    )

    stacks = Counter()

    for card in cards:
        for tech in card.get("tech_stack", []):
            if tech:
                stacks[tech] += 1

    output = [
        "PORTFOLIO SUMMARY",
        "Mode: rule-based multi-project intelligence",
        "",
        f"Projects scored: {total}",
        f"Average overall: {avg_overall}/10",
        f"Average security: {avg_security}/10",
        f"Average architecture: {avg_architecture}/10",
        f"Average maintainability: {avg_maintainability}/10",
        "",
        "Tech stack distribution:"
    ]

    if stacks:
        for tech, count in stacks.most_common(20):
            output.append(f" - {tech}: {count}")
    else:
        output.append(" - No tech stack data found.")

    output.append("")
    output.append("Project list:")

    for card in cards:
        output.append(f" - {_portfolio_line(card)}")

    return "\n".join(output)


def portfolio_roadmap():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    security_sorted = sorted(
        cards,
        key=lambda card: card["scores"]["security"]
    )

    maintainability_sorted = sorted(
        cards,
        key=lambda card: card["scores"]["maintainability"]
    )

    architecture_sorted = sorted(
        cards,
        key=lambda card: card["scores"]["architecture"]
    )

    output = [
        "PORTFOLIO ROADMAP",
        "Mode: rule-based strategic roadmap / no automatic changes",
        "",
        "Priority 1 - Security hardening:"
    ]

    for card in security_sorted[:5]:
        output.append(
            f" - {card['name']} "
            f"(security {card['scores']['security']}/10) -> "
            f"run: secure workflow {card['name']}"
        )

    output.append("")
    output.append("Priority 2 - Maintainability and tests:")

    for card in maintainability_sorted[:5]:
        output.append(
            f" - {card['name']} "
            f"(maintainability {card['scores']['maintainability']}/10) -> "
            f"run: what should i fix next {card['name']}"
        )

    output.append("")
    output.append("Priority 3 - Architecture cleanup:")

    for card in architecture_sorted[:5]:
        output.append(
            f" - {card['name']} "
            f"(architecture {card['scores']['architecture']}/10) -> "
            f"run: architect project {card['name']}"
        )

    output.append("")
    output.append("30-day portfolio plan:")
    output.append("Week 1: Fix weakest security project and re-score.")
    output.append("Week 2: Improve tests/maintainability in the two weakest projects.")
    output.append("Week 3: Generate architecture reports and sprint plans.")
    output.append("Week 4: Prepare the strongest project for production/release.")

    return "\n".join(output)


def portfolio_security_summary():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    output = [
        "PORTFOLIO SECURITY SUMMARY",
        "Mode: rule-based / no speculation",
        "",
        "Projects ranked by security:"
    ]

    cards.sort(
        key=lambda card: card["scores"]["security"],
        reverse=True
    )

    for index, card in enumerate(cards, start=1):
        output.append(
            f"{index}. {card['name']} -> "
            f"{card['scores']['security']}/10"
        )

        for reason in card["reasons"]["security"][:4]:
            output.append(f"   - {reason}")

    output.append("")
    output.append("Recommended command for weakest security:")
    weakest = cards[-1]
    output.append(f"secure workflow {weakest['name']}")

    return "\n".join(output)


def portfolio_production_readiness():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    output = [
        "PORTFOLIO PRODUCTION READINESS",
        "Mode: rule-based estimate / verify manually",
        ""
    ]

    for card in cards:
        readiness = estimate_production_readiness(card["name"])

        readiness_line = "Readiness: Not detected"

        for line in readiness.splitlines():
            if line.lower().startswith("readiness:"):
                readiness_line = line.strip()
                break

        output.append(
            f"- {card['name']} | "
            f"overall {card['scores']['overall']}/10 | "
            f"{readiness_line}"
        )

    output.append("")
    output.append("Suggested release candidate:")
    output.append(most_production_ready_project().splitlines()[0])

    return "\n".join(output)


def portfolio_next_improvements():
    cards = _portfolio_cards()

    if not cards:
        return "No remembered projects could be scored."

    output = [
        "PORTFOLIO NEXT IMPROVEMENTS",
        "Mode: rule-based / no automatic code changes",
        "",
        "Top actions across all remembered projects:"
    ]

    weakest_security = min(
        cards,
        key=lambda card: card["scores"]["security"]
    )

    weakest_maintainability = min(
        cards,
        key=lambda card: card["scores"]["maintainability"]
    )

    weakest_architecture = min(
        cards,
        key=lambda card: card["scores"]["architecture"]
    )

    best_overall = max(
        cards,
        key=lambda card: card["scores"]["overall"]
    )

    output.extend([
        f"1. Fix security first in {weakest_security['name']} "
        f"(security {weakest_security['scores']['security']}/10).",
        f"2. Improve tests/maintainability in {weakest_maintainability['name']} "
        f"(maintainability {weakest_maintainability['scores']['maintainability']}/10).",
        f"3. Review architecture in {weakest_architecture['name']} "
        f"(architecture {weakest_architecture['scores']['architecture']}/10).",
        f"4. Prepare strongest project for release: {best_overall['name']} "
        f"(overall {best_overall['scores']['overall']}/10).",
        "",
        "Suggested commands:",
        f" - secure workflow {weakest_security['name']}",
        f" - what should i fix next {weakest_maintainability['name']}",
        f" - architect project {weakest_architecture['name']}",
        f" - production workflow {best_overall['name']}",
    ])

    return "\n".join(output)


# Friendly aliases
def compare_portfolio():
    return compare_all_projects()


def show_portfolio_summary():
    return portfolio_summary()


def show_portfolio_roadmap():
    return portfolio_roadmap()


def show_portfolio_security():
    return portfolio_security_summary()


def show_portfolio_readiness():
    return portfolio_production_readiness()


def what_should_i_improve_across_all_projects():
    return portfolio_next_improvements()


def which_project_is_production_ready():
    return most_production_ready_project()


def which_project_is_most_secure():
    return most_secure_project()


def which_project_needs_most_work():
    return project_needing_most_work()

# ==========================
# STEP 18 - AUTONOMOUS REFACTORING PLANNER
# Technical debt / oversized files / modernization / migration plan.
# Safe planning only. No automatic code changes.
# ==========================
def _refactor_section(title, content, max_chars=9000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 76
        + f"\n{title}\n"
        + "=" * 76
        + "\n"
        + text
    )


def _file_line_count(content):
    return len(str(content).splitlines())


def _file_complexity_score(item):
    content = item.get("content", "")
    path = item.get("path", "")

    markers = [
        "if ",
        "elif ",
        "else:",
        "for ",
        "while ",
        "try:",
        "except ",
        "catch ",
        "switch ",
        "case ",
        "async ",
        "await ",
        "useEffect",
        "useState",
        "return ",
        "raise ",
        "throw ",
    ]

    score = 0
    lower = content.lower()

    for marker in markers:
        score += lower.count(marker.lower())

    score += _file_line_count(content) // 50

    if is_entrypoint_or_framework_file(path):
        score += 8

    if "auth" in path.lower() or "security" in path.lower() or "admin" in path.lower():
        score += 10

    return score


def find_oversized_files(project_name, limit=25):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    scored = []

    for item in files:
        lines = _file_line_count(item.get("content", ""))
        chars = len(item.get("content", ""))

        if lines >= 120 or chars >= 6000:
            scored.append(
                (
                    lines,
                    chars,
                    item
                )
            )

    scored.sort(
        key=lambda row: (
            row[0],
            row[1]
        ),
        reverse=True
    )

    output = [
        "OVERSIZED FILES REPORT",
        "Mode: rule-based / no automatic changes",
        f"Project: {project.get('name')}",
        "",
        "Files that may need splitting:"
    ]

    if not scored:
        output.append("No oversized files detected by current thresholds.")
        return "\n".join(output)

    for lines, chars, item in scored[:limit]:
        output.append(
            f"- {item['path']} | "
            f"{lines} lines | "
            f"{chars} chars | "
            f"risk: {risk_label_for_file(item['path'])}"
        )

    output.append("")
    output.append("Recommended split strategy:")
    output.append("- Move routes/controllers into route modules.")
    output.append("- Move business logic into services.")
    output.append("- Move validation into schemas/validators.")
    output.append("- Move config/constants into config modules.")
    output.append("- Add tests before splitting high-risk files.")

    return "\n".join(output)


def find_most_problematic_files(project_name, limit=20):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    scored = []

    for item in files:
        score = _file_complexity_score(item)
        lines = _file_line_count(item.get("content", ""))

        if score > 0:
            scored.append(
                (
                    score,
                    lines,
                    item
                )
            )

    scored.sort(
        key=lambda row: (
            row[0],
            row[1]
        ),
        reverse=True
    )

    output = [
        "MOST PROBLEMATIC FILES",
        "Mode: heuristic complexity scan / verify manually",
        f"Project: {project.get('name')}",
        "",
        "Files to inspect first:"
    ]

    if not scored:
        output.append("No problematic files detected by current heuristic.")
        return "\n".join(output)

    for score, lines, item in scored[:limit]:
        output.append(
            f"- {item['path']} | "
            f"complexity score {score} | "
            f"{lines} lines | "
            f"{risk_label_for_file(item['path'])}"
        )

    output.append("")
    output.append("Recommendation:")
    output.append("Start with the highest complexity + highest security relevance files.")

    return "\n".join(output)


def analyze_technical_debt(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    card, card_error = project_scorecard(project_name)

    files = get_code_files(project)

    missing_tests = []
    config_risks = []
    large_files = []
    complex_files = []

    for item in files:
        path = normalize_path(item["path"])
        content = item.get("content", "").lower()
        lines = _file_line_count(item.get("content", ""))

        if lines >= 120:
            large_files.append(item["path"])

        if _file_complexity_score(item) >= 35:
            complex_files.append(item["path"])

        if "localhost" in content or "127.0.0.1" in content:
            config_risks.append(item["path"])

    try:
        _, facts, facts_error = strict_project_facts(project_name)
    except Exception:
        facts = None
        facts_error = "Could not load strict facts."

    if facts and not facts_error:
        if not facts["test_files"]:
            missing_tests.append("No test files detected.")
        if not facts["ci_files"]:
            missing_tests.append("No CI/workflow files detected.")
        if facts["fastapi_files"] and facts["flask_files"]:
            config_risks.append("Mixed Flask/FastAPI evidence detected.")

    output = [
        "TECHNICAL DEBT ANALYSIS",
        "Mode: rule-based / safe planning only",
        f"Project: {project.get('name')}",
        f"Path: {project.get('path')}",
        ""
    ]

    if not card_error and card:
        output.append("Scorecard:")
        output.append(
            f"- Security: {card['scores']['security']}/10\n"
            f"- Architecture: {card['scores']['architecture']}/10\n"
            f"- Maintainability: {card['scores']['maintainability']}/10\n"
            f"- Overall: {card['scores']['overall']}/10"
        )

    output.append("\nDebt indicators:")

    indicators = [
        ("Large files", large_files[:15]),
        ("Complex files", complex_files[:15]),
        ("Config/local URL risks", sorted(set(config_risks))[:20]),
        ("Missing tests/CI evidence", missing_tests),
    ]

    for title, items in indicators:
        output.append(f"\n{title}:")
        if items:
            for item in items:
                output.append(f"- {item}")
        else:
            output.append("- None detected by rules.")

    output.append("")
    output.append("Recommended debt reduction order:")
    output.append("1. Add tests around risky/complex files.")
    output.append("2. Centralize config and API URLs.")
    output.append("3. Split oversized files.")
    output.append("4. Refactor duplicated logic.")
    output.append("5. Remove dead code after tests pass.")
    output.append("6. Re-run score and full workflow.")

    return "\n".join(output)


def find_refactoring_opportunities(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "REFACTORING OPPORTUNITIES\n"
        "Mode: consolidated safe analysis / no automatic changes\n\n"
        + _refactor_section(
            "1. TECHNICAL DEBT",
            analyze_technical_debt(project_name)
        )
        + _refactor_section(
            "2. OVERSIZED FILES",
            find_oversized_files(project_name)
        )
        + _refactor_section(
            "3. DUPLICATED CODE",
            find_duplicate_code(project_name)
        )
        + _refactor_section(
            "4. DEAD CODE",
            find_dead_code(project_name)
        )
        + _refactor_section(
            "5. ARCHITECTURE RISKS",
            strict_architecture_analyzer_project(project_name)
        )
    )


def what_should_i_refactor_first(project_name):
    debt = analyze_technical_debt(project_name)
    problematic = find_most_problematic_files(project_name)
    duplicates = find_duplicate_code(project_name)

    return (
        "WHAT SHOULD BE REFACTORED FIRST\n"
        "Mode: safe priority recommendation\n\n"
        "Recommended order:\n"
        "1. Files that affect auth, security, admin, uploads, config, or routing.\n"
        "2. Oversized files with high complexity.\n"
        "3. Duplicated business logic.\n"
        "4. Dead code only after tests pass.\n"
        "5. UI/component cleanup after backend/security is stable.\n"
        + _refactor_section("TECHNICAL DEBT INPUT", debt, max_chars=6000)
        + _refactor_section("PROBLEMATIC FILES INPUT", problematic, max_chars=6000)
        + _refactor_section("DUPLICATE CODE INPUT", duplicates, max_chars=6000)
    )


def generate_refactoring_roadmap(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "REFACTORING ROADMAP\n"
        "Mode: safe staged plan / no automatic changes\n"
        f"Project: {project.get('name')}\n\n"
        "Phase 0 - Safety:\n"
        "- Refresh deep project memory.\n"
        "- Create backups.\n"
        "- Run or add baseline tests.\n\n"
        "Phase 1 - Critical structure:\n"
        "- Split oversized entrypoint files.\n"
        "- Move routes/controllers into route modules.\n"
        "- Move business logic into services.\n"
        "- Move validation into schemas/validators.\n\n"
        "Phase 2 - Security-sensitive cleanup:\n"
        "- Centralize auth/JWT/config handling.\n"
        "- Centralize upload validation.\n"
        "- Replace hardcoded URLs/secrets with environment config.\n\n"
        "Phase 3 - Maintainability:\n"
        "- Remove duplicated logic.\n"
        "- Remove dead code after tests pass.\n"
        "- Add documentation for architecture decisions.\n\n"
        "Phase 4 - Verification:\n"
        "- Run tests.\n"
        "- Re-run full security audit.\n"
        "- Re-run production readiness.\n"
        "- Export project report.\n"
        + _refactor_section(
            "REFRACTORING OPPORTUNITIES INPUT",
            find_refactoring_opportunities(project_name),
            max_chars=12000
        )
    )


def generate_modernization_plan(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    try:
        _, facts, facts_error = strict_project_facts(project_name)
    except Exception:
        facts = None
        facts_error = "Could not load facts."

    modernization_items = []

    if facts and not facts_error:
        if not facts["test_files"]:
            modernization_items.append("Add automated tests.")
        if not facts["ci_files"]:
            modernization_items.append("Add CI/workflow checks.")
        if not facts["docker_files"]:
            modernization_items.append("Add Docker/deployment documentation.")
        if facts["frontend_api_calls"]:
            modernization_items.append("Centralize frontend API client configuration.")
        if facts["fastapi_files"] and facts["flask_files"]:
            modernization_items.append("Clarify or migrate mixed Flask/FastAPI backend structure.")
        if not facts["password_hashing_files"] and facts["auth_files"]:
            modernization_items.append("Confirm/modernize password hashing.")

    if not modernization_items:
        modernization_items.append("Modernization priorities should be confirmed by full project review.")

    output = [
        "MODERNIZATION PLAN",
        "Mode: rule-based / safe planning only",
        f"Project: {project.get('name')}",
        "",
        "Recommended modernization items:"
    ]

    for index, item in enumerate(modernization_items, start=1):
        output.append(f"{index}. {item}")

    output.append("")
    output.append("Modernization sequence:")
    output.append("1. Stabilize tests and CI.")
    output.append("2. Centralize config and secrets.")
    output.append("3. Modernize architecture boundaries.")
    output.append("4. Update deployment/release process.")
    output.append("5. Improve documentation and portfolio presentation.")

    return "\n".join(output)


def generate_architecture_migration_plan(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "ARCHITECTURE MIGRATION PLAN\n"
        "Mode: safe staged migration / no automatic changes\n"
        f"Project: {project.get('name')}\n\n"
        "Goal:\n"
        "Move from mixed/scattered implementation toward clear layers:\n"
        "- routes/controllers\n"
        "- services/business logic\n"
        "- models/schemas\n"
        "- config/settings\n"
        "- tests\n"
        "- docs/deployment\n\n"
        "Migration steps:\n"
        "1. Map current architecture using strict architecture analyzer.\n"
        "2. Identify oversized/high-complexity files.\n"
        "3. Choose one safe boundary to migrate first.\n"
        "4. Add tests around that boundary.\n"
        "5. Move logic gradually, keeping public behavior the same.\n"
        "6. Run tests after each file move.\n"
        "7. Re-index project and re-run architecture analyzer.\n"
        + _refactor_section(
            "CURRENT ARCHITECTURE INPUT",
            strict_architecture_analyzer_project(project_name),
            max_chars=10000
        )
        + _refactor_section(
            "OVERSIZED FILES INPUT",
            find_oversized_files(project_name),
            max_chars=8000
        )
    )


def technical_debt(project_name):
    return analyze_technical_debt(project_name)


def duplicated_code(project_name):
    return find_duplicate_code(project_name)


def oversized_files(project_name):
    return find_oversized_files(project_name)


def refactoring_opportunities(project_name):
    return find_refactoring_opportunities(project_name)


def refactoring_roadmap(project_name):
    return generate_refactoring_roadmap(project_name)


def modernization_plan(project_name):
    return generate_modernization_plan(project_name)


def architecture_migration_plan(project_name):
    return generate_architecture_migration_plan(project_name)


def most_problematic_file(project_name):
    return find_most_problematic_files(project_name)


def refactor_first(project_name):
    return what_should_i_refactor_first(project_name)


# ==========================
# STEP 21 - AUTO DOCUMENTATION GENERATOR
# README / API docs / architecture docs / developer guide / onboarding guide.
# Safe generation only. Exports Markdown docs to docs_generated/
# ==========================
DOCS_OUTPUT_DIR = "docs_generated"


def _doc_safe_filename(name):
    cleaned = re.sub(
        r"[^A-Za-z0-9_.-]+",
        "_",
        str(name).strip()
    ).strip("_")

    return cleaned or "document"


def _save_generated_doc(project_name, doc_type, content):
    os.makedirs(
        DOCS_OUTPUT_DIR,
        exist_ok=True
    )

    project_safe = _doc_safe_filename(project_name)
    doc_safe = _doc_safe_filename(doc_type)
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    path = os.path.join(
        DOCS_OUTPUT_DIR,
        f"{project_safe}_{doc_safe}_{timestamp}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _doc_header(title, project_name):
    return (
        f"# {title}\n\n"
        f"Project: `{project_name}`\n\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        "Generated by: JARVIS Engineering OS\n\n"
        "Mode: evidence-based documentation / no automatic code changes\n\n"
    )


def _doc_code_block(title, content):
    return (
        f"\n\n## {title}\n\n"
        "```text\n"
        f"{str(content).strip()}\n"
        "```\n"
    )


def _project_doc_base(project_name):
    project, error = get_project(project_name)

    if error:
        return None, None, error

    files = get_code_files(project)

    try:
        project, facts, facts_error = strict_project_facts(project_name)
    except Exception as e:
        facts = None
        facts_error = str(e)

    return project, {
        "files": files,
        "facts": facts,
        "facts_error": facts_error,
        "overview": project_overview(project),
        "manifest": build_project_manifest(project, files),
        "structure": analyze_project_structure(project_name),
        "score": score_project(project_name),
    }, None


def generate_readme(project_name):
    project, data, error = _project_doc_base(project_name)

    if error:
        return error

    project_title = project.get("name", project_name)

    content = (
        _doc_header(f"README - {project_title}", project_title)
        + "## Overview\n\n"
        + "This README was generated from the indexed project files. Review and adjust wording before publishing.\n\n"
        + _doc_code_block("Project Overview", data["overview"])
        + _doc_code_block("Tech Stack and Manifest", data["manifest"])
        + _doc_code_block("Project Score", data["score"])
        + _doc_code_block("Architecture Summary", strict_architecture_analyzer_project(project_name))
        + _doc_code_block("Security Summary", strict_security_analyzer_project(project_name))
        + _doc_code_block("Production Readiness", estimate_production_readiness(project_name))
        + "\n\n## Suggested README Sections To Keep\n\n"
        "- Project purpose\n"
        "- Features\n"
        "- Tech stack\n"
        "- Installation\n"
        "- Environment variables\n"
        "- Run commands\n"
        "- Test commands\n"
        "- Security notes\n"
        "- Deployment notes\n"
        "- Future improvements\n"
    )

    path = _save_generated_doc(project_title, "README", content)

    return f"{content}\n\nDOCUMENT EXPORTED:\n{path}"


def generate_api_docs(project_name):
    project, data, error = _project_doc_base(project_name)

    if error:
        return error

    project_title = project.get("name", project_name)
    facts = data["facts"]

    routes = []
    api_calls = []

    if facts:
        routes = facts.get("routes", [])
        api_calls = facts.get("frontend_api_calls", [])

    route_text = "\n".join(f"- {route}" for route in routes) if routes else "No backend routes detected in indexed files."
    api_text = "\n".join(f"- {call}" for call in api_calls) if api_calls else "No frontend API calls detected in indexed files."

    content = (
        _doc_header(f"API Documentation - {project_title}", project_title)
        + "## Backend Routes Detected\n\n"
        + route_text
        + "\n\n## Frontend API Calls Detected\n\n"
        + api_text
        + _doc_code_block("Strict Grounded Analyzer", strict_grounded_analyzer_project(project_name))
        + "\n\n## Notes\n\n"
        "- Only routes/API calls visible in indexed files are listed.\n"
        "- Hidden runtime routes are not included unless present in indexed files.\n"
        "- Review manually before publishing API documentation.\n"
    )

    path = _save_generated_doc(project_title, "API_DOCS", content)

    return f"{content}\n\nDOCUMENT EXPORTED:\n{path}"


def generate_architecture_docs(project_name):
    project, data, error = _project_doc_base(project_name)

    if error:
        return error

    project_title = project.get("name", project_name)

    content = (
        _doc_header(f"Architecture Documentation - {project_title}", project_title)
        + _doc_code_block("Project Structure", data["structure"])
        + _doc_code_block("Strict Architecture Analyzer", strict_architecture_analyzer_project(project_name))
        + _doc_code_block("Architecture Report", generate_architecture_report(project_name))
        + _doc_code_block("Architecture Migration Plan", generate_architecture_migration_plan(project_name))
        + "\n\n## Architecture Documentation Notes\n\n"
        "- This document is generated from indexed files only.\n"
        "- Verify diagrams, deployment topology, and runtime infrastructure manually.\n"
    )

    path = _save_generated_doc(project_title, "ARCHITECTURE_DOCS", content)

    return f"{content}\n\nDOCUMENT EXPORTED:\n{path}"


def generate_developer_guide(project_name):
    project, data, error = _project_doc_base(project_name)

    if error:
        return error

    project_title = project.get("name", project_name)

    content = (
        _doc_header(f"Developer Guide - {project_title}", project_title)
        + "## Purpose\n\n"
        "This guide helps a developer understand, run, inspect, and improve the project safely.\n"
        + _doc_code_block("Project Manifest", data["manifest"])
        + _doc_code_block("Project Structure", data["structure"])
        + _doc_code_block("Technical Debt", analyze_technical_debt(project_name))
        + _doc_code_block("Refactoring Roadmap", generate_refactoring_roadmap(project_name))
        + _doc_code_block("Testing / Maintainability Evidence", estimate_project_maturity(project_name))
        + "\n\n## Recommended Developer Workflow\n\n"
        "1. Read the README and architecture docs.\n"
        "2. Install dependencies.\n"
        "3. Configure environment variables.\n"
        "4. Run the project locally.\n"
        "5. Run tests if present.\n"
        "6. Create backups before modifying files.\n"
        "7. Apply improvements one file at a time.\n"
        "8. Re-index the project in JARVIS after major changes.\n"
    )

    path = _save_generated_doc(project_title, "DEVELOPER_GUIDE", content)

    return f"{content}\n\nDOCUMENT EXPORTED:\n{path}"


def generate_onboarding_guide(project_name):
    project, data, error = _project_doc_base(project_name)

    if error:
        return error

    project_title = project.get("name", project_name)

    content = (
        _doc_header(f"Onboarding Guide - {project_title}", project_title)
        + "## First 30 Minutes\n\n"
        "- Read the project overview.\n"
        "- Inspect the main folders and highest-priority files.\n"
        "- Review detected routes, frontend API calls, config, and security files.\n\n"
        + _doc_code_block("Project Overview", data["overview"])
        + _doc_code_block("Strict Grounded Analyzer", strict_grounded_analyzer_project(project_name))
        + _doc_code_block("Files To Inspect First", find_most_problematic_files(project_name))
        + _doc_code_block("Security Items To Understand", strict_security_analyzer_project(project_name))
        + _doc_code_block("Next Improvements", what_should_i_fix_next(project_name))
        + "\n\n## Onboarding Checklist\n\n"
        "- [ ] Understand project purpose.\n"
        "- [ ] Identify frontend/backend entry points.\n"
        "- [ ] Identify config and environment variables.\n"
        "- [ ] Identify auth/security flow.\n"
        "- [ ] Identify database/storage usage.\n"
        "- [ ] Identify test and deployment process.\n"
        "- [ ] Run JARVIS workflow before major changes.\n"
    )

    path = _save_generated_doc(project_title, "ONBOARDING_GUIDE", content)

    return f"{content}\n\nDOCUMENT EXPORTED:\n{path}"


def generate_full_documentation_pack(project_name):
    outputs = [
        generate_readme(project_name),
        generate_api_docs(project_name),
        generate_architecture_docs(project_name),
        generate_developer_guide(project_name),
        generate_onboarding_guide(project_name),
    ]

    exported = []

    for output in outputs:
        marker = "DOCUMENT EXPORTED:"
        if marker in output:
            exported.append(output.split(marker)[-1].strip())

    return (
        "FULL DOCUMENTATION PACK GENERATED\n\n"
        + "\n".join(f"- {path}" for path in exported)
    )


# Friendly aliases
def generate_readme_docs(project_name):
    return generate_readme(project_name)


def generate_api_documentation(project_name):
    return generate_api_docs(project_name)


def generate_architecture_documentation(project_name):
    return generate_architecture_docs(project_name)


def generate_dev_guide(project_name):
    return generate_developer_guide(project_name)


def generate_onboarding_docs(project_name):
    return generate_onboarding_guide(project_name)


def generate_docs_pack(project_name):
    return generate_full_documentation_pack(project_name)



# ==========================
# STEP 23 - TEST INTELLIGENCE
# Test coverage evidence / missing tests / critical code without tests / QA roadmap.
# Offline, rule-based, safe analysis only. No test execution.
# ==========================
TEST_OUTPUT_DIR = "reports"


def _test_section(title, content, max_chars=9000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 76
        + f"\n{title}\n"
        + "=" * 76
        + "\n"
        + text
    )


def _is_test_file_path(path):
    return is_test_path(path)


def _is_critical_source_file(path):
    lower = normalize_path(path)

    critical_tokens = [
        "auth",
        "login",
        "register",
        "jwt",
        "token",
        "security",
        "permission",
        "admin",
        "upload",
        "file",
        "scan",
        "database",
        "db",
        "models",
        "routes",
        "api",
        "payment",
        "config",
        "settings",
        "middleware",
    ]

    return any(token in lower for token in critical_tokens)


def _source_to_test_keywords(path):
    lower = normalize_path(path)
    base = os.path.splitext(os.path.basename(lower))[0]

    parts = [
        base,
        base.replace("_", ""),
        base.replace("-", ""),
    ]

    for token in [
        "auth",
        "login",
        "register",
        "jwt",
        "token",
        "admin",
        "upload",
        "scan",
        "dashboard",
        "api",
        "routes",
        "database",
        "db",
        "config",
        "settings",
    ]:
        if token in lower:
            parts.append(token)

    return sorted(set(part for part in parts if part))


def detect_test_frameworks(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    all_text = build_reference_text(files)

    frameworks = []

    checks = {
        "pytest": ["pytest", "def test_", "pytest.ini"],
        "unittest": ["unittest", "testcase"],
        "jest": ["jest", "describe(", "it(", "test("],
        "vitest": ["vitest", "vi.", "vite"],
        "react testing library": ["@testing-library/react", "screen.getby", "render("],
        "cypress": ["cypress", "cy."],
        "playwright": ["playwright", "@playwright/test"],
        "mocha": ["mocha", "chai", "expect("],
        "angular testing": ["testbed", "karma", "jasmine", ".spec.ts"],
    }

    for framework, markers in checks.items():
        if any(marker.lower() in all_text for marker in markers):
            frameworks.append(framework)

    if not frameworks:
        frameworks.append("No explicit test framework detected.")

    return "\n".join(
        [
            "TEST FRAMEWORKS DETECTED",
            "Mode: rule-based / indexed files only",
            f"Project: {project.get('name')}",
            "",
            *[f"- {framework}" for framework in sorted(set(frameworks))]
        ]
    )


def show_test_coverage(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    source_files = [
        item for item in files
        if item["extension"] in {".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".cs"}
        and not _is_test_file_path(item["path"])
    ]

    test_files = [
        item for item in files
        if _is_test_file_path(item["path"])
    ]

    critical_files = [
        item for item in source_files
        if _is_critical_source_file(item["path"])
    ]

    coverage_ratio = 0

    if source_files:
        coverage_ratio = round(
            len(test_files) / len(source_files) * 100,
            1
        )

    critical_ratio = 0

    if critical_files:
        critical_tests = 0
        test_text = build_reference_text(test_files)

        for item in critical_files:
            keywords = _source_to_test_keywords(item["path"])

            if any(keyword in test_text for keyword in keywords):
                critical_tests += 1

        critical_ratio = round(
            critical_tests / len(critical_files) * 100,
            1
        )

    output = [
        "TEST COVERAGE EVIDENCE REPORT",
        "Mode: heuristic / no test execution / no real coverage instrumentation",
        f"Project: {project.get('name')}",
        "",
        f"Source files detected: {len(source_files)}",
        f"Test files detected: {len(test_files)}",
        f"Approx test/source ratio: {coverage_ratio}%",
        f"Critical source files detected: {len(critical_files)}",
        f"Approx critical-test evidence ratio: {critical_ratio}%",
        "",
        "Detected test files:"
    ]

    if test_files:
        for item in test_files[:80]:
            output.append(f"- {item['path']}")
    else:
        output.append("- No test files detected.")

    output.append("")
    output.append("Important:")
    output.append("- This is not real code coverage.")
    output.append("- Run pytest --cov, npm test -- --coverage, vitest --coverage, ng test --code-coverage, or equivalent.")

    return "\n".join(output)


def missing_tests(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    source_files = [
        item for item in files
        if item["extension"] in {".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".cs"}
        and not _is_test_file_path(item["path"])
    ]

    test_files = [
        item for item in files
        if _is_test_file_path(item["path"])
    ]

    test_text = build_reference_text(test_files)
    missing = []

    for item in source_files:
        path = item["path"]

        if is_documentation_or_config_path(path):
            continue

        keywords = _source_to_test_keywords(path)

        if not any(keyword in test_text for keyword in keywords):
            priority = "HIGH" if _is_critical_source_file(path) else "MEDIUM"

            missing.append(
                f"{priority}: {path}"
            )

    output = [
        "MISSING TESTS REPORT",
        "Mode: heuristic / verify manually",
        f"Project: {project.get('name')}",
        "",
        "Source files with weak/no obvious test evidence:"
    ]

    if missing:
        output.extend(f"- {item}" for item in missing[:100])

        if len(missing) > 100:
            output.append(f"... and {len(missing) - 100} more")
    else:
        output.append("- No obvious missing test targets detected.")

    output.append("")
    output.append("Recommended priority:")
    output.append("1. Auth/login/token/admin/security files.")
    output.append("2. Upload/file/database/config routes.")
    output.append("3. Main business logic/services.")
    output.append("4. UI flows with forms/API calls.")
    output.append("5. Edge cases and error handling.")

    return "\n".join(output)


def critical_code_without_tests(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    source_files = [
        item for item in files
        if item["extension"] in {".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".cs"}
        and not _is_test_file_path(item["path"])
        and _is_critical_source_file(item["path"])
    ]

    test_files = [
        item for item in files
        if _is_test_file_path(item["path"])
    ]

    test_text = build_reference_text(test_files)
    missing = []

    for item in source_files:
        keywords = _source_to_test_keywords(item["path"])

        if not any(keyword in test_text for keyword in keywords):
            missing.append(item["path"])

    output = [
        "CRITICAL CODE WITHOUT TESTS",
        "Mode: heuristic / no test execution",
        f"Project: {project.get('name')}",
        "",
        "Critical files with weak/no obvious test evidence:"
    ]

    if missing:
        output.extend(f"- {path}" for path in missing[:80])
    else:
        output.append("- No critical untested files detected by current heuristic.")

    output.append("")
    output.append("Suggested test types:")
    output.append("- Auth/token tests.")
    output.append("- Permission/admin tests.")
    output.append("- Upload validation tests.")
    output.append("- Route/API response tests.")
    output.append("- Database/service tests.")
    output.append("- Frontend form and API error-state tests.")

    return "\n".join(output)


def generate_test_strategy(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "TEST STRATEGY\n"
        "Mode: evidence-based testing plan / no automatic test generation\n"
        f"Project: {project.get('name')}\n\n"
        "Recommended layers:\n"
        "1. Unit tests for pure functions, validators, services, and utilities.\n"
        "2. API/route tests for auth, upload, admin, dashboard, database, and error handling.\n"
        "3. Frontend component tests for forms, API states, dashboard states, and navigation.\n"
        "4. Integration tests for login -> token -> protected route flow.\n"
        "5. Security regression tests for file upload, auth bypass, invalid tokens, and bad input.\n"
        "6. Smoke tests for startup and release readiness.\n"
        + _test_section("Detected Test Frameworks", detect_test_frameworks(project_name))
        + _test_section("Current Test Evidence", show_test_coverage(project_name))
        + _test_section("Missing Tests", missing_tests(project_name))
    )


def test_roadmap(project_name):
    return (
        "TEST ROADMAP\n"
        "Mode: staged QA improvement plan\n\n"
        "Phase 1 - Baseline:\n"
        "- Identify current test framework.\n"
        "- Add one smoke test for project startup.\n"
        "- Add one test for every critical API route or auth flow.\n\n"
        "Phase 2 - Security-critical tests:\n"
        "- Invalid login.\n"
        "- Expired/invalid token.\n"
        "- Unauthorized admin access.\n"
        "- Upload invalid file type/size.\n"
        "- Dangerous payload input.\n\n"
        "Phase 3 - Frontend and workflow tests:\n"
        "- Login/register forms.\n"
        "- Dashboard loading/error states.\n"
        "- API failure handling.\n"
        "- Navigation and protected pages.\n\n"
        "Phase 4 - Release quality gate:\n"
        "- Coverage report.\n"
        "- CI test command.\n"
        "- Security regression suite.\n"
        "- Production readiness re-check.\n"
        + _test_section("Critical Untested Code", critical_code_without_tests(project_name))
    )


def generate_qa_plan(project_name):
    return (
        "QA PLAN\n"
        "Mode: practical QA checklist / no automatic execution\n\n"
        "Manual QA:\n"
        "- Run application locally.\n"
        "- Validate main user flows.\n"
        "- Validate error messages and loading states.\n"
        "- Validate responsive layout if frontend exists.\n\n"
        "Automated QA:\n"
        "- Unit tests for utilities/services.\n"
        "- API tests for routes and security flows.\n"
        "- UI tests for core screens.\n"
        "- Smoke tests before release.\n\n"
        "Security QA:\n"
        "- Invalid auth attempts.\n"
        "- Token tampering.\n"
        "- Upload validation.\n"
        "- Missing permissions.\n"
        "- Sensitive data in logs.\n\n"
        "Release QA:\n"
        "- Run dependency health report.\n"
        "- Run full security audit.\n"
        "- Run production readiness report.\n"
        "- Run tests and coverage locally.\n"
        + _test_section("Project Test Evidence", show_test_coverage(project_name))
    )


def generate_testing_checklist(project_name):
    return (
        "TESTING CHECKLIST\n\n"
        "- [ ] Test framework detected or selected.\n"
        "- [ ] Test command documented in README.\n"
        "- [ ] Auth/login tests added.\n"
        "- [ ] Token/JWT tests added.\n"
        "- [ ] Admin/permission tests added.\n"
        "- [ ] Upload/file validation tests added.\n"
        "- [ ] API route tests added.\n"
        "- [ ] Frontend form tests added.\n"
        "- [ ] Error/loading state tests added.\n"
        "- [ ] Smoke test added.\n"
        "- [ ] Coverage command documented.\n"
        "- [ ] CI test step added.\n"
        "- [ ] Tests run before release.\n"
        + _test_section("Missing Tests Input", missing_tests(project_name), max_chars=7000)
    )


def test_intelligence_report(project_name):
    return (
        "JARVIS TEST INTELLIGENCE REPORT\n"
        "Mode: offline evidence-based testing intelligence\n"
        + _test_section("1. Test Frameworks", detect_test_frameworks(project_name))
        + _test_section("2. Coverage Evidence", show_test_coverage(project_name))
        + _test_section("3. Missing Tests", missing_tests(project_name))
        + _test_section("4. Critical Code Without Tests", critical_code_without_tests(project_name))
        + _test_section("5. Test Strategy", generate_test_strategy(project_name), max_chars=9000)
        + _test_section("6. Test Roadmap", test_roadmap(project_name), max_chars=9000)
        + _test_section("7. QA Plan", generate_qa_plan(project_name), max_chars=9000)
    )


# Friendly aliases
def test_coverage(project_name):
    return show_test_coverage(project_name)


def tests_missing(project_name):
    return missing_tests(project_name)


def untested_critical_code(project_name):
    return critical_code_without_tests(project_name)


def test_strategy(project_name):
    return generate_test_strategy(project_name)


def qa_plan(project_name):
    return generate_qa_plan(project_name)


def testing_checklist(project_name):
    return generate_testing_checklist(project_name)


def test_report(project_name):
    return test_intelligence_report(project_name)



# ==========================
# STEP 24 - PERFORMANCE INTELLIGENCE
# Bottlenecks / heavy files / slow modules / frontend/backend/database audits.
# Offline, rule-based, safe analysis only. No profiling execution.
# ==========================
def _perf_section(title, content, max_chars=9000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 76
        + f"\n{title}\n"
        + "=" * 76
        + "\n"
        + text
    )


def _performance_risk_score(item):
    path = item.get("path", "")
    content = item.get("content", "")
    lower = content.lower()
    score = 0
    reasons = []

    lines = len(content.splitlines())
    chars = len(content)

    if lines >= 180:
        score += 20
        reasons.append(f"Large file: {lines} lines.")

    if chars >= 12000:
        score += 15
        reasons.append(f"Large content size: {chars} chars.")

    patterns = [
        ("nested loops / loops", ["for ", "while "], 2),
        ("sync file I/O", ["open(", "read(", "write(", "os.walk", "shutil.copy"], 5),
        ("network/API calls", ["requests.", "axios.", "fetch(", "httpx.", "urllib"], 5),
        ("database operations", ["select ", "insert ", "update ", "delete ", "query(", "execute(", "session.query"], 6),
        ("sleep/blocking calls", ["time.sleep", "thread.sleep", "sleep("], 8),
        ("large JSON handling", ["json.load", "json.loads", "json.dump", "json.dumps"], 4),
        ("image/OCR/ML processing", ["cv2.", "pytesseract", "pillow", "pil", "numpy", "pandas"], 8),
        ("render-heavy frontend", ["usestate", "useeffect", "setinterval", "settimeout", "map(", "filter("], 4),
        ("DOM/storage usage", ["localstorage", "sessionstorage", "document.", "window."], 3),
    ]

    for label, markers, weight in patterns:
        hits = sum(lower.count(marker.lower()) for marker in markers)
        if hits:
            add = min(25, hits * weight)
            score += add
            reasons.append(f"{label}: {hits} marker(s).")

    if item.get("extension") in {".tsx", ".jsx", ".ts", ".js"} and "useeffect" in lower and "set" in lower:
        score += 8
        reasons.append("Frontend state/effect logic may trigger re-render issues.")

    if item.get("extension") == ".py" and ("flask" in lower or "fastapi" in lower or "apirouter" in lower):
        score += 6
        reasons.append("Backend route/API file; request latency matters.")

    if "n+1" in lower:
        score += 15
        reasons.append("N+1 query marker/comment detected.")

    return score, reasons


def heavy_files(project_name, limit=25):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    rows = []

    for item in files:
        content = item.get("content", "")
        lines = len(content.splitlines())
        chars = len(content)

        if lines >= 120 or chars >= 8000:
            score, reasons = _performance_risk_score(item)
            rows.append((score, lines, chars, item, reasons))

    rows.sort(
        key=lambda row: (row[0], row[1], row[2]),
        reverse=True
    )

    output = [
        "HEAVY FILES REPORT",
        "Mode: offline heuristic / no runtime profiling",
        f"Project: {project.get('name')}",
        "",
        "Files that may affect performance or maintainability:"
    ]

    if not rows:
        output.append("- No heavy files detected by current thresholds.")
        return "\n".join(output)

    for score, lines, chars, item, reasons in rows[:limit]:
        output.append(
            f"- {item['path']} | score {score} | {lines} lines | {chars} chars"
        )

        for reason in reasons[:4]:
            output.append(f"  - {reason}")

    return "\n".join(output)


def find_bottlenecks(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    rows = []

    for item in files:
        score, reasons = _performance_risk_score(item)

        if score >= 12:
            rows.append((score, item, reasons))

    rows.sort(
        key=lambda row: row[0],
        reverse=True
    )

    output = [
        "POSSIBLE PERFORMANCE BOTTLENECKS",
        "Mode: static heuristic / verify with profiler and real measurements",
        f"Project: {project.get('name')}",
        "",
        "Candidates:"
    ]

    if not rows:
        output.append("- No obvious bottleneck candidates detected.")
        return "\n".join(output)

    for score, item, reasons in rows[:30]:
        output.append(f"- {item['path']} | risk score {score}")

        for reason in reasons[:5]:
            output.append(f"  - {reason}")

    output.append("")
    output.append("Recommended validation:")
    output.append("- Python: cProfile, py-spy, line_profiler.")
    output.append("- Node/React: Lighthouse, React DevTools Profiler, Chrome Performance.")
    output.append("- Backend APIs: measure route latency and database query times.")
    output.append("- Database: inspect indexes and repeated queries.")

    return "\n".join(output)


def slow_modules(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    folder_scores = defaultdict(int)
    folder_reasons = defaultdict(list)

    for item in files:
        path = item["path"]
        folder = os.path.dirname(path) or "root"
        score, reasons = _performance_risk_score(item)

        folder_scores[folder] += score

        for reason in reasons[:3]:
            folder_reasons[folder].append(f"{path}: {reason}")

    ranked = sorted(
        folder_scores.items(),
        key=lambda row: row[1],
        reverse=True
    )

    output = [
        "SLOW MODULES / HOT AREAS",
        "Mode: folder-level static heuristic",
        f"Project: {project.get('name')}",
        ""
    ]

    for folder, score in ranked[:20]:
        if score <= 0:
            continue

        output.append(f"- {folder} | accumulated risk score {score}")

        for reason in folder_reasons[folder][:5]:
            output.append(f"  - {reason}")

    if len(output) <= 4:
        output.append("- No slow module candidates detected.")

    return "\n".join(output)


def memory_hotspots(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    hotspots = []

    markers = [
        "read()",
        "readlines()",
        "json.load",
        "json.loads",
        "pandas",
        "numpy",
        "image.open",
        "cv2.imread",
        "list(",
        ".append(",
        "cache",
        "localstorage",
        "sessionstorage",
        "blob",
        "arraybuffer",
    ]

    for item in files:
        content = item.get("content", "").lower()
        hits = []

        for marker in markers:
            count = content.count(marker.lower())
            if count:
                hits.append(f"{marker}: {count}")

        if hits:
            hotspots.append((len(hits), item["path"], hits))

    hotspots.sort(
        key=lambda row: row[0],
        reverse=True
    )

    output = [
        "MEMORY HOTSPOTS",
        "Mode: static heuristic / verify with runtime memory profiler",
        f"Project: {project.get('name')}",
        ""
    ]

    if not hotspots:
        output.append("- No memory hotspot markers detected.")
        return "\n".join(output)

    for _, path, hits in hotspots[:25]:
        output.append(f"- {path}")
        for hit in hits[:6]:
            output.append(f"  - {hit}")

    output.append("")
    output.append("Recommendations:")
    output.append("- Stream large files instead of reading fully into memory.")
    output.append("- Limit cached data size.")
    output.append("- Paginate large API responses.")
    output.append("- Avoid unnecessary large array/object copies.")

    return "\n".join(output)


def cpu_hotspots(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    hotspots = []

    markers = [
        "for ",
        "while ",
        "re.findall",
        "re.search",
        "sort(",
        "sorted(",
        "json.loads",
        "json.dumps",
        "cv2.",
        "pytesseract",
        "pandas",
        "numpy",
        "map(",
        "filter(",
        "reduce(",
    ]

    for item in files:
        content = item.get("content", "").lower()
        hits = []

        for marker in markers:
            count = content.count(marker.lower())

            if count:
                hits.append((marker, count))

        score = sum(count for _, count in hits)

        if score >= 5:
            hotspots.append((score, item["path"], hits))

    hotspots.sort(
        key=lambda row: row[0],
        reverse=True
    )

    output = [
        "CPU HOTSPOTS",
        "Mode: static heuristic / verify with runtime profiler",
        f"Project: {project.get('name')}",
        ""
    ]

    if not hotspots:
        output.append("- No CPU hotspot candidates detected.")
        return "\n".join(output)

    for score, path, hits in hotspots[:25]:
        output.append(f"- {path} | CPU marker score {score}")

        for marker, count in hits[:8]:
            output.append(f"  - {marker}: {count}")

    output.append("")
    output.append("Recommendations:")
    output.append("- Profile before optimizing.")
    output.append("- Avoid repeated expensive operations inside loops.")
    output.append("- Cache safe deterministic results.")
    output.append("- Move heavy OCR/image/ML work off request threads.")
    output.append("- Use pagination/batching for large datasets.")

    return "\n".join(output)


def frontend_performance_audit(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    frontend_files = [
        item for item in files
        if item["extension"] in {".js", ".jsx", ".ts", ".tsx", ".css", ".scss", ".html"}
        and (
            "src/" in normalize_path(item["path"])
            or "components/" in normalize_path(item["path"])
            or "frontend" in normalize_path(item["path"])
            or item["extension"] in {".jsx", ".tsx"}
        )
    ]

    findings = []

    for item in frontend_files:
        path = item["path"]
        lower = item["content"].lower()

        if "useeffect" in lower and "set" in lower:
            findings.append(f"{path}: review useEffect dependencies and repeated state updates.")

        if lower.count(".map(") >= 4:
            findings.append(f"{path}: multiple map operations; verify list rendering and memoization.")

        if "localstorage" in lower or "sessionstorage" in lower:
            findings.append(f"{path}: storage access detected; avoid blocking critical render path.")

        if "axios." in lower or "fetch(" in lower:
            findings.append(f"{path}: API call detected; verify loading/error/caching behavior.")

        if "framer-motion" in lower or "animation" in lower:
            findings.append(f"{path}: animation usage detected; verify performance on low-end devices.")

        if "import *" in lower:
            findings.append(f"{path}: wildcard import may increase bundle size.")

    output = [
        "FRONTEND PERFORMANCE AUDIT",
        "Mode: static heuristic / verify with Lighthouse and React/Angular profiler",
        f"Project: {project.get('name')}",
        "",
        f"Frontend files analyzed: {len(frontend_files)}",
        "",
        "Findings:"
    ]

    if findings:
        output.extend(f"- {item}" for item in findings[:80])
    else:
        output.append("- No obvious frontend performance issues detected.")

    output.append("")
    output.append("Recommended checks:")
    output.append("- Run Lighthouse.")
    output.append("- Check bundle size.")
    output.append("- Lazy-load heavy routes/components.")
    output.append("- Memoize expensive computed values only where needed.")
    output.append("- Avoid repeated API calls from useEffect loops.")

    return "\n".join(output)


def backend_performance_audit(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    backend_files = [
        item for item in files
        if item["extension"] in {".py", ".js", ".ts"}
        and (
            "route" in normalize_path(item["path"])
            or "api" in normalize_path(item["path"])
            or "server" in normalize_path(item["path"])
            or "app.py" in normalize_path(item["path"])
            or "main.py" in normalize_path(item["path"])
            or "fastapi" in item["content"].lower()
            or "flask" in item["content"].lower()
            or "express" in item["content"].lower()
        )
    ]

    findings = []

    for item in backend_files:
        path = item["path"]
        lower = item["content"].lower()

        if "time.sleep" in lower or "sleep(" in lower:
            findings.append(f"{path}: blocking sleep detected.")

        if "os.walk" in lower:
            findings.append(f"{path}: filesystem walk detected; avoid inside request path.")

        if "read()" in lower or "readlines()" in lower:
            findings.append(f"{path}: full file read detected; stream or limit large files.")

        if "requests." in lower or "httpx." in lower:
            findings.append(f"{path}: outgoing HTTP call detected; add timeout/retry/backoff.")

        if "json.loads" in lower or "json.dumps" in lower:
            findings.append(f"{path}: JSON serialization/parsing detected; watch large payloads.")

        if "print(" in lower:
            findings.append(f"{path}: print logging detected; structured async-safe logging is better.")

        if lower.count("for ") + lower.count("while ") >= 8:
            findings.append(f"{path}: many loops detected; inspect request-time complexity.")

    output = [
        "BACKEND PERFORMANCE AUDIT",
        "Mode: static heuristic / verify with request profiling",
        f"Project: {project.get('name')}",
        "",
        f"Backend files analyzed: {len(backend_files)}",
        "",
        "Findings:"
    ]

    if findings:
        output.extend(f"- {item}" for item in findings[:80])
    else:
        output.append("- No obvious backend performance issues detected.")

    output.append("")
    output.append("Recommended checks:")
    output.append("- Add timeouts for external HTTP calls.")
    output.append("- Avoid heavy work inside request handlers.")
    output.append("- Use background jobs for OCR, scanning, ML, or big file processing.")
    output.append("- Paginate responses and avoid returning huge JSON.")
    output.append("- Measure p95 latency for important routes.")

    return "\n".join(output)


def database_performance_audit(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)

    db_files = [
        item for item in files
        if (
            "database" in normalize_path(item["path"])
            or "db" in normalize_path(item["path"])
            or "model" in normalize_path(item["path"])
            or "sql" in item["content"].lower()
            or "sqlalchemy" in item["content"].lower()
            or "sqlite" in item["content"].lower()
            or "mongoose" in item["content"].lower()
        )
    ]

    findings = []

    for item in db_files:
        path = item["path"]
        lower = item["content"].lower()

        if "select *" in lower:
            findings.append(f"{path}: SELECT * detected; select only needed columns.")

        if ".all()" in lower:
            findings.append(f"{path}: ORM .all() detected; verify pagination/limits.")

        if "limit" not in lower and ("select " in lower or ".query" in lower):
            findings.append(f"{path}: query evidence without visible limit; verify pagination.")

        if "index" not in lower and ("where" in lower or "filter(" in lower):
            findings.append(f"{path}: filters detected; verify indexes on filtered columns.")

        if "sqlite" in lower:
            findings.append(f"{path}: SQLite detected; verify concurrency and production fit.")

    output = [
        "DATABASE PERFORMANCE AUDIT",
        "Mode: static heuristic / verify with DB profiler/query plan",
        f"Project: {project.get('name')}",
        "",
        f"Database-related files analyzed: {len(db_files)}",
        "",
        "Findings:"
    ]

    if findings:
        output.extend(f"- {item}" for item in findings[:80])
    else:
        output.append("- No obvious database performance issues detected.")

    output.append("")
    output.append("Recommended checks:")
    output.append("- Add indexes for common filters.")
    output.append("- Add pagination/limits.")
    output.append("- Avoid N+1 query patterns.")
    output.append("- Measure slow queries.")
    output.append("- Consider production DB choice and connection pooling.")

    return "\n".join(output)


def optimization_roadmap(project_name):
    return (
        "OPTIMIZATION ROADMAP\n"
        "Mode: safe performance planning / no automatic code changes\n\n"
        "Phase 1 - Measure first:\n"
        "- Run app locally.\n"
        "- Measure startup time, route latency, UI load time, and memory usage.\n"
        "- Identify p95 latency for critical flows.\n\n"
        "Phase 2 - Quick wins:\n"
        "- Remove repeated API calls.\n"
        "- Add pagination and payload limits.\n"
        "- Add timeouts to external requests.\n"
        "- Avoid full file reads for large files.\n\n"
        "Phase 3 - Backend optimization:\n"
        "- Move heavy work to background tasks.\n"
        "- Cache safe deterministic results.\n"
        "- Profile slow endpoints.\n\n"
        "Phase 4 - Frontend optimization:\n"
        "- Run Lighthouse.\n"
        "- Split/lazy-load heavy routes.\n"
        "- Review useEffect loops and large renders.\n\n"
        "Phase 5 - Database optimization:\n"
        "- Add indexes.\n"
        "- Add limits/pagination.\n"
        "- Review N+1 query patterns.\n"
        + _perf_section("Bottleneck Input", find_bottlenecks(project_name), max_chars=9000)
    )


def performance_report(project_name):
    return (
        "JARVIS PERFORMANCE INTELLIGENCE REPORT\n"
        "Mode: offline static performance intelligence / verify with profiling\n"
        + _perf_section("1. Heavy Files", heavy_files(project_name))
        + _perf_section("2. Bottlenecks", find_bottlenecks(project_name))
        + _perf_section("3. Slow Modules", slow_modules(project_name))
        + _perf_section("4. Memory Hotspots", memory_hotspots(project_name))
        + _perf_section("5. CPU Hotspots", cpu_hotspots(project_name))
        + _perf_section("6. Frontend Performance", frontend_performance_audit(project_name))
        + _perf_section("7. Backend Performance", backend_performance_audit(project_name))
        + _perf_section("8. Database Performance", database_performance_audit(project_name))
        + _perf_section("9. Optimization Roadmap", optimization_roadmap(project_name), max_chars=9000)
    )


# Friendly aliases
def bottlenecks(project_name):
    return find_bottlenecks(project_name)


def performance_bottlenecks(project_name):
    return find_bottlenecks(project_name)


def performance_hotspots(project_name):
    return find_bottlenecks(project_name)


def frontend_perf(project_name):
    return frontend_performance_audit(project_name)


def backend_perf(project_name):
    return backend_performance_audit(project_name)


def database_perf(project_name):
    return database_performance_audit(project_name)


def db_performance_audit(project_name):
    return database_performance_audit(project_name)


def performance_roadmap(project_name):
    return optimization_roadmap(project_name)



# ==========================
# STEP 25 - ENTERPRISE REFACTORING INTELLIGENCE
# Enterprise readiness / clean code / maintainability / refactor commander.
# Safe planning only. No automatic code changes.
# ==========================
def _enterprise_section(title, content, max_chars=10000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 76
        + f"\n{title}\n"
        + "=" * 76
        + "\n"
        + text
    )


def _clean_code_smells_for_file(item):
    path = item.get("path", "")
    content = item.get("content", "")
    lower = content.lower()
    smells = []

    lines = content.splitlines()

    if len(lines) >= 180:
        smells.append(f"Large file ({len(lines)} lines). Consider splitting responsibilities.")

    if lower.count("print(") >= 5:
        smells.append("Many print statements. Replace production prints with structured logging.")

    if lower.count("try:") >= 6 or lower.count("except") >= 6:
        smells.append("Many try/except blocks. Review error handling boundaries.")

    if lower.count("if ") + lower.count("elif ") >= 18:
        smells.append("High conditional density. Consider smaller functions or strategy mapping.")

    if lower.count("for ") + lower.count("while ") >= 12:
        smells.append("Many loops. Review complexity and performance.")

    if "todo" in lower or "fixme" in lower:
        smells.append("TODO/FIXME markers found.")

    if "localhost" in lower or "127.0.0.1" in lower:
        smells.append("Hardcoded local URL/config marker found.")

    if "your-secret-key" in lower or "changeme" in lower or "password123" in lower:
        smells.append("Weak/demo secret marker found.")

    if item.get("extension") in {".js", ".jsx", ".ts", ".tsx"}:
        if lower.count("useeffect") >= 4:
            smells.append("Many useEffect hooks. Review side effects and re-render behavior.")

        if lower.count("usestate") >= 8:
            smells.append("Many useState hooks. Consider reducer/state grouping.")

        if lower.count("axios.") + lower.count("fetch(") >= 5:
            smells.append("Many API calls in file. Consider service layer extraction.")

    if item.get("extension") == ".py":
        if "app.route" in lower or "apirouter" in lower or "fastapi" in lower or "flask" in lower:
            if len(lines) >= 120:
                smells.append("Large backend route/API file. Split routes, services, schemas, and config.")

    return smells


def clean_code_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    findings = []

    for item in files:
        smells = _clean_code_smells_for_file(item)

        if smells:
            findings.append((len(smells), item["path"], smells))

    findings.sort(
        key=lambda row: row[0],
        reverse=True
    )

    output = [
        "CLEAN CODE REPORT",
        "Mode: offline static heuristic / no automatic changes",
        f"Project: {project.get('name')}",
        "",
        "Files with clean-code smells:"
    ]

    if not findings:
        output.append("- No major clean-code smells detected by current rules.")
        return "\n".join(output)

    for _, path, smells in findings[:40]:
        output.append(f"- {path}")

        for smell in smells[:6]:
            output.append(f"  - {smell}")

    output.append("")
    output.append("Clean code priorities:")
    output.append("1. Split large files.")
    output.append("2. Move business logic out of routes/components.")
    output.append("3. Centralize config/API clients.")
    output.append("4. Replace print/debug code with structured logging.")
    output.append("5. Add tests before refactoring behavior.")

    return "\n".join(output)


def maintainability_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    card, card_error = project_scorecard(project_name)

    score_text = card_error if card_error else format_project_scorecard(card)

    return (
        "MAINTAINABILITY REPORT\n"
        "Mode: evidence-based / safe planning\n"
        f"Project: {project.get('name')}\n"
        + _enterprise_section("1. Project Scorecard", score_text, max_chars=9000)
        + _enterprise_section("2. Technical Debt", analyze_technical_debt(project_name), max_chars=9000)
        + _enterprise_section("3. Clean Code Report", clean_code_report(project_name), max_chars=9000)
        + _enterprise_section("4. Dead Code", find_dead_code(project_name), max_chars=9000)
        + _enterprise_section("5. Duplicate Logic", find_duplicate_code(project_name), max_chars=9000)
        + _enterprise_section("6. Testing Evidence", show_test_coverage(project_name), max_chars=9000)
    )


def enterprise_readiness_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "ENTERPRISE READINESS REPORT\n"
        "Mode: static intelligence / verify manually before production\n"
        f"Project: {project.get('name')}\n"
        + _enterprise_section("1. Production Readiness", estimate_production_readiness(project_name), max_chars=9000)
        + _enterprise_section("2. Security Audit", full_security_audit(project_name), max_chars=9000)
        + _enterprise_section("3. Dependency Health", dependency_health_report(project_name) if "dependency_health_report" in globals() else "Dependency Intelligence is available in deep_project_memory.py.", max_chars=9000)
        + _enterprise_section("4. Test Intelligence", test_intelligence_report(project_name), max_chars=9000)
        + _enterprise_section("5. Performance Intelligence", performance_report(project_name), max_chars=9000)
        + _enterprise_section("6. Maintainability", maintainability_report(project_name), max_chars=12000)
        + _enterprise_section("7. Release Readiness", generate_release_checklist(project_name), max_chars=9000)
    )


def refactor_project(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "ENTERPRISE REFACTORING PLAN\n"
        "Mode: safe refactoring intelligence / no automatic code changes\n"
        f"Project: {project.get('name')}\n\n"
        "Refactoring principle:\n"
        "Do not change behavior before tests exist. Refactor one boundary at a time.\n"
        + _enterprise_section("1. What To Refactor First", what_should_i_refactor_first(project_name), max_chars=12000)
        + _enterprise_section("2. Refactoring Roadmap", generate_refactoring_roadmap(project_name), max_chars=12000)
        + _enterprise_section("3. Architecture Migration Plan", generate_architecture_migration_plan(project_name), max_chars=12000)
        + _enterprise_section("4. Clean Code Report", clean_code_report(project_name), max_chars=9000)
        + _enterprise_section("5. Performance Risks", performance_report(project_name), max_chars=9000)
        + _enterprise_section("6. Test Safety Plan", test_roadmap(project_name), max_chars=9000)
    )


def technical_debt_report(project_name):
    return (
        "TECHNICAL DEBT REPORT\n"
        "Mode: consolidated technical debt intelligence\n"
        + _enterprise_section("1. Technical Debt Analysis", analyze_technical_debt(project_name), max_chars=12000)
        + _enterprise_section("2. Problematic Files", find_most_problematic_files(project_name), max_chars=9000)
        + _enterprise_section("3. Oversized Files", find_oversized_files(project_name), max_chars=9000)
        + _enterprise_section("4. Dead Code", find_dead_code(project_name), max_chars=9000)
        + _enterprise_section("5. Duplicate Code", find_duplicate_code(project_name), max_chars=9000)
        + _enterprise_section("6. Refactoring Roadmap", generate_refactoring_roadmap(project_name), max_chars=12000)
    )


def refactoring_impact_matrix(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    rows = []

    for item in files:
        path = item["path"]
        score = 0
        reasons = []

        if is_entrypoint_or_framework_file(path):
            score += 25
            reasons.append("Entrypoint/framework/routing file.")

        if _is_critical_source_file(path):
            score += 20
            reasons.append("Critical security/API/config/business file.")

        if len(item.get("content", "").splitlines()) >= 160:
            score += 15
            reasons.append("Large file.")

        if _clean_code_smells_for_file(item):
            score += 10
            reasons.append("Clean-code smells detected.")

        if _performance_risk_score(item)[0] >= 15:
            score += 10
            reasons.append("Performance risk markers detected.")

        if is_test_path(path):
            score -= 10
            reasons.append("Test file; lower production risk.")

        if score > 0:
            if score >= 45:
                impact = "HIGH"
            elif score >= 25:
                impact = "MEDIUM"
            else:
                impact = "LOW"

            rows.append((score, impact, path, reasons))

    rows.sort(
        key=lambda row: row[0],
        reverse=True
    )

    output = [
        "REFACTORING IMPACT MATRIX",
        "Mode: heuristic / review before changing files",
        f"Project: {project.get('name')}",
        "",
        "Files ranked by refactoring impact/risk:"
    ]

    if not rows:
        output.append("- No impact candidates detected.")
        return "\n".join(output)

    for score, impact, path, reasons in rows[:60]:
        output.append(f"- {impact} | score {score} | {path}")

        for reason in reasons[:4]:
            output.append(f"  - {reason}")

    output.append("")
    output.append("Rule:")
    output.append("- HIGH impact files need tests before refactor.")
    output.append("- MEDIUM impact files should be changed one at a time.")
    output.append("- LOW impact files can be cleaned after behavior is covered.")

    return "\n".join(output)


def enterprise_fix_next(project_name):
    return (
        "ENTERPRISE NEXT FIXES",
        "Mode: priority-based / no automatic changes",
        _enterprise_section("1. Highest Risk Vulnerabilities", highest_risk_vulnerabilities(project_name), max_chars=7000),
        _enterprise_section("2. Critical Code Without Tests", critical_code_without_tests(project_name), max_chars=7000),
        _enterprise_section("3. Refactoring Impact Matrix", refactoring_impact_matrix(project_name), max_chars=10000),
        _enterprise_section("4. Performance Bottlenecks", find_bottlenecks(project_name), max_chars=7000),
        _enterprise_section("5. Final Fix Recommendation", what_should_i_fix_next(project_name), max_chars=7000),
    )


def enterprise_fix_next_report(project_name):
    parts = enterprise_fix_next(project_name)

    if isinstance(parts, tuple):
        return "\n".join(str(part) for part in parts)

    return str(parts)


def architecture_refactoring_report(project_name):
    return (
        "ARCHITECTURE REFACTORING REPORT\n"
        "Mode: architecture-focused refactor planning\n"
        + _enterprise_section("1. Architecture Analyzer", strict_architecture_analyzer_project(project_name), max_chars=9000)
        + _enterprise_section("2. Architecture Migration Plan", generate_architecture_migration_plan(project_name), max_chars=12000)
        + _enterprise_section("3. Refactoring Impact Matrix", refactoring_impact_matrix(project_name), max_chars=12000)
        + _enterprise_section("4. Modernization Plan", generate_modernization_plan(project_name), max_chars=9000)
    )


def refactoring_command_center(project_name):
    return (
        "JARVIS REFACTORING COMMAND CENTER\n"
        "Mode: full enterprise refactoring intelligence\n"
        f"Project: {project_name}\n"
        + _enterprise_section("1. Enterprise Readiness", enterprise_readiness_report(project_name), max_chars=12000)
        + _enterprise_section("2. Technical Debt", technical_debt_report(project_name), max_chars=12000)
        + _enterprise_section("3. Maintainability", maintainability_report(project_name), max_chars=12000)
        + _enterprise_section("4. Refactor Project", refactor_project(project_name), max_chars=12000)
        + _enterprise_section("5. Impact Matrix", refactoring_impact_matrix(project_name), max_chars=12000)
        + _enterprise_section("6. Next Fixes", enterprise_fix_next_report(project_name), max_chars=9000)
    )


def export_enterprise_refactoring_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    os.makedirs("reports", exist_ok=True)

    safe_name = _safe_report_filename(project.get("name", project_name))
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    path = os.path.join(
        "reports",
        f"{safe_name}_enterprise_refactoring_{timestamp}.md"
    )

    content = (
        f"# Enterprise Refactoring Report: {project.get('name')}\n\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        "Generated by: JARVIS Engineering OS\n\n"
        "Mode: safe intelligence / no automatic code changes\n\n"
        "```text\n"
        + refactoring_command_center(project_name)
        + "\n```\n"
    )

    with open(path, "w", encoding="utf-8") as f:
        f.write(content)

    return f"{content}\n\nREPORT EXPORTED:\n{path}"


# Friendly aliases
def enterprise_refactoring_report(project_name):
    return refactoring_command_center(project_name)


def refactor_roadmap(project_name):
    return generate_refactoring_roadmap(project_name)


def architecture_migration(project_name):
    return generate_architecture_migration_plan(project_name)


def technical_debt_full_report(project_name):
    return technical_debt_report(project_name)


def maintainability_full_report(project_name):
    return maintainability_report(project_name)


def clean_code(project_name):
    return clean_code_report(project_name)


def enterprise_ready(project_name):
    return enterprise_readiness_report(project_name)


def what_should_i_fix_next_enterprise(project_name):
    return enterprise_fix_next_report(project_name)


def export_refactoring_report(project_name):
    return export_enterprise_refactoring_report(project_name)





# ==========================
# STEP 26 - CI/CD & DEPLOYMENT INTELLIGENCE
# ==========================
def ci_audit(project_name):
    project, facts, error = strict_project_facts(project_name)
    if error:
        return error

    output = [
        "CI/CD AUDIT REPORT",
        f"Project: {project_name}",
        ""
    ]

    if facts["ci_files"]:
        output.append("CI/CD files detected:")
        for item in facts["ci_files"][:50]:
            output.append(f" - {item['path']}")
    else:
        output.append("No CI/CD workflow files detected.")

    output.append("")
    output.append("Recommendations:")
    output.append("- Add automated testing before deployment.")
    output.append("- Add security scanning in CI.")
    output.append("- Add dependency vulnerability checks.")
    output.append("- Add build verification.")

    return "\\n".join(output)


def pipeline_audit(project_name):
    return (
        "PIPELINE AUDIT\\n\\n"
        + ci_audit(project_name)
        + "\\n\\nReview build, test, security, package and deploy stages."
    )


def deployment_audit(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    output = [
        "DEPLOYMENT AUDIT",
        "",
        "Deployment-related files:"
    ]

    docker_files = facts["docker_files"] + facts["config_files"]

    if docker_files:
        for item in docker_files[:60]:
            output.append(f" - {item['path']}")
    else:
        output.append("No deployment files detected.")

    output.append("")
    output.append("Deployment recommendations:")
    output.append("- Environment variables for secrets.")
    output.append("- Health checks.")
    output.append("- Backup strategy.")
    output.append("- Rollback strategy.")
    output.append("- Monitoring and logging.")

    return "\\n".join(output)


def github_actions_audit(project_name):
    return ci_audit(project_name)


def docker_audit(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    output = ["DOCKER AUDIT", ""]

    if facts["docker_files"]:
        for item in facts["docker_files"]:
            output.append(f" - {item['path']}")
    else:
        output.append("No Docker files detected.")

    output.append("")
    output.append("Verify:")
    output.append("- Non-root containers")
    output.append("- Small base images")
    output.append("- Secrets not hardcoded")
    output.append("- Health checks enabled")

    return "\\n".join(output)


def devops_readiness(project_name):
    return (
        "DEVOPS READINESS REPORT\\n\\n"
        + ci_audit(project_name)
        + "\\n\\n"
        + deployment_audit(project_name)
    )



# ==========================
# STEP 27 - AI CODE EXPLANATION & TEACHING ENGINE
# Explain files, functions, classes, architecture, and project code at multiple levels.
# Evidence-based. No automatic code changes.
# ==========================
def _explain_section(title, content, max_chars=9000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 76
        + f"\n{title}\n"
        + "=" * 76
        + "\n"
        + text
    )


def _find_file_in_project(project_name, file_query):
    project, error = get_project(project_name)

    if error:
        return None, None, error

    files = get_code_files(project)
    query = normalize_path(file_query).strip()

    # Exact path match first.
    for item in files:
        if normalize_path(item["path"]) == query:
            return project, item, None

    # Endswith match.
    matches = [
        item for item in files
        if normalize_path(item["path"]).endswith(query)
    ]

    if len(matches) == 1:
        return project, matches[0], None

    if len(matches) > 1:
        output = [
            f"Multiple files matched: {file_query}",
            "Please use a more specific path:",
            ""
        ]

        for item in matches[:30]:
            output.append(f"- {item['path']}")

        return project, None, "\n".join(output)

    # Fuzzy basename/path contains match.
    matches = [
        item for item in files
        if query in normalize_path(item["path"])
    ]

    if len(matches) == 1:
        return project, matches[0], None

    if len(matches) > 1:
        output = [
            f"Multiple files matched: {file_query}",
            "Please use a more specific path:",
            ""
        ]

        for item in matches[:30]:
            output.append(f"- {item['path']}")

        return project, None, "\n".join(output)

    return project, None, f"File not found in indexed project: {file_query}"


def _extract_file_symbols(item):
    path = item.get("path", "")
    ext = item.get("extension", "")
    content = item.get("content", "")

    if ext == ".py":
        return extract_python_symbols(path, content)

    if ext in {".js", ".jsx", ".ts", ".tsx"}:
        return extract_js_symbols(content)

    return {
        "functions": [],
        "classes": [],
        "imports": []
    }


def _find_function_or_class_in_project(project_name, symbol_name):
    project, error = get_project(project_name)

    if error:
        return None, None, None, error

    files = get_code_files(project)
    target = symbol_name.strip().lower()
    matches = []

    for item in files:
        symbols = _extract_file_symbols(item)

        for func in symbols.get("functions", []):
            if func.lower() == target:
                matches.append(("function", func, item))

        for cls in symbols.get("classes", []):
            if cls.lower() == target:
                matches.append(("class", cls, item))

    if len(matches) == 1:
        kind, name, item = matches[0]
        return project, item, kind, None

    if len(matches) > 1:
        output = [
            f"Multiple symbols matched: {symbol_name}",
            "Please specify the file too:",
            ""
        ]

        for kind, name, item in matches[:30]:
            output.append(f"- {kind} {name} in {item['path']}")

        return project, None, None, "\n".join(output)

    # Fuzzy search in content.
    for item in files:
        if target in item.get("content", "").lower():
            matches.append(("content", symbol_name, item))

    if len(matches) == 1:
        kind, name, item = matches[0]
        return project, item, kind, None

    if len(matches) > 1:
        output = [
            f"Symbol/content matched multiple files: {symbol_name}",
            "Please specify the file:",
            ""
        ]

        for kind, name, item in matches[:30]:
            output.append(f"- {item['path']}")

        return project, None, None, "\n".join(output)

    return project, None, None, f"Function/class not found in indexed project: {symbol_name}"


def _extract_relevant_symbol_snippet(item, symbol_name, context_lines=35):
    content = item.get("content", "")
    lines = content.splitlines()
    target = symbol_name.strip()

    if not target:
        return content[:6000]

    patterns = [
        rf"^\s*def\s+{re.escape(target)}\s*\(",
        rf"^\s*async\s+def\s+{re.escape(target)}\s*\(",
        rf"^\s*class\s+{re.escape(target)}\b",
        rf"^\s*function\s+{re.escape(target)}\s*\(",
        rf"^\s*const\s+{re.escape(target)}\s*=",
        rf"^\s*export\s+function\s+{re.escape(target)}\s*\(",
        rf"^\s*export\s+const\s+{re.escape(target)}\s*=",
    ]

    for index, line in enumerate(lines):
        if any(re.search(pattern, line) for pattern in patterns):
            start = max(0, index - 8)
            end = min(len(lines), index + context_lines)

            return "\n".join(
                f"{i + 1}: {lines[i]}"
                for i in range(start, end)
            )

    # fallback content search
    for index, line in enumerate(lines):
        if target.lower() in line.lower():
            start = max(0, index - 10)
            end = min(len(lines), index + context_lines)

            return "\n".join(
                f"{i + 1}: {lines[i]}"
                for i in range(start, end)
            )

    return content[:6000]


def explain_file(project_name, file_query, level="intermediate"):
    project, item, error = _find_file_in_project(project_name, file_query)

    if error:
        return error

    symbols = _extract_file_symbols(item)
    related = []

    path = item["path"]
    lower_path = normalize_path(path)
    files = get_code_files(project)

    basename = file_stem(path)

    for other in files:
        if other["path"] == path:
            continue

        other_text = (
            normalize_path(other["path"])
            + "\n"
            + other.get("content", "").lower()
        )

        if basename in other_text:
            related.append(other["path"])

    prompt = f"""
You are JARVIS, a senior software mentor.

Explain this file at {level} level.

Grounding rules:
- Use only the file content and project evidence below.
- Do not invent behavior.
- Mention the exact file path.
- If something is not visible, say: Not visible in indexed files.
- Keep the explanation practical and easy to understand.

PROJECT:
{project.get("name")}

FILE PATH:
{item["path"]}

FILE EXTENSION:
{item["extension"]}

DETECTED SYMBOLS:
{symbols}

RELATED FILES BY SIMPLE REFERENCE:
{related[:25]}

FILE CONTENT:
{item["content"][:12000]}

Return:
1. What this file does
2. Why it exists in the project
3. Important functions/classes/components
4. Inputs and outputs
5. Dependencies/imports and related files
6. Security concerns if any
7. Performance concerns if any
8. How to explain this file in an interview
9. Beginner-friendly summary
"""

    return ask_llm(prompt)


def explain_function(project_name, symbol_name, level="intermediate"):
    project, item, kind, error = _find_function_or_class_in_project(
        project_name,
        symbol_name
    )

    if error:
        return error

    snippet = _extract_relevant_symbol_snippet(
        item,
        symbol_name
    )

    prompt = f"""
You are JARVIS, a senior software mentor.

Explain this {kind} at {level} level.

Grounding rules:
- Use only the indexed file content below.
- Do not invent runtime behavior.
- Mention the exact file path.
- If something is not visible, say: Not visible in indexed files.
- Explain simply, then technically.

PROJECT:
{project.get("name")}

FILE PATH:
{item["path"]}

SYMBOL:
{symbol_name}

SNIPPET:
{snippet}

Return:
1. What this {kind} does
2. Parameters / input
3. Return value / output
4. Step-by-step logic
5. Where it is likely used based on visible code
6. Risks or edge cases
7. How to test it
8. How to explain it in an interview
"""

    return ask_llm(prompt)


def explain_class(project_name, class_name, level="intermediate"):
    return explain_function(
        project_name,
        class_name,
        level=level
    )


def explain_architecture(project_name, level="intermediate"):
    project, error = get_project(project_name)

    if error:
        return error

    evidence = (
        project_overview(project)
        + "\n\n"
        + strict_grounded_analyzer_project(project_name)
        + "\n\n"
        + strict_architecture_analyzer_project(project_name)
    )

    prompt = f"""
You are JARVIS, a senior software architect and teacher.

Explain the architecture of this project at {level} level.

Grounding rules:
- Use only indexed project evidence below.
- Mention exact files when making concrete claims.
- Do not invent cloud services, databases, APIs, or features.
- If something is not visible, say: Not visible in indexed files.

PROJECT:
{project.get("name")}

EVIDENCE:
{evidence}

Return:
1. Big-picture explanation
2. Main layers/modules
3. Frontend explanation
4. Backend explanation
5. Data/config/storage explanation
6. Security layer explanation
7. How files communicate
8. What to study first
9. Interview explanation version
"""

    return ask_llm(prompt)


def teach_me_project(project_name, level="beginner"):
    project, error = get_project(project_name)

    if error:
        return error

    return (
        "PROJECT TEACHING GUIDE\n"
        f"Project: {project.get('name')}\n"
        f"Level: {level}\n"
        + _explain_section(
            "1. Project Overview",
            project_overview(project),
            max_chars=7000
        )
        + _explain_section(
            "2. Architecture Explanation",
            explain_architecture(project_name, level=level),
            max_chars=12000
        )
        + _explain_section(
            "3. Important Files",
            project_evidence_report(project_name),
            max_chars=9000
        )
        + _explain_section(
            "4. Security Concepts To Learn",
            strict_security_analyzer_project(project_name),
            max_chars=9000
        )
        + _explain_section(
            "5. What To Learn Next",
            what_should_i_fix_next(project_name),
            max_chars=9000
        )
    )


def beginner_explanation(project_name, file_query):
    return explain_file(
        project_name,
        file_query,
        level="beginner"
    )


def intermediate_explanation(project_name, file_query):
    return explain_file(
        project_name,
        file_query,
        level="intermediate"
    )


def senior_explanation(project_name, file_query):
    return explain_file(
        project_name,
        file_query,
        level="senior"
    )


def what_does_this_file_do(project_name, file_query):
    return explain_file(
        project_name,
        file_query,
        level="beginner"
    )


def explain_file_for_interview(project_name, file_query):
    explanation = explain_file(
        project_name,
        file_query,
        level="interview"
    )

    return (
        "INTERVIEW FILE EXPLANATION\n\n"
        + str(explanation)
    )


def explain_function_for_interview(project_name, symbol_name):
    explanation = explain_function(
        project_name,
        symbol_name,
        level="interview"
    )

    return (
        "INTERVIEW FUNCTION EXPLANATION\n\n"
        + str(explanation)
    )


def code_teaching_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = select_context_files(
        get_code_files(project),
        limit=12
    )

    output = [
        "CODE TEACHING REPORT",
        f"Project: {project.get('name')}",
        "",
        "Recommended learning order:"
    ]

    for index, item in enumerate(files, start=1):
        output.append(f"{index}. {item['path']}")

    output.append("")
    output.append("How to use:")
    output.append("- explain file <project> <file>")
    output.append("- explain function <project> <function>")
    output.append("- beginner explanation <project> <file>")
    output.append("- senior explanation <project> <file>")
    output.append("- explain architecture <project>")
    output.append("- teach me project <project>")

    return "\n".join(output)


# Friendly aliases
def explain_project_architecture(project_name):
    return explain_architecture(project_name)


def teach_project(project_name):
    return teach_me_project(project_name)


def teach_me_code(project_name):
    return code_teaching_report(project_name)


def explain_code(project_name):
    return code_teaching_report(project_name)



# ==========================
# STEP 28 - UML & ARCHITECTURE DIAGRAM INTELLIGENCE
# Generates Mermaid diagrams: class, component, architecture, dependency graph, data flow.
# Evidence-based. No automatic code changes.
# ==========================
DIAGRAM_OUTPUT_DIR = "diagrams_generated"


def _diagram_safe_filename(name):
    cleaned = re.sub(
        r"[^A-Za-z0-9_.-]+",
        "_",
        str(name).strip()
    ).strip("_")

    return cleaned or "diagram"


def _save_diagram(project_name, diagram_type, content):
    os.makedirs(
        DIAGRAM_OUTPUT_DIR,
        exist_ok=True
    )

    safe_project = _diagram_safe_filename(project_name)
    safe_type = _diagram_safe_filename(diagram_type)
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    path = os.path.join(
        DIAGRAM_OUTPUT_DIR,
        f"{safe_project}_{safe_type}_{timestamp}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _diagram_header(title, project_name):
    return (
        f"# {title}\n\n"
        f"Project: `{project_name}`\n\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        "Generated by: JARVIS Engineering OS\n\n"
        "Mode: evidence-based diagram generation / no automatic code changes\n\n"
    )


def _mermaid_block(code):
    return (
        "```mermaid\n"
        + code.strip()
        + "\n```\n"
    )


def _clean_node_id(value):
    cleaned = re.sub(
        r"[^A-Za-z0-9_]+",
        "_",
        str(value)
    ).strip("_")

    if not cleaned:
        cleaned = "Node"

    if cleaned[0].isdigit():
        cleaned = "N_" + cleaned

    return cleaned[:80]


def _short_label(value, max_len=55):
    text = str(value).replace('"', "'")

    if len(text) > max_len:
        text = text[:max_len - 3] + "..."

    return text


def _extract_python_class_details(item):
    content = item.get("content", "")
    classes = []

    try:
        tree = ast.parse(content)

        for node in ast.walk(tree):
            if isinstance(node, ast.ClassDef):
                methods = []

                for child in node.body:
                    if isinstance(child, (ast.FunctionDef, ast.AsyncFunctionDef)):
                        methods.append(child.name)

                bases = []

                for base in node.bases:
                    if isinstance(base, ast.Name):
                        bases.append(base.id)
                    elif isinstance(base, ast.Attribute):
                        bases.append(base.attr)

                classes.append({
                    "name": node.name,
                    "methods": methods,
                    "bases": bases,
                    "path": item.get("path", "")
                })

    except Exception:
        pass

    return classes


def _extract_ts_class_details(item):
    content = item.get("content", "")
    path = item.get("path", "")
    classes = []

    for class_match in re.finditer(
        r"(?:export\s+)?class\s+([A-Za-z_][A-Za-z0-9_]*)"
        r"(?:\s+extends\s+([A-Za-z_][A-Za-z0-9_]*))?",
        content
    ):
        name = class_match.group(1)
        base = class_match.group(2)

        # Simple method extraction from the whole file.
        methods = re.findall(
            r"^\s*(?:public\s+|private\s+|protected\s+)?"
            r"([A-Za-z_][A-Za-z0-9_]*)\s*\(",
            content,
            flags=re.MULTILINE
        )

        methods = [
            method for method in methods
            if method not in {"if", "for", "while", "switch", "catch"}
        ]

        classes.append({
            "name": name,
            "methods": sorted(set(methods))[:20],
            "bases": [base] if base else [],
            "path": path
        })

    return classes


def _collect_class_details(project_name):
    project, error = get_project(project_name)

    if error:
        return None, error

    files = get_code_files(project)
    classes = []

    for item in files:
        ext = item.get("extension", "")

        if ext == ".py":
            classes.extend(_extract_python_class_details(item))

        elif ext in {".ts", ".tsx", ".js", ".jsx"}:
            classes.extend(_extract_ts_class_details(item))

    return {
        "project": project,
        "classes": classes
    }, None


def generate_class_diagram(project_name):
    data, error = _collect_class_details(project_name)

    if error:
        return error

    project = data["project"]
    classes = data["classes"]

    lines = [
        "classDiagram"
    ]

    if not classes:
        lines.append('    class NoClassesDetected')
    else:
        for cls in classes[:60]:
            class_id = _clean_node_id(cls["name"])
            lines.append(f"    class {class_id} {{")

            for method in cls.get("methods", [])[:12]:
                lines.append(f"        +{method}()")

            lines.append("    }")

            for base in cls.get("bases", [])[:5]:
                base_id = _clean_node_id(base)
                lines.append(f"    {base_id} <|-- {class_id}")

    content = (
        _diagram_header(
            "Class Diagram",
            project.get("name", project_name)
        )
        + _mermaid_block("\n".join(lines))
        + "\n\n## Evidence\n\n"
    )

    if classes:
        for cls in classes[:80]:
            content += (
                f"- `{cls['name']}` from `{cls['path']}`"
                + (
                    f" extends `{', '.join(cls['bases'])}`"
                    if cls.get("bases")
                    else ""
                )
                + "\n"
            )
    else:
        content += "- No classes detected in indexed files.\n"

    path = _save_diagram(
        project.get("name", project_name),
        "class_diagram",
        content
    )

    return f"{content}\n\nDIAGRAM EXPORTED:\n{path}"


def _folder_component_name(path):
    normalized = normalize_path(path)
    parts = [
        part for part in normalized.split("/")
        if part and "." not in part
    ]

    if not parts:
        return "root"

    important = [
        "frontend",
        "backend",
        "src",
        "components",
        "routes",
        "api",
        "models",
        "schemas",
        "services",
        "utils",
        "config",
        "tests",
        "docs",
        ".github",
        "workflows"
    ]

    for part in parts:
        if part in important:
            return part

    return parts[0]


def generate_component_diagram(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    components = defaultdict(list)

    for item in files:
        component = _folder_component_name(item["path"])
        components[component].append(item["path"])

    lines = [
        "flowchart TD"
    ]

    project_node = _clean_node_id(project.get("name", project_name))
    lines.append(f'    {project_node}["{_short_label(project.get("name", project_name))}"]')

    for component, paths in sorted(components.items()):
        node = _clean_node_id(component)
        label = f"{component} ({len(paths)} files)"
        lines.append(f'    {project_node} --> {node}["{_short_label(label)}"]')

    # Evidence-based common relations.
    relation_pairs = [
        ("frontend", "api"),
        ("frontend", "routes"),
        ("components", "api"),
        ("routes", "services"),
        ("api", "services"),
        ("services", "models"),
        ("routes", "models"),
        ("api", "models"),
        ("services", "database"),
        ("models", "database"),
        ("routes", "schemas"),
        ("api", "schemas"),
        ("workflows", "tests"),
    ]

    existing = set(components.keys())

    for source, target in relation_pairs:
        if source in existing and target in existing:
            lines.append(f"    {_clean_node_id(source)} -.-> {_clean_node_id(target)}")

    content = (
        _diagram_header(
            "Component Diagram",
            project.get("name", project_name)
        )
        + _mermaid_block("\n".join(lines))
        + "\n\n## Components detected\n\n"
    )

    for component, paths in sorted(components.items()):
        content += f"\n### {component}\n"
        for path_item in paths[:25]:
            content += f"- `{path_item}`\n"

    path = _save_diagram(
        project.get("name", project_name),
        "component_diagram",
        content
    )

    return f"{content}\n\nDIAGRAM EXPORTED:\n{path}"


def generate_architecture_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "flowchart LR",
        '    User["User"]'
    ]

    nodes = {}

    def add_node(key, label):
        node = _clean_node_id(key)
        nodes[key] = node
        lines.append(f'    {node}["{_short_label(label)}"]')
        return node

    if facts["frontend_files"]:
        add_node("frontend", f"Frontend/UI ({len(facts['frontend_files'])} files)")

    if facts["fastapi_files"] or facts["flask_files"]:
        backend_count = len(facts["fastapi_files"]) + len(facts["flask_files"])
        add_node("backend", f"Backend/API ({backend_count} files)")

    if facts["auth_files"] or facts["jwt_files"]:
        add_node("auth", "Auth/Security Layer")

    if facts["database_files"]:
        add_node("database", f"Database/Storage ({len(facts['database_files'])} files)")

    if facts["upload_files"]:
        add_node("upload", "Upload/File Handling")

    if facts["docker_files"]:
        add_node("deployment", "Docker/Deployment")

    if facts["ci_files"]:
        add_node("ci", "CI/CD Workflows")

    if "frontend" in nodes:
        lines.append(f"    User --> {nodes['frontend']}")

    if "frontend" in nodes and "backend" in nodes:
        lines.append(f"    {nodes['frontend']} --> {nodes['backend']}")
    elif "backend" in nodes:
        lines.append(f"    User --> {nodes['backend']}")

    if "backend" in nodes and "auth" in nodes:
        lines.append(f"    {nodes['backend']} --> {nodes['auth']}")

    if "backend" in nodes and "database" in nodes:
        lines.append(f"    {nodes['backend']} --> {nodes['database']}")

    if "backend" in nodes and "upload" in nodes:
        lines.append(f"    {nodes['backend']} --> {nodes['upload']}")

    if "ci" in nodes and "deployment" in nodes:
        lines.append(f"    {nodes['ci']} --> {nodes['deployment']}")

    if "deployment" in nodes and "backend" in nodes:
        lines.append(f"    {nodes['deployment']} -. deploys .-> {nodes['backend']}")

    content = (
        _diagram_header(
            "Architecture Diagram",
            project.get("name", project_name)
        )
        + _mermaid_block("\n".join(lines))
        + "\n\n## Grounded architecture evidence\n\n"
        + "```text\n"
        + strict_architecture_analyzer_project(project_name)
        + "\n```\n"
    )

    path = _save_diagram(
        project.get("name", project_name),
        "architecture_diagram",
        content
    )

    return f"{content}\n\nDIAGRAM EXPORTED:\n{path}"


def _extract_import_edges(project_name):
    project, error = get_project(project_name)

    if error:
        return None, error

    files = get_code_files(project)
    stems_to_paths = defaultdict(list)

    for item in files:
        stems_to_paths[file_stem(item["path"])].append(item["path"])

    edges = set()

    for item in files:
        path = item["path"]
        content = item.get("content", "")
        ext = item.get("extension", "")

        imported_names = []

        if ext == ".py":
            symbols = extract_python_symbols(path, content)
            for imp in symbols.get("imports", []):
                imported_names.append(imp.split(".")[-1])

        elif ext in {".js", ".jsx", ".ts", ".tsx"}:
            imported_names.extend(
                re.findall(
                    r"import\s+.*?\s+from\s+[\"']([^\"']+)[\"']",
                    content
                )
            )
            imported_names.extend(
                re.findall(
                    r"require\(\s*[\"']([^\"']+)[\"']\s*\)",
                    content
                )
            )

        for imported in imported_names:
            imported_stem = file_stem(imported)

            if imported_stem in stems_to_paths:
                for target_path in stems_to_paths[imported_stem][:3]:
                    if target_path != path:
                        edges.add((path, target_path))

    return {
        "project": project,
        "edges": sorted(edges)
    }, None


def generate_dependency_graph(project_name):
    data, error = _extract_import_edges(project_name)

    if error:
        return error

    project = data["project"]
    edges = data["edges"]

    lines = [
        "flowchart TD"
    ]

    if not edges:
        lines.append('    NoDependencyEdges["No import dependency edges detected"]')
    else:
        used_nodes = set()

        for source, target in edges[:120]:
            source_id = _clean_node_id(source)
            target_id = _clean_node_id(target)

            if source_id not in used_nodes:
                lines.append(f'    {source_id}["{_short_label(source)}"]')
                used_nodes.add(source_id)

            if target_id not in used_nodes:
                lines.append(f'    {target_id}["{_short_label(target)}"]')
                used_nodes.add(target_id)

            lines.append(f"    {source_id} --> {target_id}")

    content = (
        _diagram_header(
            "Dependency Graph",
            project.get("name", project_name)
        )
        + _mermaid_block("\n".join(lines))
        + "\n\n## Dependency edges detected\n\n"
    )

    if edges:
        for source, target in edges[:150]:
            content += f"- `{source}` -> `{target}`\n"
    else:
        content += "- No import dependency edges detected by current rules.\n"

    path = _save_diagram(
        project.get("name", project_name),
        "dependency_graph",
        content
    )

    return f"{content}\n\nDIAGRAM EXPORTED:\n{path}"


def generate_data_flow_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "flowchart LR",
        '    User["User"]'
    ]

    if facts["frontend_files"]:
        lines.append('    Frontend["Frontend/UI"]')
        lines.append("    User --> Frontend")

    if facts["frontend_api_calls"]:
        lines.append('    ApiCalls["Frontend API Calls"]')
        if facts["frontend_files"]:
            lines.append("    Frontend --> ApiCalls")
        else:
            lines.append("    User --> ApiCalls")

    if facts["routes"]:
        lines.append('    BackendRoutes["Backend Routes/API"]')
        if facts["frontend_api_calls"]:
            lines.append("    ApiCalls --> BackendRoutes")
        elif facts["frontend_files"]:
            lines.append("    Frontend --> BackendRoutes")
        else:
            lines.append("    User --> BackendRoutes")

    if facts["auth_files"] or facts["jwt_files"]:
        lines.append('    Auth["Auth/JWT/Security"]')
        if facts["routes"]:
            lines.append("    BackendRoutes --> Auth")
        else:
            lines.append("    User --> Auth")

    if facts["upload_files"]:
        lines.append('    Uploads["Upload/File Handling"]')
        if facts["routes"]:
            lines.append("    BackendRoutes --> Uploads")
        else:
            lines.append("    User --> Uploads")

    if facts["database_files"]:
        lines.append('    Storage["Database/Storage"]')
        if facts["routes"]:
            lines.append("    BackendRoutes --> Storage")
        elif facts["auth_files"]:
            lines.append("    Auth --> Storage")

    if facts["logging_files"]:
        lines.append('    Logs["Logging/Audit"]')
        if facts["routes"]:
            lines.append("    BackendRoutes --> Logs")
        if facts["auth_files"]:
            lines.append("    Auth --> Logs")
        if facts["upload_files"]:
            lines.append("    Uploads --> Logs")

    content = (
        _diagram_header(
            "Data Flow Diagram",
            project.get("name", project_name)
        )
        + _mermaid_block("\n".join(lines))
        + "\n\n## API / route evidence\n\n"
    )

    if facts["frontend_api_calls"]:
        content += "\n### Frontend API calls\n"
        for call in facts["frontend_api_calls"][:80]:
            content += f"- `{call}`\n"

    if facts["routes"]:
        content += "\n### Backend routes\n"
        for route in facts["routes"][:80]:
            content += f"- `{route}`\n"

    if not facts["frontend_api_calls"] and not facts["routes"]:
        content += "- No explicit frontend API calls or backend routes detected.\n"

    path = _save_diagram(
        project.get("name", project_name),
        "data_flow_diagram",
        content
    )

    return f"{content}\n\nDIAGRAM EXPORTED:\n{path}"


def generate_full_diagram_pack(project_name):
    outputs = [
        generate_class_diagram(project_name),
        generate_component_diagram(project_name),
        generate_architecture_diagram(project_name),
        generate_dependency_graph(project_name),
        generate_data_flow_diagram(project_name),
    ]

    exported = []

    for output in outputs:
        marker = "DIAGRAM EXPORTED:"
        if marker in output:
            exported.append(output.split(marker)[-1].strip())

    return (
        "FULL UML / ARCHITECTURE DIAGRAM PACK GENERATED\n\n"
        + "\n".join(f"- {path}" for path in exported)
    )


# Friendly aliases
def class_diagram(project_name):
    return generate_class_diagram(project_name)


def component_diagram(project_name):
    return generate_component_diagram(project_name)


def architecture_diagram(project_name):
    return generate_architecture_diagram(project_name)


def dependency_graph(project_name):
    return generate_dependency_graph(project_name)


def data_flow_diagram(project_name):
    return generate_data_flow_diagram(project_name)


def diagram_pack(project_name):
    return generate_full_diagram_pack(project_name)



# ==========================
# STEP 29 - RELEASE READINESS INTELLIGENCE
# Release / production / deployment / enterprise go-live scoring.
# Evidence-based. No automatic code changes.
# ==========================
def _release_section(title, content, max_chars=9000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... section truncated"

    return (
        "\n\n"
        + "=" * 76
        + f"\n{title}\n"
        + "=" * 76
        + "\n"
        + text
    )


def _release_score_from_facts(facts):
    score = 0
    positives = []
    blockers = []
    warnings = []

    if facts.get("test_files"):
        score += 14
        positives.append("Tests detected.")
    else:
        blockers.append("No test files detected.")

    if facts.get("ci_files"):
        score += 12
        positives.append("CI/workflow files detected.")
    else:
        warnings.append("No CI/workflow files detected.")

    if facts.get("docker_files"):
        score += 10
        positives.append("Docker/deployment files detected.")
    else:
        warnings.append("No Docker/deployment files detected.")

    if facts.get("auth_files"):
        score += 10
        positives.append("Authentication/authorization files detected.")

    if facts.get("password_hashing_files"):
        score += 10
        positives.append("Password hashing evidence detected.")
    elif facts.get("auth_files"):
        blockers.append("Auth files exist, but no password hashing evidence detected.")

    if facts.get("secret_files"):
        score += 8
        positives.append("Config/secret handling evidence detected.")

    if facts.get("database_files"):
        score += 8
        positives.append("Database/storage evidence detected.")

    if facts.get("logging_files"):
        score += 8
        positives.append("Logging/audit evidence detected.")
    else:
        warnings.append("No logging/audit evidence detected.")

    if facts.get("packages") or facts.get("requirements"):
        score += 8
        positives.append("Dependency manifests detected.")
    else:
        warnings.append("No dependency manifest detected.")

    if facts.get("routes") or facts.get("frontend_api_calls"):
        score += 6
        positives.append("API route/API call evidence detected.")

    weak_terms = [
        "your-secret-key",
        "changeme",
        "changeme123",
        "password123",
        "secure123",
        "admin123",
    ]

    for item in facts.get("secret_files", []) + facts.get("auth_files", []):
        content = item.get("content", "").lower()

        for term in weak_terms:
            if term in content:
                blockers.append(
                    f"Weak/demo secret or credential detected in {item['path']}: {term}"
                )

    score -= min(25, len(blockers) * 5)
    score = max(0, min(100, score))

    if score >= 85 and not blockers:
        level = "GO"
    elif score >= 70 and len(blockers) <= 1:
        level = "CONDITIONAL_GO"
    elif score >= 50:
        level = "NOT_READY_YET"
    else:
        level = "BLOCKED"

    return {
        "score": score,
        "level": level,
        "positives": positives,
        "warnings": warnings,
        "blockers": blockers,
    }


def release_readiness(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    result = _release_score_from_facts(facts)

    output = [
        "RELEASE READINESS REPORT",
        "Mode: rule-based / verify manually before real release",
        f"Project: {project.get('name')}",
        f"Path: {project.get('path')}",
        "",
        f"Release score: {result['score']}/100",
        f"Release decision: {result['level']}",
        "",
        "Positive evidence:"
    ]

    if result["positives"]:
        output.extend(f"- {item}" for item in result["positives"])
    else:
        output.append("- None detected.")

    output.append("")
    output.append("Warnings:")

    if result["warnings"]:
        output.extend(f"- {item}" for item in result["warnings"])
    else:
        output.append("- No major warnings detected.")

    output.append("")
    output.append("Blockers:")

    if result["blockers"]:
        output.extend(f"- {item}" for item in result["blockers"])
    else:
        output.append("- No release blockers detected by current rules.")

    output.append("")
    output.append("Required before release:")
    output.append("- Run tests locally.")
    output.append("- Run dependency audit.")
    output.append("- Run security audit.")
    output.append("- Confirm environment variables and secrets.")
    output.append("- Confirm backup and rollback strategy.")
    output.append("- Confirm logging and monitoring.")
    output.append("- Confirm deployment documentation.")

    return "\n".join(output)


def production_readiness(project_name):
    return (
        "PRODUCTION READINESS INTELLIGENCE\n"
        "Mode: consolidated production check\n"
        + _release_section("1. Release Readiness", release_readiness(project_name), max_chars=9000)
        + _release_section("2. Production Readiness Estimate", estimate_production_readiness(project_name), max_chars=9000)
        + _release_section("3. Security Audit", full_security_audit(project_name), max_chars=9000)
        + _release_section("4. Test Intelligence", test_intelligence_report(project_name), max_chars=9000)
        + _release_section("5. Performance Intelligence", performance_report(project_name), max_chars=9000)
        + _release_section("6. CI/CD Audit", ci_audit(project_name) if "ci_audit" in globals() else "CI/CD intelligence not loaded.", max_chars=9000)
        + _release_section("7. Deployment Audit", deployment_audit(project_name) if "deployment_audit" in globals() else "Deployment intelligence not loaded.", max_chars=9000)
    )


def deployment_readiness(project_name):
    return (
        "DEPLOYMENT READINESS REPORT\n"
        "Mode: safe deployment validation checklist\n"
        + _release_section("1. Deployment Audit", deployment_audit(project_name) if "deployment_audit" in globals() else "Deployment audit not available.", max_chars=9000)
        + _release_section("2. Docker Audit", docker_audit(project_name) if "docker_audit" in globals() else "Docker audit not available.", max_chars=9000)
        + _release_section("3. CI/CD Audit", ci_audit(project_name) if "ci_audit" in globals() else "CI audit not available.", max_chars=9000)
        + _release_section("4. Release Checklist", generate_release_checklist(project_name), max_chars=9000)
        + _release_section("5. Deployment Checklist", generate_deployment_checklist(project_name), max_chars=9000)
    )


def staging_readiness(project_name):
    return (
        "STAGING READINESS REPORT\n\n"
        "Before staging, the project should have:\n"
        "- Clean local run.\n"
        "- Test command documented.\n"
        "- Environment variables separated from source code.\n"
        "- Test/staging database separated from production.\n"
        "- Logging enabled.\n"
        "- Basic rollback plan.\n"
        + _release_section("Release Readiness", release_readiness(project_name), max_chars=9000)
        + _release_section("QA Plan", generate_qa_plan(project_name), max_chars=9000)
    )


def go_live_report(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    release = _release_score_from_facts(facts)

    decision = release["level"]

    if decision == "GO":
        recommendation = "GO: Project looks ready by static rules, but still run manual production validation."
    elif decision == "CONDITIONAL_GO":
        recommendation = "CONDITIONAL GO: Fix warnings/blockers and run full validation before launch."
    elif decision == "NOT_READY_YET":
        recommendation = "NO-GO: Project needs more test/security/deployment work before go-live."
    else:
        recommendation = "BLOCKED: Critical release blockers detected."

    return (
        "GO-LIVE REPORT\n"
        "Mode: executive release decision support\n"
        f"Project: {project.get('name')}\n\n"
        f"Final release score: {release['score']}/100\n"
        f"Decision: {decision}\n"
        f"Recommendation: {recommendation}\n"
        + _release_section("1. Release Readiness", release_readiness(project_name), max_chars=9000)
        + _release_section("2. Enterprise Readiness", enterprise_readiness_report(project_name), max_chars=10000)
        + _release_section("3. Production Readiness", production_readiness(project_name), max_chars=12000)
        + _release_section("4. Final Fixes", enterprise_fix_next_report(project_name), max_chars=9000)
    )


def release_candidate_report(project_name):
    return (
        "RELEASE CANDIDATE REPORT\n"
        "Mode: release candidate validation\n"
        + _release_section("1. Go Live Report", go_live_report(project_name), max_chars=12000)
        + _release_section("2. Technical Debt", technical_debt_report(project_name), max_chars=9000)
        + _release_section("3. Dependency / Package Evidence", strict_grounded_analyzer_project(project_name), max_chars=9000)
        + _release_section("4. Diagram Pack Recommendation", "Run: generate full diagram pack <project> before final presentation.", max_chars=2000)
    )


def export_release_readiness_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    os.makedirs("reports", exist_ok=True)

    safe_name = _safe_report_filename(project.get("name", project_name))
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    path = os.path.join(
        "reports",
        f"{safe_name}_release_readiness_{timestamp}.md"
    )

    content = (
        f"# Release Readiness Report: {project.get('name')}\n\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        "Generated by: JARVIS Engineering OS\n\n"
        "```text\n"
        + release_candidate_report(project_name)
        + "\n```\n"
    )

    with open(path, "w", encoding="utf-8") as f:
        f.write(content)

    return f"{content}\n\nREPORT EXPORTED:\n{path}"


# Friendly aliases
def release_ready(project_name):
    return release_readiness(project_name)


def production_ready(project_name):
    return production_readiness(project_name)


def deployment_ready(project_name):
    return deployment_readiness(project_name)


def staging_ready(project_name):
    return staging_readiness(project_name)


def go_live(project_name):
    return go_live_report(project_name)


def release_candidate(project_name):
    return release_candidate_report(project_name)


def export_release_report(project_name):
    return export_release_readiness_report(project_name)



# ==========================
# STEP 34 - AUTO README GENERATOR
# Generates professional README.md content for any indexed project.
# Evidence-based. Safe export only. No automatic overwrite of project files.
# ==========================
README_OUTPUT_DIR = "generated_readmes"


def _readme_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _readme_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _readme_safe_filename(name):
    return _safe_report_filename(str(name))


def _readme_section(title, content):
    text = str(content).strip()

    if not text:
        text = "Not visible in indexed files."

    return f"\n\n## {title}\n\n{text}\n"


def _detect_run_commands(project, files):
    commands = []

    for item in files:
        path = normalize_path(item["path"])
        content = item["content"]

        if path.endswith("package.json"):
            try:
                data = json.loads(content)
                scripts = data.get("scripts", {})

                if isinstance(scripts, dict):
                    if "dev" in scripts:
                        commands.append("npm run dev")
                    if "start" in scripts:
                        commands.append("npm start")
                    if "build" in scripts:
                        commands.append("npm run build")
                    if "test" in scripts:
                        commands.append("npm test")
            except Exception:
                pass

        if path.endswith("requirements.txt"):
            commands.append("python -m pip install -r requirements.txt")

        if path.endswith("app.py"):
            commands.append("python app.py")

        if path.endswith("main.py"):
            commands.append("python main.py")

        if "fastapi" in content.lower() and item["extension"] == ".py":
            commands.append("uvicorn main:app --reload")

        if "flask" in content.lower() and item["extension"] == ".py":
            commands.append("python app.py")

        if path.endswith("docker-compose.yml") or path.endswith("docker-compose.yaml"):
            commands.append("docker compose up --build")

        if path.endswith("dockerfile"):
            commands.append("docker build -t <project-name> .")

    unique = []

    for command in commands:
        if command not in unique:
            unique.append(command)

    return unique


def _format_tree_from_files(files, limit=80):
    paths = sorted(item["path"] for item in files[:limit])

    output = []

    for path in paths:
        depth = path.count("/")
        name = os.path.basename(path)
        indent = "  " * min(depth, 5)
        output.append(f"{indent}- {path}")

    return "\n".join(output)


def _readme_feature_evidence(project_name):
    try:
        return strict_grounded_analyzer_project(project_name)
    except Exception as e:
        return f"Feature evidence unavailable: {e}"


def generate_readme(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    selected = select_context_files(files, limit=40)
    run_commands = _detect_run_commands(project, files)

    score_text = score_project(project_name)
    architecture_text = strict_architecture_analyzer_project(project_name)
    security_text = strict_security_analyzer_project(project_name)
    evidence_text = _readme_feature_evidence(project_name)

    prompt = f"""
You are JARVIS, a senior technical writer.

Generate a professional README.md for this project.

Grounding rules:
- Use only the indexed project evidence below.
- Do not invent features, APIs, screenshots, cloud services, teams, or deployment platforms.
- If something is not visible, write: Not visible in indexed files.
- Keep the README useful for GitHub and portfolio presentation.
- Mention exact files only when helpful.
- Use clean Markdown.

PROJECT:
{project.get("name")}

PROJECT PATH:
{project.get("path")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

FILES INDEXED:
{project.get("files_count")}

RUN COMMANDS DETECTED:
{run_commands}

PROJECT STRUCTURE:
{_format_tree_from_files(selected)}

SCORE:
{score_text}

ARCHITECTURE EVIDENCE:
{architecture_text}

SECURITY EVIDENCE:
{security_text}

GENERAL EVIDENCE:
{evidence_text}

Return a complete README.md with these sections:

# Project Name
1. Overview
2. Main Features
3. Tech Stack
4. Project Structure
5. Installation
6. How to Run
7. Available Scripts / Commands
8. Architecture Overview
9. Security Notes
10. Testing
11. Deployment
12. Future Improvements
13. Author / Notes

Do not use fake badges or fake screenshots.
"""

    return ask_llm(prompt)


def generate_readme_strict(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    run_commands = _detect_run_commands(project, files)

    facts = strict_grounded_analyzer_project(project_name)
    security = strict_security_analyzer_project(project_name)
    architecture = strict_architecture_analyzer_project(project_name)

    output = [
        f"# {project.get('name')}",
        "",
        "## Overview",
        "",
        "This README was generated from indexed project files only.",
        "",
        f"- Project path: `{project.get('path')}`",
        f"- Files indexed: `{project.get('files_count')}`",
        f"- Tech stack: `{', '.join(project.get('tech_stack', [])) or 'Not visible in indexed files.'}`",
        "",
        "## Project Structure",
        "",
        "```text",
        _format_tree_from_files(select_context_files(files, limit=80)),
        "```",
        "",
        "## Installation",
        "",
    ]

    install_commands = [
        command for command in run_commands
        if "install" in command
    ]

    if install_commands:
        output.extend("```bash\n" + "\n".join(install_commands) + "\n```")
    else:
        output.append("Not visible in indexed files.")

    output.extend([
        "",
        "## How to Run",
        "",
    ])

    runnable = [
        command for command in run_commands
        if "install" not in command
    ]

    if runnable:
        output.extend("```bash\n" + "\n".join(runnable) + "\n```")
    else:
        output.append("Not visible in indexed files.")

    output.extend([
        "",
        "## Architecture Overview",
        "",
        "```text",
        architecture,
        "```",
        "",
        "## Security Notes",
        "",
        "```text",
        security,
        "```",
        "",
        "## Testing",
        "",
        "Check the indexed test evidence below:",
        "",
        "```text",
        facts,
        "```",
        "",
        "## Future Improvements",
        "",
        "- Add or improve tests for critical modules.",
        "- Keep configuration and secrets outside source code.",
        "- Add deployment documentation if not already present.",
        "- Re-index the project after major changes.",
    ])

    return "\n".join(output)


def export_readme(project_name, strict=False):
    project, error = get_project(project_name)

    if error:
        return error

    os.makedirs(
        README_OUTPUT_DIR,
        exist_ok=True
    )

    safe_name = _readme_safe_filename(
        project.get("name", project_name)
    )

    filename = f"{safe_name}_README_{_readme_timestamp()}.md"
    path = os.path.join(
        README_OUTPUT_DIR,
        filename
    )

    content = (
        generate_readme_strict(project_name)
        if strict
        else generate_readme(project_name)
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return f"{content}\n\nREADME EXPORTED:\n{path}"


def export_readme_strict(project_name):
    return export_readme(
        project_name,
        strict=True
    )


def readme_preview(project_name):
    content = generate_readme(project_name)

    if len(content) > 6000:
        return content[:6000] + "\n\n... README preview truncated."

    return content


def readme_checklist(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    text = build_reference_text(files)

    checks = [
        ("Overview", True),
        ("Tech stack", bool(project.get("tech_stack"))),
        ("Installation commands", "requirements.txt" in text or "package.json" in text),
        ("Run commands", "scripts" in text or "app.py" in text or "main.py" in text),
        ("Testing section", "test" in text or "pytest" in text or "jest" in text),
        ("Docker/deployment section", "dockerfile" in text or "docker-compose" in text or "deploy" in text),
        ("Security notes", "auth" in text or "jwt" in text or "security" in text),
        ("Architecture overview", True),
        ("Future improvements", True),
    ]

    output = [
        "README CHECKLIST",
        f"Project: {project.get('name')}",
        ""
    ]

    for title, ok in checks:
        output.append(
            f"- [{'x' if ok else ' '}] {title}"
        )

    return "\n".join(output)


# Friendly aliases
def auto_readme(project_name):
    return generate_readme(project_name)


def generate_project_readme(project_name):
    return generate_readme(project_name)


def create_readme(project_name):
    return generate_readme(project_name)


def readme(project_name):
    return generate_readme(project_name)


def readme_strict(project_name):
    return generate_readme_strict(project_name)


def export_project_readme(project_name):
    return export_readme(project_name)


def export_project_readme_strict(project_name):
    return export_readme_strict(project_name)



# ==========================
# STEP 35 - ARCHITECTURE DIAGRAM GENERATOR
# Generates Mermaid architecture diagrams and Markdown reports.
# Evidence-based. Safe export only.
# ==========================
ARCHITECTURE_DIAGRAM_DIR = "architecture_diagrams"


def _arch35_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _arch35_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _arch35_safe_filename(name):
    return _safe_report_filename(str(name))


def _arch35_clean_node_id(value):
    cleaned = re.sub(
        r"[^A-Za-z0-9_]+",
        "_",
        str(value)
    ).strip("_")

    if not cleaned:
        cleaned = "Node"

    if cleaned[0].isdigit():
        cleaned = "N_" + cleaned

    return cleaned[:80]


def _arch35_label(value, max_len=55):
    text = str(value).replace('"', "'")

    if len(text) > max_len:
        text = text[:max_len - 3] + "..."

    return text


def _arch35_mermaid_block(code):
    return "```mermaid\n" + code.strip() + "\n```"


def _arch35_save_file(project_name, filename_suffix, content):
    os.makedirs(
        ARCHITECTURE_DIAGRAM_DIR,
        exist_ok=True
    )

    safe_project = _arch35_safe_filename(project_name)
    path = os.path.join(
        ARCHITECTURE_DIAGRAM_DIR,
        f"{safe_project}_{filename_suffix}_{_arch35_timestamp()}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _arch35_nodes_from_facts(facts):
    nodes = {}

    if facts.get("frontend_files"):
        nodes["frontend"] = f"Frontend/UI ({len(facts['frontend_files'])} files)"

    if facts.get("fastapi_files") or facts.get("flask_files") or facts.get("routes"):
        backend_count = (
            len(facts.get("fastapi_files", []))
            + len(facts.get("flask_files", []))
            + len(facts.get("routes", []))
        )
        nodes["backend"] = f"Backend/API ({backend_count} evidence items)"

    if facts.get("auth_files") or facts.get("jwt_files") or facts.get("password_hashing_files"):
        nodes["security"] = "Auth / JWT / Security"

    if facts.get("database_files"):
        nodes["database"] = f"Database/Storage ({len(facts['database_files'])} files)"

    if facts.get("upload_files"):
        nodes["upload"] = "Upload / File Processing"

    if facts.get("logging_files"):
        nodes["logging"] = "Logging / Audit"

    if facts.get("docker_files"):
        nodes["deployment"] = "Docker / Deployment"

    if facts.get("ci_files"):
        nodes["cicd"] = "CI/CD Workflows"

    if facts.get("test_files"):
        nodes["tests"] = "Tests"

    return nodes


def generate_architecture_overview_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    nodes = _arch35_nodes_from_facts(facts)

    lines = [
        "flowchart LR",
        '    User["User"]'
    ]

    for key, label in nodes.items():
        lines.append(
            f'    {_arch35_clean_node_id(key)}["{_arch35_label(label)}"]'
        )

    if "frontend" in nodes:
        lines.append("    User --> frontend")
    elif "backend" in nodes:
        lines.append("    User --> backend")

    if "frontend" in nodes and "backend" in nodes:
        lines.append("    frontend --> backend")

    if "backend" in nodes and "security" in nodes:
        lines.append("    backend --> security")

    if "backend" in nodes and "database" in nodes:
        lines.append("    backend --> database")

    if "backend" in nodes and "upload" in nodes:
        lines.append("    backend --> upload")

    if "backend" in nodes and "logging" in nodes:
        lines.append("    backend --> logging")

    if "security" in nodes and "logging" in nodes:
        lines.append("    security --> logging")

    if "upload" in nodes and "logging" in nodes:
        lines.append("    upload --> logging")

    if "cicd" in nodes and "tests" in nodes:
        lines.append("    cicd --> tests")

    if "cicd" in nodes and "deployment" in nodes:
        lines.append("    cicd --> deployment")

    if "deployment" in nodes and "backend" in nodes:
        lines.append("    deployment -. deploys .-> backend")

    if "deployment" in nodes and "frontend" in nodes:
        lines.append("    deployment -. deploys .-> frontend")

    return "\n".join(lines)


def generate_backend_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "flowchart TD",
        '    Backend["Backend/API Layer"]'
    ]

    if facts["fastapi_files"]:
        lines.append('    FastAPI["FastAPI Evidence"]')
        lines.append("    Backend --> FastAPI")

    if facts["flask_files"]:
        lines.append('    Flask["Flask Evidence"]')
        lines.append("    Backend --> Flask")

    if facts["routes"]:
        lines.append('    Routes["Routes / Endpoints"]')
        lines.append("    Backend --> Routes")

    if facts["auth_files"]:
        lines.append('    Auth["Auth / Permissions"]')
        lines.append("    Routes --> Auth")

    if facts["jwt_files"]:
        lines.append('    JWT["JWT / Token Handling"]')
        lines.append("    Auth --> JWT")

    if facts["database_files"]:
        lines.append('    Database["Database / Storage"]')
        lines.append("    Routes --> Database")

    if facts["upload_files"]:
        lines.append('    Upload["Upload / File Handling"]')
        lines.append("    Routes --> Upload")

    if facts["logging_files"]:
        lines.append('    Logs["Logging / Audit"]')
        lines.append("    Routes --> Logs")

    if len(lines) == 2:
        lines.append('    NoBackendEvidence["No backend evidence detected"]')
        lines.append("    Backend --> NoBackendEvidence")

    return "\n".join(lines)


def generate_frontend_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "flowchart TD",
        '    Frontend["Frontend/UI Layer"]'
    ]

    if facts["frontend_files"]:
        lines.append('    Components["Components / Pages / UI Files"]')
        lines.append("    Frontend --> Components")

    if facts["frontend_api_calls"]:
        lines.append('    ApiClient["API Calls"]')
        lines.append("    Components --> ApiClient")

    if facts["routes"]:
        lines.append('    Backend["Backend Routes/API"]')
        if facts["frontend_api_calls"]:
            lines.append("    ApiClient --> Backend")
        else:
            lines.append("    Components --> Backend")

    if facts["upload_files"]:
        lines.append('    UploadUi["Upload/File UI"]')
        lines.append("    Components --> UploadUi")

    if len(lines) == 2:
        lines.append('    NoFrontendEvidence["No frontend evidence detected"]')
        lines.append("    Frontend --> NoFrontendEvidence")

    return "\n".join(lines)


def generate_api_flow_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "sequenceDiagram",
        "    participant U as User"
    ]

    if facts["frontend_files"]:
        lines.append("    participant F as Frontend")

    if facts["routes"] or facts["fastapi_files"] or facts["flask_files"]:
        lines.append("    participant B as Backend/API")

    if facts["auth_files"] or facts["jwt_files"]:
        lines.append("    participant A as Auth/Security")

    if facts["database_files"]:
        lines.append("    participant D as Database/Storage")

    if facts["logging_files"]:
        lines.append("    participant L as Logs/Audit")

    if facts["frontend_files"]:
        lines.append("    U->>F: User action")
        if facts["routes"] or facts["frontend_api_calls"]:
            lines.append("    F->>B: API request")
    elif facts["routes"]:
        lines.append("    U->>B: Request")

    if facts["auth_files"] or facts["jwt_files"]:
        lines.append("    B->>A: Validate auth/token")
        lines.append("    A-->>B: Auth result")

    if facts["database_files"]:
        lines.append("    B->>D: Read/write data")
        lines.append("    D-->>B: Data result")

    if facts["logging_files"]:
        lines.append("    B->>L: Write audit/log event")

    if facts["frontend_files"] and (facts["routes"] or facts["frontend_api_calls"]):
        lines.append("    B-->>F: API response")
        lines.append("    F-->>U: UI update")
    elif facts["routes"]:
        lines.append("    B-->>U: Response")

    if len(lines) <= 3:
        lines.append("    U->>U: No API flow evidence detected")

    return "\n".join(lines)


def generate_security_layer_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "flowchart TD",
        '    Entry["User / Request"]'
    ]

    if facts["auth_files"]:
        lines.append('    Auth["Authentication / Authorization"]')
        lines.append("    Entry --> Auth")

    if facts["jwt_files"]:
        lines.append('    JWT["JWT / Token Validation"]')
        if facts["auth_files"]:
            lines.append("    Auth --> JWT")
        else:
            lines.append("    Entry --> JWT")

    if facts["password_hashing_files"]:
        lines.append('    Hash["Password Hashing"]')
        if facts["auth_files"]:
            lines.append("    Auth --> Hash")
        else:
            lines.append("    Entry --> Hash")

    if facts["secret_files"]:
        lines.append('    Config["Secrets / Environment Config"]')
        if facts["auth_files"]:
            lines.append("    Auth --> Config")
        else:
            lines.append("    Entry --> Config")

    if facts["upload_files"]:
        lines.append('    UploadSecurity["Upload Validation"]')
        lines.append("    Entry --> UploadSecurity")

    if facts["logging_files"]:
        lines.append('    Audit["Security Logs / Audit"]')
        if facts["auth_files"]:
            lines.append("    Auth --> Audit")
        if facts["upload_files"]:
            lines.append("    UploadSecurity --> Audit")
        if not facts["auth_files"] and not facts["upload_files"]:
            lines.append("    Entry --> Audit")

    if len(lines) == 2:
        lines.append('    NoSecurityEvidence["No security layer evidence detected"]')
        lines.append("    Entry --> NoSecurityEvidence")

    return "\n".join(lines)


def generate_deployment_diagram(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    lines = [
        "flowchart LR",
        '    Dev["Developer"]'
    ]

    if facts["test_files"]:
        lines.append('    Tests["Tests"]')
        lines.append("    Dev --> Tests")

    if facts["ci_files"]:
        lines.append('    CI["CI/CD Workflow"]')
        if facts["test_files"]:
            lines.append("    Tests --> CI")
        else:
            lines.append("    Dev --> CI")

    if facts["docker_files"]:
        lines.append('    Docker["Docker Build/Compose"]')
        if facts["ci_files"]:
            lines.append("    CI --> Docker")
        else:
            lines.append("    Dev --> Docker")

    if facts["docker_files"] or facts["ci_files"]:
        lines.append('    Runtime["Runtime / Deployment Target"]')
        if facts["docker_files"]:
            lines.append("    Docker --> Runtime")
        else:
            lines.append("    CI --> Runtime")

    if not facts["docker_files"] and not facts["ci_files"] and not facts["test_files"]:
        lines.append('    Manual["Manual run / deployment not visible"]')
        lines.append("    Dev --> Manual")

    return "\n".join(lines)


def generate_architecture_diagram_pack(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    diagrams = {
        "Architecture Overview": generate_architecture_overview_diagram(project_name),
        "Backend Diagram": generate_backend_diagram(project_name),
        "Frontend Diagram": generate_frontend_diagram(project_name),
        "API Flow Diagram": generate_api_flow_diagram(project_name),
        "Security Layer Diagram": generate_security_layer_diagram(project_name),
        "Deployment Diagram": generate_deployment_diagram(project_name),
    }

    content = [
        f"# Architecture Diagram Pack: {project.get('name')}",
        "",
        f"Generated: {_arch35_now()}",
        "",
        f"Project path: `{project.get('path')}`",
        f"Files indexed: `{project.get('files_count')}`",
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`",
        "",
        "Generated by: JARVIS Engineering OS",
        "",
        "Mode: evidence-based Mermaid diagrams / verify manually.",
    ]

    for title, diagram in diagrams.items():
        content.append(f"\n\n## {title}\n")
        content.append(_arch35_mermaid_block(diagram))

    content.append("\n\n## Architecture Evidence\n")
    content.append("```text")
    content.append(strict_architecture_analyzer_project(project_name))
    content.append("```")

    return "\n".join(content)


def export_architecture_diagram(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_architecture_diagram_pack(project_name)

    report_path = _arch35_save_file(
        project.get("name", project_name),
        "architecture_report",
        content
    )

    mmd_content = generate_architecture_overview_diagram(project_name)

    mmd_path = _arch35_save_file(
        project.get("name", project_name),
        "architecture_diagram_mmd",
        mmd_content
    )

    return (
        content
        + "\n\nARCHITECTURE REPORT EXPORTED:\n"
        + report_path
        + "\n\nMERMAID DIAGRAM EXPORTED:\n"
        + mmd_path
    )


# Friendly aliases
def architecture_overview_diagram(project_name):
    return generate_architecture_overview_diagram(project_name)


def backend_diagram(project_name):
    return generate_backend_diagram(project_name)


def frontend_diagram(project_name):
    return generate_frontend_diagram(project_name)


def api_flow_diagram(project_name):
    return generate_api_flow_diagram(project_name)


def security_layer_diagram(project_name):
    return generate_security_layer_diagram(project_name)


def deployment_diagram(project_name):
    return generate_deployment_diagram(project_name)


def architecture_diagram_pack(project_name):
    return generate_architecture_diagram_pack(project_name)


def export_architecture_report(project_name):
    return export_architecture_diagram(project_name)



# ==========================
# STEP 36 - API DOCUMENTATION GENERATOR
# Generates API documentation, endpoint catalog, auth flow docs, request/response notes.
# Evidence-based. Safe export only.
# ==========================
API_DOCS_OUTPUT_DIR = "api_docs_generated"


def _api36_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _api36_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _api36_safe_filename(name):
    return _safe_report_filename(str(name))


def _api36_save_file(project_name, suffix, content):
    os.makedirs(
        API_DOCS_OUTPUT_DIR,
        exist_ok=True
    )

    safe_project = _api36_safe_filename(project_name)
    path = os.path.join(
        API_DOCS_OUTPUT_DIR,
        f"{safe_project}_{suffix}_{_api36_timestamp()}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _api36_extract_routes_with_context(files):
    endpoints = []

    for item in files:
        path = item.get("path", "")
        content = item.get("content", "")
        ext = item.get("extension", "")

        if ext != ".py":
            continue

        lines = content.splitlines()

        for index, line in enumerate(lines):
            stripped = line.strip()

            route_patterns = [
                r"@router\.(get|post|put|delete|patch)\(\s*[\"']([^\"']+)[\"']",
                r"@app\.(get|post|put|delete|patch|route)\(\s*[\"']([^\"']+)[\"']",
                r"@[\w_]+\.route\(\s*[\"']([^\"']+)[\"']",
            ]

            for pattern in route_patterns:
                match = re.search(pattern, stripped)

                if not match:
                    continue

                if len(match.groups()) >= 2:
                    method = match.group(1).upper()
                    route = match.group(2)
                else:
                    method = "ROUTE"
                    route = match.group(1)

                function_name = "Not visible in indexed files."

                for next_line in lines[index + 1:index + 8]:
                    func_match = re.search(
                        r"^\s*(?:async\s+def|def)\s+([A-Za-z_][A-Za-z0-9_]*)\s*\(",
                        next_line
                    )

                    if func_match:
                        function_name = func_match.group(1)
                        break

                snippet_start = max(0, index - 3)
                snippet_end = min(len(lines), index + 18)

                snippet = "\n".join(
                    f"{line_no + 1}: {lines[line_no]}"
                    for line_no in range(snippet_start, snippet_end)
                )

                endpoints.append({
                    "method": method,
                    "route": route,
                    "file": path,
                    "function": function_name,
                    "snippet": snippet,
                })

    return endpoints


def _api36_extract_frontend_calls(files):
    calls = []

    for item in files:
        path = item.get("path", "")
        content = item.get("content", "")
        ext = item.get("extension", "")

        if ext not in {".js", ".jsx", ".ts", ".tsx"}:
            continue

        lines = content.splitlines()

        for index, line in enumerate(lines):
            patterns = [
                (r"axios\.(get|post|put|delete|patch)\(\s*[\"']([^\"']+)[\"']", "AXIOS"),
                (r"fetch\(\s*[\"']([^\"']+)[\"']", "FETCH"),
            ]

            for pattern, call_type in patterns:
                match = re.search(pattern, line)

                if not match:
                    continue

                if call_type == "AXIOS":
                    method = match.group(1).upper()
                    url = match.group(2)
                else:
                    method = "FETCH"
                    url = match.group(1)

                snippet_start = max(0, index - 3)
                snippet_end = min(len(lines), index + 10)

                snippet = "\n".join(
                    f"{line_no + 1}: {lines[line_no]}"
                    for line_no in range(snippet_start, snippet_end)
                )

                calls.append({
                    "method": method,
                    "url": url,
                    "file": path,
                    "snippet": snippet,
                })

    return calls


def show_api_endpoints(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    endpoints = _api36_extract_routes_with_context(files)

    output = [
        "API ENDPOINT CATALOG",
        "Mode: rule-based / no speculation",
        f"Project: {project.get('name')}",
        "",
    ]

    if not endpoints:
        output.append("No backend API endpoints detected in indexed files.")
        return "\n".join(output)

    for index, endpoint in enumerate(endpoints, start=1):
        output.append(
            f"{index}. {endpoint['method']} {endpoint['route']}"
        )
        output.append(f"   File: {endpoint['file']}")
        output.append(f"   Function: {endpoint['function']}")

    return "\n".join(output)


def show_frontend_api_calls(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    calls = _api36_extract_frontend_calls(files)

    output = [
        "FRONTEND API CALLS CATALOG",
        "Mode: rule-based / no speculation",
        f"Project: {project.get('name')}",
        "",
    ]

    if not calls:
        output.append("No frontend API calls detected in indexed files.")
        return "\n".join(output)

    for index, call in enumerate(calls, start=1):
        output.append(
            f"{index}. {call['method']} {call['url']}"
        )
        output.append(f"   File: {call['file']}")

    return "\n".join(output)


def show_authentication_flow(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return error

    output = [
        "AUTHENTICATION FLOW DOCUMENTATION",
        "Mode: evidence-based / no speculation",
        f"Project: {project.get('name')}",
        "",
        "Authentication evidence files:"
    ]

    auth_related = (
        facts["auth_files"]
        + facts["jwt_files"]
        + facts["password_hashing_files"]
        + facts["secret_files"]
    )

    if not auth_related:
        output.append("- No authentication/JWT/password hashing evidence detected.")
    else:
        seen = set()

        for item in auth_related:
            path = item["path"]

            if path in seen:
                continue

            seen.add(path)
            output.append(f"- {path}")

    output.append("")
    output.append("Detected authentication concepts:")

    if facts["auth_files"]:
        output.append("- Authentication/authorization files are present.")
    else:
        output.append("- Authentication files not detected by rules.")

    if facts["jwt_files"]:
        output.append("- JWT/token handling evidence is present.")
    else:
        output.append("- JWT/token handling not detected by rules.")

    if facts["password_hashing_files"]:
        output.append("- Password hashing evidence is present.")
    else:
        output.append("- Password hashing evidence not detected by rules.")

    if facts["secret_files"]:
        output.append("- Secret/config handling evidence is present.")
    else:
        output.append("- Secret/config handling not detected by rules.")

    output.append("")
    output.append("Security requirements to document manually:")
    output.append("- Required auth header format if present.")
    output.append("- Token expiration policy.")
    output.append("- Roles/permissions if present.")
    output.append("- Password policy.")
    output.append("- Error format for unauthorized/forbidden requests.")

    return "\n".join(output)


def generate_api_docs(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    endpoints = _api36_extract_routes_with_context(files)
    calls = _api36_extract_frontend_calls(files)

    endpoint_catalog = show_api_endpoints(project_name)
    frontend_calls = show_frontend_api_calls(project_name)
    auth_flow = show_authentication_flow(project_name)
    security = strict_security_analyzer_project(project_name)
    architecture = strict_architecture_analyzer_project(project_name)

    prompt = f"""
You are JARVIS, a senior API technical writer.

Generate API documentation for this project.

Grounding rules:
- Use only the indexed project evidence below.
- Do not invent endpoints, parameters, responses, auth schemes, databases, or rate limits.
- If something is not visible, write: Not visible in indexed files.
- Mention exact file paths for concrete claims.
- Keep it useful for GitHub/portfolio and developer onboarding.

PROJECT:
{project.get("name")}

PROJECT PATH:
{project.get("path")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

BACKEND ENDPOINTS:
{endpoint_catalog}

FRONTEND API CALLS:
{frontend_calls}

AUTH FLOW:
{auth_flow}

SECURITY EVIDENCE:
{security}

ARCHITECTURE EVIDENCE:
{architecture}

ROUTE SNIPPETS:
{endpoints[:30]}

FRONTEND CALL SNIPPETS:
{calls[:30]}

Return Markdown with:

# API Documentation
1. Overview
2. API Base URL
3. Authentication
4. Endpoint Catalog
5. Endpoint Details
6. Request Examples
7. Response Examples
8. Error Handling
9. Security Requirements
10. Rate Limiting
11. Frontend API Usage
12. Testing Recommendations
13. Missing Information
"""

    return ask_llm(prompt)


def generate_api_docs_strict(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    endpoint_catalog = show_api_endpoints(project_name)
    frontend_calls = show_frontend_api_calls(project_name)
    auth_flow = show_authentication_flow(project_name)

    return (
        f"# API Documentation: {project.get('name')}\n\n"
        f"Generated: {_api36_now()}\n\n"
        f"Project path: `{project.get('path')}`\n\n"
        f"Files indexed: `{project.get('files_count')}`\n\n"
        "Mode: strict rule-based API documentation. No invented endpoints.\n\n"
        "## Endpoint Catalog\n\n"
        "```text\n"
        + endpoint_catalog
        + "\n```\n\n"
        "## Frontend API Calls\n\n"
        "```text\n"
        + frontend_calls
        + "\n```\n\n"
        "## Authentication Flow\n\n"
        "```text\n"
        + auth_flow
        + "\n```\n\n"
        "## Security Requirements\n\n"
        "- Confirm authentication requirements manually.\n"
        "- Confirm request/response schemas manually.\n"
        "- Confirm rate limiting manually.\n"
        "- Confirm error response format manually.\n"
        "- Confirm CORS policy manually.\n\n"
        "## Missing Information\n\n"
        "Any endpoint parameter, response body, rate limit, or auth rule not visible in indexed files must be verified manually.\n"
    )


def export_api_docs(project_name, strict=False):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        generate_api_docs_strict(project_name)
        if strict
        else generate_api_docs(project_name)
    )

    path = _api36_save_file(
        project.get("name", project_name),
        "api_documentation",
        content
    )

    return f"{content}\n\nAPI DOCUMENTATION EXPORTED:\n{path}"


def export_api_docs_strict(project_name):
    return export_api_docs(
        project_name,
        strict=True
    )


def export_api_endpoints_catalog(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        f"# API Endpoints Catalog: {project.get('name')}\n\n"
        f"Generated: {_api36_now()}\n\n"
        "```text\n"
        + show_api_endpoints(project_name)
        + "\n\n"
        + show_frontend_api_calls(project_name)
        + "\n```"
    )

    path = _api36_save_file(
        project.get("name", project_name),
        "api_endpoints_catalog",
        content
    )

    return f"{content}\n\nAPI ENDPOINTS CATALOG EXPORTED:\n{path}"


def api_documentation_report(project_name):
    return generate_api_docs(project_name)


# Friendly aliases
def api_docs(project_name):
    return generate_api_docs(project_name)


def generate_api_documentation(project_name):
    return generate_api_docs(project_name)


def export_api_documentation(project_name):
    return export_api_docs(project_name)


def api_endpoints(project_name):
    return show_api_endpoints(project_name)


def frontend_api_calls(project_name):
    return show_frontend_api_calls(project_name)


def authentication_flow(project_name):
    return show_authentication_flow(project_name)



# ==========================
# STEP 37 - CHANGELOG & RELEASE NOTES GENERATOR
# Generates CHANGELOG.md, RELEASE_NOTES.md and VERSION_HISTORY.md from indexed evidence.
# Evidence-based. Safe export only.
# ==========================
CHANGELOG_OUTPUT_DIR = "changelog_generated"


def _changelog37_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _changelog37_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _changelog37_safe_filename(name):
    return _safe_report_filename(str(name))


def _changelog37_save_file(project_name, suffix, content):
    os.makedirs(
        CHANGELOG_OUTPUT_DIR,
        exist_ok=True
    )

    safe_project = _changelog37_safe_filename(project_name)

    path = os.path.join(
        CHANGELOG_OUTPUT_DIR,
        f"{safe_project}_{suffix}_{_changelog37_timestamp()}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _changelog37_detect_version(project, files):
    version_candidates = []

    for item in files:
        path = normalize_path(item["path"])
        content = item["content"]

        if path.endswith("package.json"):
            try:
                data = json.loads(content)
                version = data.get("version")

                if version:
                    version_candidates.append(
                        f"{version} from {item['path']}"
                    )
            except Exception:
                pass

        if path.endswith("pyproject.toml"):
            match = re.search(
                r"version\s*=\s*[\"']([^\"']+)[\"']",
                content
            )

            if match:
                version_candidates.append(
                    f"{match.group(1)} from {item['path']}"
                )

        if path.endswith("__init__.py"):
            match = re.search(
                r"__version__\s*=\s*[\"']([^\"']+)[\"']",
                content
            )

            if match:
                version_candidates.append(
                    f"{match.group(1)} from {item['path']}"
                )

    if version_candidates:
        return version_candidates[0]

    return "0.1.0 generated / version not visible in indexed files"


def _changelog37_detect_existing_changelog(files):
    matches = []

    for item in files:
        path = normalize_path(item["path"])

        if (
            "changelog" in path
            or "release_notes" in path
            or "version_history" in path
        ):
            matches.append(item["path"])

    return matches


def _changelog37_feature_buckets(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return None, error

    buckets = {
        "Added": [],
        "Changed": [],
        "Security": [],
        "Testing": [],
        "DevOps": [],
        "Documentation": [],
        "Known Gaps": [],
    }

    if facts["frontend_files"]:
        buckets["Added"].append(
            f"Frontend/UI layer detected in {len(facts['frontend_files'])} indexed file(s)."
        )

    if facts["fastapi_files"]:
        buckets["Added"].append(
            f"FastAPI backend evidence detected in {len(facts['fastapi_files'])} file(s)."
        )

    if facts["flask_files"]:
        buckets["Added"].append(
            f"Flask backend evidence detected in {len(facts['flask_files'])} file(s)."
        )

    if facts["routes"]:
        buckets["Added"].append(
            f"Backend routes/endpoints detected: {len(facts['routes'])}."
        )

    if facts["frontend_api_calls"]:
        buckets["Added"].append(
            f"Frontend API calls detected: {len(facts['frontend_api_calls'])}."
        )

    if facts["database_files"]:
        buckets["Added"].append(
            f"Database/storage layer evidence detected in {len(facts['database_files'])} file(s)."
        )

    if facts["upload_files"]:
        buckets["Added"].append(
            f"Upload/file-handling evidence detected in {len(facts['upload_files'])} file(s)."
        )

    if facts["auth_files"]:
        buckets["Security"].append(
            f"Authentication/authorization evidence detected in {len(facts['auth_files'])} file(s)."
        )

    if facts["jwt_files"]:
        buckets["Security"].append(
            f"JWT/token handling evidence detected in {len(facts['jwt_files'])} file(s)."
        )

    if facts["password_hashing_files"]:
        buckets["Security"].append(
            f"Password hashing evidence detected in {len(facts['password_hashing_files'])} file(s)."
        )

    if facts["logging_files"]:
        buckets["Security"].append(
            f"Logging/audit evidence detected in {len(facts['logging_files'])} file(s)."
        )

    if facts["test_files"]:
        buckets["Testing"].append(
            f"Test files detected: {len(facts['test_files'])}."
        )
    else:
        buckets["Known Gaps"].append(
            "No test files detected by current rules."
        )

    if facts["docker_files"]:
        buckets["DevOps"].append(
            f"Docker/deployment evidence detected in {len(facts['docker_files'])} file(s)."
        )
    else:
        buckets["Known Gaps"].append(
            "No Docker/deployment files detected by current rules."
        )

    if facts["ci_files"]:
        buckets["DevOps"].append(
            f"CI/workflow evidence detected in {len(facts['ci_files'])} file(s)."
        )
    else:
        buckets["Known Gaps"].append(
            "No CI/workflow files detected by current rules."
        )

    if facts["packages"] or facts["requirements"]:
        buckets["Changed"].append(
            "Dependency manifest evidence detected."
        )

    existing_docs = []

    for item in get_code_files(project):
        path = normalize_path(item["path"])

        if path.endswith(".md") or "docs/" in path:
            existing_docs.append(item["path"])

    if existing_docs:
        buckets["Documentation"].append(
            f"Documentation/Markdown files detected: {len(existing_docs)}."
        )
    else:
        buckets["Known Gaps"].append(
            "No Markdown/docs files detected by current rules."
        )

    return buckets, None


def generate_changelog(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    version = _changelog37_detect_version(project, files)
    existing_changelog = _changelog37_detect_existing_changelog(files)

    buckets, bucket_error = _changelog37_feature_buckets(project_name)

    if bucket_error:
        return bucket_error

    content = [
        f"# Changelog - {project.get('name')}",
        "",
        "All notable changes are generated from indexed project evidence.",
        "",
        "Important: This generated changelog is evidence-based and should be reviewed manually before publishing.",
        "",
        "## [Unreleased]",
        "",
        "### Notes",
        "",
        "- Generated by JARVIS Engineering OS.",
        "- This file does not replace Git commit history.",
        "- Items below are inferred from indexed files, not from Git diffs.",
        "",
        f"## [{version}] - {_changelog37_now().split()[0]}",
        "",
    ]

    for section in [
        "Added",
        "Changed",
        "Security",
        "Testing",
        "DevOps",
        "Documentation",
        "Known Gaps",
    ]:
        content.append(f"### {section}")
        content.append("")

        items = buckets.get(section, [])

        if items:
            for item in items:
                content.append(f"- {item}")
        else:
            content.append("- Not visible in indexed files.")

        content.append("")

    content.append("## Existing changelog/release files detected")
    content.append("")

    if existing_changelog:
        for path in existing_changelog:
            content.append(f"- `{path}`")
    else:
        content.append("- None detected.")

    return "\n".join(content)


def generate_release_notes(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    version = _changelog37_detect_version(project, files)
    score = score_project(project_name)
    release = release_readiness(project_name) if "release_readiness" in globals() else "Release readiness module not loaded."
    architecture = strict_architecture_analyzer_project(project_name)
    security = strict_security_analyzer_project(project_name)

    prompt = f"""
You are JARVIS, a senior release manager.

Generate professional release notes for this project.

Grounding rules:
- Use only the indexed project evidence below.
- Do not invent features, customers, deadlines, production incidents, cloud providers, or release dates.
- If something is not visible, write: Not visible in indexed files.
- Keep it practical for GitHub/portfolio.

PROJECT:
{project.get("name")}

VERSION:
{version}

PROJECT PATH:
{project.get("path")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

SCORE:
{score}

RELEASE READINESS:
{release}

ARCHITECTURE:
{architecture}

SECURITY:
{security}

CHANGELOG:
{generate_changelog(project_name)}

Return Markdown with:
# Release Notes
1. Release Summary
2. Highlights
3. Added
4. Changed
5. Security Notes
6. Testing Notes
7. Deployment Notes
8. Known Limitations
9. Upgrade / Migration Notes
10. Final Release Recommendation
"""

    return ask_llm(prompt)


def generate_release_notes_strict(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    version = _changelog37_detect_version(project, files)

    return (
        f"# Release Notes - {project.get('name')}\n\n"
        f"Generated: {_changelog37_now()}\n\n"
        f"Version: `{version}`\n\n"
        f"Project path: `{project.get('path')}`\n\n"
        f"Files indexed: `{project.get('files_count')}`\n\n"
        "## Summary\n\n"
        "These release notes were generated from indexed project evidence only.\n\n"
        "## Changelog\n\n"
        + generate_changelog(project_name)
        + "\n\n"
        "## Release Readiness\n\n"
        "```text\n"
        + (release_readiness(project_name) if "release_readiness" in globals() else "Release readiness module not loaded.")
        + "\n```\n\n"
        "## Manual Verification Required\n\n"
        "- Confirm version number.\n"
        "- Confirm completed features against Git commits.\n"
        "- Confirm tests pass.\n"
        "- Confirm deployment procedure.\n"
        "- Confirm secrets and environment variables are safe.\n"
    )


def show_version_history(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    files = get_code_files(project)
    version = _changelog37_detect_version(project, files)
    existing = _changelog37_detect_existing_changelog(files)

    output = [
        "VERSION HISTORY",
        "Mode: evidence-based / no Git history access",
        f"Project: {project.get('name')}",
        "",
        f"Detected/current version: {version}",
        "",
        "Existing version/changelog/release files:"
    ]

    if existing:
        output.extend(f"- {path}" for path in existing)
    else:
        output.append("- None detected in indexed files.")

    output.append("")
    output.append("Note:")
    output.append("- Full historical version data requires Git commit/tag history.")
    output.append("- Current report uses only indexed project files.")

    return "\n".join(output)


def export_changelog(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_changelog(project_name)

    path = _changelog37_save_file(
        project.get("name", project_name),
        "CHANGELOG",
        content
    )

    return f"{content}\n\nCHANGELOG EXPORTED:\n{path}"


def export_release_notes(project_name, strict=False):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        generate_release_notes_strict(project_name)
        if strict
        else generate_release_notes(project_name)
    )

    path = _changelog37_save_file(
        project.get("name", project_name),
        "RELEASE_NOTES",
        content
    )

    return f"{content}\n\nRELEASE NOTES EXPORTED:\n{path}"


def export_release_notes_strict(project_name):
    return export_release_notes(
        project_name,
        strict=True
    )


def export_version_history(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        f"# Version History - {project.get('name')}\n\n"
        f"Generated: {_changelog37_now()}\n\n"
        "```text\n"
        + show_version_history(project_name)
        + "\n```"
    )

    path = _changelog37_save_file(
        project.get("name", project_name),
        "VERSION_HISTORY",
        content
    )

    return f"{content}\n\nVERSION HISTORY EXPORTED:\n{path}"


# Friendly aliases
def changelog(project_name):
    return generate_changelog(project_name)


def generate_project_changelog(project_name):
    return generate_changelog(project_name)


def release_notes(project_name):
    return generate_release_notes(project_name)


def generate_project_release_notes(project_name):
    return generate_release_notes(project_name)


def version_history(project_name):
    return show_version_history(project_name)


def export_project_changelog(project_name):
    return export_changelog(project_name)


def export_project_release_notes(project_name):
    return export_release_notes(project_name)



# ==========================
# STEP 38 - TEST COVERAGE INTELLIGENCE
# Detects untested files, critical code without tests, endpoint/component test gaps,
# estimated coverage, testing roadmap, and quality report.
# Evidence-based. Safe reporting only. No automatic code changes.
# ==========================
TEST_COVERAGE_OUTPUT_DIR = "test_coverage_generated"


def _test38_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _test38_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _test38_safe_filename(name):
    return _safe_report_filename(str(name))


def _test38_save_file(project_name, suffix, content):
    os.makedirs(
        TEST_COVERAGE_OUTPUT_DIR,
        exist_ok=True
    )

    safe_project = _test38_safe_filename(project_name)

    path = os.path.join(
        TEST_COVERAGE_OUTPUT_DIR,
        f"{safe_project}_{suffix}_{_test38_timestamp()}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _test38_is_source_file(item):
    path = normalize_path(item.get("path", ""))
    ext = item.get("extension", "").lower()

    if is_test_path(path):
        return False

    if is_documentation_or_config_path(path):
        return False

    if ext not in {
        ".py",
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".java",
        ".cs",
        ".cpp",
        ".c",
        ".h",
        ".hpp",
    }:
        return False

    return True


def _test38_related_test_names(path):
    normalized = normalize_path(path)
    base = file_stem(normalized)

    return {
        f"test_{base}",
        f"{base}_test",
        f"{base}.test",
        f"{base}.spec",
        f"{base}-test",
        f"{base}-spec",
        f"{base}test",
        f"{base}spec",
    }


def _test38_find_related_tests(source_item, test_files):
    source_path = normalize_path(source_item.get("path", ""))
    source_stem = file_stem(source_path)
    candidates = _test38_related_test_names(source_path)

    matches = []

    for test in test_files:
        test_path = normalize_path(test.get("path", ""))
        test_content = test.get("content", "").lower()

        if source_stem in test_path:
            matches.append(test["path"])
            continue

        if any(candidate in test_path for candidate in candidates):
            matches.append(test["path"])
            continue

        if source_stem in test_content:
            matches.append(test["path"])
            continue

    return sorted(set(matches))


def _test38_criticality(item):
    path = normalize_path(item.get("path", ""))
    content = item.get("content", "").lower()
    score = 0
    reasons = []

    critical_markers = [
        ("auth", 25, "auth/security path or content"),
        ("login", 20, "login/authentication logic"),
        ("jwt", 20, "JWT/token logic"),
        ("token", 15, "token logic"),
        ("password", 15, "password handling"),
        ("permission", 18, "permission/authorization logic"),
        ("admin", 15, "admin/privileged logic"),
        ("upload", 18, "upload/file handling"),
        ("file", 8, "file handling"),
        ("database", 15, "database/storage logic"),
        ("sqlalchemy", 14, "ORM/database logic"),
        ("sqlite", 12, "SQLite/storage logic"),
        ("payment", 25, "payment/business-critical logic"),
        ("email", 10, "email/notification logic"),
        ("api", 10, "API logic"),
        ("route", 10, "route/controller logic"),
        ("scan", 12, "scanner/security scan logic"),
        ("audit", 12, "audit/logging logic"),
        ("logger", 8, "logging logic"),
        ("delete", 10, "delete/destructive operation"),
        ("remove", 8, "remove/destructive operation"),
        ("backup", 12, "backup/recovery logic"),
    ]

    haystack = path + "\n" + content[:6000]

    for marker, points, reason in critical_markers:
        if marker in haystack:
            score += points
            reasons.append(reason)

    if is_entrypoint_or_framework_file(path):
        score += 15
        reasons.append("entrypoint/framework file")

    if item.get("extension") == ".py":
        symbols = extract_python_symbols(
            item.get("path", ""),
            item.get("content", "")
        )

        function_count = len(symbols.get("functions", []))
        class_count = len(symbols.get("classes", []))

        if function_count >= 5:
            score += 8
            reasons.append("many Python functions")

        if class_count >= 2:
            score += 8
            reasons.append("multiple Python classes")

    elif item.get("extension") in {".js", ".jsx", ".ts", ".tsx"}:
        symbols = extract_js_symbols(item.get("content", ""))

        if len(symbols.get("functions", [])) >= 4:
            score += 8
            reasons.append("many JS/TS functions/components")

        if len(symbols.get("classes", [])) >= 1:
            score += 5
            reasons.append("JS/TS class evidence")

    score = max(0, min(100, score))

    if score >= 70:
        level = "CRITICAL"
    elif score >= 45:
        level = "HIGH"
    elif score >= 25:
        level = "MEDIUM"
    else:
        level = "LOW"

    return {
        "score": score,
        "level": level,
        "reasons": sorted(set(reasons)),
    }


def _test38_extract_critical_symbols(item):
    ext = item.get("extension", "")
    content = item.get("content", "")
    path = item.get("path", "")

    symbols = []

    critical_words = [
        "auth",
        "login",
        "jwt",
        "token",
        "password",
        "permission",
        "admin",
        "upload",
        "delete",
        "remove",
        "scan",
        "backup",
        "encrypt",
        "decrypt",
        "validate",
        "verify",
        "create",
        "update",
    ]

    if ext == ".py":
        extracted = extract_python_symbols(path, content)

        for func in extracted.get("functions", []):
            lower = func.lower()

            if any(word in lower for word in critical_words):
                symbols.append(f"function {func}")

        for cls in extracted.get("classes", []):
            lower = cls.lower()

            if any(word in lower for word in critical_words):
                symbols.append(f"class {cls}")

    elif ext in {".js", ".jsx", ".ts", ".tsx"}:
        extracted = extract_js_symbols(content)

        for func in extracted.get("functions", []):
            lower = func.lower()

            if any(word in lower for word in critical_words):
                symbols.append(f"function/component {func}")

        for cls in extracted.get("classes", []):
            lower = cls.lower()

            if any(word in lower for word in critical_words):
                symbols.append(f"class {cls}")

    return sorted(set(symbols))


def _test38_analyze(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return None, error

    files = get_code_files(project)
    source_files = [
        item for item in files
        if _test38_is_source_file(item)
    ]

    test_files = facts.get("test_files", [])

    covered = []
    uncovered = []
    critical_uncovered = []

    for item in source_files:
        related_tests = _test38_find_related_tests(
            item,
            test_files
        )
        criticality = _test38_criticality(item)
        critical_symbols = _test38_extract_critical_symbols(item)

        record = {
            "path": item["path"],
            "extension": item["extension"],
            "related_tests": related_tests,
            "criticality": criticality,
            "critical_symbols": critical_symbols,
        }

        if related_tests:
            covered.append(record)
        else:
            uncovered.append(record)

            if criticality["level"] in {"CRITICAL", "HIGH"} or critical_symbols:
                critical_uncovered.append(record)

    endpoint_gaps = []

    for route in facts.get("routes", []):
        route_lower = route.lower()
        route_has_test = False

        for test in test_files:
            test_haystack = (
                normalize_path(test.get("path", ""))
                + "\n"
                + test.get("content", "").lower()
            )

            route_parts = re.findall(r"/[A-Za-z0-9_/{}/.-]+", route_lower)

            for part in route_parts:
                clean_part = part.strip()

                if clean_part and clean_part in test_haystack:
                    route_has_test = True
                    break

            if route_has_test:
                break

        if not route_has_test:
            endpoint_gaps.append(route)

    frontend_gaps = []

    for call in facts.get("frontend_api_calls", []):
        call_lower = call.lower()
        has_test = False

        for test in test_files:
            test_haystack = (
                normalize_path(test.get("path", ""))
                + "\n"
                + test.get("content", "").lower()
            )

            parts = re.findall(r"/[A-Za-z0-9_/{}/.-]+", call_lower)

            for part in parts:
                if part and part in test_haystack:
                    has_test = True
                    break

            if has_test:
                break

        if not has_test:
            frontend_gaps.append(call)

    total = len(source_files)
    covered_count = len(covered)

    estimated_coverage = round(
        (covered_count / total) * 100,
        1
    ) if total else 0.0

    if estimated_coverage >= 80 and not critical_uncovered:
        level = "STRONG"
    elif estimated_coverage >= 60:
        level = "OK_BUT_NEEDS_IMPROVEMENT"
    elif estimated_coverage >= 35:
        level = "WEAK"
    else:
        level = "VERY_WEAK"

    return {
        "project": project,
        "facts": facts,
        "source_files": source_files,
        "test_files": test_files,
        "covered": covered,
        "uncovered": uncovered,
        "critical_uncovered": sorted(
            critical_uncovered,
            key=lambda item: item["criticality"]["score"],
            reverse=True
        ),
        "endpoint_gaps": endpoint_gaps,
        "frontend_gaps": frontend_gaps,
        "estimated_coverage": estimated_coverage,
        "coverage_level": level,
    }, None


def show_untested_files(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    output = [
        "UNTESTED FILES REPORT",
        "Mode: heuristic / related test name and content matching",
        f"Project: {data['project'].get('name')}",
        "",
        f"Source files analyzed: {len(data['source_files'])}",
        f"Test files detected: {len(data['test_files'])}",
        f"Untested source files: {len(data['uncovered'])}",
        "",
    ]

    if not data["uncovered"]:
        output.append("No untested source files detected by current rules.")
        return "\n".join(output)

    for item in sorted(
        data["uncovered"],
        key=lambda row: row["criticality"]["score"],
        reverse=True
    )[:120]:
        output.append(
            f"- {item['path']} | "
            f"{item['criticality']['level']} "
            f"({item['criticality']['score']}/100)"
        )

        reasons = item["criticality"]["reasons"]

        if reasons:
            output.append(
                "  Reasons: " + ", ".join(reasons[:6])
            )

    return "\n".join(output)


def show_critical_code_without_tests(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    output = [
        "CRITICAL CODE WITHOUT TESTS",
        "Mode: heuristic / verify manually",
        f"Project: {data['project'].get('name')}",
        "",
    ]

    if not data["critical_uncovered"]:
        output.append("No critical untested source files detected by current rules.")
        return "\n".join(output)

    for item in data["critical_uncovered"][:80]:
        output.append(
            f"- {item['path']} | "
            f"{item['criticality']['level']} "
            f"({item['criticality']['score']}/100)"
        )

        if item["criticality"]["reasons"]:
            output.append(
                "  Reasons: " + ", ".join(item["criticality"]["reasons"][:8])
            )

        if item["critical_symbols"]:
            output.append(
                "  Critical symbols: "
                + ", ".join(item["critical_symbols"][:10])
            )

    return "\n".join(output)


def show_untested_endpoints(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    output = [
        "UNTESTED ENDPOINTS / API FLOWS",
        "Mode: heuristic / route text search in tests",
        f"Project: {data['project'].get('name')}",
        "",
    ]

    if not data["endpoint_gaps"]:
        output.append("No untested backend endpoints detected by current rules.")
    else:
        output.append("Backend routes without detected tests:")
        for route in data["endpoint_gaps"][:100]:
            output.append(f"- {route}")

    output.append("")

    if not data["frontend_gaps"]:
        output.append("No untested frontend API calls detected by current rules.")
    else:
        output.append("Frontend API calls without detected tests:")
        for call in data["frontend_gaps"][:100]:
            output.append(f"- {call}")

    return "\n".join(output)


def analyze_test_coverage(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    output = [
        "TEST COVERAGE INTELLIGENCE REPORT",
        "Mode: heuristic / static analysis / no real test runner coverage",
        f"Generated: {_test38_now()}",
        "",
        f"Project: {data['project'].get('name')}",
        f"Path: {data['project'].get('path')}",
        "",
        f"Source files analyzed: {len(data['source_files'])}",
        f"Test files detected: {len(data['test_files'])}",
        f"Files with related tests: {len(data['covered'])}",
        f"Files without detected tests: {len(data['uncovered'])}",
        f"Estimated test coverage: {data['estimated_coverage']}%",
        f"Coverage level: {data['coverage_level']}",
        "",
        "Important note:",
        "- This is not line/branch coverage. It estimates coverage from indexed files, naming, and test references.",
        "- Run pytest/jest/vitest/coverage tools for real coverage.",
        "",
        "Top critical gaps:"
    ]

    if data["critical_uncovered"]:
        for item in data["critical_uncovered"][:15]:
            output.append(
                f"- {item['path']} | "
                f"{item['criticality']['level']} "
                f"({item['criticality']['score']}/100)"
            )
    else:
        output.append("- None detected by current rules.")

    output.append("")
    output.append("Endpoint/API gaps:")
    if data["endpoint_gaps"]:
        for route in data["endpoint_gaps"][:20]:
            output.append(f"- {route}")
    else:
        output.append("- None detected by current rules.")

    output.append("")
    output.append("Recommended next tests:")
    output.extend(_test38_recommended_tests(data)[:30])

    return "\n".join(output)


def _test38_recommended_tests(data):
    recommendations = []

    for item in data["critical_uncovered"][:20]:
        path = item["path"]
        ext = item["extension"]

        if "auth" in normalize_path(path) or "token" in normalize_path(path) or "jwt" in normalize_path(path):
            recommendations.append(
                f"- Add auth/security tests for `{path}`: valid login, invalid login, expired token, missing token, unauthorized access."
            )
        elif "upload" in normalize_path(path) or "file" in normalize_path(path):
            recommendations.append(
                f"- Add upload/file tests for `{path}`: valid file, invalid type, oversized file, empty file, malicious filename."
            )
        elif "database" in normalize_path(path) or "db" in normalize_path(path) or "model" in normalize_path(path):
            recommendations.append(
                f"- Add database tests for `{path}`: create/read/update/delete, validation errors, connection failure."
            )
        elif ext in {".jsx", ".tsx", ".js", ".ts"}:
            recommendations.append(
                f"- Add frontend/component tests for `{path}`: render, user actions, loading/error states, API failure."
            )
        else:
            recommendations.append(
                f"- Add unit tests for `{path}` covering main functions, edge cases, and error handling."
            )

    for route in data["endpoint_gaps"][:20]:
        recommendations.append(
            f"- Add API endpoint test for `{route}`: success, validation error, unauthorized/forbidden, failure path."
        )

    if not recommendations:
        recommendations.append(
            "- Test coverage looks acceptable by current rules. Run real coverage tools for line/branch coverage."
        )

    return recommendations


def generate_testing_roadmap(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    output = [
        "TESTING ROADMAP",
        f"Project: {data['project'].get('name')}",
        f"Generated: {_test38_now()}",
        "",
        "Priority 1 - Critical safety tests"
    ]

    critical = data["critical_uncovered"][:15]

    if critical:
        for item in critical:
            output.append(
                f"- Test `{item['path']}` first "
                f"({item['criticality']['level']} / {item['criticality']['score']}/100)."
            )
    else:
        output.append("- No critical untested files detected by current rules.")

    output.extend([
        "",
        "Priority 2 - API and integration tests"
    ])

    if data["endpoint_gaps"]:
        for route in data["endpoint_gaps"][:25]:
            output.append(f"- Add integration test for `{route}`.")
    else:
        output.append("- No endpoint gaps detected by current rules.")

    output.extend([
        "",
        "Priority 3 - Frontend tests"
    ])

    if data["frontend_gaps"]:
        for call in data["frontend_gaps"][:25]:
            output.append(f"- Add frontend test around API call `{call}`.")
    else:
        output.append("- No frontend API call gaps detected by current rules.")

    output.extend([
        "",
        "Priority 4 - Regression and release tests",
        "- Add smoke tests for app startup.",
        "- Add test for invalid configuration/environment variables.",
        "- Add security regression tests for known risks.",
        "- Add CI command to run all tests before release.",
        "",
        "Suggested test commands:",
        "```bash",
        "pytest",
        "npm test",
        "npm run test",
        "npm run coverage",
        "```",
    ])

    return "\n".join(output)


def generate_quality_report(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    release = release_readiness(project_name) if "release_readiness" in globals() else "Release readiness module not loaded."
    score = score_project(project_name)

    output = [
        f"# Quality Report - {data['project'].get('name')}",
        "",
        f"Generated: {_test38_now()}",
        "",
        "## Test Coverage Intelligence",
        "",
        "```text",
        analyze_test_coverage(project_name),
        "```",
        "",
        "## Project Score",
        "",
        "```text",
        score,
        "```",
        "",
        "## Release Readiness",
        "",
        "```text",
        release,
        "```",
        "",
        "## Critical Code Without Tests",
        "",
        "```text",
        show_critical_code_without_tests(project_name),
        "```",
        "",
        "## Testing Roadmap",
        "",
        "```text",
        generate_testing_roadmap(project_name),
        "```",
    ]

    return "\n".join(output)


def export_test_coverage(project_name):
    data, error = _test38_analyze(project_name)

    if error:
        return error

    content = (
        f"# Test Coverage Report - {data['project'].get('name')}\n\n"
        f"Generated: {_test38_now()}\n\n"
        "```text\n"
        + analyze_test_coverage(project_name)
        + "\n```"
    )

    path = _test38_save_file(
        data["project"].get("name", project_name),
        "test_coverage_report",
        content
    )

    return f"{content}\n\nTEST COVERAGE REPORT EXPORTED:\n{path}"


def export_testing_roadmap(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        f"# Testing Roadmap - {project.get('name')}\n\n"
        + generate_testing_roadmap(project_name)
    )

    path = _test38_save_file(
        project.get("name", project_name),
        "testing_roadmap",
        content
    )

    return f"{content}\n\nTESTING ROADMAP EXPORTED:\n{path}"


def export_quality_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_quality_report(project_name)

    path = _test38_save_file(
        project.get("name", project_name),
        "quality_report",
        content
    )

    return f"{content}\n\nQUALITY REPORT EXPORTED:\n{path}"


# Friendly aliases
def test_coverage(project_name):
    return analyze_test_coverage(project_name)


def coverage_report(project_name):
    return analyze_test_coverage(project_name)


def testing_roadmap(project_name):
    return generate_testing_roadmap(project_name)


def untested_files(project_name):
    return show_untested_files(project_name)


def critical_code_without_tests(project_name):
    return show_critical_code_without_tests(project_name)


def untested_endpoints(project_name):
    return show_untested_endpoints(project_name)


def export_coverage_report(project_name):
    return export_test_coverage(project_name)


def export_project_quality_report(project_name):
    return export_quality_report(project_name)



# ==========================
# STEP 39 - PORTFOLIO PACK GENERATOR
# Generates portfolio-ready summaries: GitHub showcase, LinkedIn description,
# CV summary, recruiter summary, technical highlights, business value, demo script.
# Evidence-based. Safe export only.
# ==========================
PORTFOLIO_PACK_OUTPUT_DIR = "portfolio_pack_generated"


def _portfolio39_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _portfolio39_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _portfolio39_safe_filename(name):
    return _safe_report_filename(str(name))


def _portfolio39_save_file(project_name, suffix, content):
    os.makedirs(
        PORTFOLIO_PACK_OUTPUT_DIR,
        exist_ok=True
    )

    safe_project = _portfolio39_safe_filename(project_name)

    path = os.path.join(
        PORTFOLIO_PACK_OUTPUT_DIR,
        f"{safe_project}_{suffix}_{_portfolio39_timestamp()}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _portfolio39_evidence(project_name):
    project, error = get_project(project_name)

    if error:
        return None, error

    evidence = {
        "overview": project_overview(project),
        "score": score_project(project_name),
        "grounded": strict_grounded_analyzer_project(project_name),
        "architecture": strict_architecture_analyzer_project(project_name),
        "security": strict_security_analyzer_project(project_name),
        "readme_checklist": readme_checklist(project_name) if "readme_checklist" in globals() else "README generator not loaded.",
        "api_endpoints": show_api_endpoints(project_name) if "show_api_endpoints" in globals() else "API documentation generator not loaded.",
        "test_coverage": analyze_test_coverage(project_name) if "analyze_test_coverage" in globals() else "Test coverage intelligence not loaded.",
        "release": release_readiness(project_name) if "release_readiness" in globals() else "Release readiness module not loaded.",
        "changelog": generate_changelog(project_name) if "generate_changelog" in globals() else "Changelog generator not loaded.",
    }

    return evidence, None


def generate_cv_summary(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a senior career coach and technical recruiter.

Generate a CV-ready project summary for this project.

Grounding rules:
- Use only the indexed project evidence below.
- Do not invent metrics, users, production usage, companies, cloud services, awards, or live deployment.
- If something is not visible, write: Not visible in indexed files.
- Make it sound professional, but still honest.
- Keep it suitable for a junior/mid frontend, full-stack, or cybersecurity-oriented developer CV.

PROJECT:
{project.get("name")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

EVIDENCE:
{evidence}

Return:
1. One-line CV project title
2. 3 bullet points for CV
3. Short recruiter-friendly paragraph
4. Technical keywords
5. What to say in an interview
"""

    return ask_llm(prompt)


def generate_github_showcase(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a senior GitHub portfolio writer.

Generate a GitHub showcase page for this project.

Grounding rules:
- Use only indexed project evidence.
- Do not invent fake badges, screenshots, deployment links, stars, contributors, production status, or cloud services.
- If something is not visible, write: Not visible in indexed files.
- Make it polished and recruiter-friendly.
- Use Markdown.

PROJECT:
{project.get("name")}

PROJECT PATH:
{project.get("path")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

EVIDENCE:
{evidence}

Return Markdown with:
# Project Showcase
1. Short description
2. Why this project matters
3. Main features
4. Tech stack
5. Architecture highlights
6. Security highlights
7. API / backend highlights
8. Testing / quality notes
9. Demo script
10. Future improvements
11. Interview pitch
"""

    return ask_llm(prompt)


def generate_linkedin_showcase(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a LinkedIn personal branding assistant.

Generate a LinkedIn post/description for this project.

Grounding rules:
- Use only indexed project evidence.
- Do not invent fake production results, clients, users, funding, awards, or deployment.
- Be professional and human.
- Keep it suitable for a student/developer portfolio.
- Mention technologies only when supported by evidence.

PROJECT:
{project.get("name")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

EVIDENCE:
{evidence}

Return:
1. Short LinkedIn project description
2. Longer LinkedIn post
3. 5 recruiter-friendly bullets
4. Suggested hashtags
5. Short version for About/Featured section
"""

    return ask_llm(prompt)


def generate_recruiter_summary(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a technical recruiter.

Create a recruiter-friendly summary for this project.

Grounding rules:
- Use only indexed project evidence.
- Do not invent business impact, revenue, users, or production claims.
- Explain why the project is impressive for interviews.
- Keep it clear and concise.

PROJECT:
{project.get("name")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

EVIDENCE:
{evidence}

Return:
1. Recruiter summary
2. Why it is relevant for frontend roles
3. Why it is relevant for cybersecurity roles
4. Why it is relevant for full-stack roles
5. Interview questions this project can answer
"""

    return ask_llm(prompt)


def generate_technical_highlights(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a senior software engineer.

Generate technical highlights for this project.

Grounding rules:
- Use only indexed project evidence.
- Mention exact capabilities only when visible.
- Do not invent cloud architecture, ML modules, databases, APIs, or security features.
- Keep it technical and useful for an interview.

PROJECT:
{project.get("name")}

EVIDENCE:
{evidence}

Return:
1. Architecture highlights
2. Backend highlights
3. Frontend highlights
4. Security highlights
5. Testing/quality highlights
6. DevOps/deployment highlights
7. Technical risks to be honest about
8. Best interview explanation
"""

    return ask_llm(prompt)


def generate_business_value_summary(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a product-minded engineering assistant.

Generate a business value summary for this project.

Grounding rules:
- Use only indexed project evidence.
- Do not invent customers, users, revenue, production deployment, or quantified impact.
- Explain value in realistic, portfolio-friendly language.

PROJECT:
{project.get("name")}

EVIDENCE:
{evidence}

Return:
1. Problem this project appears to address
2. Practical value
3. User/developer value
4. Security/quality value
5. Business-style pitch
6. Limitations / not visible in indexed files
"""

    return ask_llm(prompt)


def generate_demo_script(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    prompt = f"""
You are JARVIS, a demo coach.

Create a 2-3 minute demo script for presenting this project.

Grounding rules:
- Use only indexed project evidence.
- Do not invent screens, features, flows, users, production data, or deployment.
- If something is not visible, say it is not visible.
- Make it natural for an interview or portfolio presentation.

PROJECT:
{project.get("name")}

EVIDENCE:
{evidence}

Return:
1. 30-second elevator pitch
2. 2-minute demo script
3. Technical deep-dive script
4. Questions interviewers may ask
5. Best answers based on visible evidence
"""

    return ask_llm(prompt)


def generate_portfolio_pack(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = [
        f"# Portfolio Pack - {project.get('name')}",
        "",
        f"Generated: {_portfolio39_now()}",
        "",
        f"Project path: `{project.get('path')}`",
        f"Files indexed: `{project.get('files_count')}`",
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`",
        "",
        "Generated by: JARVIS Engineering OS",
        "",
        "Mode: evidence-based portfolio generation / verify manually before publishing.",
        "",
        "## CV Summary",
        "",
        generate_cv_summary(project_name),
        "",
        "## GitHub Showcase",
        "",
        generate_github_showcase(project_name),
        "",
        "## LinkedIn Showcase",
        "",
        generate_linkedin_showcase(project_name),
        "",
        "## Recruiter Summary",
        "",
        generate_recruiter_summary(project_name),
        "",
        "## Technical Highlights",
        "",
        generate_technical_highlights(project_name),
        "",
        "## Business Value Summary",
        "",
        generate_business_value_summary(project_name),
        "",
        "## Demo Script",
        "",
        generate_demo_script(project_name),
    ]

    return "\n".join(content)


def generate_portfolio_pack_strict(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    evidence, evidence_error = _portfolio39_evidence(project_name)

    if evidence_error:
        return evidence_error

    return (
        f"# Portfolio Pack - {project.get('name')}\n\n"
        f"Generated: {_portfolio39_now()}\n\n"
        f"Project path: `{project.get('path')}`\n\n"
        f"Files indexed: `{project.get('files_count')}`\n\n"
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`\n\n"
        "## Portfolio Summary\n\n"
        "This project can be presented as a portfolio project based on the indexed technical evidence below.\n\n"
        "## Technical Evidence\n\n"
        "```text\n"
        + evidence["grounded"]
        + "\n```\n\n"
        "## Architecture Evidence\n\n"
        "```text\n"
        + evidence["architecture"]
        + "\n```\n\n"
        "## Security Evidence\n\n"
        "```text\n"
        + evidence["security"]
        + "\n```\n\n"
        "## Testing / Quality Evidence\n\n"
        "```text\n"
        + evidence["test_coverage"]
        + "\n```\n\n"
        "## Honest Portfolio Notes\n\n"
        "- Do not claim production deployment unless deployment evidence exists.\n"
        "- Do not claim real users or business impact unless externally verified.\n"
        "- Emphasize architecture, security thinking, documentation, and testing roadmap.\n"
    )


def export_portfolio_pack(project_name, strict=False):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        generate_portfolio_pack_strict(project_name)
        if strict
        else generate_portfolio_pack(project_name)
    )

    path = _portfolio39_save_file(
        project.get("name", project_name),
        "portfolio_pack",
        content
    )

    return f"{content}\n\nPORTFOLIO PACK EXPORTED:\n{path}"


def export_portfolio_pack_strict(project_name):
    return export_portfolio_pack(
        project_name,
        strict=True
    )


def export_github_showcase(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_github_showcase(project_name)

    path = _portfolio39_save_file(
        project.get("name", project_name),
        "github_showcase",
        content
    )

    return f"{content}\n\nGITHUB SHOWCASE EXPORTED:\n{path}"


def export_linkedin_showcase(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_linkedin_showcase(project_name)

    path = _portfolio39_save_file(
        project.get("name", project_name),
        "linkedin_showcase",
        content
    )

    return f"{content}\n\nLINKEDIN SHOWCASE EXPORTED:\n{path}"


def export_cv_project_summary(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_cv_summary(project_name)

    path = _portfolio39_save_file(
        project.get("name", project_name),
        "cv_project_summary",
        content
    )

    return f"{content}\n\nCV PROJECT SUMMARY EXPORTED:\n{path}"


def portfolio_readiness(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    checks = {
        "README generator": "generate_readme" in globals(),
        "Architecture diagrams": "generate_architecture_diagram_pack" in globals() or "generate_architecture_overview_diagram" in globals(),
        "API docs": "generate_api_docs" in globals(),
        "Changelog": "generate_changelog" in globals(),
        "Test coverage": "analyze_test_coverage" in globals(),
        "Release readiness": "release_readiness" in globals(),
    }

    output = [
        "PORTFOLIO READINESS CHECK",
        f"Project: {project.get('name')}",
        "",
    ]

    for name, ok in checks.items():
        output.append(f"- [{'x' if ok else ' '}] {name}")

    output.append("")
    output.append("Recommendation:")
    if all(checks.values()):
        output.append("- Portfolio pack can combine all major generated reports.")
    else:
        output.append("- Some generators are not loaded. Portfolio pack will still work, but may be less complete.")

    return "\n".join(output)


# Friendly aliases
def portfolio_pack(project_name):
    return generate_portfolio_pack(project_name)


def generate_project_portfolio_pack(project_name):
    return generate_portfolio_pack(project_name)


def github_showcase(project_name):
    return generate_github_showcase(project_name)


def linkedin_showcase(project_name):
    return generate_linkedin_showcase(project_name)


def cv_summary(project_name):
    return generate_cv_summary(project_name)


def recruiter_summary(project_name):
    return generate_recruiter_summary(project_name)


def technical_highlights(project_name):
    return generate_technical_highlights(project_name)


def business_value(project_name):
    return generate_business_value_summary(project_name)


def demo_script(project_name):
    return generate_demo_script(project_name)


def export_project_portfolio_pack(project_name):
    return export_portfolio_pack(project_name)



# ==========================
# STEP 40 - FINAL ENTERPRISE RELEASE PACK
# Combines README, architecture, API docs, changelog, release notes, test coverage,
# quality report, portfolio pack, security audit, release readiness and executive summary.
# Evidence-based. Safe export only.
# ==========================
ENTERPRISE_RELEASE_PACK_OUTPUT_DIR = "enterprise_release_pack_generated"


def _release40_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _release40_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _release40_safe_filename(name):
    return _safe_report_filename(str(name))


def _release40_save_file(project_name, suffix, content):
    os.makedirs(
        ENTERPRISE_RELEASE_PACK_OUTPUT_DIR,
        exist_ok=True
    )

    safe_project = _release40_safe_filename(project_name)

    path = os.path.join(
        ENTERPRISE_RELEASE_PACK_OUTPUT_DIR,
        f"{safe_project}_{suffix}_{_release40_timestamp()}.md"
    )

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as f:
        f.write(content)

    return path


def _release40_section(title, content, max_chars=18000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n... section truncated for final pack readability."

    return (
        f"\n\n## {title}\n\n"
        "```text\n"
        + text
        + "\n```"
    )


def _release40_markdown_section(title, content, max_chars=22000):
    text = str(content).strip()

    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n... section truncated for final pack readability."

    return f"\n\n## {title}\n\n{text}\n"


def _release40_module_status():
    checks = {
        "README Generator": "generate_readme_strict" in globals(),
        "Architecture Diagram Generator": "generate_architecture_diagram_pack" in globals(),
        "API Documentation Generator": "generate_api_docs_strict" in globals(),
        "Changelog Generator": "generate_changelog" in globals(),
        "Release Notes Generator": "generate_release_notes_strict" in globals(),
        "Test Coverage Intelligence": "analyze_test_coverage" in globals(),
        "Quality Report": "generate_quality_report" in globals(),
        "Portfolio Pack Generator": "generate_portfolio_pack_strict" in globals(),
        "Release Readiness": "release_readiness" in globals(),
        "Security Analyzer": "strict_security_analyzer_project" in globals(),
        "Architecture Analyzer": "strict_architecture_analyzer_project" in globals(),
        "Grounded Analyzer": "strict_grounded_analyzer_project" in globals(),
    }

    output = [
        "MODULE STATUS",
        "Mode: loaded function availability",
        ""
    ]

    for name, loaded in checks.items():
        output.append(f"- [{'x' if loaded else ' '}] {name}")

    return "\n".join(output)


def generate_executive_summary(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    score = score_project(project_name)
    release = release_readiness(project_name) if "release_readiness" in globals() else "Release readiness module not loaded."
    coverage = analyze_test_coverage(project_name) if "analyze_test_coverage" in globals() else "Test coverage intelligence not loaded."
    security = strict_security_analyzer_project(project_name)
    architecture = strict_architecture_analyzer_project(project_name)

    prompt = f"""
You are JARVIS, an enterprise software release advisor.

Generate an executive summary for this project.

Grounding rules:
- Use only the indexed project evidence below.
- Do not invent production deployment, real customers, revenue, users, cloud infrastructure, or business impact.
- If something is not visible, write: Not visible in indexed files.
- Be clear, honest, and suitable for a final release/portfolio report.

PROJECT:
{project.get("name")}

PROJECT PATH:
{project.get("path")}

TECH STACK:
{", ".join(project.get("tech_stack", []))}

PROJECT SCORE:
{score}

RELEASE READINESS:
{release}

TEST COVERAGE:
{coverage}

SECURITY:
{security}

ARCHITECTURE:
{architecture}

Return Markdown with:
1. Executive overview
2. Project value
3. Technical maturity
4. Security maturity
5. Testing and quality maturity
6. Release readiness
7. Main risks
8. Final recommendation
"""

    return ask_llm(prompt)


def generate_project_audit(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = [
        f"# Project Audit - {project.get('name')}",
        "",
        f"Generated: {_release40_now()}",
        "",
        f"Project path: `{project.get('path')}`",
        f"Files indexed: `{project.get('files_count')}`",
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`",
        "",
        "Generated by: JARVIS Engineering OS",
        "",
        "Mode: evidence-based audit / verify manually before production.",
    ]

    content.append(
        _release40_section(
            "Module Status",
            _release40_module_status()
        )
    )

    content.append(
        _release40_section(
            "Project Score",
            score_project(project_name)
        )
    )

    content.append(
        _release40_section(
            "Grounded Analyzer",
            strict_grounded_analyzer_project(project_name)
        )
    )

    content.append(
        _release40_section(
            "Security Audit",
            strict_security_analyzer_project(project_name)
        )
    )

    content.append(
        _release40_section(
            "Architecture Audit",
            strict_architecture_analyzer_project(project_name)
        )
    )

    content.append(
        _release40_section(
            "Dead Code Scan",
            find_dead_code(project_name)
        )
    )

    content.append(
        _release40_section(
            "Duplicate Logic Scan",
            find_duplicate_code(project_name)
        )
    )

    if "analyze_test_coverage" in globals():
        content.append(
            _release40_section(
                "Test Coverage Intelligence",
                analyze_test_coverage(project_name)
            )
        )

    if "release_readiness" in globals():
        content.append(
            _release40_section(
                "Release Readiness",
                release_readiness(project_name)
            )
        )

    if "production_readiness" in globals():
        content.append(
            _release40_section(
                "Production Readiness",
                production_readiness(project_name)
            )
        )

    return "\n".join(content)


def generate_final_release_report(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = [
        f"# Final Release Report - {project.get('name')}",
        "",
        f"Generated: {_release40_now()}",
        "",
        f"Project path: `{project.get('path')}`",
        f"Files indexed: `{project.get('files_count')}`",
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`",
        "",
        "Generated by: JARVIS Engineering OS",
        "",
        "Mode: final release review / evidence-based / manual verification required.",
    ]

    content.append(
        _release40_markdown_section(
            "Executive Summary",
            generate_executive_summary(project_name)
        )
    )

    if "release_readiness" in globals():
        content.append(
            _release40_section(
                "Release Readiness",
                release_readiness(project_name)
            )
        )

    if "go_live_report" in globals():
        content.append(
            _release40_section(
                "Go-Live Report",
                go_live_report(project_name)
            )
        )

    if "generate_quality_report" in globals():
        content.append(
            _release40_markdown_section(
                "Quality Report",
                generate_quality_report(project_name)
            )
        )

    if "generate_changelog" in globals():
        content.append(
            _release40_markdown_section(
                "Changelog",
                generate_changelog(project_name)
            )
        )

    if "generate_release_notes_strict" in globals():
        content.append(
            _release40_markdown_section(
                "Release Notes",
                generate_release_notes_strict(project_name)
            )
        )

    content.append(
        "\n\n## Final Manual Checklist\n\n"
        "- [ ] Run all tests locally.\n"
        "- [ ] Run dependency/security audit.\n"
        "- [ ] Confirm environment variables and secrets.\n"
        "- [ ] Confirm database/storage migration plan if needed.\n"
        "- [ ] Confirm logging and monitoring strategy.\n"
        "- [ ] Confirm backup and rollback strategy.\n"
        "- [ ] Confirm README/setup instructions.\n"
        "- [ ] Confirm demo script and portfolio claims are honest.\n"
        "- [ ] Confirm no production claims are made without evidence.\n"
    )

    return "\n".join(content)


def generate_enterprise_release_pack(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = [
        f"# Enterprise Release Pack - {project.get('name')}",
        "",
        f"Generated: {_release40_now()}",
        "",
        f"Project path: `{project.get('path')}`",
        f"Files indexed: `{project.get('files_count')}`",
        f"Tech stack: `{', '.join(project.get('tech_stack', []))}`",
        "",
        "Generated by: JARVIS Engineering OS",
        "",
        "Mode: enterprise pack / evidence-based / manual verification required.",
        "",
        "> This pack combines all major JARVIS analysis modules into one release-ready documentation set.",
    ]

    content.append(
        _release40_section(
            "Module Status",
            _release40_module_status()
        )
    )

    content.append(
        _release40_markdown_section(
            "Executive Summary",
            generate_executive_summary(project_name)
        )
    )

    if "generate_readme_strict" in globals():
        content.append(
            _release40_markdown_section(
                "Generated README",
                generate_readme_strict(project_name)
            )
        )

    if "generate_architecture_diagram_pack" in globals():
        content.append(
            _release40_markdown_section(
                "Architecture Diagram Pack",
                generate_architecture_diagram_pack(project_name)
            )
        )
    elif "generate_architecture_report" in globals():
        content.append(
            _release40_markdown_section(
                "Architecture Report",
                generate_architecture_report(project_name)
            )
        )

    if "generate_api_docs_strict" in globals():
        content.append(
            _release40_markdown_section(
                "API Documentation",
                generate_api_docs_strict(project_name)
            )
        )

    if "generate_changelog" in globals():
        content.append(
            _release40_markdown_section(
                "Changelog",
                generate_changelog(project_name)
            )
        )

    if "generate_release_notes_strict" in globals():
        content.append(
            _release40_markdown_section(
                "Release Notes",
                generate_release_notes_strict(project_name)
            )
        )

    if "analyze_test_coverage" in globals():
        content.append(
            _release40_section(
                "Test Coverage Intelligence",
                analyze_test_coverage(project_name)
            )
        )

    if "generate_testing_roadmap" in globals():
        content.append(
            _release40_section(
                "Testing Roadmap",
                generate_testing_roadmap(project_name)
            )
        )

    if "generate_quality_report" in globals():
        content.append(
            _release40_markdown_section(
                "Quality Report",
                generate_quality_report(project_name)
            )
        )

    content.append(
        _release40_markdown_section(
            "Project Audit",
            generate_project_audit(project_name)
        )
    )

    if "generate_portfolio_pack_strict" in globals():
        content.append(
            _release40_markdown_section(
                "Portfolio Pack",
                generate_portfolio_pack_strict(project_name)
            )
        )

    content.append(
        "\n\n## Final Enterprise Recommendation\n\n"
        "Use this pack as a final documentation and review artifact. "
        "Before publishing or deploying the project, manually verify all generated findings, "
        "run the real test suite, confirm secrets/configuration, and check that every public claim "
        "is supported by the actual repository.\n"
    )

    return "\n".join(content)


def export_executive_summary(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = (
        f"# Executive Summary - {project.get('name')}\n\n"
        f"Generated: {_release40_now()}\n\n"
        + generate_executive_summary(project_name)
    )

    path = _release40_save_file(
        project.get("name", project_name),
        "EXECUTIVE_SUMMARY",
        content
    )

    return f"{content}\n\nEXECUTIVE SUMMARY EXPORTED:\n{path}"


def export_project_audit(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_project_audit(project_name)

    path = _release40_save_file(
        project.get("name", project_name),
        "PROJECT_AUDIT",
        content
    )

    return f"{content}\n\nPROJECT AUDIT EXPORTED:\n{path}"


def export_final_release(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_final_release_report(project_name)

    path = _release40_save_file(
        project.get("name", project_name),
        "FINAL_RELEASE_REPORT",
        content
    )

    return f"{content}\n\nFINAL RELEASE REPORT EXPORTED:\n{path}"


def export_enterprise_release_pack(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    content = generate_enterprise_release_pack(project_name)

    path = _release40_save_file(
        project.get("name", project_name),
        "ENTERPRISE_RELEASE_PACK",
        content
    )

    return f"{content}\n\nENTERPRISE RELEASE PACK EXPORTED:\n{path}"


def export_complete_release_bundle(project_name):
    outputs = [
        export_executive_summary(project_name),
        export_project_audit(project_name),
        export_final_release(project_name),
        export_enterprise_release_pack(project_name),
    ]

    paths = []

    for output in outputs:
        for marker in [
            "EXECUTIVE SUMMARY EXPORTED:",
            "PROJECT AUDIT EXPORTED:",
            "FINAL RELEASE REPORT EXPORTED:",
            "ENTERPRISE RELEASE PACK EXPORTED:",
        ]:
            if marker in output:
                paths.append(output.split(marker)[-1].strip())

    return (
        "COMPLETE ENTERPRISE RELEASE BUNDLE EXPORTED\n\n"
        + "\n".join(f"- {path}" for path in paths)
    )


def enterprise_release_readiness(project_name):
    project, error = get_project(project_name)

    if error:
        return error

    checks = {
        "README": "generate_readme_strict" in globals(),
        "Architecture": "generate_architecture_diagram_pack" in globals(),
        "API Docs": "generate_api_docs_strict" in globals(),
        "Changelog": "generate_changelog" in globals(),
        "Release Notes": "generate_release_notes_strict" in globals(),
        "Test Coverage": "analyze_test_coverage" in globals(),
        "Portfolio Pack": "generate_portfolio_pack_strict" in globals(),
        "Release Readiness": "release_readiness" in globals(),
        "Quality Report": "generate_quality_report" in globals(),
    }

    loaded = sum(1 for ok in checks.values() if ok)
    total = len(checks)
    percent = round((loaded / total) * 100, 1)

    output = [
        "ENTERPRISE RELEASE PACK READINESS",
        f"Project: {project.get('name')}",
        "",
        f"Modules loaded: {loaded}/{total} ({percent}%)",
        "",
    ]

    for name, ok in checks.items():
        output.append(f"- [{'x' if ok else ' '}] {name}")

    output.append("")
    if percent == 100:
        output.append("Status: COMPLETE - all major release pack modules are available.")
    elif percent >= 75:
        output.append("Status: STRONG - most major modules are available.")
    else:
        output.append("Status: PARTIAL - some major modules are missing.")

    return "\n".join(output)


# Friendly aliases
def final_release(project_name):
    return generate_final_release_report(project_name)


def enterprise_release_pack(project_name):
    return generate_enterprise_release_pack(project_name)


def final_enterprise_pack(project_name):
    return generate_enterprise_release_pack(project_name)


def executive_summary(project_name):
    return generate_executive_summary(project_name)


def project_audit(project_name):
    return generate_project_audit(project_name)


def export_final_enterprise_pack(project_name):
    return export_enterprise_release_pack(project_name)


def export_release_bundle(project_name):
    return export_complete_release_bundle(project_name)


def release_pack_readiness(project_name):
    return enterprise_release_readiness(project_name)

# ==========================================================
# J.A.R.V.I.S LOCAL ENTERPRISE REPORT ENGINE
# No dependency on developer_assistant.py.
# Supports: Markdown, HTML, TXT, JSON, CSV, DOCX/DOC, PDF, PPTX/PPT, XLSX/XLS.
# ==========================================================
REPORTS_DIR = "reports"


def _jarvis_ensure_reports_dir():
    os.makedirs(REPORTS_DIR, exist_ok=True)
    return REPORTS_DIR


def _jarvis_report_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _jarvis_safe_report_name(value):
    value = str(value).strip()
    value = re.sub(r"[^A-Za-z0-9_.-]+", "_", value)
    return value.strip("_") or "jarvis_report"


def _jarvis_open_file(path):
    try:
        os.startfile(os.path.abspath(path))
        return True
    except Exception:
        pass

    try:
        subprocess.Popen(
            ["cmd", "/c", "start", "", os.path.abspath(path)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            shell=False,
        )
        return True
    except Exception:
        return False


def _jarvis_write_text(path, content):
    with open(path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(str(content))


def _jarvis_content_lines(content):
    return str(content).replace("\r\n", "\n").replace("\r", "\n").split("\n")


def _jarvis_shorten(text, limit=900):
    text = str(text).strip()
    if len(text) <= limit:
        return text
    return text[:limit - 3] + "..."


def _jarvis_report_path(project_name, report_kind, extension):
    _jarvis_ensure_reports_dir()

    safe_project = _jarvis_safe_report_name(project_name)
    safe_kind = _jarvis_safe_report_name(report_kind)

    return os.path.join(
        REPORTS_DIR,
        f"{safe_project}_{safe_kind}_{_jarvis_report_timestamp()}.{extension}"
    )


def build_enterprise_project_review(project_name):
    sections = []

    def add_section(title, func):
        try:
            content = func(project_name)
        except Exception as e:
            content = f"Section failed safely: {e}"

        sections.append(
            "\n\n"
            + "=" * 78
            + f"\n{title}\n"
            + "=" * 78
            + "\n"
            + str(content)
        )

    sections.append(
        "J.A.R.V.I.S ENTERPRISE PROJECT REVIEW\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"Project: {project_name}\n"
        "Mode: safe analysis / no automatic code changes\n"
        "All recommendations must be verified with tests before being applied."
    )

    add_section("1. PROJECT SCORECARD", score_project)
    add_section("2. STRICT GROUNDED ANALYZER", strict_grounded_analyzer_project)
    add_section("3. STRICT SECURITY ANALYZER", strict_security_analyzer_project)
    add_section("4. STRICT ARCHITECTURE ANALYZER", strict_architecture_analyzer_project)
    add_section("5. DEAD CODE SCAN", find_dead_code)
    add_section("6. DUPLICATE CODE SCAN", find_duplicate_code)
    add_section("7. PROJECT STRUCTURE", analyze_project_structure)

    # LLM-based sections are kept, but they fail safely if Ollama is not ready.
    add_section("8. IMPROVEMENT ROADMAP", generate_improvement_roadmap)
    add_section("9. PROJECT DOCUMENTATION", generate_project_documentation)

    return "\n".join(sections)


def build_security_project_review(project_name):
    sections = []

    def add_section(title, func):
        try:
            content = func(project_name)
        except Exception as e:
            content = f"Section failed safely: {e}"

        sections.append(
            "\n\n"
            + "=" * 78
            + f"\n{title}\n"
            + "=" * 78
            + "\n"
            + str(content)
        )

    sections.append(
        "J.A.R.V.I.S SECURITY REVIEW\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"Project: {project_name}\n"
        "Mode: safe security analysis / no automatic code changes"
    )

    add_section("1. PROJECT SCORECARD", score_project)
    add_section("2. STRICT SECURITY ANALYZER", strict_security_analyzer_project)
    add_section("3. FULL SECURITY AUDIT", full_security_audit)

    try:
        add_section("4. HIGHEST RISK VULNERABILITIES", highest_risk_vulnerabilities)
    except Exception:
        pass

    return "\n".join(sections)


def build_architecture_project_review(project_name):
    sections = []

    def add_section(title, func):
        try:
            content = func(project_name)
        except Exception as e:
            content = f"Section failed safely: {e}"

        sections.append(
            "\n\n"
            + "=" * 78
            + f"\n{title}\n"
            + "=" * 78
            + "\n"
            + str(content)
        )

    sections.append(
        "J.A.R.V.I.S ARCHITECTURE REVIEW\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"Project: {project_name}\n"
        "Mode: safe architecture analysis / no automatic code changes"
    )

    add_section("1. PROJECT SCORECARD", score_project)
    add_section("2. STRICT ARCHITECTURE ANALYZER", strict_architecture_analyzer_project)
    add_section("3. PROJECT STRUCTURE", analyze_project_structure)
    add_section("4. ARCHITECTURE REPORT", generate_architecture_report)

    return "\n".join(sections)


def build_score_project_review(project_name):
    return score_project(project_name)


def build_report_content_by_kind(project_name, report_kind="project_review"):
    kind = str(report_kind).lower().strip()

    if kind in {"security", "security_review", "audit", "full_security_audit"}:
        return build_security_project_review(project_name)

    if kind in {"architecture", "architecture_review"}:
        return build_architecture_project_review(project_name)

    if kind in {"score", "score_report"}:
        return build_score_project_review(project_name)

    return build_enterprise_project_review(project_name)


def export_markdown_report(project_name, content, report_kind="project_review", open_after=True):
    path = _jarvis_report_path(project_name, report_kind, "md")

    body = (
        f"# J.A.R.V.I.S {str(report_kind).replace('_', ' ').title()}\n\n"
        f"**Project:** {project_name}\n\n"
        f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        "```text\n"
        f"{content}\n"
        "```\n"
    )

    _jarvis_write_text(path, body)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_txt_report(project_name, content, report_kind="project_review", open_after=True):
    path = _jarvis_report_path(project_name, report_kind, "txt")
    _jarvis_write_text(path, content)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_html_report(project_name, content, report_kind="project_review", open_after=True):
    import html

    path = _jarvis_report_path(project_name, report_kind, "html")

    doc = f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>J.A.R.V.I.S {html.escape(str(report_kind))}</title>
<style>
body {{
    font-family: Arial, sans-serif;
    margin: 40px;
    background: #0b1020;
    color: #e7f6ff;
}}
h1 {{ color: #58d8ff; }}
pre {{
    white-space: pre-wrap;
    background: #111a33;
    border: 1px solid #263a70;
    border-radius: 12px;
    padding: 20px;
    line-height: 1.45;
}}
.meta {{ color: #9fb3c8; }}
</style>
</head>
<body>
<h1>J.A.R.V.I.S {html.escape(str(report_kind).replace('_', ' ').title())}</h1>
<p class="meta"><b>Project:</b> {html.escape(str(project_name))}</p>
<p class="meta"><b>Generated:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
<pre>{html.escape(str(content))}</pre>
</body>
</html>"""

    _jarvis_write_text(path, doc)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_json_report(project_name, content, report_kind="project_review", open_after=True):
    path = _jarvis_report_path(project_name, report_kind, "json")

    data = {
        "project": project_name,
        "report_kind": report_kind,
        "generated_at": datetime.now().isoformat(),
        "content": str(content),
    }

    with open(path, "w", encoding="utf-8", errors="ignore") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_csv_report(project_name, content, report_kind="project_review", open_after=True):
    path = _jarvis_report_path(project_name, report_kind, "csv")

    with open(path, "w", encoding="utf-8", errors="ignore", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["project", "report_kind", "line_number", "text"])

        for index, line in enumerate(_jarvis_content_lines(content), start=1):
            writer.writerow([project_name, report_kind, index, line])

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_docx_report(project_name, content, report_kind="project_review", open_after=True):
    try:
        from docx import Document
    except Exception:
        return export_markdown_report(
            project_name,
            "python-docx is not installed. Markdown fallback created.\n\n" + str(content),
            report_kind + "_docx_fallback",
            open_after=open_after
        )

    path = _jarvis_report_path(project_name, report_kind, "docx")

    doc = Document()
    doc.add_heading(f"J.A.R.V.I.S {str(report_kind).replace('_', ' ').title()}", level=1)
    doc.add_paragraph(f"Project: {project_name}")
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_paragraph("Mode: safe analysis / no automatic code changes")

    for block in str(content).split("\n\n"):
        block = block.strip()

        if not block:
            continue

        if len(block) <= 90 and block.isupper():
            doc.add_heading(block.title(), level=2)
        else:
            doc.add_paragraph(block)

    doc.save(path)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_pdf_report(project_name, content, report_kind="project_review", open_after=True):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception:
        return export_markdown_report(
            project_name,
            "reportlab is not installed. Markdown fallback created.\n\n" + str(content),
            report_kind + "_pdf_fallback",
            open_after=open_after
        )

    path = _jarvis_report_path(project_name, report_kind, "pdf")

    c = canvas.Canvas(path, pagesize=A4)
    width, height = A4
    x = 42
    y = height - 50

    c.setFont("Helvetica-Bold", 15)
    c.drawString(x, y, f"J.A.R.V.I.S {str(report_kind).replace('_', ' ').title()}"[:90])
    y -= 22

    c.setFont("Helvetica", 9)
    c.drawString(x, y, f"Project: {project_name}"[:110])
    y -= 14
    c.drawString(x, y, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    y -= 24

    c.setFont("Helvetica", 8)

    for line in _jarvis_content_lines(content):
        chunks = [str(line)[i:i + 115] for i in range(0, len(str(line)), 115)] or [""]

        for chunk in chunks:
            if y < 45:
                c.showPage()
                c.setFont("Helvetica", 8)
                y = height - 45

            c.drawString(x, y, chunk)
            y -= 10

    c.save()

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_pptx_report(project_name, content, report_kind="project_review", open_after=True):
    try:
        from pptx import Presentation
    except Exception:
        return export_markdown_report(
            project_name,
            "python-pptx is not installed. Markdown fallback created.\n\n" + str(content),
            report_kind + "_pptx_fallback",
            open_after=open_after
        )

    path = _jarvis_report_path(project_name, report_kind, "pptx")

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"J.A.R.V.I.S {str(report_kind).replace('_', ' ').title()}"[:80]
    slide.placeholders[1].text = (
        f"Project: {project_name}\n"
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    )

    sections = re.split(r"\n={10,}\n", str(content))

    for section in sections[:18]:
        section = section.strip()

        if not section:
            continue

        lines = section.splitlines()
        title = _jarvis_shorten(lines[0] if lines else "Report Section", 60)
        body = _jarvis_shorten("\n".join(lines[1:]) if len(lines) > 1 else section, 900)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    prs.save(path)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_xlsx_report(project_name, content, report_kind="project_review", open_after=True):
    path = _jarvis_report_path(project_name, report_kind, "xlsx")

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
    except Exception:
        return export_csv_report(project_name, content, report_kind + "_xlsx_fallback", open_after=open_after)

    wb = Workbook()

    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = f"J.A.R.V.I.S {str(report_kind).replace('_', ' ').title()}"
    ws["A1"].font = Font(bold=True, size=16)
    ws["A3"] = "Project"
    ws["B3"] = project_name
    ws["A4"] = "Report type"
    ws["B4"] = report_kind
    ws["A5"] = "Generated"
    ws["B5"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    lines = _jarvis_content_lines(content)

    ws["A7"] = "Report preview"
    ws["A7"].font = Font(bold=True)

    for row_index, line in enumerate(lines[:500], start=8):
        ws.cell(row=row_index, column=1).value = line

    ws.column_dimensions["A"].width = 120
    ws.column_dimensions["B"].width = 40

    ws2 = wb.create_sheet("Report Lines")
    ws2.append(["Line", "Text"])

    for cell in ws2[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="D9EAF7")

    for index, line in enumerate(lines, start=1):
        ws2.append([index, line])

    ws2.column_dimensions["A"].width = 12
    ws2.column_dimensions["B"].width = 140

    ws3 = wb.create_sheet("Risk Matrix")
    ws3.append(["Severity", "Finding"])

    for cell in ws3[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="FCE4D6")

    risk_keywords = ["critical", "high", "medium", "low", "risk", "vulnerability", "security", "weak", "secret", "password"]

    for line in lines:
        lower = line.lower()

        if any(keyword in lower for keyword in risk_keywords):
            if "critical" in lower:
                severity = "CRITICAL"
            elif "high" in lower:
                severity = "HIGH"
            elif "medium" in lower:
                severity = "MEDIUM"
            elif "low" in lower:
                severity = "LOW"
            elif "security" in lower:
                severity = "SECURITY"
            else:
                severity = "INFO"

            ws3.append([severity, line])

    ws3.column_dimensions["A"].width = 18
    ws3.column_dimensions["B"].width = 140

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

    wb.save(path)

    if open_after:
        _jarvis_open_file(path)

    return os.path.abspath(path)


def export_report_content(project_name, content, format_type="md", report_kind="project_review", open_after=True):
    fmt = str(format_type).lower().strip()

    if fmt in {"md", "markdown"}:
        path = export_markdown_report(project_name, content, report_kind, open_after)
    elif fmt in {"txt", "text"}:
        path = export_txt_report(project_name, content, report_kind, open_after)
    elif fmt in {"html", "web"}:
        path = export_html_report(project_name, content, report_kind, open_after)
    elif fmt == "json":
        path = export_json_report(project_name, content, report_kind, open_after)
    elif fmt == "csv":
        path = export_csv_report(project_name, content, report_kind, open_after)
    elif fmt in {"docx", "doc", "word"}:
        path = export_docx_report(project_name, content, report_kind, open_after)
    elif fmt == "pdf":
        path = export_pdf_report(project_name, content, report_kind, open_after)
    elif fmt in {"pptx", "ppt", "powerpoint", "presentation"}:
        path = export_pptx_report(project_name, content, report_kind, open_after)
    elif fmt in {"xlsx", "xls", "excel", "spreadsheet"}:
        path = export_xlsx_report(project_name, content, report_kind, open_after)
    else:
        path = export_markdown_report(project_name, content, report_kind, open_after)

    return f"Report created:\n{path}"


def export_project_review(project_name, format_type="md", open_after=True):
    content = build_enterprise_project_review(project_name)

    return export_report_content(
        project_name,
        content,
        format_type=format_type,
        report_kind="project_review",
        open_after=open_after
    )


def export_security_review(project_name, format_type="md", open_after=True):
    content = build_security_project_review(project_name)

    return export_report_content(
        project_name,
        content,
        format_type=format_type,
        report_kind="security_review",
        open_after=open_after
    )


def export_architecture_review(project_name, format_type="md", open_after=True):
    content = build_architecture_project_review(project_name)

    return export_report_content(
        project_name,
        content,
        format_type=format_type,
        report_kind="architecture_review",
        open_after=open_after
    )


def export_score_report(project_name, format_type="md", open_after=True):
    content = build_score_project_review(project_name)

    return export_report_content(
        project_name,
        content,
        format_type=format_type,
        report_kind="score_report",
        open_after=open_after
    )


def export_all_project_reports(project_name, open_after=True):
    content = build_enterprise_project_review(project_name)

    paths = []

    for fmt in ["md", "html", "docx", "pdf", "pptx", "xlsx", "json", "csv"]:
        result = export_report_content(
            project_name,
            content,
            format_type=fmt,
            report_kind="project_review",
            open_after=False
        )
        paths.append(result.replace("Report created:\n", ""))

    if open_after and paths:
        _jarvis_open_file(paths[0])

    return "All project reports created:\n" + "\n".join(paths)


# Friendly / compatibility aliases
enterprise_project_report = build_enterprise_project_review
project_review_to_word = lambda project_name: export_project_review(project_name, "docx", True)
project_review_to_pdf = lambda project_name: export_project_review(project_name, "pdf", True)
project_review_to_ppt = lambda project_name: export_project_review(project_name, "pptx", True)
project_review_to_excel = lambda project_name: export_project_review(project_name, "xlsx", True)
project_review_to_html = lambda project_name: export_project_review(project_name, "html", True)
project_review_to_markdown = lambda project_name: export_project_review(project_name, "md", True)
project_review_to_json = lambda project_name: export_project_review(project_name, "json", True)
project_review_to_csv = lambda project_name: export_project_review(project_name, "csv", True)

export_report_docx = project_review_to_word
export_report_pdf = project_review_to_pdf
export_report_pptx = project_review_to_ppt
export_report_xlsx = project_review_to_excel
export_project_word_report = project_review_to_word
export_project_pdf_report = project_review_to_pdf
export_project_ppt_report = project_review_to_ppt
export_project_excel_report = project_review_to_excel
export_project_html_report = project_review_to_html
export_project_json_report = project_review_to_json
export_project_csv_report = project_review_to_csv

# ==========================================================
# FINAL EXPORT OVERRIDES
# These keep old command compatibility while supporting format selection.
# ==========================================================
def export_project_report(project_name, format_type="md", open_after=True):
    return export_project_review(project_name, format_type=format_type, open_after=open_after)


def export_report(project_name, format_type="md", open_after=True):
    return export_project_review(project_name, format_type=format_type, open_after=open_after)


def export_project_markdown_report(project_name):
    return export_project_review(project_name, format_type="md", open_after=True)


def export_project_excel_report(project_name):
    return export_project_review(project_name, format_type="xlsx", open_after=True)



# ==========================================================
# JARVIS ENTERPRISE REPORT EXPORTER
# Adds professional exports:
# - Word / DOCX
# - PDF
# - PowerPoint / PPTX
# - Excel / XLSX with charts
# - Markdown / HTML / JSON / CSV
#
# These functions are intentionally appended at the end so they override
# older/simple report export functions without breaking existing imports.
# ==========================================================

ENTERPRISE_REPORT_DIR = "reports"


def _enterprise_ensure_dir(path=ENTERPRISE_REPORT_DIR):
    os.makedirs(path, exist_ok=True)
    return path


def _enterprise_timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _enterprise_safe_filename(name):
    cleaned = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(name).strip())
    cleaned = cleaned.strip("_")
    return cleaned or "jarvis_report"


def _enterprise_open_file(path):
    try:
        os.startfile(os.path.abspath(path))
        return True
    except Exception:
        pass

    try:
        subprocess.Popen(
            ["cmd", "/c", "start", "", os.path.abspath(path)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            shell=False,
        )
        return True
    except Exception:
        return False


def _enterprise_normalize_format(format_type):
    fmt = str(format_type or "pdf").lower().strip()

    aliases = {
        "doc": "word",
        "docx": "word",
        "word": "word",
        "pdf": "pdf",
        "ppt": "powerpoint",
        "pptx": "powerpoint",
        "presentation": "powerpoint",
        "powerpoint": "powerpoint",
        "excel": "excel",
        "xls": "excel",
        "xlsx": "excel",
        "spreadsheet": "excel",
        "md": "markdown",
        "markdown": "markdown",
        "html": "html",
        "json": "json",
        "csv": "csv",
    }

    return aliases.get(fmt, fmt)


def _enterprise_extract_scores(content):
    text = str(content or "")

    result = {
        "security": None,
        "architecture": None,
        "maintainability": None,
        "overall": None,
    }

    patterns = {
        "security": [
            r"security[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*/\s*10",
            r"security[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*%",
        ],
        "architecture": [
            r"architecture[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*/\s*10",
            r"architecture[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*%",
        ],
        "maintainability": [
            r"maintainability[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*/\s*10",
            r"maintainability[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*%",
        ],
        "overall": [
            r"overall[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*/\s*10",
            r"overall[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*%",
            r"code quality[^0-9]{0,40}([0-9]+(?:\.[0-9]+)?)\s*/\s*10",
        ],
    }

    for key, key_patterns in patterns.items():
        for pattern in key_patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)

            if match:
                try:
                    value = float(match.group(1))

                    if value > 10:
                        value = value / 10

                    result[key] = round(max(0, min(10, value)), 1)
                    break
                except Exception:
                    pass

    available = [
        value for value in result.values()
        if isinstance(value, (int, float))
    ]

    if result["overall"] is None and available:
        result["overall"] = round(sum(available) / len(available), 1)

    return result


def _enterprise_extract_findings(content, limit=40):
    lines = [line.strip() for line in str(content or "").splitlines()]
    findings = []

    keywords = [
        "high",
        "medium",
        "low",
        "risk",
        "vulnerability",
        "security",
        "issue",
        "fix",
        "improvement",
        "recommend",
        "warning",
        "critical",
        "positive",
    ]

    for line in lines:
        if not line:
            continue

        lower = line.lower()

        if any(keyword in lower for keyword in keywords):
            findings.append(line[:500])

        if len(findings) >= limit:
            break

    return findings


def _enterprise_split_sections(content):
    text = str(content or "").strip()

    if not text:
        return [("Report", "No content generated.")]

    sections = []
    current_title = "Executive Summary"
    current_lines = []

    for raw_line in text.splitlines():
        line = raw_line.strip()

        is_heading = False

        if line.startswith("#"):
            is_heading = True
            title = line.strip("# ").strip() or "Section"
        elif re.match(r"^\d+\.\s+.+", line):
            is_heading = True
            title = re.sub(r"^\d+\.\s+", "", line).strip()
        elif line.isupper() and len(line) <= 80 and len(line.split()) <= 8:
            is_heading = True
            title = line.title()

        if is_heading:
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
                current_lines = []

            current_title = title
        else:
            current_lines.append(raw_line)

    if current_lines:
        sections.append((current_title, "\n".join(current_lines).strip()))

    return sections[:40] or [("Report", text)]


def _enterprise_report_content(project_name, report_kind="project_review"):
    kind = str(report_kind or "project_review").lower().strip()

    if kind in {"security", "security_review", "audit", "security audit"}:
        return full_security_audit(project_name)

    if kind in {"architecture", "architecture_review", "architect"}:
        return strict_architecture_analyzer_project(project_name)

    if kind in {"score", "score_report", "scorecard"}:
        return score_project(project_name)

    if kind in {"documentation", "docs"}:
        return generate_project_documentation(project_name)

    if kind in {"roadmap", "improvement", "improvements"}:
        return generate_improvement_roadmap(project_name)

    if kind in {"dead_code", "dead code", "unused"}:
        return find_dead_code(project_name)

    if kind in {"duplicates", "duplicate_code"}:
        return find_duplicate_code(project_name)

    # Enterprise default: combine deterministic sections and grounded review.
    parts = [
        "ENTERPRISE PROJECT REVIEW",
        "",
        "PROJECT SCORECARD",
        score_project(project_name),
        "",
        "STRICT GROUNDED ANALYZER",
        strict_grounded_analyzer_project(project_name),
        "",
        "STRICT SECURITY ANALYZER",
        strict_security_analyzer_project(project_name),
        "",
        "STRICT ARCHITECTURE ANALYZER",
        strict_architecture_analyzer_project(project_name),
        "",
        "DEAD CODE SCAN",
        find_dead_code(project_name),
        "",
        "DUPLICATE CODE SCAN",
        find_duplicate_code(project_name),
    ]

    return "\n".join(str(part) for part in parts)


def build_enterprise_project_review(project_name):
    return _enterprise_report_content(project_name, report_kind="project_review")


def _enterprise_report_metadata(project_name, content, report_kind):
    project, error = get_project(project_name)

    if error:
        project = {
            "name": project_name,
            "path": "Unknown",
            "files_count": "Unknown",
            "tech_stack": [],
        }

    scores = _enterprise_extract_scores(content)

    return {
        "project_name": project.get("name", project_name),
        "project_path": project.get("path", "Unknown"),
        "files_count": project.get("files_count", "Unknown"),
        "tech_stack": project.get("tech_stack", []),
        "generated": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "report_kind": report_kind,
        "scores": scores,
        "findings": _enterprise_extract_findings(content),
    }


def _enterprise_write_markdown(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.md"
    )

    metadata = _enterprise_report_metadata(project_name, content, report_kind)

    text = (
        f"# J.A.R.V.I.S Enterprise Report - {metadata['project_name']}\n\n"
        f"Generated: {metadata['generated']}\n\n"
        f"Report type: `{metadata['report_kind']}`\n\n"
        f"Files indexed: `{metadata['files_count']}`\n\n"
        f"Tech stack: `{', '.join(metadata['tech_stack'])}`\n\n"
        "---\n\n"
        f"{content}\n"
    )

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(text)

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_html(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.html"
    )

    metadata = _enterprise_report_metadata(project_name, content, report_kind)

    escaped = (
        str(content)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )

    html = f"""<!doctype html>
<html>
<head>
<meta charset="utf-8">
<title>J.A.R.V.I.S Enterprise Report - {metadata['project_name']}</title>
<style>
body {{ font-family: Arial, sans-serif; margin: 40px; background: #f7f9fc; color: #1f2937; }}
.card {{ background: white; border-radius: 14px; padding: 24px; margin-bottom: 20px; box-shadow: 0 4px 18px rgba(0,0,0,.08); }}
h1 {{ color: #111827; }}
pre {{ white-space: pre-wrap; font-family: Consolas, monospace; background: #111827; color: #e5e7eb; padding: 20px; border-radius: 10px; }}
.badge {{ display: inline-block; background: #dbeafe; padding: 6px 10px; border-radius: 999px; margin: 4px; }}
</style>
</head>
<body>
<div class="card">
<h1>J.A.R.V.I.S Enterprise Report</h1>
<p><strong>Project:</strong> {metadata['project_name']}</p>
<p><strong>Generated:</strong> {metadata['generated']}</p>
<p><strong>Report type:</strong> {metadata['report_kind']}</p>
<p><strong>Files indexed:</strong> {metadata['files_count']}</p>
<p><strong>Tech stack:</strong> {', '.join(metadata['tech_stack'])}</p>
</div>
<div class="card">
<pre>{escaped}</pre>
</div>
</body>
</html>"""

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(html)

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_json(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.json"
    )

    metadata = _enterprise_report_metadata(project_name, content, report_kind)

    payload = {
        "metadata": metadata,
        "content": str(content),
        "sections": [
            {"title": title, "content": body}
            for title, body in _enterprise_split_sections(content)
        ],
    }

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        json.dump(payload, file, indent=2)

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_csv(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.csv"
    )

    import csv

    with open(path, "w", encoding="utf-8", errors="ignore", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["project", "report_kind", "line", "text"])

        for index, line in enumerate(str(content).splitlines(), start=1):
            writer.writerow([project_name, report_kind, index, line])

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_docx(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.docx"
    )

    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.section import WD_SECTION
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
    except Exception:
        return _enterprise_write_markdown(project_name, content, report_kind, open_after=open_after)

    metadata = _enterprise_report_metadata(project_name, content, report_kind)
    scores = metadata["scores"]

    doc = Document()

    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = doc.add_heading("J.A.R.V.I.S Enterprise Report", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(str(metadata["project_name"]))
    run.bold = True
    run.font.size = Pt(16)

    doc.add_paragraph(f"Generated: {metadata['generated']}")
    doc.add_paragraph(f"Report type: {metadata['report_kind']}")
    doc.add_paragraph(f"Files indexed: {metadata['files_count']}")
    doc.add_paragraph(f"Tech stack: {', '.join(metadata['tech_stack'])}")

    doc.add_heading("Executive Scorecard", level=1)

    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Score"

    for metric in ["security", "architecture", "maintainability", "overall"]:
        row = table.add_row().cells
        row[0].text = metric.title()
        row[1].text = "N/A" if scores.get(metric) is None else f"{scores.get(metric)}/10"

    findings = metadata["findings"]

    doc.add_heading("Key Findings", level=1)

    if findings:
        for finding in findings[:12]:
            doc.add_paragraph(finding, style="List Bullet")
    else:
        doc.add_paragraph("No key findings extracted automatically.")

    doc.add_heading("Full Technical Report", level=1)

    for title_text, body in _enterprise_split_sections(content):
        doc.add_heading(title_text[:80], level=2)

        for paragraph in str(body).splitlines():
            paragraph = paragraph.strip()

            if not paragraph:
                continue

            if paragraph.startswith("- "):
                doc.add_paragraph(paragraph[2:], style="List Bullet")
            else:
                doc.add_paragraph(paragraph)

    doc.save(path)

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_pdf(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.pdf"
    )

    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, KeepTogether
        )
    except Exception:
        return _enterprise_write_markdown(project_name, content, report_kind, open_after=open_after)

    metadata = _enterprise_report_metadata(project_name, content, report_kind)
    scores = metadata["scores"]

    doc = SimpleDocTemplate(
        path,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="JarvisTitle",
        parent=styles["Title"],
        fontSize=22,
        leading=26,
        spaceAfter=18,
    ))
    styles.add(ParagraphStyle(
        name="JarvisBody",
        parent=styles["BodyText"],
        fontSize=9,
        leading=12,
        spaceAfter=6,
    ))

    story = []

    story.append(Paragraph("J.A.R.V.I.S Enterprise Report", styles["JarvisTitle"]))
    story.append(Paragraph(f"<b>Project:</b> {metadata['project_name']}", styles["JarvisBody"]))
    story.append(Paragraph(f"<b>Generated:</b> {metadata['generated']}", styles["JarvisBody"]))
    story.append(Paragraph(f"<b>Report type:</b> {metadata['report_kind']}", styles["JarvisBody"]))
    story.append(Paragraph(f"<b>Files indexed:</b> {metadata['files_count']}", styles["JarvisBody"]))
    story.append(Spacer(1, 0.2 * inch))

    score_rows = [["Metric", "Score"]]

    for metric in ["security", "architecture", "maintainability", "overall"]:
        value = scores.get(metric)
        score_rows.append([metric.title(), "N/A" if value is None else f"{value}/10"])

    score_table = Table(score_rows, colWidths=[2.3 * inch, 2.0 * inch])
    score_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ALIGN", (1, 1), (1, -1), "CENTER"),
    ]))
    story.append(score_table)
    story.append(Spacer(1, 0.25 * inch))

    story.append(Paragraph("<b>Key Findings</b>", styles["Heading2"]))

    findings = metadata["findings"]

    if findings:
        for finding in findings[:15]:
            safe = str(finding).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph("• " + safe, styles["JarvisBody"]))
    else:
        story.append(Paragraph("No key findings extracted automatically.", styles["JarvisBody"]))

    story.append(PageBreak())
    story.append(Paragraph("<b>Full Technical Report</b>", styles["Heading1"]))

    for title_text, body in _enterprise_split_sections(content):
        safe_title = str(title_text[:80]).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        story.append(Paragraph(safe_title, styles["Heading2"]))

        for line in str(body).splitlines():
            line = line.strip()

            if not line:
                continue

            safe_line = line[:900].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe_line, styles["JarvisBody"]))

    doc.build(story)

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_excel(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.xlsx"
    )

    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference
        from openpyxl.styles import Font, PatternFill, Alignment
    except Exception:
        return _enterprise_write_csv(project_name, content, report_kind, open_after=open_after)

    metadata = _enterprise_report_metadata(project_name, content, report_kind)
    scores = metadata["scores"]

    wb = Workbook()

    ws = wb.active
    ws.title = "Executive Summary"

    ws["A1"] = "J.A.R.V.I.S Enterprise Report"
    ws["A1"].font = Font(bold=True, size=16)

    summary_rows = [
        ("Project", metadata["project_name"]),
        ("Generated", metadata["generated"]),
        ("Report type", metadata["report_kind"]),
        ("Files indexed", metadata["files_count"]),
        ("Tech stack", ", ".join(metadata["tech_stack"])),
    ]

    start_row = 3

    for offset, (key, value) in enumerate(summary_rows):
        ws.cell(row=start_row + offset, column=1).value = key
        ws.cell(row=start_row + offset, column=2).value = value
        ws.cell(row=start_row + offset, column=1).font = Font(bold=True)

    ws["A10"] = "Metric"
    ws["B10"] = "Score / 10"
    ws["A10"].font = Font(bold=True)
    ws["B10"].font = Font(bold=True)

    score_items = [
        ("Security", scores.get("security")),
        ("Architecture", scores.get("architecture")),
        ("Maintainability", scores.get("maintainability")),
        ("Overall", scores.get("overall")),
    ]

    for index, (metric, value) in enumerate(score_items, start=11):
        ws.cell(row=index, column=1).value = metric
        ws.cell(row=index, column=2).value = value if value is not None else 0

    chart = BarChart()
    chart.title = "Project Scorecard"
    chart.y_axis.title = "Score / 10"
    chart.x_axis.title = "Metric"

    data = Reference(ws, min_col=2, min_row=10, max_row=14)
    cats = Reference(ws, min_col=1, min_row=11, max_row=14)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 8
    chart.width = 14
    ws.add_chart(chart, "D10")

    ws2 = wb.create_sheet("Findings")
    ws2.append(["Index", "Finding"])

    for cell in ws2[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="D9EAF7")

    findings = metadata["findings"]

    for index, finding in enumerate(findings, start=1):
        ws2.append([index, finding])

    ws3 = wb.create_sheet("Full Report")
    ws3.append(["Line", "Text"])

    for cell in ws3[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="E2F0D9")

    for index, line in enumerate(str(content).splitlines(), start=1):
        ws3.append([index, line])

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        sheet.column_dimensions["A"].width = 18
        sheet.column_dimensions["B"].width = 120

    wb.save(path)

    if open_after:
        _enterprise_open_file(path)

    return path


def _enterprise_write_powerpoint(project_name, content, report_kind, open_after=True):
    _enterprise_ensure_dir()

    safe_name = _enterprise_safe_filename(project_name)
    safe_kind = _enterprise_safe_filename(report_kind)
    path = os.path.join(
        ENTERPRISE_REPORT_DIR,
        f"{safe_name}_{safe_kind}_{_enterprise_timestamp()}.pptx"
    )

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return _enterprise_write_markdown(project_name, content, report_kind, open_after=open_after)

    metadata = _enterprise_report_metadata(project_name, content, report_kind)
    scores = metadata["scores"]
    findings = metadata["findings"]
    sections = _enterprise_split_sections(content)

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "J.A.R.V.I.S Enterprise Report"
    slide.placeholders[1].text = (
        f"{metadata['project_name']}\n"
        f"{metadata['report_kind']} | {metadata['generated']}"
    )

    # Scorecard slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Scorecard"

    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)

    table = slide.shapes.add_table(5, 2, left, top, width, height).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Score"

    for row_index, metric in enumerate(["security", "architecture", "maintainability", "overall"], start=1):
        table.cell(row_index, 0).text = metric.title()
        value = scores.get(metric)
        table.cell(row_index, 1).text = "N/A" if value is None else f"{value}/10"

    # Findings slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Findings"
    body = slide.placeholders[1]
    body.text = "\n".join(f"• {finding[:150]}" for finding in findings[:8]) or "No key findings extracted automatically."

    # Section slides
    for title_text, body_text in sections[:10]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(title_text)[:70]

        lines = [
            line.strip()
            for line in str(body_text).splitlines()
            if line.strip()
        ]

        slide.placeholders[1].text = "\n".join(lines[:10])[:1200] or "No details."

    prs.save(path)

    if open_after:
        _enterprise_open_file(path)

    return path


def export_report_content(project_name, content, format_type="pdf", report_kind="project_review", open_after=True):
    fmt = _enterprise_normalize_format(format_type)

    if fmt == "word":
        path = _enterprise_write_docx(project_name, content, report_kind, open_after=open_after)
    elif fmt == "pdf":
        path = _enterprise_write_pdf(project_name, content, report_kind, open_after=open_after)
    elif fmt == "powerpoint":
        path = _enterprise_write_powerpoint(project_name, content, report_kind, open_after=open_after)
    elif fmt == "excel":
        path = _enterprise_write_excel(project_name, content, report_kind, open_after=open_after)
    elif fmt == "html":
        path = _enterprise_write_html(project_name, content, report_kind, open_after=open_after)
    elif fmt == "json":
        path = _enterprise_write_json(project_name, content, report_kind, open_after=open_after)
    elif fmt == "csv":
        path = _enterprise_write_csv(project_name, content, report_kind, open_after=open_after)
    else:
        path = _enterprise_write_markdown(project_name, content, report_kind, open_after=open_after)

    return f"Report created:\n{path}"


def export_report(project_name, format_type="pdf", report_kind="project_review", open_after=True):
    content = _enterprise_report_content(project_name, report_kind=report_kind)

    return export_report_content(
        project_name,
        content,
        format_type=format_type,
        report_kind=report_kind,
        open_after=open_after
    )


def export_project_report(project_name, format_type="markdown", report_kind="project_review", open_after=True):
    return export_report(
        project_name,
        format_type=format_type,
        report_kind=report_kind,
        open_after=open_after
    )


def export_project_markdown_report(project_name):
    return export_report(project_name, format_type="markdown", report_kind="project_review", open_after=True)


def create_project_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="project_review", open_after=True)


def create_security_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="security", open_after=True)


def create_architecture_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="architecture", open_after=True)


def create_score_report(project_name, format_type="excel"):
    return export_report(project_name, format_type=format_type, report_kind="score", open_after=True)


def generate_enterprise_report(project_name, format_type="pdf", report_kind="project_review"):
    return export_report(project_name, format_type=format_type, report_kind=report_kind, open_after=True)


def parse_and_generate_report(command):
    text = str(command or "").strip()
    lower = text.lower()

    fmt = _enterprise_normalize_format("pdf")

    for candidate in [
        "word", "docx", "doc", "pdf", "ppt", "pptx", "powerpoint",
        "presentation", "excel", "xlsx", "xls", "spreadsheet",
        "markdown", "md", "html", "json", "csv"
    ]:
        if re.search(rf"\b{re.escape(candidate)}\b", lower):
            fmt = _enterprise_normalize_format(candidate)
            break

    kind = "project_review"

    if "security" in lower or "audit" in lower:
        kind = "security"
    elif "architecture" in lower or "architect" in lower:
        kind = "architecture"
    elif "score" in lower or "scorecard" in lower:
        kind = "score"
    elif "documentation" in lower or "docs" in lower:
        kind = "documentation"
    elif "roadmap" in lower:
        kind = "roadmap"

    match = re.search(
        r"\b(?:for|from|of)\s+(?:project\s+)?(.+)$",
        text,
        flags=re.IGNORECASE
    )

    if match:
        project_name = match.group(1).strip(" .,:;")
        project_name = re.sub(
            r"\b(word|docx|doc|pdf|ppt|pptx|powerpoint|presentation|excel|xlsx|xls|spreadsheet|markdown|md|html|json|csv|report|review|security|architecture|score|documentation|roadmap)\b",
            "",
            project_name,
            flags=re.IGNORECASE
        ).strip()

        if not project_name:
            project_name = "CyberShield AI"
    else:
        project_name = "CyberShield AI"

    aliases = {
        "cyber": "CyberShield AI",
        "cyber shield": "CyberShield AI",
        "cyber shield ai": "CyberShield AI",
        "cybershield": "CyberShield AI",
        "cybershield ai": "CyberShield AI",
        "jarvis": "J.A.R.V.I.S",
        "jervis": "J.A.R.V.I.S",
        "manager app": "ManagerApp",
        "managerapp": "ManagerApp",
    }

    project_name = aliases.get(project_name.lower(), project_name)

    return export_report(project_name, format_type=fmt, report_kind=kind, open_after=True)



# ==========================================================
# J.A.R.V.I.S MARK XLVII PROJECT REVIEW UPGRADE
# Appended safely at the end of project_review_assistant.py.
#
# Adds:
# - stronger enterprise analysis orchestration
# - project health dashboard
# - executive summaries
# - risk matrix
# - action plans
# - autonomous review bundle
# - robust report exports in DOCX / PDF / PPTX / XLSX / HTML / JSON / CSV / MD
# - self-test helpers for J.A.R.V.I.S Mark XLVII
# ==========================================================

MARK47_REPORT_DIR = "reports"
MARK47_REVIEW_VERSION = "J.A.R.V.I.S Mark XLVII Project Review Engine"


def _m47_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _m47_stamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _m47_safe_name(name):
    cleaned = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(name or "").strip())
    cleaned = cleaned.strip("_")
    return cleaned or "jarvis_report"


def _m47_ensure_reports_dir():
    os.makedirs(MARK47_REPORT_DIR, exist_ok=True)
    return MARK47_REPORT_DIR


def _m47_open_file(path):
    try:
        os.startfile(os.path.abspath(path))
        return True
    except Exception:
        pass

    try:
        subprocess.Popen(
            ["cmd", "/c", "start", "", os.path.abspath(path)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            shell=False,
        )
        return True
    except Exception:
        return False


def _m47_normalize_format(format_type):
    value = str(format_type or "pdf").lower().strip()

    aliases = {
        "doc": "word",
        "docx": "word",
        "word": "word",
        "pdf": "pdf",
        "ppt": "powerpoint",
        "pptx": "powerpoint",
        "presentation": "powerpoint",
        "powerpoint": "powerpoint",
        "excel": "excel",
        "xls": "excel",
        "xlsx": "excel",
        "spreadsheet": "excel",
        "md": "markdown",
        "markdown": "markdown",
        "html": "html",
        "json": "json",
        "csv": "csv",
        "txt": "text",
        "text": "text",
    }

    return aliases.get(value, value)


def _m47_project_alias(project_name):
    lower = str(project_name or "").lower().strip()

    aliases = {
        "cyber": "CyberShield AI",
        "cyber shield": "CyberShield AI",
        "cyber shield ai": "CyberShield AI",
        "cybershield": "CyberShield AI",
        "cybershield ai": "CyberShield AI",
        "cybers in the": "CyberShield AI",
        "jarvis": "J.A.R.V.I.S",
        "jervis": "J.A.R.V.I.S",
        "j a r v i s": "J.A.R.V.I.S",
        "manager app": "ManagerApp",
        "managerapp": "ManagerApp",
    }

    return aliases.get(lower, str(project_name or "").strip())


def _m47_project_or_error(project_name):
    project_name = _m47_project_alias(project_name)
    project, error = get_project(project_name)

    if error:
        return None, error

    return project, None


def _m47_lines(text):
    return [line.strip() for line in str(text or "").splitlines() if line.strip()]


def _m47_limit(text, limit=8000):
    text = str(text or "")

    if len(text) <= limit:
        return text

    return text[:limit] + "\n\n[Output truncated for report section.]"


def _m47_extract_score_from_text(text, label=None):
    text = str(text or "")
    label_part = re.escape(label) if label else r"(?:overall|score|security|architecture|maintainability)"

    patterns = [
        rf"{label_part}[^0-9]{{0,45}}([0-9]+(?:\.[0-9]+)?)\s*/\s*10",
        rf"{label_part}[^0-9]{{0,45}}([0-9]+(?:\.[0-9]+)?)\s*%",
    ]

    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)

        if match:
            try:
                value = float(match.group(1))
                if value > 10:
                    value = value / 10
                return round(max(0, min(10, value)), 1)
            except Exception:
                pass

    return None


def _m47_scorecard_dict(project_name):
    try:
        card, error = project_scorecard(project_name)

        if not error and card:
            return card
    except Exception:
        pass

    try:
        score_text = score_project(project_name)
    except Exception as error:
        score_text = str(error)

    scores = {
        "security": _m47_extract_score_from_text(score_text, "security"),
        "architecture": _m47_extract_score_from_text(score_text, "architecture"),
        "maintainability": _m47_extract_score_from_text(score_text, "maintainability"),
        "overall": _m47_extract_score_from_text(score_text, "overall"),
    }

    values = [v for v in scores.values() if isinstance(v, (int, float))]

    if scores["overall"] is None and values:
        scores["overall"] = round(sum(values) / len(values), 1)

    return {
        "name": project_name,
        "path": "",
        "files_count": "",
        "tech_stack": [],
        "scores": scores,
        "reasons": {},
        "evidence": {},
    }


def _m47_collect_evidence_counts(project_name):
    project, facts, error = strict_project_facts(project_name)

    if error:
        return {}, error

    counts = {
        "files_indexed": facts.get("files_count", 0),
        "routes": len(facts.get("routes", [])),
        "frontend_api_calls": len(facts.get("frontend_api_calls", [])),
        "auth_files": len(facts.get("auth_files", [])),
        "jwt_files": len(facts.get("jwt_files", [])),
        "database_files": len(facts.get("database_files", [])),
        "frontend_files": len(facts.get("frontend_files", [])),
        "test_files": len(facts.get("test_files", [])),
        "docker_files": len(facts.get("docker_files", [])),
        "ci_files": len(facts.get("ci_files", [])),
        "config_files": len(facts.get("config_files", [])),
    }

    return counts, None


def _m47_risk_matrix(project_name):
    risks = []

    try:
        security_text = strict_security_analyzer_project(project_name)
    except Exception as error:
        security_text = f"Security analyzer unavailable: {error}"

    for line in _m47_lines(security_text):
        lower = line.lower()

        if "high" in lower or "critical" in lower:
            severity = "HIGH"
        elif "medium" in lower:
            severity = "MEDIUM"
        elif "low" in lower:
            severity = "LOW"
        elif "positive" in lower:
            severity = "POSITIVE"
        elif "info" in lower:
            severity = "INFO"
        else:
            continue

        risks.append({
            "severity": severity,
            "finding": line[:600],
        })

    if not risks:
        risks.append({
            "severity": "INFO",
            "finding": "No explicit risk lines detected by the rule-based analyzer.",
        })

    return risks


def _m47_priority_action_plan(project_name):
    risks = _m47_risk_matrix(project_name)
    actions = []

    high_count = sum(1 for item in risks if item["severity"] in {"HIGH", "CRITICAL"})
    medium_count = sum(1 for item in risks if item["severity"] == "MEDIUM")

    if high_count:
        actions.append("Fix HIGH severity findings first and verify secrets/authentication/configuration.")
    if medium_count:
        actions.append("Review MEDIUM findings and create tracked tasks for each affected file.")
    if not high_count and not medium_count:
        actions.append("No major rule-based security blockers detected; continue with tests and documentation.")

    try:
        counts, error = _m47_collect_evidence_counts(project_name)
        if not error:
            if counts.get("test_files", 0) == 0:
                actions.append("Add tests for core routes, authentication, uploads and project workflows.")
            if counts.get("docker_files", 0) == 0:
                actions.append("Add Docker/deployment documentation if this project must run outside the local machine.")
            if counts.get("ci_files", 0) == 0:
                actions.append("Add CI workflow for linting, tests and security checks.")
    except Exception:
        pass

    actions.extend([
        "Run the application locally after changes and record the exact command used.",
        "Generate a fresh PDF/Word report after fixes to compare progress.",
        "Keep backups before applying automatic patches.",
    ])

    return actions


def mark47_project_health(project_name):
    project_name = _m47_project_alias(project_name)
    project, error = _m47_project_or_error(project_name)

    if error:
        return error

    card = _m47_scorecard_dict(project_name)
    counts, counts_error = _m47_collect_evidence_counts(project_name)
    risks = _m47_risk_matrix(project_name)
    actions = _m47_priority_action_plan(project_name)

    scores = card.get("scores", {})

    output = [
        "J.A.R.V.I.S MARK XLVII PROJECT HEALTH",
        "",
        f"Project: {project.get('name', project_name)}",
        f"Path: {project.get('path', 'Unknown')}",
        f"Files indexed: {project.get('files_count', 'Unknown')}",
        f"Tech stack: {', '.join(project.get('tech_stack', []))}",
        "",
        "Scores:",
        f" - Security: {scores.get('security', 'N/A')}/10",
        f" - Architecture: {scores.get('architecture', 'N/A')}/10",
        f" - Maintainability: {scores.get('maintainability', 'N/A')}/10",
        f" - Overall: {scores.get('overall', 'N/A')}/10",
        "",
        "Evidence counts:"
    ]

    if counts_error:
        output.append(f" - {counts_error}")
    else:
        for key, value in counts.items():
            output.append(f" - {key}: {value}")

    output.append("\nRisk matrix:")
    for item in risks[:20]:
        output.append(f" - {item['severity']}: {item['finding']}")

    output.append("\nPriority action plan:")
    for index, action in enumerate(actions, start=1):
        output.append(f"{index}. {action}")

    return "\n".join(output)


def mark47_executive_summary(project_name):
    project_name = _m47_project_alias(project_name)
    project, error = _m47_project_or_error(project_name)

    if error:
        return error

    card = _m47_scorecard_dict(project_name)
    scores = card.get("scores", {})
    counts, _ = _m47_collect_evidence_counts(project_name)
    risks = _m47_risk_matrix(project_name)

    high = sum(1 for item in risks if item["severity"] in {"HIGH", "CRITICAL"})
    medium = sum(1 for item in risks if item["severity"] == "MEDIUM")

    return (
        "EXECUTIVE SUMMARY\n\n"
        f"{project.get('name', project_name)} was analyzed by {MARK47_REVIEW_VERSION} using indexed project evidence.\n"
        f"The current overall score is {scores.get('overall', 'N/A')}/10, "
        f"with security {scores.get('security', 'N/A')}/10, architecture {scores.get('architecture', 'N/A')}/10, "
        f"and maintainability {scores.get('maintainability', 'N/A')}/10.\n\n"
        f"Evidence detected: {counts.get('files_indexed', project.get('files_count', 'N/A'))} indexed files, "
        f"{counts.get('routes', 0)} backend routes, {counts.get('frontend_api_calls', 0)} frontend API calls, "
        f"{counts.get('auth_files', 0)} auth-related files, and {counts.get('test_files', 0)} test files.\n\n"
        f"Risk summary: {high} high/critical findings and {medium} medium findings were detected by the rule-based analyzer.\n"
        "The recommended next step is to fix the highest-risk security/configuration issues first, then strengthen tests, CI, and documentation."
    )


def mark47_autonomous_review_bundle(project_name):
    project_name = _m47_project_alias(project_name)

    sections = [
        ("Executive Summary", mark47_executive_summary(project_name)),
        ("Project Health", mark47_project_health(project_name)),
        ("Scorecard", score_project(project_name)),
        ("Strict Grounded Analyzer", strict_grounded_analyzer_project(project_name)),
        ("Strict Security Analyzer", strict_security_analyzer_project(project_name)),
        ("Strict Architecture Analyzer", strict_architecture_analyzer_project(project_name)),
        ("Dead Code Scan", find_dead_code(project_name)),
        ("Duplicate Code Scan", find_duplicate_code(project_name)),
        ("Improvement Roadmap", generate_project_roadmap(project_name)),
        ("Production Readiness", production_readiness(project_name)),
        ("Release Checklist", release_checklist(project_name)),
        ("Deployment Checklist", deployment_checklist(project_name)),
    ]

    output = [
        f"{MARK47_REVIEW_VERSION}",
        f"Generated: {_m47_now()}",
        f"Project: {project_name}",
        "",
    ]

    for title, body in sections:
        output.append("\n" + "=" * 80)
        output.append(title.upper())
        output.append("=" * 80)
        output.append(str(body))

    return "\n".join(output)


def mark47_security_bundle(project_name):
    project_name = _m47_project_alias(project_name)

    sections = [
        ("Executive Security Summary", mark47_executive_summary(project_name)),
        ("Strict Security Analyzer", strict_security_analyzer_project(project_name)),
        ("Full Security Audit", full_security_audit(project_name)),
        ("Hardcoded Secrets", find_hardcoded_secrets(project_name)),
        ("API Keys", find_api_keys(project_name)),
        ("Passwords", find_passwords(project_name)),
        ("SQL Injection", find_sql_injection(project_name)),
        ("XSS Risks", find_xss_risks(project_name)),
        ("Dangerous Imports", find_dangerous_imports(project_name)),
        ("Security Roadmap", generate_security_roadmap(project_name)),
    ]

    output = [
        f"{MARK47_REVIEW_VERSION} - Security Bundle",
        f"Generated: {_m47_now()}",
        f"Project: {project_name}",
        "",
    ]

    for title, body in sections:
        output.append("\n" + "=" * 80)
        output.append(title.upper())
        output.append("=" * 80)
        output.append(str(body))

    return "\n".join(output)


def mark47_architecture_bundle(project_name):
    project_name = _m47_project_alias(project_name)

    sections = [
        ("Executive Architecture Summary", mark47_executive_summary(project_name)),
        ("Strict Architecture Analyzer", strict_architecture_analyzer_project(project_name)),
        ("Architecture Report", generate_architecture_report(project_name)),
        ("Project Structure", analyze_project_structure(project_name)),
        ("Duplicate Code Scan", find_duplicate_code(project_name)),
        ("Dead Code Scan", find_dead_code(project_name)),
        ("Architecture Roadmap", generate_project_roadmap(project_name)),
    ]

    output = [
        f"{MARK47_REVIEW_VERSION} - Architecture Bundle",
        f"Generated: {_m47_now()}",
        f"Project: {project_name}",
        "",
    ]

    for title, body in sections:
        output.append("\n" + "=" * 80)
        output.append(title.upper())
        output.append("=" * 80)
        output.append(str(body))

    return "\n".join(output)


def _m47_report_content(project_name, report_kind="project_review"):
    kind = str(report_kind or "project_review").lower().strip()
    project_name = _m47_project_alias(project_name)

    if kind in {"project_review", "project", "full", "enterprise", "autonomous"}:
        return mark47_autonomous_review_bundle(project_name)

    if kind in {"security", "security_review", "audit", "security audit"}:
        return mark47_security_bundle(project_name)

    if kind in {"architecture", "architecture_review", "architect"}:
        return mark47_architecture_bundle(project_name)

    if kind in {"health", "dashboard"}:
        return mark47_project_health(project_name)

    if kind in {"executive", "summary", "executive_summary"}:
        return mark47_executive_summary(project_name)

    if kind in {"score", "scorecard"}:
        return score_project(project_name)

    if kind in {"roadmap", "improvement"}:
        return generate_project_roadmap(project_name)

    return mark47_autonomous_review_bundle(project_name)


def _m47_split_sections(content):
    text = str(content or "")
    sections = []
    current_title = "Executive Summary"
    current_lines = []

    for raw in text.splitlines():
        line = raw.strip()

        is_title = False
        title = ""

        if line and set(line) == {"="}:
            continue

        if line.isupper() and 4 <= len(line) <= 90:
            is_title = True
            title = line.title()
        elif line.startswith("#"):
            is_title = True
            title = line.strip("# ").strip() or "Section"

        if is_title:
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
                current_lines = []
            current_title = title
        else:
            current_lines.append(raw)

    if current_lines:
        sections.append((current_title, "\n".join(current_lines).strip()))

    return sections or [("Report", text)]


def _m47_write_markdown(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.md"
    )

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(f"# {MARK47_REVIEW_VERSION}\n\n")
        file.write(f"Generated: {_m47_now()}\n\n")
        file.write(f"Project: {project_name}\n\n")
        file.write(f"Report kind: {report_kind}\n\n")
        file.write("---\n\n")
        file.write(str(content))

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_text(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.txt"
    )

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(str(content))

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_json(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.json"
    )

    payload = {
        "engine": MARK47_REVIEW_VERSION,
        "generated": _m47_now(),
        "project": project_name,
        "report_kind": report_kind,
        "health": mark47_project_health(project_name),
        "sections": [
            {"title": title, "content": body}
            for title, body in _m47_split_sections(content)
        ],
        "content": str(content),
    }

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        json.dump(payload, file, indent=2, ensure_ascii=False)

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_csv(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.csv"
    )

    import csv

    with open(path, "w", encoding="utf-8", errors="ignore", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["project", "report_kind", "line", "text"])

        for index, line in enumerate(str(content).splitlines(), start=1):
            writer.writerow([project_name, report_kind, index, line])

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_html(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.html"
    )

    def esc(value):
        return str(value).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    sections_html = []

    for title, body in _m47_split_sections(content):
        sections_html.append(
            f"<section class='card'><h2>{esc(title)}</h2><pre>{esc(_m47_limit(body, 12000))}</pre></section>"
        )

    html = f"""<!doctype html>
<html>
<head>
<meta charset="utf-8">
<title>{MARK47_REVIEW_VERSION} - {esc(project_name)}</title>
<style>
body {{ margin: 0; background: #0f172a; color: #e5e7eb; font-family: Arial, sans-serif; }}
header {{ background: linear-gradient(135deg, #111827, #1e3a8a); padding: 32px 48px; }}
h1 {{ margin: 0; font-size: 32px; }}
.card {{ background: #111827; border: 1px solid #334155; border-radius: 16px; margin: 24px 48px; padding: 24px; box-shadow: 0 10px 30px rgba(0,0,0,.25); }}
pre {{ white-space: pre-wrap; line-height: 1.45; color: #d1d5db; }}
.badge {{ display: inline-block; background: #2563eb; padding: 6px 12px; border-radius: 999px; margin-top: 12px; }}
</style>
</head>
<body>
<header>
<h1>{MARK47_REVIEW_VERSION}</h1>
<div class="badge">Project: {esc(project_name)}</div>
<div class="badge">Generated: {esc(_m47_now())}</div>
<div class="badge">Kind: {esc(report_kind)}</div>
</header>
{''.join(sections_html)}
</body>
</html>"""

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(html)

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_docx(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.docx"
    )

    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except Exception:
        return _m47_write_markdown(project_name, content, report_kind, open_after=open_after)

    doc = Document()

    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = doc.add_heading(MARK47_REVIEW_VERSION, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(project_name)
    run.bold = True
    run.font.size = Pt(15)

    doc.add_paragraph(f"Generated: {_m47_now()}")
    doc.add_paragraph(f"Report kind: {report_kind}")

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(mark47_executive_summary(project_name))

    doc.add_heading("Project Health", level=1)
    for line in mark47_project_health(project_name).splitlines():
        if line.strip().startswith("- "):
            doc.add_paragraph(line.strip()[2:], style="List Bullet")
        elif line.strip():
            doc.add_paragraph(line)

    doc.add_page_break()
    doc.add_heading("Full Report", level=1)

    for title_text, body in _m47_split_sections(content):
        doc.add_heading(title_text[:80], level=2)

        for line in str(body).splitlines():
            clean = line.strip()
            if not clean:
                continue
            if clean.startswith("- "):
                doc.add_paragraph(clean[2:], style="List Bullet")
            elif re.match(r"^\d+\.\s+", clean):
                doc.add_paragraph(clean)
            else:
                doc.add_paragraph(clean)

    doc.save(path)

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_pdf(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.pdf"
    )

    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    except Exception:
        return _m47_write_markdown(project_name, content, report_kind, open_after=open_after)

    def safe(value):
        return str(value).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    doc = SimpleDocTemplate(
        path,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="JarvisSmall",
        parent=styles["BodyText"],
        fontSize=8.5,
        leading=11,
        spaceAfter=5,
    ))

    story = []
    story.append(Paragraph(MARK47_REVIEW_VERSION, styles["Title"]))
    story.append(Paragraph(f"<b>Project:</b> {safe(project_name)}", styles["BodyText"]))
    story.append(Paragraph(f"<b>Generated:</b> {safe(_m47_now())}", styles["BodyText"]))
    story.append(Paragraph(f"<b>Report kind:</b> {safe(report_kind)}", styles["BodyText"]))
    story.append(Spacer(1, 0.2 * inch))

    card = _m47_scorecard_dict(project_name)
    scores = card.get("scores", {})

    rows = [["Metric", "Score / 10"]]
    for metric in ["security", "architecture", "maintainability", "overall"]:
        rows.append([metric.title(), str(scores.get(metric, "N/A"))])

    table = Table(rows, colWidths=[2.8 * inch, 1.8 * inch])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    story.append(table)
    story.append(Spacer(1, 0.25 * inch))

    story.append(Paragraph("Executive Summary", styles["Heading1"]))
    for line in mark47_executive_summary(project_name).splitlines():
        if line.strip():
            story.append(Paragraph(safe(line[:1000]), styles["JarvisSmall"]))

    story.append(PageBreak())

    for title_text, body in _m47_split_sections(content):
        story.append(Paragraph(safe(title_text[:80]), styles["Heading2"]))

        for line in str(body).splitlines():
            line = line.strip()
            if not line:
                continue
            story.append(Paragraph(safe(line[:1100]), styles["JarvisSmall"]))

    doc.build(story)

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_excel(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.xlsx"
    )

    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference
        from openpyxl.styles import Font, PatternFill, Alignment
    except Exception:
        return _m47_write_csv(project_name, content, report_kind, open_after=open_after)

    wb = Workbook()
    card = _m47_scorecard_dict(project_name)
    scores = card.get("scores", {})
    counts, _ = _m47_collect_evidence_counts(project_name)
    risks = _m47_risk_matrix(project_name)
    actions = _m47_priority_action_plan(project_name)

    ws = wb.active
    ws.title = "Dashboard"

    ws["A1"] = MARK47_REVIEW_VERSION
    ws["A1"].font = Font(bold=True, size=16)
    ws["A3"] = "Project"
    ws["B3"] = project_name
    ws["A4"] = "Generated"
    ws["B4"] = _m47_now()
    ws["A5"] = "Report kind"
    ws["B5"] = report_kind

    ws["A8"] = "Metric"
    ws["B8"] = "Score"

    for cell in ws[8]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="D9EAF7")

    score_rows = [
        ("Security", scores.get("security") or 0),
        ("Architecture", scores.get("architecture") or 0),
        ("Maintainability", scores.get("maintainability") or 0),
        ("Overall", scores.get("overall") or 0),
    ]

    for index, (metric, score) in enumerate(score_rows, start=9):
        ws.cell(row=index, column=1).value = metric
        ws.cell(row=index, column=2).value = score

    chart = BarChart()
    chart.title = "Project Scorecard"
    chart.y_axis.title = "Score / 10"
    chart.x_axis.title = "Metric"
    data = Reference(ws, min_col=2, min_row=8, max_row=12)
    cats = Reference(ws, min_col=1, min_row=9, max_row=12)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 8
    chart.width = 14
    ws.add_chart(chart, "D8")

    ws_counts = wb.create_sheet("Evidence Counts")
    ws_counts.append(["Metric", "Value"])
    for key, value in counts.items():
        ws_counts.append([key, value])

    ws_risks = wb.create_sheet("Risk Matrix")
    ws_risks.append(["Severity", "Finding"])
    for item in risks:
        ws_risks.append([item.get("severity"), item.get("finding")])

    ws_actions = wb.create_sheet("Action Plan")
    ws_actions.append(["Priority", "Action"])
    for index, action in enumerate(actions, start=1):
        ws_actions.append([index, action])

    ws_full = wb.create_sheet("Full Report")
    ws_full.append(["Line", "Text"])
    for index, line in enumerate(str(content).splitlines(), start=1):
        ws_full.append([index, line])

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        sheet.column_dimensions["A"].width = 24
        sheet.column_dimensions["B"].width = 120

    wb.save(path)

    if open_after:
        _m47_open_file(path)

    return path


def _m47_write_powerpoint(project_name, content, report_kind, open_after=True):
    _m47_ensure_reports_dir()
    path = os.path.join(
        MARK47_REPORT_DIR,
        f"{_m47_safe_name(project_name)}_{_m47_safe_name(report_kind)}_{_m47_stamp()}.pptx"
    )

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        return _m47_write_markdown(project_name, content, report_kind, open_after=open_after)

    prs = Presentation()
    card = _m47_scorecard_dict(project_name)
    scores = card.get("scores", {})
    risks = _m47_risk_matrix(project_name)
    actions = _m47_priority_action_plan(project_name)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = MARK47_REVIEW_VERSION
    slide.placeholders[1].text = f"{project_name}\n{report_kind} | {_m47_now()}"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Scorecard"

    table = slide.shapes.add_table(5, 2, Inches(0.8), Inches(1.5), Inches(8.2), Inches(3.2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Score"

    for row, metric in enumerate(["security", "architecture", "maintainability", "overall"], start=1):
        table.cell(row, 0).text = metric.title()
        table.cell(row, 1).text = str(scores.get(metric, "N/A"))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk Matrix"
    slide.placeholders[1].text = "\n".join(
        f"• {item['severity']}: {item['finding'][:130]}"
        for item in risks[:8]
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Priority Action Plan"
    slide.placeholders[1].text = "\n".join(
        f"{index}. {action[:150]}"
        for index, action in enumerate(actions[:8], start=1)
    )

    for title_text, body in _m47_split_sections(content)[:10]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_text[:65]

        lines = [line.strip() for line in str(body).splitlines() if line.strip()]
        slide.placeholders[1].text = "\n".join(lines[:10])[:1200] or "No details."

    prs.save(path)

    if open_after:
        _m47_open_file(path)

    return path


def export_report_content(project_name, content, format_type="pdf", report_kind="project_review", open_after=True):
    project_name = _m47_project_alias(project_name)
    fmt = _m47_normalize_format(format_type)

    if fmt == "word":
        path = _m47_write_docx(project_name, content, report_kind, open_after=open_after)
    elif fmt == "pdf":
        path = _m47_write_pdf(project_name, content, report_kind, open_after=open_after)
    elif fmt == "powerpoint":
        path = _m47_write_powerpoint(project_name, content, report_kind, open_after=open_after)
    elif fmt == "excel":
        path = _m47_write_excel(project_name, content, report_kind, open_after=open_after)
    elif fmt == "html":
        path = _m47_write_html(project_name, content, report_kind, open_after=open_after)
    elif fmt == "json":
        path = _m47_write_json(project_name, content, report_kind, open_after=open_after)
    elif fmt == "csv":
        path = _m47_write_csv(project_name, content, report_kind, open_after=open_after)
    elif fmt == "text":
        path = _m47_write_text(project_name, content, report_kind, open_after=open_after)
    else:
        path = _m47_write_markdown(project_name, content, report_kind, open_after=open_after)

    return f"Report created:\n{path}"


def export_report(project_name, format_type="pdf", report_kind="project_review", open_after=True):
    project_name = _m47_project_alias(project_name)
    content = _m47_report_content(project_name, report_kind=report_kind)

    return export_report_content(
        project_name,
        content,
        format_type=format_type,
        report_kind=report_kind,
        open_after=open_after
    )


def export_project_report(project_name, format_type="markdown", report_kind="project_review", open_after=True):
    return export_report(
        project_name,
        format_type=format_type,
        report_kind=report_kind,
        open_after=open_after
    )


def export_project_markdown_report(project_name):
    return export_report(project_name, format_type="markdown", report_kind="project_review", open_after=True)


def create_project_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="project_review", open_after=True)


def create_security_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="security", open_after=True)


def create_architecture_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="architecture", open_after=True)


def create_health_report(project_name, format_type="pdf"):
    return export_report(project_name, format_type=format_type, report_kind="health", open_after=True)


def generate_enterprise_report(project_name, format_type="pdf", report_kind="project_review"):
    return export_report(project_name, format_type=format_type, report_kind=report_kind, open_after=True)


def build_enterprise_project_review(project_name):
    return mark47_autonomous_review_bundle(project_name)


def parse_and_generate_report(command):
    text = str(command or "").strip()
    lower = text.lower()

    fmt = "pdf"
    for candidate in [
        "word", "docx", "doc", "pdf", "ppt", "pptx", "powerpoint",
        "presentation", "excel", "xlsx", "xls", "spreadsheet",
        "markdown", "md", "html", "json", "csv", "txt"
    ]:
        if re.search(rf"\b{re.escape(candidate)}\b", lower):
            fmt = _m47_normalize_format(candidate)
            break

    kind = "project_review"

    if "security" in lower or "audit" in lower:
        kind = "security"
    elif "architecture" in lower or "architect" in lower:
        kind = "architecture"
    elif "health" in lower or "dashboard" in lower:
        kind = "health"
    elif "executive" in lower or "summary" in lower:
        kind = "executive"
    elif "score" in lower or "scorecard" in lower:
        kind = "score"
    elif "roadmap" in lower:
        kind = "roadmap"

    match = re.search(
        r"\b(?:for|from|of)\s+(?:project\s+)?(.+)$",
        text,
        flags=re.IGNORECASE
    )

    if match:
        project_name = match.group(1).strip(" .,:;")
        project_name = re.sub(
            r"\b(word|docx|doc|pdf|ppt|pptx|powerpoint|presentation|excel|xlsx|xls|spreadsheet|markdown|md|html|json|csv|txt|report|review|security|architecture|score|health|dashboard|executive|summary|roadmap)\b",
            "",
            project_name,
            flags=re.IGNORECASE
        ).strip()
    else:
        project_name = "CyberShield AI"

    project_name = _m47_project_alias(project_name or "CyberShield AI")

    return export_report(project_name, format_type=fmt, report_kind=kind, open_after=True)


def mark47_review_self_test(project_name="CyberShield AI"):
    project_name = _m47_project_alias(project_name)

    checks = []

    def check(name, func):
        try:
            result = func()
            ok = bool(result)
            checks.append((name, "OK" if ok else "EMPTY", str(result)[:300]))
        except Exception as error:
            checks.append((name, "FAILED", str(error)))

    check("Project health", lambda: mark47_project_health(project_name))
    check("Executive summary", lambda: mark47_executive_summary(project_name))
    check("Scorecard", lambda: score_project(project_name))
    check("Risk matrix", lambda: str(_m47_risk_matrix(project_name)))
    check("Action plan", lambda: str(_m47_priority_action_plan(project_name)))

    output = [
        "MARK XLVII PROJECT REVIEW SELF TEST",
        f"Project: {project_name}",
        "",
    ]

    for name, status, detail in checks:
        output.append(f"- {name}: {status}")
        if status != "OK":
            output.append(f"  Detail: {detail}")

    return "\n".join(output)


# Friendly aliases used by jarvis_agent.py
def project_health(project_name):
    return mark47_project_health(project_name)


def executive_summary(project_name):
    return mark47_executive_summary(project_name)


def autonomous_review_bundle(project_name):
    return mark47_autonomous_review_bundle(project_name)


def security_bundle(project_name):
    return mark47_security_bundle(project_name)


def architecture_bundle(project_name):
    return mark47_architecture_bundle(project_name)


def review_self_test(project_name="CyberShield AI"):
    return mark47_review_self_test(project_name)



# ==========================================================
# J.A.R.V.I.S FAST REPORT MODE
# Goal:
# - Generate reports fast, usually under 1 minute.
# - Avoid slow full LLM analysis for normal report exports.
# - Use rule-based project evidence already indexed in deep memory.
# - Add quick scorecard / risk matrix / action plan.
# - Add charts for Excel and PowerPoint.
# - Keep PDF/Word lightweight and reliable.
# ==========================================================

FAST_REPORT_VERSION = "J.A.R.V.I.S Fast Report Mode"
FAST_REPORT_DIR = "reports"
FAST_REPORT_CACHE_FILE = os.path.join(FAST_REPORT_DIR, "fast_report_cache.json")
FAST_REPORT_CACHE_SECONDS = 10 * 60
FAST_REPORT_OPEN_AFTER_CREATE = True


def _fast_now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def _fast_stamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def _fast_safe_name(name):
    value = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(name or "").strip())
    value = value.strip("_")
    return value or "jarvis_report"


def _fast_ensure_dir():
    os.makedirs(FAST_REPORT_DIR, exist_ok=True)
    return FAST_REPORT_DIR


def _fast_load_json(path, default=None):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as file:
            return json.load(file)
    except Exception:
        return default if default is not None else {}


def _fast_save_json(path, data):
    try:
        parent = os.path.dirname(path)
        if parent:
            os.makedirs(parent, exist_ok=True)

        with open(path, "w", encoding="utf-8") as file:
            json.dump(data, file, indent=2, ensure_ascii=False)

        return True
    except Exception:
        return False


def _fast_open_file(path):
    try:
        os.startfile(os.path.abspath(path))
        return True
    except Exception:
        pass

    try:
        subprocess.Popen(
            ["cmd", "/c", "start", "", os.path.abspath(path)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            shell=False,
        )
        return True
    except Exception:
        return False


def _fast_format_type(format_type):
    fmt = str(format_type or "pdf").lower().strip()

    aliases = {
        "doc": "word",
        "docx": "word",
        "word": "word",
        "pdf": "pdf",
        "ppt": "powerpoint",
        "pptx": "powerpoint",
        "presentation": "powerpoint",
        "powerpoint": "powerpoint",
        "excel": "excel",
        "xls": "excel",
        "xlsx": "excel",
        "spreadsheet": "excel",
        "html": "html",
        "json": "json",
        "csv": "csv",
        "md": "markdown",
        "markdown": "markdown",
        "txt": "text",
        "text": "text",
    }

    return aliases.get(fmt, fmt)


def _fast_project_alias(project_name):
    lower = str(project_name or "").lower().strip()

    aliases = {
        "cyber": "CyberShield AI",
        "cyber shield": "CyberShield AI",
        "cyber shield ai": "CyberShield AI",
        "cybershield": "CyberShield AI",
        "cybershield ai": "CyberShield AI",
        "cybers in the": "CyberShield AI",
        "this project": "CyberShield AI",
        "current project": "CyberShield AI",
        "jarvis": "J.A.R.V.I.S",
        "jervis": "J.A.R.V.I.S",
        "j a r v i s": "J.A.R.V.I.S",
        "manager app": "ManagerApp",
        "managerapp": "ManagerApp",
    }

    return aliases.get(lower, str(project_name or "").strip() or "CyberShield AI")


def _fast_cache_key(project_name, report_kind):
    project_name = _fast_project_alias(project_name)
    return f"{project_name.lower()}::{str(report_kind or 'project_review').lower()}"


def _fast_get_cached_analysis(project_name, report_kind):
    cache = _fast_load_json(FAST_REPORT_CACHE_FILE, {})

    if not isinstance(cache, dict):
        return None

    key = _fast_cache_key(project_name, report_kind)
    item = cache.get(key)

    if not item:
        return None

    created = int(item.get("created", 0) or 0)

    if datetime.now().timestamp() - created > FAST_REPORT_CACHE_SECONDS:
        return None

    return item.get("analysis")


def _fast_set_cached_analysis(project_name, report_kind, analysis):
    cache = _fast_load_json(FAST_REPORT_CACHE_FILE, {})

    if not isinstance(cache, dict):
        cache = {}

    key = _fast_cache_key(project_name, report_kind)
    cache[key] = {
        "created": int(datetime.now().timestamp()),
        "project": _fast_project_alias(project_name),
        "report_kind": report_kind,
        "analysis": analysis,
    }

    _fast_save_json(FAST_REPORT_CACHE_FILE, cache)


def _fast_get_scorecard(project_name):
    try:
        card, error = project_scorecard(project_name)

        if not error and card:
            return card
    except Exception:
        pass

    return {
        "name": project_name,
        "path": "",
        "files_count": 0,
        "tech_stack": [],
        "scores": {
            "security": 0,
            "architecture": 0,
            "maintainability": 0,
            "overall": 0,
        },
        "reasons": {
            "security": [],
            "architecture": [],
            "maintainability": [],
        },
        "evidence": {
            "routes": [],
            "frontend_api_calls": [],
            "auth_files": [],
            "jwt_files": [],
            "frontend_files": [],
            "test_files": [],
            "docker_files": [],
            "ci_files": [],
        }
    }


def _fast_evidence_counts(card):
    evidence = card.get("evidence", {})

    return {
        "Backend routes": len(evidence.get("routes", [])),
        "Frontend API calls": len(evidence.get("frontend_api_calls", [])),
        "Auth files": len(evidence.get("auth_files", [])),
        "JWT files": len(evidence.get("jwt_files", [])),
        "Frontend files": len(evidence.get("frontend_files", [])),
        "Test files": len(evidence.get("test_files", [])),
        "Docker files": len(evidence.get("docker_files", [])),
        "CI files": len(evidence.get("ci_files", [])),
    }


def _fast_risk_matrix(project_name):
    risks = []

    try:
        security_text = strict_security_analyzer_project(project_name)
    except Exception as error:
        security_text = f"Security analyzer unavailable: {error}"

    for line in str(security_text).splitlines():
        clean = line.strip()

        if not clean:
            continue

        lower = clean.lower()

        if "high" in lower or "critical" in lower:
            severity = "HIGH"
        elif "medium" in lower:
            severity = "MEDIUM"
        elif "positive" in lower:
            severity = "POSITIVE"
        elif "info" in lower:
            severity = "INFO"
        elif "low" in lower:
            severity = "LOW"
        else:
            continue

        risks.append({
            "severity": severity,
            "finding": clean[:600],
        })

    if not risks:
        risks.append({
            "severity": "INFO",
            "finding": "No explicit rule-based risks detected.",
        })

    return risks[:40]


def _fast_action_plan(card, risks):
    evidence = card.get("evidence", {})
    scores = card.get("scores", {})
    actions = []

    high_count = sum(1 for item in risks if item.get("severity") == "HIGH")
    medium_count = sum(1 for item in risks if item.get("severity") == "MEDIUM")

    if high_count:
        actions.append("Fix HIGH security findings first, especially secrets, auth, uploads and tokens.")
    if medium_count:
        actions.append("Create tracked tasks for MEDIUM findings and verify affected files manually.")

    if scores.get("security", 0) < 7:
        actions.append("Improve security score with stronger secret handling, token validation and production configuration.")
    if scores.get("architecture", 0) < 7:
        actions.append("Improve architecture by centralizing config, API calls and project structure.")
    if scores.get("maintainability", 0) < 7:
        actions.append("Improve maintainability with tests, CI, documentation and cleanup.")

    if not evidence.get("test_files"):
        actions.append("Add tests for critical flows: auth, project operations, uploads and reports.")
    if not evidence.get("ci_files"):
        actions.append("Add CI workflow for linting, tests and security checks.")
    if not evidence.get("docker_files"):
        actions.append("Add Docker or deployment documentation if the project must run outside local development.")

    actions.append("Generate a fresh report after fixes and compare score progression.")
    actions.append("Do not delete files automatically; run tests before removing dead code.")

    # Stable, short action plan for speed.
    return actions[:10]


def _fast_analysis(project_name, report_kind="project_review", use_cache=True):
    project_name = _fast_project_alias(project_name)
    report_kind = str(report_kind or "project_review").lower().strip()

    if use_cache:
        cached = _fast_get_cached_analysis(project_name, report_kind)
        if cached:
            return cached

    project, error = get_project(project_name)

    if error:
        return {
            "error": error,
            "project_name": project_name,
            "generated": _fast_now(),
        }

    card = _fast_get_scorecard(project_name)
    counts = _fast_evidence_counts(card)
    risks = _fast_risk_matrix(project_name)
    actions = _fast_action_plan(card, risks)

    try:
        grounded = strict_grounded_analyzer_project(project_name)
    except Exception as error:
        grounded = f"Strict grounded analyzer unavailable: {error}"

    try:
        architecture = strict_architecture_analyzer_project(project_name)
    except Exception as error:
        architecture = f"Strict architecture analyzer unavailable: {error}"

    try:
        security = strict_security_analyzer_project(project_name)
    except Exception as error:
        security = f"Strict security analyzer unavailable: {error}"

    # Keep slow sections short and rule-based. No LLM here.
    analysis = {
        "version": FAST_REPORT_VERSION,
        "generated": _fast_now(),
        "project_name": card.get("name", project_name),
        "project_path": card.get("path", project.get("path", "")),
        "files_count": card.get("files_count", project.get("files_count", 0)),
        "tech_stack": card.get("tech_stack", project.get("tech_stack", [])),
        "report_kind": report_kind,
        "scores": card.get("scores", {}),
        "reasons": card.get("reasons", {}),
        "evidence_counts": counts,
        "risk_matrix": risks,
        "action_plan": actions,
        "sections": {
            "scorecard": format_project_scorecard(card),
            "grounded": grounded[:12000],
            "security": security[:12000],
            "architecture": architecture[:12000],
        },
    }

    _fast_set_cached_analysis(project_name, report_kind, analysis)

    return analysis


def _fast_summary_text(analysis):
    if analysis.get("error"):
        return analysis["error"]

    scores = analysis.get("scores", {})
    risks = analysis.get("risk_matrix", [])
    high_count = sum(1 for item in risks if item.get("severity") == "HIGH")
    medium_count = sum(1 for item in risks if item.get("severity") == "MEDIUM")

    return (
        "EXECUTIVE SUMMARY\n\n"
        f"Project: {analysis.get('project_name')}\n"
        f"Generated: {analysis.get('generated')}\n"
        f"Files indexed: {analysis.get('files_count')}\n"
        f"Tech stack: {', '.join(analysis.get('tech_stack', []))}\n\n"
        f"Scores: Security {scores.get('security', 'N/A')}/10, "
        f"Architecture {scores.get('architecture', 'N/A')}/10, "
        f"Maintainability {scores.get('maintainability', 'N/A')}/10, "
        f"Overall {scores.get('overall', 'N/A')}/10.\n\n"
        f"Risk summary: {high_count} HIGH findings and {medium_count} MEDIUM findings detected by rule-based analysis.\n"
        "This fast report avoids slow full-project LLM analysis and uses indexed evidence plus deterministic analyzers."
    )


def _fast_report_text(analysis):
    if analysis.get("error"):
        return analysis["error"]

    lines = [
        FAST_REPORT_VERSION,
        "",
        _fast_summary_text(analysis),
        "",
        "SCORECARD",
        "-" * 60,
        analysis.get("sections", {}).get("scorecard", ""),
        "",
        "EVIDENCE COUNTS",
        "-" * 60,
    ]

    for key, value in analysis.get("evidence_counts", {}).items():
        lines.append(f"- {key}: {value}")

    lines.extend(["", "RISK MATRIX", "-" * 60])

    for item in analysis.get("risk_matrix", []):
        lines.append(f"- {item.get('severity')}: {item.get('finding')}")

    lines.extend(["", "ACTION PLAN", "-" * 60])

    for index, action in enumerate(analysis.get("action_plan", []), start=1):
        lines.append(f"{index}. {action}")

    lines.extend(["", "STRICT SECURITY ANALYZER", "-" * 60])
    lines.append(analysis.get("sections", {}).get("security", ""))

    lines.extend(["", "STRICT ARCHITECTURE ANALYZER", "-" * 60])
    lines.append(analysis.get("sections", {}).get("architecture", ""))

    lines.extend(["", "STRICT GROUNDED ANALYZER", "-" * 60])
    lines.append(analysis.get("sections", {}).get("grounded", ""))

    return "\n".join(str(item) for item in lines)


def _fast_write_markdown(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.md"
    )

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write("# J.A.R.V.I.S Fast Project Report\n\n")
        file.write(_fast_report_text(analysis))

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_text(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.txt"
    )

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(_fast_report_text(analysis))

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_json(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.json"
    )

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        json.dump(analysis, file, indent=2, ensure_ascii=False)

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_csv(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.csv"
    )

    import csv

    with open(path, "w", encoding="utf-8", errors="ignore", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["section", "key", "value"])

        for key, value in analysis.get("scores", {}).items():
            writer.writerow(["scores", key, value])

        for key, value in analysis.get("evidence_counts", {}).items():
            writer.writerow(["evidence_counts", key, value])

        for item in analysis.get("risk_matrix", []):
            writer.writerow(["risk_matrix", item.get("severity"), item.get("finding")])

        for index, action in enumerate(analysis.get("action_plan", []), start=1):
            writer.writerow(["action_plan", index, action])

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_html(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.html"
    )

    def esc(value):
        return str(value).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    scores = analysis.get("scores", {})
    counts = analysis.get("evidence_counts", {})

    score_cards = "".join(
        f"<div class='metric'><strong>{esc(k.title())}</strong><span>{esc(v)}/10</span></div>"
        for k, v in scores.items()
    )

    count_cards = "".join(
        f"<div class='metric'><strong>{esc(k)}</strong><span>{esc(v)}</span></div>"
        for k, v in counts.items()
    )

    risks = "".join(
        f"<li><strong>{esc(item.get('severity'))}</strong>: {esc(item.get('finding'))}</li>"
        for item in analysis.get("risk_matrix", [])
    )

    actions = "".join(
        f"<li>{esc(action)}</li>"
        for action in analysis.get("action_plan", [])
    )

    html = f"""<!doctype html>
<html>
<head>
<meta charset="utf-8">
<title>J.A.R.V.I.S Fast Report - {esc(analysis.get('project_name'))}</title>
<style>
body {{ font-family: Arial, sans-serif; background:#0f172a; color:#e5e7eb; margin:0; }}
header {{ padding:32px 48px; background:linear-gradient(135deg,#111827,#1d4ed8); }}
.card {{ background:#111827; margin:24px 48px; padding:24px; border-radius:16px; border:1px solid #334155; }}
.grid {{ display:grid; grid-template-columns:repeat(auto-fit,minmax(180px,1fr)); gap:12px; }}
.metric {{ background:#1f2937; border-radius:12px; padding:14px; display:flex; justify-content:space-between; }}
pre {{ white-space:pre-wrap; color:#cbd5e1; }}
</style>
</head>
<body>
<header>
<h1>J.A.R.V.I.S Fast Project Report</h1>
<p>{esc(analysis.get('project_name'))} | Generated {esc(analysis.get('generated'))}</p>
</header>
<div class="card"><h2>Executive Summary</h2><pre>{esc(_fast_summary_text(analysis))}</pre></div>
<div class="card"><h2>Scorecard</h2><div class="grid">{score_cards}</div></div>
<div class="card"><h2>Evidence Counts</h2><div class="grid">{count_cards}</div></div>
<div class="card"><h2>Risk Matrix</h2><ul>{risks}</ul></div>
<div class="card"><h2>Action Plan</h2><ol>{actions}</ol></div>
<div class="card"><h2>Full Report</h2><pre>{esc(_fast_report_text(analysis))}</pre></div>
</body>
</html>"""

    with open(path, "w", encoding="utf-8", errors="ignore") as file:
        file.write(html)

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_docx(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.docx"
    )

    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except Exception:
        return _fast_write_markdown(analysis, open_after=open_after)

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = doc.add_heading("J.A.R.V.I.S Fast Project Report", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(str(analysis.get("project_name")))
    run.bold = True
    run.font.size = Pt(15)

    doc.add_paragraph(f"Generated: {analysis.get('generated')}")
    doc.add_paragraph(f"Mode: {FAST_REPORT_VERSION}")

    doc.add_heading("Executive Summary", level=1)
    for line in _fast_summary_text(analysis).splitlines():
        if line.strip():
            doc.add_paragraph(line)

    doc.add_heading("Scorecard", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Score / 10"

    for metric, score in analysis.get("scores", {}).items():
        row = table.add_row().cells
        row[0].text = metric.title()
        row[1].text = str(score)

    doc.add_heading("Evidence Counts", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Evidence"
    table.rows[0].cells[1].text = "Count"

    for key, value in analysis.get("evidence_counts", {}).items():
        row = table.add_row().cells
        row[0].text = key
        row[1].text = str(value)

    doc.add_heading("Risk Matrix", level=1)
    for item in analysis.get("risk_matrix", []):
        doc.add_paragraph(f"{item.get('severity')}: {item.get('finding')}", style="List Bullet")

    doc.add_heading("Action Plan", level=1)
    for action in analysis.get("action_plan", []):
        doc.add_paragraph(action, style="List Number")

    doc.add_page_break()
    doc.add_heading("Full Rule-Based Report", level=1)
    for line in _fast_report_text(analysis).splitlines():
        if line.strip():
            doc.add_paragraph(line[:1200])

    doc.save(path)

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_pdf(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.pdf"
    )

    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    except Exception:
        return _fast_write_markdown(analysis, open_after=open_after)

    def safe(value):
        return str(value).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    doc = SimpleDocTemplate(
        path,
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="FastBody",
        parent=styles["BodyText"],
        fontSize=8.5,
        leading=11,
        spaceAfter=5,
    ))

    story = []
    story.append(Paragraph("J.A.R.V.I.S Fast Project Report", styles["Title"]))
    story.append(Paragraph(f"<b>Project:</b> {safe(analysis.get('project_name'))}", styles["BodyText"]))
    story.append(Paragraph(f"<b>Generated:</b> {safe(analysis.get('generated'))}", styles["BodyText"]))
    story.append(Paragraph(f"<b>Mode:</b> {safe(FAST_REPORT_VERSION)}", styles["BodyText"]))
    story.append(Spacer(1, 0.2 * inch))

    score_rows = [["Metric", "Score / 10"]]
    for metric, score in analysis.get("scores", {}).items():
        score_rows.append([metric.title(), str(score)])

    score_table = Table(score_rows, colWidths=[2.6 * inch, 1.8 * inch])
    score_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    story.append(score_table)
    story.append(Spacer(1, 0.25 * inch))

    count_rows = [["Evidence", "Count"]]
    for key, value in analysis.get("evidence_counts", {}).items():
        count_rows.append([key, str(value)])

    count_table = Table(count_rows, colWidths=[3.0 * inch, 1.4 * inch])
    count_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0F766E")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.grey),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ]))
    story.append(count_table)
    story.append(Spacer(1, 0.25 * inch))

    story.append(Paragraph("Executive Summary", styles["Heading1"]))
    for line in _fast_summary_text(analysis).splitlines():
        if line.strip():
            story.append(Paragraph(safe(line[:1000]), styles["FastBody"]))

    story.append(Paragraph("Risk Matrix", styles["Heading1"]))
    for item in analysis.get("risk_matrix", [])[:18]:
        story.append(Paragraph(f"<b>{safe(item.get('severity'))}</b>: {safe(item.get('finding')[:900])}", styles["FastBody"]))

    story.append(Paragraph("Action Plan", styles["Heading1"]))
    for index, action in enumerate(analysis.get("action_plan", []), start=1):
        story.append(Paragraph(f"{index}. {safe(action)}", styles["FastBody"]))

    story.append(PageBreak())
    story.append(Paragraph("Full Rule-Based Report", styles["Heading1"]))

    # Limit PDF body to keep it very fast.
    for line in _fast_report_text(analysis).splitlines()[:450]:
        line = line.strip()
        if line:
            story.append(Paragraph(safe(line[:1000]), styles["FastBody"]))

    doc.build(story)

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_excel(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.xlsx"
    )

    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, PieChart, Reference
        from openpyxl.styles import Font, PatternFill, Alignment
    except Exception:
        return _fast_write_csv(analysis, open_after=open_after)

    wb = Workbook()
    ws = wb.active
    ws.title = "Dashboard"

    ws["A1"] = "J.A.R.V.I.S Fast Project Report"
    ws["A1"].font = Font(bold=True, size=16)
    ws["A3"] = "Project"
    ws["B3"] = analysis.get("project_name")
    ws["A4"] = "Generated"
    ws["B4"] = analysis.get("generated")
    ws["A5"] = "Mode"
    ws["B5"] = FAST_REPORT_VERSION

    ws["A8"] = "Metric"
    ws["B8"] = "Score"
    ws["A8"].font = Font(bold=True)
    ws["B8"].font = Font(bold=True)
    ws["A8"].fill = PatternFill("solid", fgColor="D9EAF7")
    ws["B8"].fill = PatternFill("solid", fgColor="D9EAF7")

    score_items = list(analysis.get("scores", {}).items())
    for row_index, (metric, score) in enumerate(score_items, start=9):
        ws.cell(row=row_index, column=1).value = metric.title()
        ws.cell(row=row_index, column=2).value = score or 0

    if score_items:
        chart = BarChart()
        chart.title = "Project Scorecard"
        chart.y_axis.title = "Score / 10"
        chart.x_axis.title = "Metric"
        data = Reference(ws, min_col=2, min_row=8, max_row=8 + len(score_items))
        cats = Reference(ws, min_col=1, min_row=9, max_row=8 + len(score_items))
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 8
        chart.width = 14
        ws.add_chart(chart, "D8")

    ws_counts = wb.create_sheet("Evidence Counts")
    ws_counts.append(["Evidence", "Count"])
    for key, value in analysis.get("evidence_counts", {}).items():
        ws_counts.append([key, value])

    if analysis.get("evidence_counts"):
        pie = PieChart()
        pie.title = "Evidence Distribution"
        data = Reference(ws_counts, min_col=2, min_row=1, max_row=1 + len(analysis.get("evidence_counts", {})))
        labels = Reference(ws_counts, min_col=1, min_row=2, max_row=1 + len(analysis.get("evidence_counts", {})))
        pie.add_data(data, titles_from_data=True)
        pie.set_categories(labels)
        pie.height = 8
        pie.width = 12
        ws_counts.add_chart(pie, "D2")

    ws_risks = wb.create_sheet("Risk Matrix")
    ws_risks.append(["Severity", "Finding"])
    for item in analysis.get("risk_matrix", []):
        ws_risks.append([item.get("severity"), item.get("finding")])

    ws_actions = wb.create_sheet("Action Plan")
    ws_actions.append(["Priority", "Action"])
    for index, action in enumerate(analysis.get("action_plan", []), start=1):
        ws_actions.append([index, action])

    ws_full = wb.create_sheet("Full Report")
    ws_full.append(["Line", "Text"])
    for index, line in enumerate(_fast_report_text(analysis).splitlines(), start=1):
        ws_full.append([index, line])

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
        sheet.column_dimensions["A"].width = 26
        sheet.column_dimensions["B"].width = 100

    wb.save(path)

    if open_after:
        _fast_open_file(path)

    return path


def _fast_write_powerpoint(analysis, open_after=True):
    _fast_ensure_dir()
    path = os.path.join(
        FAST_REPORT_DIR,
        f"{_fast_safe_name(analysis.get('project_name'))}_fast_{_fast_stamp()}.pptx"
    )

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        return _fast_write_markdown(analysis, open_after=open_after)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "J.A.R.V.I.S Fast Project Report"
    slide.placeholders[1].text = f"{analysis.get('project_name')}\n{analysis.get('generated')}"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Scorecard"

    scores = analysis.get("scores", {})
    rows = max(2, len(scores) + 1)
    table = slide.shapes.add_table(rows, 2, Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Score / 10"
    for row_index, (metric, score) in enumerate(scores.items(), start=1):
        table.cell(row_index, 0).text = metric.title()
        table.cell(row_index, 1).text = str(score)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Evidence Counts"
    slide.placeholders[1].text = "\n".join(
        f"• {key}: {value}"
        for key, value in analysis.get("evidence_counts", {}).items()
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk Matrix"
    slide.placeholders[1].text = "\n".join(
        f"• {item.get('severity')}: {item.get('finding')[:140]}"
        for item in analysis.get("risk_matrix", [])[:8]
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Priority Action Plan"
    slide.placeholders[1].text = "\n".join(
        f"{index}. {action[:150]}"
        for index, action in enumerate(analysis.get("action_plan", [])[:8], start=1)
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = _fast_summary_text(analysis)[:1100]

    prs.save(path)

    if open_after:
        _fast_open_file(path)

    return path


def export_fast_report(project_name="CyberShield AI", format_type="pdf", report_kind="project_review", open_after=True):
    project_name = _fast_project_alias(project_name)
    fmt = _fast_format_type(format_type)
    start = datetime.now().timestamp()

    analysis = _fast_analysis(project_name, report_kind=report_kind, use_cache=True)

    if fmt == "word":
        path = _fast_write_docx(analysis, open_after=open_after)
    elif fmt == "pdf":
        path = _fast_write_pdf(analysis, open_after=open_after)
    elif fmt == "powerpoint":
        path = _fast_write_powerpoint(analysis, open_after=open_after)
    elif fmt == "excel":
        path = _fast_write_excel(analysis, open_after=open_after)
    elif fmt == "html":
        path = _fast_write_html(analysis, open_after=open_after)
    elif fmt == "json":
        path = _fast_write_json(analysis, open_after=open_after)
    elif fmt == "csv":
        path = _fast_write_csv(analysis, open_after=open_after)
    elif fmt == "text":
        path = _fast_write_text(analysis, open_after=open_after)
    else:
        path = _fast_write_markdown(analysis, open_after=open_after)

    elapsed = round(datetime.now().timestamp() - start, 2)

    return (
        f"Fast report created in {elapsed}s:\n{path}\n\n"
        "Charts/diagrams: Excel includes scorecard bar chart and evidence pie chart. "
        "PowerPoint includes scorecard/risk/action slides. "
        "PDF/Word include tables and structured scorecards."
    )


def export_report(project_name, format_type="pdf", report_kind="project_review", open_after=True):
    # Fast default for voice commands and daily use.
    return export_fast_report(
        project_name=project_name,
        format_type=format_type,
        report_kind=report_kind,
        open_after=open_after
    )


def export_project_report(project_name, format_type="markdown", report_kind="project_review", open_after=True):
    return export_fast_report(
        project_name=project_name,
        format_type=format_type,
        report_kind=report_kind,
        open_after=open_after
    )


def create_project_report(project_name, format_type="pdf"):
    return export_fast_report(project_name, format_type=format_type, report_kind="project_review", open_after=True)


def create_security_report(project_name, format_type="pdf"):
    return export_fast_report(project_name, format_type=format_type, report_kind="security", open_after=True)


def create_architecture_report(project_name, format_type="pdf"):
    return export_fast_report(project_name, format_type=format_type, report_kind="architecture", open_after=True)


def generate_enterprise_report(project_name, format_type="pdf", report_kind="project_review"):
    return export_fast_report(project_name, format_type=format_type, report_kind=report_kind, open_after=True)


def parse_and_generate_report(command):
    text = str(command or "").strip()
    lower = text.lower()

    fmt = "pdf"

    for candidate in [
        "word", "docx", "doc", "pdf", "ppt", "pptx", "powerpoint",
        "presentation", "excel", "xlsx", "xls", "spreadsheet",
        "markdown", "md", "html", "json", "csv", "txt"
    ]:
        if re.search(rf"\b{re.escape(candidate)}\b", lower):
            fmt = _fast_format_type(candidate)
            break

    kind = "project_review"

    if "security" in lower or "audit" in lower:
        kind = "security"
    elif "architecture" in lower or "architect" in lower:
        kind = "architecture"
    elif "health" in lower or "dashboard" in lower:
        kind = "health"
    elif "score" in lower:
        kind = "score"

    match = re.search(
        r"\b(?:for|from|of|about)\s+(?:this\s+)?(?:current\s+)?(?:project\s+)?(.+)$",
        text,
        flags=re.IGNORECASE
    )

    if match:
        project_name = match.group(1).strip(" .,:;")
        project_name = re.sub(
            r"\b(word|docx|doc|pdf|ppt|pptx|powerpoint|presentation|excel|xlsx|xls|spreadsheet|markdown|md|html|json|csv|txt|report|review|security|architecture|score|health|dashboard|project)\b",
            "",
            project_name,
            flags=re.IGNORECASE
        ).strip()
    else:
        project_name = "CyberShield AI"

    if not project_name or project_name.lower() in {"this", "current", "it"}:
        project_name = "CyberShield AI"

    return export_fast_report(
        project_name=_fast_project_alias(project_name),
        format_type=fmt,
        report_kind=kind,
        open_after=True
    )


def fast_report_status(project_name="CyberShield AI"):
    analysis = _fast_analysis(project_name, report_kind="project_review", use_cache=True)
    return _fast_summary_text(analysis)


def clear_fast_report_cache():
    _fast_save_json(FAST_REPORT_CACHE_FILE, {})
    return "Fast report cache cleared."


# Legacy-friendly names
def fast_pdf_report(project_name="CyberShield AI"):
    return export_fast_report(project_name, "pdf")


def fast_word_report(project_name="CyberShield AI"):
    return export_fast_report(project_name, "word")


def fast_excel_report(project_name="CyberShield AI"):
    return export_fast_report(project_name, "excel")


def fast_powerpoint_report(project_name="CyberShield AI"):
    return export_fast_report(project_name, "powerpoint")


def report_self_test(project_name="CyberShield AI"):
    analysis = _fast_analysis(project_name, use_cache=True)
    return (
        "FAST REPORT SELF TEST\n"
        f"Project: {analysis.get('project_name')}\n"
        f"Generated: {analysis.get('generated')}\n"
        f"Scores: {analysis.get('scores')}\n"
        f"Evidence counts: {analysis.get('evidence_counts')}\n"
        f"Risks: {len(analysis.get('risk_matrix', []))}\n"
        f"Actions: {len(analysis.get('action_plan', []))}\n"
    )

