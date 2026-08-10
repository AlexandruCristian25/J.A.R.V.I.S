import os
import re
import sys
import json
import shutil
import subprocess
import webbrowser
from pathlib import Path
from datetime import datetime

# ==========================================================
# J.A.R.V.I.S Developer Assistant
# Steps 1-6:
# 1. IDE Manager
# 2. Code Navigator
# 3. Code Copier / Line Replacement
# 4. AI Code Reviewer / Improvements
# 5. Report Generator: MD / DOCX / PDF / PPTX
# 6. Refactoring Helpers
# ==========================================================

SKIP_DIRS = {
    "node_modules", "venv", ".venv", "jarvis-env", "__pycache__",
    ".git", ".idea", ".vscode", "dist", "build", ".next",
    ".cache", "site-packages", "file_backups",
}

TEXT_EXTENSIONS = {
    ".py", ".js", ".jsx", ".ts", ".tsx", ".html", ".css", ".scss",
    ".sass", ".json", ".md", ".txt", ".yml", ".yaml", ".env",
    ".ini", ".cfg", ".toml", ".xml", ".java", ".c", ".cpp",
    ".h", ".hpp", ".cs", ".php", ".rb", ".go", ".rs", ".sql",
    ".bat", ".ps1", ".sh",
}

REPORT_DIR = "developer_reports"
BACKUP_DIR = "developer_backups"

# ==========================================================
# BASIC HELPERS
# ==========================================================
def clean_text(text):
    text = str(text).strip()
    return re.sub(r"\s+", " ", text)


def normalize_lower(text):
    return clean_text(text).lower()


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def timestamp():
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def safe_read_text(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[READ ERROR] {e}"


def safe_write_text(path, content):
    with open(path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(str(content))


def numbered_lines(text, start_line=1):
    lines = str(text).splitlines()
    width = len(str(start_line + len(lines)))
    return "\n".join(f"{str(index).rjust(width)} | {line}" for index, line in enumerate(lines, start=start_line))


def short(text, limit=900):
    text = str(text).strip()
    return text if len(text) <= limit else text[:limit - 3] + "..."


def open_path(path):
    try:
        os.startfile(str(path))
        return True
    except Exception:
        pass
    try:
        subprocess.Popen(["cmd", "/c", "start", "", str(path)], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, shell=False)
        return True
    except Exception:
        return False


def create_backup(file_path):
    ensure_dir(BACKUP_DIR)
    file_path = Path(file_path)
    if not file_path.exists():
        return None
    backup_path = Path(BACKUP_DIR) / f"{file_path.stem}_{timestamp()}{file_path.suffix}.bak"
    shutil.copy2(file_path, backup_path)
    return str(backup_path)

# ==========================================================
# PROJECT RESOLUTION
# ==========================================================
def find_project(project_name):
    try:
        from tools import find_project as tools_find_project
        project = tools_find_project(project_name)
        if project and isinstance(project, dict):
            path = project.get("path")
            if path and os.path.exists(path):
                return project
    except Exception:
        pass

    query = normalize_lower(project_name).replace(" ", "").replace("_", "").replace("-", "")
    roots = [os.getcwd(), str(Path.home())]
    for letter in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
        drive = f"{letter}:\\"
        if os.path.exists(drive):
            roots.append(drive)

    for root_dir in roots:
        try:
            for root, dirs, files in os.walk(root_dir):
                dirs[:] = [d for d in dirs if d.lower() not in SKIP_DIRS]
                base = os.path.basename(root).lower().replace(" ", "").replace("_", "").replace("-", "")
                if query and (query in base or base in query):
                    return {"name": os.path.basename(root), "path": root, "type": "Detected"}
        except Exception:
            continue
    return None


def resolve_project_path(project_name):
    project = find_project(project_name)
    if not project:
        return None, f"Project not found: {project_name}"
    path = project.get("path")
    if not path or not os.path.exists(path):
        return None, f"Project path not found: {path}"
    return path, None


def iter_project_files(project_path):
    for root, dirs, files in os.walk(project_path):
        dirs[:] = [d for d in dirs if d.lower() not in SKIP_DIRS]
        for filename in files:
            full = os.path.join(root, filename)
            if os.path.splitext(filename)[1].lower() in TEXT_EXTENSIONS:
                yield full


def find_file_in_project(project_name, file_query):
    project_path, error = resolve_project_path(project_name)
    if error:
        return None, error

    file_query_clean = normalize_lower(file_query).replace("\\", "/")
    file_query_base = os.path.basename(file_query_clean)
    candidates = []
    for path in iter_project_files(project_path):
        rel = os.path.relpath(path, project_path)
        rel_lower = rel.lower().replace("\\", "/")
        base_lower = os.path.basename(path).lower()
        score = 0
        if rel_lower == file_query_clean:
            score = 100
        elif base_lower == file_query_base:
            score = 95
        elif file_query_clean in rel_lower:
            score = 85
        elif file_query_base and file_query_base in base_lower:
            score = 75
        if score:
            candidates.append((score, path, rel))
    if not candidates:
        return None, f"File not found in project {project_name}: {file_query}"
    candidates.sort(key=lambda item: (-item[0], item[2].lower()))
    return candidates[0][1], None

# ==========================================================
# STEP 1 - IDE MANAGER
# ==========================================================
IDE_ALIASES = {
    "vscode": "vscode", "vs code": "vscode", "visual studio code": "vscode", "code": "vscode",
    "visual studio": "visualstudio", "visual studio community": "visualstudio", "vs community": "visualstudio", "vscommunity": "visualstudio",
    "intellij": "intellij", "intellij idea": "intellij", "idea": "intellij",
    "pycharm": "pycharm", "android studio": "androidstudio", "androidstudio": "androidstudio", "eclipse": "eclipse",
}

IDE_CANDIDATES = {
    "vscode": [os.path.expandvars(r"%LOCALAPPDATA%\Programs\Microsoft VS Code\Code.exe"), r"C:\Program Files\Microsoft VS Code\Code.exe", "code", "code.cmd"],
    "visualstudio": [r"C:\Program Files\Microsoft Visual Studio\2022\Community\Common7\IDE\devenv.exe", r"C:\Program Files\Microsoft Visual Studio\2022\Professional\Common7\IDE\devenv.exe", r"C:\Program Files\Microsoft Visual Studio\2022\Enterprise\Common7\IDE\devenv.exe", "devenv.exe"],
    "intellij": [os.path.expandvars(r"%LOCALAPPDATA%\JetBrains\Toolbox\scripts\idea.cmd"), r"C:\Program Files\JetBrains\IntelliJ IDEA Community Edition 2024.3\bin\idea64.exe", r"C:\Program Files\JetBrains\IntelliJ IDEA 2024.3\bin\idea64.exe", "idea64.exe", "idea"],
    "pycharm": [os.path.expandvars(r"%LOCALAPPDATA%\JetBrains\Toolbox\scripts\pycharm.cmd"), r"C:\Program Files\JetBrains\PyCharm Community Edition 2024.3\bin\pycharm64.exe", r"C:\Program Files\JetBrains\PyCharm 2024.3\bin\pycharm64.exe", "pycharm64.exe", "pycharm"],
    "androidstudio": [r"C:\Program Files\Android\Android Studio\bin\studio64.exe", "studio64.exe"],
    "eclipse": ["eclipse.exe"],
}


def resolve_ide(ide_name):
    key = IDE_ALIASES.get(normalize_lower(ide_name), normalize_lower(ide_name))
    for candidate in IDE_CANDIDATES.get(key, []):
        expanded = os.path.expandvars(candidate)
        if os.path.exists(expanded):
            return key, expanded
        found = shutil.which(expanded) or shutil.which(candidate)
        if found:
            return key, found
    return key, None


def open_project_in_ide(project_name, ide_name):
    project_path, error = resolve_project_path(project_name)
    if error:
        return error
    ide_key, ide_path = resolve_ide(ide_name)
    if not ide_path:
        return f"Could not find IDE: {ide_name}\nProject found at: {project_path}\nInstall the IDE or add it to PATH, then try again."
    try:
        subprocess.Popen([ide_path, project_path], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, shell=False)
        return f"Opening project '{project_name}' in {ide_name}: {project_path}"
    except Exception as e:
        return f"Could not open project in {ide_name}: {e}"


def open_file_in_ide(project_name, file_query, ide_name):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    ide_key, ide_path = resolve_ide(ide_name)
    if not ide_path:
        return f"Could not find IDE: {ide_name}"
    try:
        subprocess.Popen([ide_path, file_path], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, shell=False)
        return f"Opening file in {ide_name}: {file_path}"
    except Exception as e:
        return f"Could not open file in {ide_name}: {e}"

# ==========================================================
# STEP 2 - CODE NAVIGATOR
# ==========================================================
def read_file_lines(project_name, file_query, start_line=None, end_line=None):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    text = safe_read_text(file_path)
    if text.startswith("[READ ERROR]"):
        return text
    lines = text.splitlines()
    start_line = 1 if start_line is None else max(1, int(start_line))
    end_line = len(lines) if end_line is None else min(len(lines), int(end_line))
    if start_line > end_line:
        return "Invalid line range."
    selected = "\n".join(lines[start_line - 1:end_line])
    return f"File: {file_path}\nLines: {start_line}-{end_line}\n\n```text\n{numbered_lines(selected, start_line)}\n```"


def find_function_in_file(project_name, file_query, function_name):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    lines = safe_read_text(file_path).splitlines()
    pattern = re.compile(rf"^\s*(def|function|const|let|var|class|public|private|protected|static)\s+{re.escape(function_name)}\b|^\s*{re.escape(function_name)}\s*[:=]", re.IGNORECASE)
    matches = [(i, line) for i, line in enumerate(lines, start=1) if pattern.search(line)]
    if not matches:
        return f"Function/symbol not found: {function_name}"
    output = [f"Matches for '{function_name}' in {file_path}:", ""]
    output.extend(f"- line {line_no}: {line.strip()}" for line_no, line in matches[:20])
    return "\n".join(output)

# ==========================================================
# STEP 3 - CODE COPIER / SAFE LINE REPLACEMENT
# ==========================================================
def copy_lines_between_files(source_project, source_file, source_start, source_end, target_project, target_file, target_start, target_end, mode="replace"):
    source_path, error = find_file_in_project(source_project, source_file)
    if error:
        return error
    target_path, error = find_file_in_project(target_project, target_file)
    if error:
        return error
    source_lines = safe_read_text(source_path).splitlines()
    target_lines = safe_read_text(target_path).splitlines()
    source_start = max(1, int(source_start)); source_end = min(len(source_lines), int(source_end))
    target_start = max(1, int(target_start)); target_end = min(len(target_lines), int(target_end))
    if source_start > source_end or target_start > target_end:
        return "Invalid source or target line range."
    extracted = source_lines[source_start - 1:source_end]
    backup = create_backup(target_path)
    if mode == "insert":
        new_lines = target_lines[:target_start - 1] + extracted + target_lines[target_start - 1:]
    else:
        new_lines = target_lines[:target_start - 1] + extracted + target_lines[target_end:]
    safe_write_text(target_path, "\n".join(new_lines) + "\n")
    return f"Code transfer completed safely.\nSource: {source_path} lines {source_start}-{source_end}\nTarget: {target_path} lines {target_start}-{target_end}\nMode: {mode}\nBackup: {backup}"


def replace_lines_with_text(project_name, file_query, start_line, end_line, new_code):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    lines = safe_read_text(file_path).splitlines()
    start_line = max(1, int(start_line)); end_line = min(len(lines), int(end_line))
    if start_line > end_line:
        return "Invalid line range."
    backup = create_backup(file_path)
    new_lines = lines[:start_line - 1] + str(new_code).splitlines() + lines[end_line:]
    safe_write_text(file_path, "\n".join(new_lines) + "\n")
    return f"Lines replaced successfully.\nFile: {file_path}\nLines: {start_line}-{end_line}\nBackup: {backup}"

# ==========================================================
# STEP 4 - AI REVIEWER / IMPROVER
# ==========================================================
def ask_ai(prompt):
    try:
        from llm_local import ask_llm
        return ask_llm(prompt)
    except Exception as e:
        return "AI engine unavailable.\nReason: " + str(e) + "\n\nFallback: check readability, naming, duplication, error handling, security and edge cases."


def review_code_text(code_text, context="code"):
    prompt = f"""
You are J.A.R.V.I.S, a senior software engineering assistant.
Review the following {context}.
Return:
1. Short summary
2. Bugs or risks
3. Security issues
4. Performance issues
5. Maintainability improvements
6. Improved version of the code if possible

CODE:
```text
{code_text}
```
"""
    return ask_ai(prompt)


def improve_code_text(code_text, context="code"):
    prompt = f"""
You are J.A.R.V.I.S, a senior software engineer.
Improve the following {context}.
Keep behavior the same unless there is an obvious bug.
Return:
1. What you changed
2. Improved code only in a code block
3. Notes

CODE:
```text
{code_text}
```
"""
    return ask_ai(prompt)


def review_file(project_name, file_query, start_line=None, end_line=None):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    text = safe_read_text(file_path)
    lines = text.splitlines()
    if start_line is not None and end_line is not None:
        start_line = max(1, int(start_line)); end_line = min(len(lines), int(end_line))
        code = "\n".join(lines[start_line - 1:end_line])
        context = f"{file_path}, lines {start_line}-{end_line}"
    else:
        code = text
        context = file_path
    return review_code_text(code, context=context)


def improve_file(project_name, file_query, start_line=None, end_line=None):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    text = safe_read_text(file_path)
    lines = text.splitlines()
    if start_line is not None and end_line is not None:
        start_line = max(1, int(start_line)); end_line = min(len(lines), int(end_line))
        code = "\n".join(lines[start_line - 1:end_line])
        context = f"{file_path}, lines {start_line}-{end_line}"
    else:
        code = text
        context = file_path
    return improve_code_text(code, context=context)

# ==========================================================
# STEP 5 - REPORT GENERATOR
# ==========================================================
def create_markdown_report(title, content, open_after=True):
    ensure_dir(REPORT_DIR)
    path = Path(REPORT_DIR) / f"{title.replace(' ', '_')}_{timestamp()}.md"
    safe_write_text(path, "\n".join([f"# {title}", "", f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", "", str(content), ""]))
    if open_after:
        open_path(path)
    return str(path)


def create_docx_report(title, content, open_after=True):
    ensure_dir(REPORT_DIR)
    try:
        from docx import Document
    except Exception:
        md_path = create_markdown_report(title, content, open_after=open_after)
        return f"python-docx is not installed. Markdown report created instead:\n{md_path}"
    path = Path(REPORT_DIR) / f"{title.replace(' ', '_')}_{timestamp()}.docx"
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    for block in str(content).split("\n\n"):
        doc.add_paragraph(block)
    doc.save(path)
    if open_after:
        open_path(path)
    return str(path)


def create_pdf_report(title, content, open_after=True):
    ensure_dir(REPORT_DIR)
    path = Path(REPORT_DIR) / f"{title.replace(' ', '_')}_{timestamp()}.pdf"
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception:
        md_path = create_markdown_report(title, content, open_after=open_after)
        return f"reportlab is not installed. Markdown report created instead:\n{md_path}"
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4
    x, y = 42, height - 50
    c.setFont("Helvetica-Bold", 16); c.drawString(x, y, title[:80]); y -= 28
    c.setFont("Helvetica", 9); c.drawString(x, y, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"); y -= 28
    for raw_line in str(content).splitlines():
        if y < 50:
            c.showPage(); c.setFont("Helvetica", 9); y = height - 50
        c.drawString(x, y, raw_line[:115]); y -= 13
    c.save()
    if open_after:
        open_path(path)
    return str(path)


def create_pptx_report(title, content, open_after=True):
    ensure_dir(REPORT_DIR)
    try:
        from pptx import Presentation
    except Exception:
        md_path = create_markdown_report(title, content, open_after=open_after)
        return f"python-pptx is not installed. Markdown report created instead:\n{md_path}"
    path = Path(REPORT_DIR) / f"{title.replace(' ', '_')}_{timestamp()}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    for section in str(content).split("\n\n")[:12]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = section.splitlines()
        slide.shapes.title.text = short(lines[0] if lines else "Report", 60)
        slide.placeholders[1].text = short("\n".join(lines[1:]) if len(lines) > 1 else section, 900)
    prs.save(path)
    if open_after:
        open_path(path)
    return str(path)


def create_report(title, content, format_type="md", open_after=True):
    fmt = normalize_lower(format_type)
    if fmt in {"word", "docx", "doc"}:
        path = create_docx_report(title, content, open_after=open_after)
    elif fmt == "pdf":
        path = create_pdf_report(title, content, open_after=open_after)
    elif fmt in {"ppt", "pptx", "powerpoint", "presentation"}:
        path = create_pptx_report(title, content, open_after=open_after)
    else:
        path = create_markdown_report(title, content, open_after=open_after)
    return f"Report created:\n{path}"


def review_project_report(project_name, format_type="md"):
    try:
        from project_review_assistant import review_project
        content = review_project(project_name)
    except Exception:
        project_path, error = resolve_project_path(project_name)
        if error:
            return error
        content = f"Project review fallback for {project_name}\nPath: {project_path}\n\nReview engine unavailable."
    return create_report(f"JARVIS Project Review - {project_name}", content, format_type=format_type, open_after=True)


def review_file_report(project_name, file_query, format_type="md"):
    return create_report(f"JARVIS File Review - {file_query}", review_file(project_name, file_query), format_type=format_type, open_after=True)

# ==========================================================
# STEP 6 - REFACTORING HELPERS
# ==========================================================
def remove_trailing_whitespace(project_name, file_query):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    lines = safe_read_text(file_path).splitlines()
    backup = create_backup(file_path)
    safe_write_text(file_path, "\n".join(line.rstrip() for line in lines) + "\n")
    return f"Trailing whitespace removed.\nFile: {file_path}\nBackup: {backup}"


def add_header_comment(project_name, file_query, comment):
    file_path, error = find_file_in_project(project_name, file_query)
    if error:
        return error
    text = safe_read_text(file_path)
    backup = create_backup(file_path)
    ext = Path(file_path).suffix.lower()
    if ext in {".py", ".sh", ".ps1"}:
        header = "\n".join(f"# {line}" for line in str(comment).splitlines())
    elif ext in {".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".cpp", ".cs", ".go", ".rs", ".php"}:
        header = "/*\n" + "\n".join(str(comment).splitlines()) + "\n*/"
    else:
        header = str(comment)
    safe_write_text(file_path, header + "\n\n" + text)
    return f"Header comment added.\nFile: {file_path}\nBackup: {backup}"


def format_light(project_name, file_query):
    return remove_trailing_whitespace(project_name, file_query)

# ==========================================================
# COMMAND ROUTER
# ==========================================================
def handle_developer_command(command):
    original = str(command).strip()
    lower = normalize_lower(original)

    match = re.match(r"^open project (.+?) in (vs code|vscode|visual studio code|visual studio community|vs community|visual studio|intellij|intellij idea|idea|pycharm|android studio|eclipse)$", lower)
    if match:
        project_name = original[13: original.lower().rfind(" in ")].strip()
        ide_name = original[original.lower().rfind(" in ") + 4:].strip()
        return open_project_in_ide(project_name, ide_name)

    match = re.match(r"^open file (.+?) from (.+?) in (.+)$", original, flags=re.IGNORECASE)
    if match:
        return open_file_in_ide(match.group(2).strip(), match.group(1).strip(), match.group(3).strip())

    match = re.match(r"^(show|read|get|give me) lines? (\d+) (to|-|until) (\d+) from file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return read_file_lines(match.group(7).strip(), match.group(5).strip(), int(match.group(2)), int(match.group(4)))

    match = re.match(r"^(give me|show|get|read) (the )?(code )?(from )?line (\d+) (to|-|until) (line )?(\d+) from (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return read_file_lines(match.group(11).strip(), match.group(9).strip(), int(match.group(5)), int(match.group(8)))

    match = re.match(r"^find (function|class|symbol) (.+?) in file (.+?) (from|in) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return find_function_in_file(match.group(5).strip(), match.group(3).strip(), match.group(2).strip())

    match = re.match(r"^(copy|replace|move) lines? (\d+) (to|-|until) (\d+) from file (.+?) (in|from) project (.+?) to file (.+?) (in|from) project (.+?) lines? (\d+) (to|-|until) (\d+)$", original, flags=re.IGNORECASE)
    if match:
        return copy_lines_between_files(match.group(7).strip(), match.group(5).strip(), int(match.group(2)), int(match.group(4)), match.group(10).strip(), match.group(8).strip(), int(match.group(11)), int(match.group(13)), mode="replace")

    match = re.match(r"^review lines? (\d+) (to|-|until) (\d+) from file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return review_file(match.group(6).strip(), match.group(4).strip(), int(match.group(1)), int(match.group(3)))

    match = re.match(r"^review file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return review_file(match.group(3).strip(), match.group(1).strip())

    match = re.match(r"^(improve|optimize|refactor) lines? (\d+) (to|-|until) (\d+) from file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return improve_file(match.group(7).strip(), match.group(5).strip(), int(match.group(2)), int(match.group(4)))

    match = re.match(r"^(improve|optimize|refactor) file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return improve_file(match.group(4).strip(), match.group(2).strip())

    match = re.match(r"^generate (word|docx|pdf|ppt|pptx|powerpoint|markdown|md) (report|review) for project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return review_project_report(match.group(3).strip(), match.group(1).strip())

    match = re.match(r"^generate (word|docx|pdf|ppt|pptx|powerpoint|markdown|md) (report|review) for file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return review_file_report(match.group(5).strip(), match.group(3).strip(), match.group(1).strip())

    match = re.match(r"^(format|clean) file (.+?) (in|from) project (.+)$", original, flags=re.IGNORECASE)
    if match:
        return format_light(match.group(4).strip(), match.group(2).strip())

    return None


if __name__ == "__main__":
    print("J.A.R.V.I.S Developer Assistant")
    print("Type a developer command or 'exit'.")
    while True:
        cmd = input("\nDeveloper command: ").strip()
        if cmd.lower() in {"exit", "quit"}:
            break
        result = handle_developer_command(cmd)
        print(result if result is not None else "Command not handled by developer_assistant.py")
